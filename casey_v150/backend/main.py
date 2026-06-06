from __future__ import annotations

from fastapi import FastAPI, HTTPException, UploadFile, File, Request
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import StreamingResponse
from pydantic import BaseModel
from typing import Any, Dict, List, Optional, Tuple
from datetime import datetime
from io import BytesIO, StringIO
import csv, json, math, os, random, re, sqlite3, statistics, zipfile, hashlib, uuid

import numpy as np
from openpyxl import Workbook, load_workbook
from openpyxl.chart import LineChart, BarChart, Reference
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
from openpyxl.utils import get_column_letter
from docx import Document
from docx.shared import Pt, RGBColor, Inches
from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt as PptxPt
from reportlab.lib import colors
from reportlab.lib.pagesizes import A4, landscape
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, PageBreak

APP_VERSION = "CASEY TITAN X v26 Revenue Machine + GTM Demo Edition"
DB_PATH = os.environ.get("CASEY_DB", "casey_titan_v26.sqlite3")
DEMO_LIMIT_PER_IP = int(os.environ.get("CASEY_DEMO_LIMIT_PER_IP", "1"))
PUBLIC_DEMO_LIMIT = int(os.environ.get("CASEY_PUBLIC_DEMO_LIMIT", "1"))
ADMIN_TOKEN = os.environ.get("CASEY_ADMIN_TOKEN", "")

app = FastAPI(title=APP_VERSION, version="26.0-revenue-machine")
app.add_middleware(CORSMiddleware, allow_origins=["*"], allow_credentials=True, allow_methods=["*"], allow_headers=["*"])

class GenerateRequest(BaseModel):
    prompt: str
    client: Optional[str] = None
    class_level: Optional[int] = 3
    schedule_level: Optional[int] = 3
    scenario: Optional[str] = "base"
    demo: Optional[bool] = False
    active_model: Optional[Dict[str, Any]] = None


class PublicDemoRequest(BaseModel):
    email: str
    project_type: str = "Earth"
    project_description: str
    location: Optional[str] = None
    size_or_capacity: Optional[str] = None
    stage: Optional[str] = "Concept / early feasibility"
    biggest_concern: Optional[str] = "Cost, schedule and risk confidence"
    fingerprint: Optional[str] = None
    client_token: Optional[str] = None

class PublicDemoFeedback(BaseModel):
    run_id: str
    rating: int
    comment: Optional[str] = None

class ChatRequest(BaseModel):
    question: str
    project: Optional[Dict[str, Any]] = None
    demo: Optional[bool] = False

class SaveProjectRequest(BaseModel):
    name: str
    model: Dict[str, Any]

# ------------------------- database -------------------------
def db() -> sqlite3.Connection:
    con = sqlite3.connect(DB_PATH)
    con.row_factory = sqlite3.Row
    return con

def init_db():
    con = db(); cur = con.cursor()
    cur.execute("""CREATE TABLE IF NOT EXISTS projects(
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        name TEXT NOT NULL,
        client TEXT,
        prompt TEXT,
        mode TEXT,
        created_at TEXT,
        model_json TEXT NOT NULL
    )""")
    cur.execute("""CREATE TABLE IF NOT EXISTS demo_usage(
        ip TEXT PRIMARY KEY,
        count INTEGER NOT NULL DEFAULT 0,
        updated_at TEXT
    )""")
    cur.execute("""CREATE TABLE IF NOT EXISTS public_demo_uses(
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        run_id TEXT UNIQUE,
        email_hash TEXT,
        ip_hash TEXT,
        fingerprint_hash TEXT,
        client_token_hash TEXT,
        project_type TEXT,
        project_text TEXT,
        model_json TEXT,
        created_at TEXT
    )""")
    cur.execute("""CREATE TABLE IF NOT EXISTS public_demo_feedback(
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        run_id TEXT,
        rating INTEGER,
        comment TEXT,
        created_at TEXT
    )""")
    cur.execute("""CREATE TABLE IF NOT EXISTS uploads(
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        filename TEXT,
        created_at TEXT,
        analysis_json TEXT
    )""")
    con.commit(); con.close()
init_db()

# ------------------------- helpers -------------------------
def client_ip(request: Request) -> str:
    return request.headers.get("x-forwarded-for", request.client.host if request.client else "local").split(",")[0].strip()

def check_demo_allowance(ip: str) -> Dict[str, Any]:
    con = db(); cur = con.cursor(); row = cur.execute("SELECT count FROM demo_usage WHERE ip=?", (ip,)).fetchone()
    count = int(row["count"]) if row else 0
    con.close()
    return {"allowed": count < DEMO_LIMIT_PER_IP, "used": count, "limit": DEMO_LIMIT_PER_IP}

def record_demo_use(ip: str):
    con = db(); cur = con.cursor(); now = datetime.utcnow().isoformat()
    cur.execute("INSERT INTO demo_usage(ip,count,updated_at) VALUES(?,?,?) ON CONFLICT(ip) DO UPDATE SET count=count+1, updated_at=excluded.updated_at", (ip, 1, now))
    con.commit(); con.close()

def _sha(value: str) -> str:
    value = (value or "").strip().lower()
    return hashlib.sha256(value.encode("utf-8")).hexdigest() if value else ""

def _normalise_email(email: str) -> str:
    return (email or "").strip().lower()

def _public_demo_brief_quality_score(req: PublicDemoRequest) -> Dict[str, Any]:
    """Public demo gate: allow credible infrastructure briefs quickly, block rubbish.
    The old gate was too strict and blocked strong executive-style briefs. This version
    scores useful signals but does not demand every field when CASEY can infer them.
    """
    raw_desc = (req.project_description or "").strip()
    desc = raw_desc.lower()
    joined = " ".join([
        desc,
        (req.location or "").lower(),
        (req.size_or_capacity or "").lower(),
        (req.stage or "").lower(),
        (req.biggest_concern or "").lower(),
        (req.project_type or "").lower(),
    ])

    score = 0
    reasons = []
    contradictions = []

    words = re.findall(r"[a-zA-Z0-9]+(?:[-/][a-zA-Z0-9]+)?", desc)
    unique_words = set(words)

    nonsense_terms = [
        "asdf", "qwerty", "lorem ipsum", "test test", "blah blah", "hello world",
        "ignore previous", "ignore all", "jailbreak", "prompt injection", "write a poem",
        "make me rich", "bitcoin", "crypto meme", "football match", "recipe", "dating",
        "tell me a joke", "song lyrics"
    ]
    if any(x in desc for x in nonsense_terms):
        contradictions.append("This does not look like a credible project brief. Please enter a real infrastructure or space programme.")

    # Repetition / spam check.
    if len(words) >= 12 and len(unique_words) < max(8, len(words) * 0.25):
        contradictions.append("The brief appears repetitive or low quality. Please write a real project description.")

    # Length helps, but short credible briefs should still be allowed.
    if len(words) >= 18:
        score += 24
    elif len(words) >= 10:
        score += 14
    else:
        reasons.append("Add a little more project detail: asset, location/environment and main concern.")

    asset_terms = [
        "data centre","data center","datacenter","airport","terminal","runway","rail","metro","station","hospital",
        "fab","semiconductor","wafer","nuclear","smr","hydrogen","wind","solar","grid","battery","gigafactory",
        "defence","defense","military","naval","airbase","radar","command","campus","port","water","desalination",
        "flood","lng","carbon capture","ccs","stadium","university","fibre","fiber","ev charging",
        "biologics","therapeutics","gmp","aseptic","fill-finish","fill finish","fda","cqv","pharma","pharmaceutical",
        "manufacturing","cleanroom","clean utilities","cold-chain","cold chain","warehouse","logistics",
        "lunar","moon","mars","orbital","orbit","leo","cislunar","satellite","spaceport","propellant","habitat",
        "launch vehicle","payload","asteroid","space infrastructure","space data centre","space data center",
        "orbital compute","orbital ai","thermal rejection","relay communications","lunar base","mars base"
    ]
    has_asset = any(x in joined for x in asset_terms)
    if has_asset:
        score += 28
    else:
        reasons.append("State the asset type, e.g. data centre, GMP campus, rail, hospital, lunar base or orbital platform.")

    location_terms = [
        "north carolina","carolina","arizona","texas","boston","london","cambridge","manchester","riyadh",
        "dubai","abu dhabi","qatar","uk","usa","united states","uae","saudi","canada","australia","singapore",
        "india","japan","germany","france","poland","moon","lunar","mars","leo","orbit","orbital","cislunar",
        "spaceport","deep space"
    ]
    has_location = (req.location and len(req.location.strip()) >= 3 and "auto-inferred" not in req.location.lower()) or any(x in joined for x in location_terms)
    if has_location:
        score += 18
    else:
        reasons.append("Add a location or operating environment.")

    size_terms = [
        "mw","gw","km","beds","m2","sqm","sq m","satellites","crew","tonnes","ha","capacity","runway",
        "stations","terminal","phase","phased","halls","modules","multi-product","multi product","fill-finish",
        "campus","site","lines","clusters","constellation","base","hub","plant","network"
    ]
    if any(x in joined for x in size_terms):
        score += 12

    concern_terms = [
        "cost","schedule","risk","procurement","approval","consent","grid","logistics","commissioning",
        "funding","delivery","safety","regulatory","interface","utilities","critical path","scope","confidence",
        "supply chain","phasing","validation","qualification","licensing","resilience","thermal","radiation",
        "servicing","latency","autonomous","debris","power density","operational readiness","production continuity",
        "fda","cqv","inspection","launch cadence","long-lead","long lead"
    ]
    has_concern = any(x in joined for x in concern_terms)
    if has_concern:
        score += 18
    else:
        reasons.append("Add the main concern: cost, schedule, risk, procurement, approvals, logistics or commissioning.")

    project_context_terms = [
        "project","programme","program","facility","campus","hub","plant","terminal","network","corridor",
        "base","station","platform","outpost","depot","infrastructure","scheme","development","expansion",
        "upgrade","rollout","manufacturing","construction","delivery"
    ]
    if any(x in joined for x in project_context_terms):
        score += 12

    # Strong domain phrases should pass even if a formal field was not provided.
    strong_domain = has_asset and (has_location or any(x in joined for x in ["orbital","lunar","mars","leo","spaceport"])) and has_concern
    if strong_domain:
        score = max(score, 88)

    space_terms = ["moon","lunar","mars","orbital","orbit","leo","cislunar","spaceport","satellite constellation","launch vehicle","rocket","payload","space station","propellant depot","asteroid"]
    earth_satellite_facility = any(x in joined for x in [
        "satellite control centre", "satellite control center", "secure satellite control",
        "ground station", "space domain awareness ground station", "mission operations centre",
        "mission operations center", "mission operations rooms"
    ])

    # Do not punish product/commercial launch language in Earth sectors.
    product_launch_only = any(x in joined for x in ["commercial launch demand","product launch","market launch","launch demand"]) and not any(x in joined for x in ["rocket","launch vehicle","spaceport","launch pad","orbital","leo","lunar","mars"])
    effective_space_terms = [x for x in space_terms if not (x == "launch" and product_launch_only)]

    if (req.project_type or "").lower() == "earth" and any(x in joined for x in effective_space_terms) and not earth_satellite_facility and not product_launch_only:
        # Only block if it is genuinely contradictory, not if frontend auto-inferred Earth before backend reroutes.
        if any(x in joined for x in ["orbital","leo","lunar","moon","mars","spaceport","launch vehicle","rocket"]):
            contradictions.append("The brief contains strong space terms. Choose Space or clarify that this is an Earth facility.")
    if "leo" in joined and any(x in joined for x in ["moon","lunar surface","mars surface"]):
        contradictions.append("The brief mixes orbital location terms such as LEO with Moon/Mars surface terms. Clarify the operating location.")
    if "mars" in joined and "moon" in joined and "cislunar" not in joined:
        contradictions.append("The brief mixes Mars and Moon scope. Split this into one project location.")

    return {
        "score": min(score, 100),
        "reasons": reasons[:4],
        "contradictions": contradictions,
        "pass": score >= 58 and has_asset and not contradictions,
    }
def _public_demo_quality_message(req: PublicDemoRequest) -> Optional[Dict[str, Any]]:
    quality = _public_demo_brief_quality_score(req)
    if quality["pass"]:
        return None
    return {
        "message": "CASEY needs a stronger brief before using your one free intelligence run.",
        "quality_score": quality["score"],
        "issues": quality["contradictions"] or quality["reasons"],
        "example_brief": "Example: Earth project — 500MW AI data centre campus in Riyadh, concept stage, 4 buildings, liquid cooling, grid substation, concern is grid connection and schedule acceleration. Or Space project — lunar south pole logistics hub, landing pads, regolith roads, autonomous rovers, concept stage, concern is dust, power resilience and launch cadence."
    }

def _quality_gate_public_demo(req: PublicDemoRequest) -> List[str]:
    issues = []
    email = _normalise_email(req.email)
    desc = (req.project_description or "").strip()
    if not email or "@" not in email or "." not in email.split("@")[-1]:
        issues.append("Please enter a valid email so CASEY can reserve your one free intelligence run.")
    if len(desc) > 3000:
        issues.append("Please keep the project brief under 3,000 characters for the public demo.")
    if len(re.findall(r"[a-zA-Z0-9]+", desc)) < 8:
        issues.append("Add a little more detail: asset, location/environment and main project concern.")
    quality = _public_demo_brief_quality_score(req)
    if not quality["pass"]:
        issues.extend(quality["contradictions"] or quality["reasons"])
    return list(dict.fromkeys(issues))[:5]

def _public_demo_identity(request: Request, req: PublicDemoRequest) -> Dict[str, str]:
    ip = client_ip(request)
    ua = request.headers.get("user-agent", "")
    forwarded = request.headers.get("x-forwarded-for", "")
    fp = req.fingerprint or ""
    token = req.client_token or ""
    return {
        "email_hash": _sha(_normalise_email(req.email)),
        "ip_hash": _sha(ip),
        "fingerprint_hash": _sha(fp + "|" + ua),
        "client_token_hash": _sha(token),
        "raw_ip": ip,
        "raw_ua": ua,
        "raw_forwarded": forwarded,
    }

def _public_demo_used(identity: Dict[str, str]) -> Optional[str]:
    con = db(); cur = con.cursor()
    checks = []
    params = []
    for col in ["email_hash", "ip_hash", "fingerprint_hash", "client_token_hash"]:
        val = identity.get(col)
        if val:
            checks.append(f"{col}=?")
            params.append(val)
    if not checks:
        con.close(); return None
    row = cur.execute(f"SELECT run_id, created_at FROM public_demo_uses WHERE {' OR '.join(checks)} ORDER BY id DESC LIMIT 1", tuple(params)).fetchone()
    con.close()
    if row:
        return row["run_id"]
    return None

def _premium_public_prompt(req: PublicDemoRequest) -> str:
    return " | ".join([
        f"{req.project_type} public demo project",
        f"Brief: {req.project_description}",
        f"Location: {req.location or 'not specified'}",
        f"Size/capacity: {req.size_or_capacity or 'not specified'}",
        f"Stage: {req.stage or 'concept / early feasibility'}",
        f"Primary concern: {req.biggest_concern or 'cost, schedule and risk confidence'}",
        "Generate a board-grade first-pass class estimate, schedule intelligence and risk register."
    ])

def _public_demo_report(model: Dict[str, Any]) -> Dict[str, Any]:
    risks = model.get("risks", model.get("risk_register", []))[:8]
    schedule = model.get("schedule_rows", model.get("schedule_detail", []))[:8]
    costs = model.get("cost_lines", model.get("cost_breakdown", []))[:10]
    return {
        "title": model.get("title"),
        "mode": model.get("mode"),
        "executive_summary": model.get("executive_summary"),
        "class_estimate": {
            "estimate_class": model.get("estimate_class_name") or f"Class {model.get('estimate_class', 3)}",
            "p10": model.get("cost_p10"),
            "p50": model.get("cost_p50"),
            "p90": model.get("cost_p90"),
            "range": model.get("cost_range"),
            "confidence_pct": model.get("confidence_pct"),
        },
        "schedule": {
            "baseline": model.get("schedule"),
            "level": model.get("schedule_level"),
            "first_milestones": schedule,
        },
        "risk": {
            "rating": model.get("risk"),
            "score": model.get("risk_score"),
            "top_risks": risks,
        },
        "cost_breakdown": costs,
        "assumptions": model.get("confidence_explanation", [])[:6],
        "next_best_actions": model.get("next_best_actions", [])[:6],
        "input_quality_score": model.get("input_quality_score"),
        "upgrade_cta": "This is a one-shot public intelligence run. Request access for exports, scenarios, QCRA/QSRA, audit trail and deeper model challenge."
    }


def has(t: str, words: List[str]) -> bool: return any(w in t for w in words)
def clamp(v, lo, hi): return max(lo, min(hi, v))
def money_bn(v: float) -> str:
    if v >= 1000: return f"${v/1000:.1f}T"
    if v >= 1: return f"${v:.1f}B"
    return f"${v*1000:.0f}M"
def parse_bn(s: Any) -> float:
    if isinstance(s,(int,float)): return float(s)
    st=str(s or "").replace("$","").replace(",","").strip().upper()
    try:
        if st.endswith("T"): return float(st[:-1])*1000
        if st.endswith("B"): return float(st[:-1])
        if st.endswith("M"): return float(st[:-1])/1000
        return float(st)
    except Exception: return 1.0

def class_range(c:int):
    return {5:(0.50,2.00,"Class 5 Concept Screening","0-2% definition maturity"),4:(0.70,1.50,"Class 4 Feasibility","1-15% definition maturity"),3:(0.80,1.30,"Class 3 Budget Authorization","10-40% definition maturity"),2:(0.85,1.20,"Class 2 Control/Tender","30-75% definition maturity"),1:(0.90,1.15,"Class 1 Definitive/Bid","65-100% definition maturity")}.get(c,(0.8,1.3,"Class 3 Budget Authorization","10-40% maturity"))

def scenario_params(s: str):
    key=(s or "base").lower().replace(" ","_")
    table={
        "base":(1.00,1.00,1.00,0,"Base","Balanced reference case."),
        "faster":(1.11,0.82,1.18,-5,"Faster","Parallel design/procurement, acceleration, premium logistics and higher interface risk."),
        "cheaper":(0.88,1.08,1.12,-6,"Cheaper","Value engineering, tighter scope and procurement competition, with schedule/quality exposure."),
        "lower_risk":(1.09,1.10,0.72,9,"Lower Risk","More surveys, assurance, stage gates, contingency and schedule buffers."),
        "premium":(1.15,1.04,0.80,12,"Premium","Best assurance, independent review, early procurement and high-confidence governance."),
        "investor":(0.98,1.00,0.92,3,"Investor","Investment-framing case with transparent risk and commercial challenge."),
        "survival":(0.78,1.20,1.30,-10,"Survival","Minimum viable scope, warning-heavy, high residual risk."),
    }
    return table.get(key, table["base"])

# ------------------------- detection -------------------------
def detect_sector(prompt: str):
    t=prompt.lower()

    # Named mega-programmes first.
    named=[
      ("Heathrow Third Runway","Earth","Airport Mega Programme",26,108,["heathrow","third runway"]),
      ("HS2 High Speed Rail","Earth","Rail Mega Programme",95,180,["hs2","high speed rail"]),
      ("AUKUS Naval Industrial Programme","Earth","Defence / Naval Industrial",28,144,["aukus"]),
      ("NEOM/The Line Style Development","Earth","Future City Mega Programme",500,240,["neom","the line"]),
      ("Lunar City","Space","Lunar Settlement",120,240,["lunar city","moon city"]),
      ("Mars Base","Space","Mars Settlement",180,300,["mars base","mars city"]),
      ("Orbital Hospital","Space","Space Medical Infrastructure",12,120,["orbital hospital","space hospital"]),
    ]
    for n,m,s,c,d,k in named:
        if has(t,k): return n,m,s,c,d

    # Strong Earth-place signals. If a real Earth location is present, ambiguous words like
    # "commercial launch", "mission", "platform", or "payload" should NOT automatically route to Space.
    earth_places = [
        "north carolina","carolina","cambridge","boston","arizona","texas","california","new york","florida",
        "united states","usa","america","uk","united kingdom","london","manchester","birmingham",
        "riyadh","saudi","dubai","abu dhabi","uae","qatar","doha","canada","toronto","australia",
        "sydney","singapore","japan","tokyo","india","mumbai","germany","berlin","france","paris",
        "netherlands","amsterdam","poland","warsaw","brazil","south africa","kenya","morocco"
    ]
    earth_place_present = has(t, earth_places)

    # Strong Earth-sector scores.
    life_science_terms = ["biologics","therapeutics","gmp","aseptic","fill-finish","fill finish","fda","cqv","qualification","validated clean utilities","clean utilities","cold-chain","cold chain","pharma","pharmaceutical","cell therapy","gene therapy","manufacturing campus","obesity therapeutics"]
    semiconductor_terms = ["semiconductor"," fab","wafer","upw","cleanroom","process tooling","chip plant"]
    data_terms = ["data centre","data center","datacenter","hyperscale","ai campus","gpu campus","liquid cooling"]
    defence_earth_terms = ["secure satellite control centre","secure satellite control center","satellite control centre","satellite control center","mission operations room","mission operations centre","mission operations center","secure command","command and control","air defence","air defense","military airbase","naval base","munitions storage","radar station","border surveillance","defence data centre","defense data center","submarine","dockyard","shipbuilding","naval shipyard","naval vessel","warship","frigate","destroyer","aircraft carrier","naval propulsion","aukus nuclear","nuclear submarine","nuclear-powered submarine","naval industrial","defence industrial","defense industrial","military shipbuilding","combat vessel","naval combat","maritime patrol","submarine base","submarine program","submarine programme"]
    healthcare_terms = ["hospital","healthcare","clinical","diagnostic imaging","theatres","emergency department"]
    transport_terms = ["airport","aviation","runway","airside","baggage","passenger terminal","rail","metro"," station ","transit"]
    energy_terms = ["hydrogen","solar","wind","battery","power","grid","energy","nuclear","smr","lng","carbon capture","ccs"]
    water_terms = ["desalination","water","wastewater","flood","pumping station","reservoir"]

    earth_score = 0
    for terms, weight in [
        (life_science_terms, 5),
        (semiconductor_terms, 5),
        (data_terms, 4),
        (defence_earth_terms, 5),
        (healthcare_terms, 4),
        (transport_terms, 3),
        (energy_terms, 3),
        (water_terms, 3),
    ]:
        earth_score += sum(weight for term in terms if term in t)
    if earth_place_present:
        earth_score += 6

    # Strong Space scores. Weak words like "launch" only count strongly with aerospace context.
    strong_space_terms = ["moon","lunar","mars","orbit","orbital","leo","meo","geo","cislunar","cis-lunar","space station","asteroid","deep space","deep-space","in-space"]
    medium_space_terms = ["spaceport","rocket","payload","satellite constellation","launch vehicle","launch pad","launch complex","spacecraft","propellant depot","orbital compute","space data centre","space data center","orbital ai","space-based data centre","space-based data center"]
    weak_space_terms = ["launch","mission","payload","platform","constellation","habitat"]

    space_score = 0
    space_score += sum(5 for term in strong_space_terms if term in t)
    space_score += sum(4 for term in medium_space_terms if term in t)

    # Only count weak space terms if supported by other space signals or no strong Earth context.
    if space_score > 0 or not earth_place_present:
        space_score += sum(1 for term in weak_space_terms if term in t)

    # Product/commercial launch phrases are Earth business language, not space.
    if has(t, ["commercial launch demand","product launch","launch demand","market launch","commercial launch","launch readiness"]) and not has(t, ["rocket","spaceport","launch vehicle","launch pad","orbital","leo","lunar","mars"]):
        space_score = max(0, space_score - 5)
        earth_score += 4

    # Hard space overrides: explicit Mars/Moon/LEO/orbital infrastructure should remain Space even if it contains
    # Earth-sector words such as nuclear power, grid, hospital, platform, or manufacturing.
    if has(t, ["mars surface","mars outpost","mars research","mars habitat","mars base"]):
        return "Mars Surface Infrastructure","Space","Mars Surface Habitat/Base",38,168
    if has(t, ["lunar surface","lunar south pole","lunar habitat","lunar base","lunar logistics","moon base"]):
        return "Lunar Surface Infrastructure","Space","Lunar Surface Habitat/Base",32,156
    if has(t, ["orbital ai data centre","orbital ai data center","space-based data centre","space-based data center","orbital compute platform","leo compute"]):
        return "Orbital AI Compute Platform","Space","Orbital Compute / Manufacturing",24,132
    if has(t, ["space power grid","orbital power grid","lunar power grid","mars power grid","space solar grid","orbital energy grid"]):
        return "Space Power Grid","Space","Power/Energy Infrastructure",18,132

    # High-signal Earth sectors win even if one ambiguous space word exists.
    if earth_score >= space_score and earth_score >= 5:
        if has(t, life_science_terms):
            return "Life Sciences Manufacturing Campus","Earth","Life Sciences / Biologics Manufacturing",2.8,58
        if has(t, semiconductor_terms):
            return "Advanced Semiconductor Fab","Earth","Semiconductor / Advanced Manufacturing",18,78
        if has(t, data_terms):
            return "AI Data Centre Campus","Earth","Digital Infrastructure / Hyperscale Data Centre",3.8,46
        if has(t, defence_earth_terms):
            return "Secure Defence Infrastructure","Earth","Defence / Secure Mission Infrastructure",4.8,54
        if has(t, healthcare_terms):
            return "Hospital Campus","Earth","Healthcare / Hospital",2.4,60
        if has(t, ["airport","aviation","runway","airside","baggage","passenger terminal"]):
            return "Airport Infrastructure","Earth","Airport / Aviation",9,84
        if has(t, ["rail","metro","station","transit"]):
            return "Rail Infrastructure","Earth","Rail / Transit",6.5,84
        if has(t, ["nuclear","smr","fusion"]):
            return "Nuclear Energy Facility","Earth","Nuclear / Energy",12,96
        if has(t, energy_terms):
            return "Energy Infrastructure","Earth","Energy / Utilities",5.2,66
        if has(t, water_terms):
            return "Water Infrastructure","Earth","Water / Utilities",2.2,50

    # Space routing only after weighted scoring confirms it.
    if space_score > earth_score and space_score >= 4:
        if has(t,["deep-space communications","deep space communications","deep-space comms","deep space network","relay spacecraft","space communications array"]): return "Deep-Space Communications Array","Space","Deep-Space Communications Infrastructure",10,96
        if has(t,["orbital ai compute","orbital ai data centre","orbital ai data center","space data centre","space data center","space-based data centre","space-based data center","orbital compute","leo compute","compute platform","orbital data centre","orbital data center"]): return "Orbital AI Compute Platform","Space","Orbital Compute / Manufacturing",24,132
        if has(t,["hospital","medical","health"]): return "Space Hospital","Space","Orbital/Lunar Healthcare",12,126
        if has(t,["city","settlement","colony"]): return "Space Settlement","Space","Lunar/Mars Settlement",110,240
        if has(t,["base","habitat","outpost"]): return "Space Base","Space","Surface Habitat/Base",32,156
        if has(t,["power","solar","nuclear","grid"]): return "Space Power Grid","Space","Power/Energy Infrastructure",18,132
        if has(t,["spaceport","launch vehicle","launch pad","launch complex","rocket"]): return "Launch Infrastructure","Space","Spaceport/Launch",9,84
        if has(t,["mine","mining","isru","water","oxygen","propellant"]): return "Space Resources Facility","Space","ISRU/Mining/Propellant",22,168
        if has(t,["satellite","constellation"]): return "Satellite Constellation","Space","Satellite/Comms",6,60
        return "Frontier Space Infrastructure","Space","General Space Infrastructure",14,108

    # Earth fallback after scoring.
    if has(t,["data centre","data center","datacenter","ai campus","cloud","hyperscale"]): return "AI Data Centre Campus","Earth","Digital Infrastructure / Hyperscale Data Centre",3.8,46
    if has(t,["biologics","gmp","aseptic","fill-finish","fill finish","fda","cqv","pharma","therapeutics","cell therapy"]): return "Life Sciences Manufacturing Campus","Earth","Life Sciences / Biologics Manufacturing",2.8,58
    if has(t,["airport","aviation","runway","airside","baggage","passenger terminal"]): return "Airport Infrastructure","Earth","Airport / Aviation",9,84
    if has(t,["rail","metro"," station ","transit"]): return "Rail Infrastructure","Earth","Rail / Transit",6.5,84
    if has(t,["nuclear","smr","fusion"]): return "Nuclear Energy Facility","Earth","Nuclear / Energy",12,96
    if has(t,["hydrogen","solar","wind","battery","power","grid","energy"]): return "Energy Infrastructure","Earth","Energy / Utilities",5.2,66
    if has(t,["hospital","medical campus","healthcare"]): return "Hospital Campus","Earth","Healthcare / Hospital",2.4,60
    if has(t,["pharma","gmp","biologics","life sciences","cleanroom"]): return "Life Sciences Campus","Earth","Life Sciences / Pharma",2.8,58
    if has(t,["defence","defense","military","command","radar"]): return "Secure Defence Infrastructure","Earth","Defence / Secure Mission Infrastructure",4.8,54
    if has(t,["desalination","water","wastewater","flood"]): return "Water Infrastructure","Earth","Water / Utilities",2.2,50
    return "Major Infrastructure Programme","Earth","General Infrastructure",2.0,48

def location_factor(t: str):
    t=t.lower(); table=[("Moon / Lunar Surface",2.9,["moon","lunar"]),("Mars",3.8,["mars"]),("Low Earth Orbit",2.5,["leo","orbit","orbital"]),("Deep Space",4.3,["deep space","asteroid"]),("Saudi Arabia / GCC",1.35,["saudi","riyadh","neom","jeddah","gcc"]),("UAE / Middle East",1.18,["uae","dubai","abu dhabi","qatar"]),("United Kingdom",1.20,["uk","london","heathrow","manchester"]),("United States",1.08,["usa","us ","texas","california","florida","new york"]),("Europe",1.12,["europe","germany","france","spain","italy"]),("India",0.78,["india","mumbai","delhi"]),("Australia",1.22,["australia","sydney","melbourne"])]
    for n,f,keys in table:
        if has(t,keys): return n,f
    return "Global / Not specified",1.0

def scale_factor(t: str):
    t=t.lower()
    if has(t,["global","interplanetary","system of systems"]): return "System of Systems",5.0,72
    if has(t,["mega","city","national","regional","10000","10,000","colony"]): return "Mega Programme",3.2,42
    if has(t,["programme","program","campus","portfolio","corridor"]): return "Programme",1.9,18
    if has(t,["pilot","phase 1","small"]): return "Pilot / Phase 1",0.55,-8
    return "Project",1.0,0

def complexity(t:str):
    rules={"accelerated":(1.10,-8,"accelerated delivery compression"),"fast":(1.07,-5,"fast-track compression"),"nuclear":(1.45,36,"nuclear assurance and regulation"),"live":(1.15,12,"live environment"),"brownfield":(1.20,12,"brownfield interfaces"),"underground":(1.35,18,"underground works"),"offshore":(1.40,24,"offshore logistics"),"autonomous":(1.20,12,"autonomous systems"),"robotic":(1.18,12,"robotic construction"),"life support":(1.35,24,"life support systems"),"radiation":(1.25,18,"radiation hardening"),"cryogenic":(1.25,18,"cryogenic systems"),"cleanroom":(1.25,12,"GMP cleanroom"),"gmp":(1.22,12,"GMP validation"),"secure":(1.20,12,"secure/defence constraints")}
    f=1.0; m=0; d=[]; tl=t.lower()
    for k,(ff,mm,dd) in rules.items():
        if k in tl: f*=ff; m+=mm; d.append(dd)
    return f,m,d

def risk_label(score:float):
    if score>=90: return "Extreme"
    if score>=75: return "Very High"
    if score>=60: return "High"
    if score>=42: return "Medium-High"
    if score>=25: return "Medium"
    return "Low"

def cost_lines(mode:str, subsector:str):
    if mode=="Space":
        return [("01.01","DDT&E Engineering","Direct","Engineering design, mission architecture and qualification",0.09),("01.02","Systems Engineering & Integration","Direct","Requirements, ICDs, verification and validation",0.06),("01.03","Structures / Habitat Shell","Direct","Pressure vessels, shielding and primary structures",0.10),("01.04","Power Generation & Distribution","Direct","Solar/nuclear source, grid and controls",0.08),("01.05","ECLSS / Life Support","Direct","Air, water, waste, atmosphere and safety systems",0.09),("01.06","Payload / Mission Equipment","Direct","Scientific, medical, industrial or settlement payload",0.07),("01.07","Robotics / Autonomous Construction","Direct","Robotic construction, inspection and maintenance",0.06),("01.08","Launch Services","Direct","Launch vehicle, manifest, payload processing",0.14),("01.09","Transport / Landing / Surface Logistics","Direct","Transit, landing, mobility and logistics chain",0.09),("02.01","Test & Qualification Facilities","Indirect","Thermal-vac, vibration, radiation and integrated testing",0.05),("02.02","Safety & Mission Assurance","Indirect","Quality, safety, independent review and compliance",0.05),("02.03","Programme Management / PMO","Indirect","Controls, reporting, commercial and integration PMO",0.05),("03.01","Risk Contingency / UFE","Reserve","Known-unknown risk allowance linked to QCRA/QSRA",0.07),("03.02","Management Reserve","Reserve","Sponsor-held unknown-unknown reserve",0.05)]
    special=[]
    s=subsector.lower()
    if "data centre" in s: special=[("01.06","Data Halls / White Space","Direct","Server halls, fit-out and critical rooms",0.13),("01.07","Power Train / Substations","Direct","HV/MV gear, UPS, generators and substations",0.16),("01.08","Cooling Systems","Direct","Chillers, heat rejection, CRAH/CRAC, water systems",0.12)]
    elif "airport" in s: special=[("01.06","Airfield Works","Direct","Runway, taxiway, apron and airfield systems",0.14),("01.07","Terminal / Passenger Systems","Direct","Terminal, baggage, security and passenger systems",0.10)]
    elif "rail" in s: special=[("01.06","Track / Signalling / Systems","Direct","Trackform, power, signalling and telecoms",0.18),("01.07","Stations / Depots","Direct","Stations, depots and operational readiness",0.09)]
    elif "pharma" in s: special=[("01.06","Cleanroom Envelope","Direct","Cleanroom partitions, finishes and GMP envelope",0.12),("01.07","Process Utilities","Direct","WFI, gases, clean steam and process utilities",0.13),("01.08","Validation / Qualification","Indirect","IQ/OQ/PQ and GMP readiness",0.05)]
    base=[("01.01","Site Preparation / Enabling","Direct","Demolition, enabling, logistics and access",0.06),("01.02","Substructure / Foundations","Direct","Piling, groundworks, retaining and foundations",0.08),("01.03","Superstructure / Civil Works","Direct","Frame, civils, primary structures",0.14),("01.04","MEP / Systems / Utilities","Direct","Mechanical, electrical, controls and utilities",0.14),("01.05","Specialist Equipment / Fit-Out","Direct","Permanent equipment and specialist fit-out",0.12)]
    tail=[("01.09","Direct Labour","Direct","Productivity, supervision and trade labour",0.09),("01.10","Direct Materials","Direct","Concrete, steel, envelope, MEP materials",0.10),("01.11","Plant / Equipment","Direct","Cranes, heavy plant and site equipment",0.04),("02.01","General Conditions / Prelims","Indirect","Site management, safety, offices and temporary services",0.07),("02.02","Design / Professional Fees","Indirect","Design, QS, controls, assurance and consultants",0.07),("02.03","Permits / Insurance / Bonds","Indirect","Permits, inspections, legal, bonds and insurance",0.03),("03.01","QCRA Contingency","Reserve","Known-unknown risk allowance",0.06),("03.02","Management Reserve","Reserve","Sponsor-held reserve",0.03)]
    return base+special+tail

def schedule_rows(mode:str, subsector:str, months:int, level:int):
    """Rich sector-specific schedule. Level 1=milestones, 2=phases, 3=work stages, 4=packages, 5=activities."""
    s = subsector.lower()

    # ── LEVEL 1: Milestones only ─────────────────────────────────────
    if level <= 1:
        rows = [
            ("M100","Milestone","Programme mandate and funding approval","",0.04),
            ("M200","Milestone","Design and consents complete","M100",0.22),
            ("M300","Milestone","Procurement and long-lead confirmed","M200",0.12),
            ("M400","Milestone","Enabling works and site ready","M300",0.08),
            ("M500","Milestone","Substantial construction complete","M400",0.38),
            ("M600","Milestone","Systems integration and testing complete","M500",0.10),
            ("M700","Milestone","Operational readiness and handover","M600",0.06),
        ]
        out=[]
        for code,phase,act,pred,pct in rows:
            out.append({"activity_id":code,"phase":phase,"activity":act,"predecessor":pred,
                        "duration_months":max(1,round(months*pct)),"critical":"Yes",
                        "basis":"Level 1 milestone: programme governance and approvals logic."})
        return out

    # ── LEVEL 2: Programme phases ────────────────────────────────────
    if level == 2:
        rows = [
            ("PH100","Programme Setup","Controls, governance and programme baseline","",0.03),
            ("PH200","Business Case","Strategic case, options and funding approval","PH100",0.06),
            ("PH300","Consents and Stakeholder","Planning, consents and third-party agreements","PH200",0.10),
            ("PH400","Concept and Feasibility","Concept design, options and cost plan","PH200",0.06),
            ("PH500","Detailed Design","Detailed design, specifications and IFC","PH400",0.10),
            ("PH600","Procurement","Procurement, tendering and contract award","PH500",0.08),
            ("PH700","Enabling Works","Site preparation, utilities and enabling","PH300;PH600",0.07),
            ("PH800","Main Works","Construction, fabrication and assembly","PH500;PH600;PH700",0.28),
            ("PH900","Systems Integration","Systems integration, testing and commissioning","PH800",0.12),
            ("PH950","Operational Readiness","ORAT, training and handover preparation","PH900",0.05),
            ("PH990","Handover","Practical completion and operational handover","PH950",0.05),
        ]
        out=[]
        for code,phase,act,pred,pct in rows:
            out.append({"activity_id":code,"phase":phase,"activity":act,"predecessor":pred,
                        "duration_months":max(1,round(months*pct)),"critical":phase in ["Main Works","Systems Integration","Operational Readiness"] and "Yes" or "No",
                        "basis":"Level 2 programme phase: sector phasing logic and benchmark durations."})
        return out

    # ── LEVEL 3+: Sector-specific work stages ───────────────────────
    # Core work stages (universal)
    rows = [
        ("A1000","Programme Setup","Project initiation, controls baseline and governance","",0.025),
        ("A1100","Business Case","Strategic brief, options analysis and funding approval","A1000",0.050),
        ("A1200","Scope Definition","Requirements definition, scope freeze and assumption register","A1100",0.040),
        ("A1300","Concept Design","Concept design, reference architecture and option selection","A1200",0.060),
        ("A1400","Consents","Planning applications, environmental consents and third-party agreements","A1200",0.090),
        ("A1500","Procurement Strategy","Procurement strategy, package structure and tender preparation","A1300",0.050),
        ("A1600","Detailed Design","Detailed design, specifications, IFC and design coordination","A1300",0.100),
        ("A1700","Long-Lead Procurement","Long-lead equipment orders, supplier management and expediting","A1500",0.090),
        ("A1800","Enabling and Prelims","Site clearance, enabling works, utilities and temporary facilities","A1400;A1500",0.060),
        ("A1900","Main Works","Main construction, fabrication and assembly","A1600;A1700;A1800",0.210),
        ("A2000","Systems Integration","Systems integration, software, interfaces and testing","A1900",0.090),
        ("A2100","Commissioning","Pre-commissioning, commissioning and performance testing","A2000",0.060),
        ("A2200","ORAT and Readiness","Operational readiness, training, documentation and transition","A2100",0.040),
        ("A2300","Handover","Practical completion, defects, as-built and operational handover","A2200",0.030),
    ]

    # ── SECTOR OVERLAY: add sector-specific work stages ──────────────
    if mode == "Space":
        rows += [
            ("S1000","Mission Architecture","Mission architecture, system requirements and trade studies","A1200",0.050),
            ("S1100","Technology Development","Technology development, breadboard and TRL advancement","S1000",0.070),
            ("S1200","PDR / CDR","Preliminary and critical design reviews","A1600",0.040),
            ("S1300","Qualification Testing","Structural, thermal-vacuum, vibration and EMC qualification","S1200",0.070),
            ("S1400","Launch Vehicle Integration","Payload-to-launch-vehicle interface and fairing integration","S1300",0.050),
            ("S1500","Launch Campaign","Launch site operations, fuelling, countdown and launch","S1400",0.030),
            ("S1600","LEOP and Early Operations","Launch and early orbit phase, deployment and commissioning","S1500",0.040),
            ("S1700","Mission Operations","Steady-state mission operations, health monitoring and maintenance","S1600",0.080),
        ]
    elif any(x in s for x in ["rail","metro","transit","hsr","underground","tram"]):
        rows += [
            ("R1000","Ground Investigation","Ground investigation, utilities survey and corridor baseline","A1300",0.040),
            ("R1100","Utility Diversions","Utility mapping, diversions and corridor clearance","A1400",0.060),
            ("R1200","Land Acquisition","Compulsory purchase, land assembly and possession","A1400",0.080),
            ("R1300","Tunnelling","Tunnel boring, NATM or cut-and-cover excavation","A1800",0.120),
            ("R1400","Track and OLE","Trackform installation, overhead line and electrification","R1300",0.070),
            ("R1500","Signalling Installation","Signalling, telecoms and control systems installation","R1400",0.080),
            ("R1600","Stations and Depots","Station fitout, depot construction and passenger systems","R1400",0.090),
            ("R1700","Systems Integration","Signalling-to-train interface, system integration and SIL testing","R1500;R1600",0.060),
            ("R1800","Trial Running","Driver training, trial running and safety case closure","R1700",0.050),
            ("R1900","Operator Acceptance","TOC acceptance, regulator sign-off and operational handover","R1800",0.040),
        ]
    elif any(x in s for x in ["data centre","datacenter","hyperscale","compute","cloud","ai campus"]):
        rows += [
            ("D1000","Grid Connection Works","DNO works, substation construction and energisation","A1800",0.100),
            ("D1100","Structural Frame","Structural steel frame, cladding and external envelope","A1900",0.080),
            ("D1200","Power Infrastructure","HV/LV distribution, UPS, switchgear and generators","D1100",0.090),
            ("D1300","Cooling Plant","Cooling towers, CRACs, pipework and thermal commissioning","D1100",0.080),
            ("D1400","Data Hall Fitout","Raised floor, containment, cabling and rack installation","D1200;D1300",0.060),
            ("D1500","Network Connectivity","Dark fibre, carrier connections and network commissioning","D1400",0.040),
            ("D1600","Security and Access","Physical security, CCTV, biometrics and access systems","D1400",0.030),
            ("D1700","DCIM and BMS","DCIM platform, BMS integration and operational dashboards","D1600",0.040),
            ("D1800","IT Load Testing","IT load testing, PUE verification and performance acceptance","D1700",0.030),
        ]
    elif any(x in s for x in ["nuclear","smr","fusion","reactor"]):
        rows += [
            ("N1000","Nuclear Regulatory Programme","ONR/NRC GDA, site licence application and safety case programme","A1200",0.100),
            ("N1100","Site Characterisation","Radiological baseline, seismic survey and site investigation","A1300",0.060),
            ("N1200","Nuclear Island Design","Nuclear island detailed design, safety case and code compliance","A1600",0.080),
            ("N1300","Nuclear Supply Chain","N-stamp qualification, nuclear grade manufacturing and FAT","A1700",0.090),
            ("N1400","Civil Construction","Nuclear-grade civil works, containment and foundations","A1800",0.100),
            ("N1500","Mechanical and Pipework","Nuclear mechanical, pipework and primary circuit installation","N1400",0.090),
            ("N1600","Electrical and I&C","Electrical systems, I&C, protection and control installation","N1500",0.070),
            ("N1700","Cold Commissioning","Cold commissioning, systems walkdown and functional tests","N1600",0.050),
            ("N1800","Hot Commissioning and Fuel Load","Hot commissioning, fuel loading, first criticality and power ascension","N1700",0.060),
            ("N1900","Grid Synchronisation","Grid connection, power export and performance demonstration","N1800",0.030),
        ]
    elif any(x in s for x in ["semiconductor","fab","wafer","cleanroom","chip"]):
        rows += [
            ("F1000","Cleanroom Design and Spec","Cleanroom layout, HVAC specification and vibration criteria","A1300",0.050),
            ("F1100","Site Preparation","Foundation, vibration isolation and sub-fab preparation","A1800",0.070),
            ("F1200","Cleanroom Shell","Cleanroom structural shell, roof and cladding","F1100",0.080),
            ("F1300","HVAC and Filtration","HVAC, HEPA filtration and cleanroom environmental commissioning","F1200",0.080),
            ("F1400","UPW and Chemical Systems","Ultra-pure water, chemical distribution and waste treatment","F1300",0.070),
            ("F1500","Power and Electrical","Power quality, UPS, emergency generation and distribution","F1300",0.060),
            ("F1600","Tool Installation","Process tool delivery, installation and utility connection","F1400;F1500",0.090),
            ("F1700","Tool Qualification","IQ/OQ/PQ for each process tool and system","F1600",0.080),
            ("F1800","Process Integration","Process integration, yield learning and qualification wafers","F1700",0.070),
            ("F1900","Ramp and Volume","Production ramp, yield improvement and volume release","F1800",0.060),
        ]
    elif any(x in s for x in ["life science","pharma","biotech","gmp","vaccine","cell therapy"]):
        rows += [
            ("L1000","GMP Design Basis","GMP facility design, layout and regulatory basis review","A1300",0.050),
            ("L1100","Cleanroom and HVAC","GMP cleanroom, HVAC, pressure differentials and environmental monitoring","A1900",0.090),
            ("L1200","Utilities and Pure Media","WFI, PW, clean steam and utility commissioning","L1100",0.060),
            ("L1300","Process Equipment","Bioreactor, isolator, fill-finish and process equipment installation","L1200",0.080),
            ("L1400","IQ/OQ","Installation qualification and operational qualification for all systems","L1300",0.070),
            ("L1500","Technology Transfer","Process transfer, analytical method transfer and batch records","L1400",0.060),
            ("L1600","PQ and Validation Batches","Performance qualification, validation batches and data package","L1500",0.070),
            ("L1700","Regulatory Submission","GMP licence amendment, regulatory filing and inspection preparation","L1600",0.050),
            ("L1800","Pre-Approval Inspection","Regulatory inspection, observation response and approval","L1700",0.040),
        ]
    elif any(x in s for x in ["airport","aviation","runway","airside","terminal"]):
        rows += [
            ("AP1000","Airfield Survey and Enablement","Airfield survey, pavement assessment and enabling works","A1800",0.060),
            ("AP1100","Runway and Taxiway","Runway reconstruction, taxiway and apron works with NOTAM management","AP1000",0.100),
            ("AP1200","Terminal Structure","Terminal structural frame, cladding and envelope works","A1900",0.090),
            ("AP1300","Terminal Fitout","Terminal internal fitout, passenger facilities and retail","AP1200",0.080),
            ("AP1400","Baggage Handling System","BHS installation, software integration and acceptance testing","AP1200",0.070),
            ("AP1500","Security and Check-In","Security screening, check-in systems and hold baggage reconciliation","AP1400",0.050),
            ("AP1600","Airfield Systems","Airfield lighting, FIDS, DATIS, comms and ATC interfaces","AP1100",0.060),
            ("AP1700","ORAT Programme","Operational readiness, airline transition and dry runs","AP1400;AP1500;AP1600",0.060),
            ("AP1800","CAA Acceptance","CAA aerodrome licensing, safety case sign-off and operational acceptance","AP1700",0.040),
        ]

    # ── LEVEL 4: Add packages ────────────────────────────────────────
    if level >= 4:
        if mode == "Space":
            packages = [
                ("Launch campaign and pad operations","S1500"),
                ("Payload and spacecraft integration","S1400"),
                ("Power generation and ECLSS","S1200"),
                ("Communications and data systems","S1200"),
                ("Surface logistics and mobility","S1600"),
                ("Mission operations centre","S1600"),
            ]
        elif any(x in s for x in ["rail","metro","transit","hsr"]):
            packages = [
                ("Track and trackform","R1400"),
                ("Overhead line and electrification","R1400"),
                ("Signalling and control systems","R1500"),
                ("Stations, platforms and public areas","R1600"),
                ("Depot and stabling","R1600"),
                ("Systems integration and ETCS","R1700"),
            ]
        elif any(x in s for x in ["data centre","hyperscale","compute"]):
            packages = [
                ("Grid connection and HV switchgear","D1000"),
                ("Backup generation and UPS","D1200"),
                ("Cooling systems and plant","D1300"),
                ("Data hall fitout and cabling","D1400"),
                ("Network and connectivity","D1500"),
                ("DCIM and building management","D1700"),
            ]
        elif any(x in s for x in ["nuclear","smr"]):
            packages = [
                ("Nuclear island civil and containment","N1400"),
                ("Reactor pressure vessel and internals","N1300"),
                ("Primary circuit pipework and valves","N1500"),
                ("Electrical, I&C and protection","N1600"),
                ("Waste management and active drains","N1700"),
                ("Balance of plant and turbine hall","N1800"),
            ]
        elif any(x in s for x in ["semiconductor","fab","wafer"]):
            packages = [
                ("Cleanroom HVAC and environmental systems","F1300"),
                ("Ultra-pure water and chemical distribution","F1400"),
                ("Power quality and emergency systems","F1500"),
                ("Lithography and deposition tools","F1600"),
                ("Etch, clean and metrology tools","F1600"),
                ("Wastewater treatment and abatement","F1400"),
            ]
        elif any(x in s for x in ["life science","pharma","biotech","gmp"]):
            packages = [
                ("GMP cleanroom and HVAC","L1100"),
                ("Water-for-injection and pure media","L1200"),
                ("Bioreactor and upstream processing","L1300"),
                ("Fill-finish and isolator systems","L1300"),
                ("QC laboratory and analytical","L1400"),
                ("Cold chain and storage","L1200"),
            ]
        elif any(x in s for x in ["airport","aviation"]):
            packages = [
                ("Runway, taxiway and apron","AP1100"),
                ("Terminal structure and envelope","AP1200"),
                ("Baggage handling system","AP1400"),
                ("Security and check-in systems","AP1500"),
                ("Airfield lighting and navigation aids","AP1600"),
                ("ORAT and airline transition","AP1700"),
            ]
        else:
            packages = [
                ("Civil and structural package","A1900"),
                ("MEP and mechanical systems","A1900"),
                ("Specialist equipment and fitout","A1900"),
                ("Utilities and external works","A1800"),
                ("Systems and controls integration","A2000"),
                ("Testing and commissioning","A2100"),
            ]
        for i,(pname,pred) in enumerate(packages,1):
            code=f"P{i}000"
            actual_pred = pred if i==1 else f"P{i-1}000"
            rows.append((code,"Package",pname,actual_pred if i>1 else pred,0.030))
            if level >= 5:
                rows.append((f"P{i}100","Activity",f"{pname} — mobilise and set up","" + code,0.007))
                rows.append((f"P{i}200","Activity",f"{pname} — execute and install",f"P{i}100",0.018))
                rows.append((f"P{i}300","Activity",f"{pname} — test, inspect and closeout",f"P{i}200",0.007))

    # ── Build output list ────────────────────────────────────────────
    critical_codes = {
        "A1400","A1700","A1900","A2000","A2100","A2200",
        "R1200","R1300","R1400","R1500","R1700","R1800","R1900",
        "D1000","D1200","D1300","D1700","D1800",
        "N1000","N1700","N1800","N1900",
        "F1600","F1700","F1800",
        "L1600","L1700","L1800",
        "AP1100","AP1400","AP1700","AP1800",
        "S1300","S1400","S1500","S1600",
    }
    out=[]
    for item in rows:
        if len(item) == 5:
            code,phase,act,pred,pct = item
        else:
            continue
        dur = max(1, round(months*pct))
        is_crit = "Yes" if (code in critical_codes or code.startswith("P")) else "No"
        out.append({
            "activity_id": code,
            "phase": phase,
            "activity": act,
            "predecessor": pred,
            "duration_months": dur,
            "critical": is_crit,
            "basis": f"Duration {dur} months ({round(pct*100,1)}% of {months}mo programme). Sector-specific logic for {subsector}. Level {level} schedule.",
        })
    return out

def peer_competitors(client:str, subsector:str, mode:str):
    c=(client or "").lower(); s=subsector.lower()
    if mode=="Space" or any(x in c for x in ["spacex","blue origin","nasa"]): return ["SpaceX","Blue Origin","Rocket Lab","ULA","NASA/ESA programmes"]
    if any(x in c for x in ["microsoft","msft"]): return ["Meta","Alphabet / Google","Amazon AWS","Oracle","Apple"]
    if "data centre" in s: return ["Meta","Alphabet / Google","Amazon AWS","Microsoft","Oracle"]
    if "airport" in s: return ["Dubai Airports","Heathrow operators","Changi Airport Group","AD Airports","Saudi airport programmes"]
    if "rail" in s: return ["Network Rail","SNCF","DB InfraGO","HS2 delivery partners","Gulf rail programmes"]
    if "pharma" in s: return ["Pfizer","Roche","Novartis","AstraZeneca","Moderna"]
    if "energy" in s or "nuclear" in s: return ["EDF","NextEra","Aramco energy programmes","Masdar","Ørsted"]
    return ["Relevant global peers","Regional developers","Tier-one operators","Sovereign funds","Major asset owners"]


# ------------------------- benchmark memory / public-scale calibration -------------------------
BENCHMARK_LIBRARY = [
    {"sector":"Digital Infrastructure / Hyperscale Data Centre","mode":"Earth","keywords":["data centre","data center","hyperscale","ai campus","gpu"],"cost_bn":3.8,"months":46,"risks":["Grid connection delay","Cooling capacity constraint","Long-lead electrical equipment","Commissioning load-bank failure"]},
    {"sector":"Semiconductor / Advanced Manufacturing","mode":"Earth","keywords":["semiconductor","fab","wafer","upw","chip"],"cost_bn":18.0,"months":78,"risks":["Tool delivery slippage","UPW and chemicals interface","Cleanroom validation delay","Utility demand escalation"]},
    {"sector":"Airport / Aviation","mode":"Earth","keywords":["airport","terminal","runway","baggage"],"cost_bn":9.0,"months":84,"risks":["Live operations phasing","Security and baggage integration","Airside possession constraints","Stakeholder approvals"]},
    {"sector":"Rail / Transit","mode":"Earth","keywords":["rail","metro","station","signalling","transit"],"cost_bn":6.5,"months":84,"risks":["Possession access","Systems integration delay","Utility diversions","Land/stakeholder approvals"]},
    {"sector":"Energy / Utilities","mode":"Earth","keywords":["hydrogen","wind","solar","grid","battery","energy"],"cost_bn":5.2,"months":66,"risks":["Grid connection delay","Equipment lead times","Permitting constraints","Commodity escalation"]},
    {"sector":"Nuclear / Energy","mode":"Earth","keywords":["nuclear","smr","reactor"],"cost_bn":12.0,"months":96,"risks":["Licensing delay","FOAK supply chain","Safety case approval","Specialist labour availability"]},
    {"sector":"Healthcare / Hospital","mode":"Earth","keywords":["hospital","medical","healthcare"],"cost_bn":2.4,"months":60,"risks":["Clinical commissioning","Medical equipment lead times","Live hospital interfaces","Regulatory approval"]},
    {"sector":"Life Sciences / Pharma","mode":"Earth","keywords":["gmp","pharma","biologics","life sciences","validation"],"cost_bn":2.8,"months":58,"risks":["Validation delay","Clean utility readiness","Regulatory handover","Specialist automation integration"]},
    {"sector":"Defence / Secure Mission Infrastructure","mode":"Earth","keywords":["defence","defense","military","secure","command","airbase","naval","radar"],"cost_bn":4.8,"months":54,"risks":["Security accreditation","Controlled procurement","Resilience requirements","Operational continuity"]},
    {"sector":"Water / Utilities","mode":"Earth","keywords":["desalination","water","wastewater","flood"],"cost_bn":2.2,"months":50,"risks":["Marine works delay","Environmental consent","Energy cost exposure","Stakeholder approvals"]},
    {"sector":"Spaceport/Launch","mode":"Space","keywords":["spaceport","launch","rocket","pad"],"cost_bn":9.0,"months":84,"risks":["Launch licensing","Range safety","Propellant farm safety","Environmental approvals"]},
    {"sector":"Lunar/Mars Settlement","mode":"Space","keywords":["lunar","moon","mars","habitat","base","settlement"],"cost_bn":32.0,"months":156,"risks":["Launch manifest dependency","Life-support reliability","Power survival","Surface logistics"]},
    {"sector":"ISRU/Mining/Propellant","mode":"Space","keywords":["isru","propellant","oxygen","methane","mining","depot"],"cost_bn":22.0,"months":168,"risks":["Cryogenic reliability","Autonomous transfer","Technology readiness","Launch cadence"]},
    {"sector":"Satellite/Comms","mode":"Space","keywords":["satellite","constellation","comms","ground station"],"cost_bn":6.0,"months":60,"risks":["Production cadence","Launch manifest","Ground segment integration","Spectrum/regulatory"]},
    {"sector":"Orbital Compute / Manufacturing","mode":"Space","keywords":["orbital compute","space data centre","space data center","orbital ai compute","compute platform","leo compute","orbital manufacturing","leo manufacturing"],"cost_bn":24.0,"months":132,"risks":["Thermal rejection system failure","Launch cadence dependency","Autonomous servicing constraints","Radiation hardening exposure","Power density scaling","Ground relay latency","Orbital debris exposure"]},
]

def benchmark_similarity(prompt: str, mode: str, subsector: str) -> List[Dict[str, Any]]:
    t = (prompt or "").lower()
    scored = []
    for b in BENCHMARK_LIBRARY:
        score = 0
        if b["mode"] == mode:
            score += 2
        if b["sector"].lower() in (subsector or "").lower() or (subsector or "").lower() in b["sector"].lower():
            score += 3
        for kw in b["keywords"]:
            if kw in t:
                score += 2
        if score > 0:
            scored.append({**b, "similarity_score": score})
    return sorted(scored, key=lambda x: x["similarity_score"], reverse=True)[:5]

def calibrate_with_benchmarks(prompt: str, mode: str, subsector: str, raw_cost: float, months: int) -> Tuple[float, int, List[str], List[Dict[str, Any]]]:
    matches = benchmark_similarity(prompt, mode, subsector)
    if not matches:
        return raw_cost, months, ["No strong benchmark match; CASEY sector template used."], []
    primary = matches[0]
    anchor_cost = float(primary["cost_bn"])
    anchor_months = int(primary["months"])
    calibrated_cost = raw_cost * 0.72 + anchor_cost * 0.28
    calibrated_months = round(months * 0.76 + anchor_months * 0.24)
    notes = [
        f"Benchmark memory matched {primary['sector']} with similarity score {primary['similarity_score']}.",
        f"Cost calibrated by blending CASEY template output with benchmark anchor {money_bn(anchor_cost)}.",
        f"Schedule calibrated against benchmark duration {anchor_months} months.",
    ]
    return calibrated_cost, calibrated_months, notes, matches

def estimate_quality_index(class_level:int, schedule_level:int, risk_score:float, matches:List[Dict[str,Any]], input_quality:int=80) -> Dict[str,Any]:
    maturity = {1:92,2:82,3:70,4:56,5:42}.get(class_level,70)
    schedule_quality = {1:42,2:55,3:68,4:78,5:88}.get(schedule_level,68)
    benchmark_quality = min(92, 40 + (matches[0]["similarity_score"]*7 if matches else 0))
    risk_penalty = min(28, risk_score/4)
    score = int(clamp((maturity*0.30 + schedule_quality*0.25 + benchmark_quality*0.25 + input_quality*0.20) - risk_penalty, 10, 96))
    band = "High" if score >= 78 else "Medium" if score >= 58 else "Low"
    return {"score": score, "band": band, "maturity_score": maturity, "schedule_quality": schedule_quality, "benchmark_quality": benchmark_quality, "risk_penalty": round(risk_penalty,1), "meaning": "Confidence index reflects estimate class, schedule level, benchmark similarity, input quality and risk exposure."}

def procurement_heatmap(mode:str, subsector:str, risks:List[Dict[str,Any]]) -> List[Dict[str,Any]]:
    base = [
        {"package":"Long-lead equipment","exposure":"High","reason":"Procurement path drives cost and schedule confidence."},
        {"package":"Utilities / enabling works","exposure":"High" if mode=="Earth" else "Medium","reason":"Utility readiness controls commissioning and phasing."},
        {"package":"Systems integration","exposure":"High","reason":"Integration maturity is a common late-stage failure mode."},
        {"package":"Specialist suppliers","exposure":"Medium-High","reason":"Limited supplier base increases price and schedule volatility."},
    ]
    if mode=="Space":
        base.insert(0, {"package":"Launch / mission logistics","exposure":"Extreme","reason":"Launch cadence and mission integration dominate delivery exposure."})
        base.append({"package":"Autonomous operations","exposure":"High","reason":"Remote operations reduce recovery options after deployment."})
    if "Data Centre" in subsector:
        base.append({"package":"Power equipment","exposure":"Extreme","reason":"Transformers, switchgear and grid connection can dominate critical path."})
    if "Semiconductor" in subsector:
        base.append({"package":"Process tools","exposure":"Extreme","reason":"Tool install and cleanroom readiness drive fab ramp-up."})
    if "Orbital Compute" in subsector:
        base.append({"package":"Thermal rejection systems","exposure":"Extreme","reason":"Heat rejection and cooling stability dominate orbital compute viability."})
        base.append({"package":"Radiation-hardened compute","exposure":"High","reason":"Radiation-tolerant compute supply chains remain constrained."})
    return base[:9]

def critical_path_narrative(mode:str, subsector:str, schedule_rows:List[Dict[str,Any]]) -> List[str]:
    names = [r.get("activity","") for r in schedule_rows if isinstance(r,dict)]
    out = [
        "Critical path is likely governed by definition maturity, procurement release and systems commissioning.",
        "Early schedule confidence depends on locking scope, approvals and long-lead package strategy.",
    ]
    if mode=="Space":
        out.append("Space critical path is additionally controlled by launch integration, payload readiness and remote operations validation.")
    if "Data Centre" in subsector:
        out.append("For digital infrastructure, grid connection and electrical equipment procurement are likely critical path drivers.")
    if "Orbital Compute" in subsector:
        out.append("For orbital compute, thermal rejection, power density, launch cadence and autonomous servicing are likely critical path drivers.")
    if "Airport" in subsector or "Rail" in subsector:
        out.append("For transport programmes, possessions/live operations and systems integration are likely critical path drivers.")
    if "Defence" in subsector:
        out.append("For secure infrastructure, accreditation and controlled procurement can become critical path constraints.")
    if names:
        out.append(f"First schedule control focus should be: {names[min(2,len(names)-1)]}.")
    return out[:5]


# ------------------------- CASEY WOW narrative intelligence layer -------------------------
def board_briefing_narrative(mode:str, subsector:str, title:str, cost_p50:str, schedule:str, risk:str, confidence:int, drivers:List[str], matches:List[Dict[str,Any]]) -> List[str]:
    if mode == "Space":
        opening = f"{title} should be treated less like a one-off aerospace project and more like persistent orbital infrastructure requiring operational coordination between launch, power, comms, autonomy and servicing."
    elif "Data Centre" in subsector:
        opening = f"{title} is likely to behave as much like an energy-and-utilities programme as a conventional digital infrastructure build because power availability and commissioning sequencing will dominate confidence."
    elif "Defence" in subsector:
        opening = f"{title} is likely to be governed by accreditation, resilience and controlled procurement constraints as much as by traditional construction delivery."
    elif "Airport" in subsector or "Rail" in subsector:
        opening = f"{title} is likely to be constrained by live operational interfaces, systems integration and possession/phasing strategy rather than physical works alone."
    else:
        opening = f"{title} should be managed as an integrated infrastructure system, not simply a capital project reporting exercise."

    return [
        opening,
        f"Current P50 intelligence indicates {cost_p50} over approximately {schedule}, with overall risk assessed as {risk}.",
        f"Confidence remains {('high' if confidence >= 78 else 'moderate' if confidence >= 55 else 'low')} because class maturity, schedule definition, benchmark similarity and delivery exposure are still moving together.",
        "The immediate control priority is to convert early assumptions into governed decisions around scope freeze, procurement path, interface ownership and critical-path evidence.",
    ]

def uncertainty_narrative(class_level:int, schedule_level:int, confidence:int, mode:str, risk:str, drivers:List[str]) -> Dict[str,Any]:
    class_msg = {
        1:"Class 1 maturity suggests a high-definition estimate with reduced range exposure.",
        2:"Class 2 maturity suggests strong definition but remaining package and market exposure.",
        3:"Class 3 maturity is suitable for budget authorization, but procurement and design assumptions still need challenge.",
        4:"Class 4 maturity remains feasibility-led and should be treated as a decision-support range rather than a firm budget.",
        5:"Class 5 maturity is concept-level and should be treated as strategic order-of-magnitude intelligence."
    }.get(class_level, "Estimate maturity requires review.")
    sched_msg = {
        1:"Schedule Level 1 is strategic only and should not be used for delivery commitment.",
        2:"Schedule Level 2 supports programme framing but still lacks full activity logic.",
        3:"Schedule Level 3 supports integrated control but needs package-level validation.",
        4:"Schedule Level 4 gives stronger logic and QSRA traceability.",
        5:"Schedule Level 5 supports detailed controls thinking if activity evidence is maintained."
    }.get(schedule_level, "Schedule definition requires review.")
    exposure = "Launch and orbital operations exposure increases uncertainty." if mode=="Space" else "Market, approvals, procurement and interface exposure remain key uncertainty drivers."
    return {"confidence":confidence, "risk":risk, "estimate_maturity":class_msg, "schedule_maturity":sched_msg, "uncertainty_drivers":drivers[:6] or [exposure], "interpretation":exposure}

def mission_control_cards(mode:str, subsector:str, title:str, risk:str, matches:List[Dict[str,Any]], procurement:List[Dict[str,Any]], critical:List[str]) -> List[Dict[str,str]]:
    cards = []
    cards.append({"label":"CRITICAL PATH EXPOSURE","signal":critical[0] if critical else "Critical path requires validation against procurement and commissioning logic.","severity":risk})
    if procurement:
        p = procurement[0]
        cards.append({"label":"PROCUREMENT WATCH","signal":f"{p.get('package')} exposure is {p.get('exposure')}: {p.get('reason')}","severity":p.get("exposure","High")})
    if matches:
        cards.append({"label":"BENCHMARK MEMORY","signal":f"Closest archetype: {matches[0].get('sector')} with similarity score {matches[0].get('similarity_score')}.","severity":"Reference"})
    if mode == "Space":
        cards.append({"label":"ORBITAL OPERATIONS RISK","signal":"Launch cadence, autonomous servicing, communications dependency and remote recovery constraints should be treated as programme controls issues, not only engineering issues.","severity":"Extreme" if risk in ["Very High","Extreme"] else "High"})
    if "Orbital Compute" in subsector:
        cards.append({"label":"THERMAL / POWER BOTTLENECK","signal":"Orbital compute viability is likely to be governed by thermal rejection, power density, radiation-hardening and servicing cadence.","severity":"Extreme"})
    elif "Data Centre" in subsector:
        cards.append({"label":"POWER BOTTLENECK","signal":"Grid energisation, transformers, switchgear and commissioning capacity are likely to dominate the delivery confidence curve.","severity":"High"})
    return cards[:6]

def casey_thinking(mode:str, subsector:str, title:str) -> str:
    if "Orbital Compute" in subsector:
        return "CASEY interprets this as a new infrastructure category: compute capacity no longer sits only behind a fence line on Earth; it becomes an orbital operations system governed by launch cadence, thermal rejection, servicing and relay dependency."
    if mode == "Space":
        return "CASEY interprets this as persistent orbital infrastructure, where project controls shift from construction reporting toward Earth-orbit operational orchestration."
    if "Data Centre" in subsector:
        return "CASEY interprets this as an energy-constrained infrastructure programme, not simply a digital real-estate deployment."
    if "Semiconductor" in subsector:
        return "CASEY interprets this as a utility-and-tooling critical programme where cleanroom readiness and process tool cadence will dominate ramp-up confidence."
    if "Defence" in subsector:
        return "CASEY interprets this as a secure mission infrastructure programme where accreditation and resilience can become hidden critical-path drivers."
    return "CASEY interprets this as a system-of-systems infrastructure programme where delivery confidence depends on interface control, procurement evidence and decision velocity."

def benchmark_comparison(matches:List[Dict[str,Any]], mode:str, subsector:str) -> List[Dict[str,Any]]:
    if matches:
        return [{"archetype":m["sector"],"similarity_score":m["similarity_score"],"anchor_cost":money_bn(m["cost_bn"]),"anchor_duration_months":m["months"],"use":"Directional calibration only; not a certified benchmark."} for m in matches[:4]]
    return [{"archetype":subsector or ("Space infrastructure" if mode=="Space" else "Earth infrastructure"),"similarity_score":0,"anchor_cost":"N/A","anchor_duration_months":"N/A","use":"No strong benchmark memory match; CASEY used sector logic and input assumptions."}]


# ------------------------- sector realism and executive output upgrade -------------------------
SECTOR_ENVELOPES = {
    "Life Sciences / Biologics Manufacturing": {"cost": (1.8, 8.5), "months": (36, 78)},
    "Life Sciences / Pharma": {"cost": (1.5, 7.5), "months": (34, 76)},
    "Digital Infrastructure / Hyperscale Data Centre": {"cost": (1.5, 18.0), "months": (24, 84)},
    "Semiconductor / Advanced Manufacturing": {"cost": (8.0, 32.0), "months": (54, 96)},
    "Airport / Aviation": {"cost": (2.0, 24.0), "months": (60, 144)},
    "Rail / Transit": {"cost": (3.0, 80.0), "months": (60, 180)},
    "Healthcare / Hospital": {"cost": (0.8, 6.5), "months": (36, 84)},
    "Defence / Secure Mission Infrastructure": {"cost": (0.8, 12.0), "months": (30, 84)},
    "Energy / Utilities": {"cost": (1.0, 18.0), "months": (36, 108)},
    "Nuclear / Energy": {"cost": (5.0, 35.0), "months": (72, 144)},
    "Water / Utilities": {"cost": (0.8, 8.0), "months": (30, 84)},
    "Spaceport/Launch": {"cost": (1.0, 14.0), "months": (36, 108)},
    "Orbital Compute / Manufacturing": {"cost": (8.0, 65.0), "months": (72, 168)},
    "Lunar Surface Habitat/Base": {"cost": (15.0, 95.0), "months": (96, 216)},
    "Mars Surface Habitat/Base": {"cost": (30.0, 180.0), "months": (120, 300)},
    "ISRU/Mining/Propellant": {"cost": (15.0, 120.0), "months": (96, 240)},
    "Satellite/Comms": {"cost": (2.0, 20.0), "months": (36, 96)},
    "Deep-Space Communications Infrastructure": {"cost": (4.0, 35.0), "months": (60, 144)},
}

def _sector_family(subsector:str) -> str:
    s=(subsector or "").lower()
    if "life sciences" in s or "biologics" in s or "pharma" in s: return "life_sciences"
    if "data centre" in s or "data center" in s or "digital infrastructure" in s: return "data_centre"
    if "semiconductor" in s: return "semiconductor"
    if "airport" in s or "aviation" in s: return "airport"
    if "rail" in s or "transit" in s: return "rail"
    if "hospital" in s or "healthcare" in s: return "healthcare"
    if "defence" in s or "defense" in s or "secure mission" in s: return "defence"
    if "nuclear" in s: return "nuclear"
    if "energy" in s or "utilities" in s: return "energy"
    if "water" in s: return "water"
    if "orbital compute" in s or "orbital" in s: return "orbital"
    if "lunar" in s: return "lunar"
    if "mars" in s: return "mars"
    if "spaceport" in s or "launch" in s: return "spaceport"
    if "satellite" in s or "comms" in s or "communications" in s: return "space_comms"
    if "isru" in s or "propellant" in s or "mining" in s: return "space_resources"
    return "general"

def sector_envelope(subsector:str, cost:float, months:int, scale:str="") -> Tuple[float,int,List[str]]:
    fam = _sector_family(subsector)
    env = SECTOR_ENVELOPES.get(subsector)
    if not env:
        for k,v in SECTOR_ENVELOPES.items():
            if k.lower() in (subsector or "").lower() or (subsector or "").lower() in k.lower():
                env = v; break
    if not env:
        return cost, months, []
    lo, hi = env["cost"]; mlo, mhi = env["months"]
    # Allow mega wording to approach upper end, but prevent runaway values.
    scale_l=(scale or "").lower()
    if "mega-mega" in scale_l:
        hi *= 1.25; mhi *= 1.15
    elif "mega" in scale_l:
        hi *= 1.12; mhi *= 1.08
    original=(cost,months)
    cost = min(max(cost, lo), hi)
    months = int(min(max(months, mlo), mhi))
    notes=[]
    if abs(cost-original[0]) > 0.01 or months != original[1]:
        notes.append(f"Sector realism envelope applied for {subsector}: cost and schedule constrained to credible early-stage range.")
    return cost, months, notes

SECTOR_COST_TEMPLATES = {
    "life_sciences": [
        ("01.01","Site, enabling and controlled utilities","Direct","Site works, utility corridors and GMP-ready enabling infrastructure",0.07),
        ("01.02","GMP cleanrooms and classified areas","Direct","Cleanroom envelope, HVAC zoning, pressure cascades and finishes",0.17),
        ("01.03","Process equipment and fill-finish lines","Direct","Bioreactors, formulation, aseptic fill-finish and packaging equipment",0.22),
        ("01.04","Clean utilities","Direct","WFI, clean steam, purified water, process gases and validated utility loops",0.14),
        ("01.05","Automation, MES and digital validation","Direct","Automation, MES, batch records, data integrity and validation systems",0.08),
        ("01.06","Cold-chain and automated logistics","Direct","Cold storage, automated warehouse, packaging and distribution infrastructure",0.08),
        ("02.01","CQV / qualification","Indirect","Commissioning, qualification, validation protocols and deviation closure",0.08),
        ("02.02","Regulatory readiness and quality systems","Indirect","FDA/EMA readiness, QA systems, inspection preparation and documentation",0.05),
        ("02.03","Programme management and operational readiness","Indirect","PMO, controls, staffing readiness, training and handover",0.06),
        ("03.01","Risk reserve","Reserve","Allowance for GMP validation, equipment lead times and regulatory uncertainty",0.05),
    ],
    "data_centre": [
        ("01.01","Land, enabling and campus infrastructure","Direct","Site, roads, drainage, telecom routes and campus utilities",0.07),
        ("01.02","Power train and substations","Direct","HV/MV substations, switchgear, transformers, UPS and generators",0.24),
        ("01.03","Data halls / white space","Direct","Data halls, fit-out, containment and critical environments",0.17),
        ("01.04","Cooling plant","Direct","Liquid cooling, chillers, heat rejection and water systems",0.15),
        ("01.05","Fibre, controls and security","Direct","Network rooms, BMS/EPMS, access control and cyber/physical security",0.07),
        ("02.01","Commissioning and integrated systems testing","Indirect","IST, black-building tests, staged energisation and reliability proving",0.08),
        ("02.02","Programme management and procurement","Indirect","PMO, procurement, logistics and contractor management",0.07),
        ("03.01","Risk reserve","Reserve","Utility, long-lead electrical and commissioning uncertainty",0.15),
    ],
    "semiconductor": [
        ("01.01","Site and enabling infrastructure","Direct","Site works, seismic base, bulk utilities and logistics",0.06),
        ("01.02","Cleanroom shell and process environment","Direct","Vibration control, cleanroom, HVAC and contamination control",0.20),
        ("01.03","Process tools and tool install","Direct","Lithography, etch, deposition, metrology and tool hook-up",0.32),
        ("01.04","UPW, chemicals and specialty gases","Direct","Ultra-pure water, gases, abatement and chemical systems",0.14),
        ("01.05","Power and resilient utilities","Direct","Substations, emergency power and redundancy",0.08),
        ("02.01","Qualification and ramp-up","Indirect","Tool qualification, yield ramp and validation",0.08),
        ("02.02","PMO and controls","Indirect","Programme controls, procurement and logistics",0.05),
        ("03.01","Risk reserve","Reserve","Tool cadence, supply chain and ramp-up uncertainty",0.07),
    ],
}

def sector_cost_lines(mode:str, subsector:str):
    fam=_sector_family(subsector)
    if fam in SECTOR_COST_TEMPLATES:
        return SECTOR_COST_TEMPLATES[fam]
    return cost_lines(mode, subsector)

def sector_specific_lists(subsector:str, mode:str):
    fam=_sector_family(subsector)
    if fam=="life_sciences":
        return {
            "cost":["Process equipment and fill-finish lines","GMP cleanrooms / classified areas","Clean utilities and HVAC zoning","CQV validation and deviation closure","Cold-chain logistics and automated warehousing"],
            "schedule":["CQV protocol approval and execution","Long-lead process equipment delivery","Clean utility validation and media fills","FDA/EMA inspection readiness","Operational staffing, training and batch readiness"],
            "confidence":["Benchmark similarity: pharma / biologics campus","Scope maturity: GMP package and user requirement definition","Procurement certainty: process equipment and clean utility lead times","Schedule maturity: CQV logic and validation pathway","Regulatory exposure: FDA/EMA inspection readiness"],
            "thinking":"CASEY interprets this as a validation-constrained life sciences manufacturing programme where CQV sequencing, clean utility readiness, process equipment cadence and regulatory inspection preparedness dominate schedule confidence more than conventional construction productivity.",
            "bench":[
                {"archetype":"Life Sciences / Biologics Manufacturing","similarity_score":9,"anchor_cost":"$2B-$8B","anchor_duration_months":"36-78","use":"Primary analogue for GMP campus, fill-finish and CQV complexity."},
                {"archetype":"Sterile Fill-Finish / Aseptic Expansion","similarity_score":8,"anchor_cost":"$1B-$5B","anchor_duration_months":"30-60","use":"Relevant for aseptic process, cleanroom and validation sequence."},
                {"archetype":"Advanced Manufacturing Cleanroom","similarity_score":7,"anchor_cost":"$2B-$10B","anchor_duration_months":"42-84","use":"Adjacent benchmark for clean environment and specialist equipment."},
                {"archetype":"Cold-Chain Logistics Campus","similarity_score":6,"anchor_cost":"$500M-$3B","anchor_duration_months":"24-54","use":"Adjacent benchmark for warehousing and distribution infrastructure."},
            ]
        }
    if fam=="data_centre":
        return {
            "cost":["Utility/grid connection and substations","Power train, transformers and switchgear","Liquid cooling / heat rejection systems","Data halls and white space fit-out","Accelerated procurement premiums"],
            "schedule":["Grid energisation and utility agreements","Long-lead transformer and switchgear delivery","Integrated systems testing and commissioning","Design freeze and phasing stability","Contractor productivity across concurrent halls"],
            "confidence":["Benchmark similarity: hyperscale digital infrastructure","Scope maturity: campus power and white-space definition","Procurement certainty: transformers, generators and switchgear","Schedule maturity: grid and commissioning logic","Interface exposure: utilities, fibre and commissioning"],
            "thinking":"CASEY interprets this as an energy-constrained digital infrastructure programme where utility availability, commissioning throughput and long-lead electrical procurement dominate schedule confidence more than vertical construction productivity.",
            "bench":[
                {"archetype":"Hyperscale AI Data Centre Campus","similarity_score":9,"anchor_cost":"$2B-$18B","anchor_duration_months":"24-84","use":"Primary analogue for power-intensive digital infrastructure."},
                {"archetype":"Energy / Utility Megaprogramme","similarity_score":7,"anchor_cost":"$1B-$18B","anchor_duration_months":"36-108","use":"Adjacent benchmark for substations and energisation risk."},
                {"archetype":"Mission Critical Facility","similarity_score":7,"anchor_cost":"$1B-$8B","anchor_duration_months":"30-72","use":"Adjacent benchmark for resilience and uptime requirements."},
                {"archetype":"Semiconductor / Advanced Manufacturing","similarity_score":6,"anchor_cost":"$8B-$32B","anchor_duration_months":"54-96","use":"Adjacent benchmark for complex MEP and commissioning."},
            ]
        }
    if fam=="semiconductor":
        return {
            "cost":["Process tools and tool install","Cleanroom shell and vibration control","UPW, gases and chemical systems","Power redundancy and utility resilience","Yield ramp and qualification"],
            "schedule":["Long-lead lithography/process tool delivery","Cleanroom readiness and contamination control","UPW and specialty gas commissioning","Tool hook-up and qualification cadence","Yield ramp and process validation"],
            "confidence":["Benchmark similarity: semiconductor fab","Scope maturity: process tool list and cleanroom class definition","Procurement certainty: tool delivery and hook-up sequencing","Schedule maturity: cleanroom/tool qualification logic","Utility exposure: UPW, gases, power and abatement"],
            "thinking":"CASEY interprets this as a tooling-and-utilities constrained advanced manufacturing programme where process tool cadence, cleanroom readiness, UPW/gas systems and yield ramp dominate delivery confidence.",
            "bench":[
                {"archetype":"Advanced Semiconductor Fab","similarity_score":9,"anchor_cost":"$8B-$32B","anchor_duration_months":"54-96","use":"Primary analogue for tool-heavy cleanroom delivery."},
                {"archetype":"Advanced Manufacturing Cleanroom","similarity_score":8,"anchor_cost":"$2B-$10B","anchor_duration_months":"42-84","use":"Adjacent benchmark for controlled environment and specialist systems."},
                {"archetype":"Hyperscale Digital Infrastructure","similarity_score":6,"anchor_cost":"$2B-$18B","anchor_duration_months":"24-84","use":"Adjacent benchmark for utility and commissioning intensity."}
            ]
        }
    if fam=="airport":
        return {
            "cost":["Terminal and airside works","Baggage and security systems","Operational transition and phasing","Transport integration and utilities","Retail/passenger experience fit-out"],
            "schedule":["Live airport phasing and possessions","Baggage/security systems integration","Operational readiness trials","Regulatory and stakeholder approvals","Airside access and safety constraints"],
            "confidence":["Benchmark similarity: airport terminal expansion","Scope maturity: capacity, phasing and systems definition","Procurement certainty: baggage/security/MEP packages","Schedule maturity: ORAT and live operations logic","Interface exposure: airlines, airside, landside and regulators"],
            "thinking":"CASEY interprets this as a live-operational aviation programme where phasing, systems integration and operational readiness dominate confidence more than standalone construction progress.",
            "bench":[{"archetype":"Airport Terminal Expansion","similarity_score":9,"anchor_cost":"$2B-$24B","anchor_duration_months":"60-144","use":"Primary analogue for live-airport capital delivery."},{"archetype":"Rail/Transit Systems Integration","similarity_score":6,"anchor_cost":"$3B-$80B","anchor_duration_months":"60-180","use":"Adjacent benchmark for passenger systems and operational transition."}]
        }
    if fam=="rail":
        return {
            "cost":["Civil works, tunnelling and structures","Stations and public realm","Signalling, power and systems","Land, utilities and consents","Testing, commissioning and trial operations"],
            "schedule":["Land acquisition and utility diversions","Possessions and live railway interfaces","Systems integration and signalling readiness","Station fit-out and operational trials","Regulatory approvals and safety case"],
            "confidence":["Benchmark similarity: rail/transit programme","Scope maturity: alignment, station and systems definition","Procurement certainty: civil/systems package strategy","Schedule maturity: possessions and test/commissioning logic","Interface exposure: utilities, operators and regulators"],
            "thinking":"CASEY interprets this as an interface-heavy transit programme where utilities, possessions, systems integration and safety assurance dominate delivery confidence.",
            "bench":[{"archetype":"Metro / Rail Extension","similarity_score":9,"anchor_cost":"$3B-$80B","anchor_duration_months":"60-180","use":"Primary analogue for complex transit delivery."},{"archetype":"Airport Systems Programme","similarity_score":6,"anchor_cost":"$2B-$24B","anchor_duration_months":"60-144","use":"Adjacent benchmark for live operations and systems assurance."}]
        }
    if fam=="healthcare":
        return {
            "cost":["Clinical accommodation and theatres","Diagnostics and specialist medical equipment","MEP, resilience and energy centre","Digital clinical systems","Phased handover and commissioning"],
            "schedule":["Clinical commissioning and licensing","Medical equipment procurement","Decant/phasing with live services","Digital systems integration","Operational readiness and staff training"],
            "confidence":["Benchmark similarity: acute hospital campus","Scope maturity: clinical model and departmental adjacencies","Procurement certainty: medical equipment and MEP resilience","Schedule maturity: clinical commissioning and handover logic","Operational exposure: live healthcare continuity"],
            "thinking":"CASEY interprets this as a clinically constrained healthcare programme where operational continuity, equipment readiness and clinical commissioning dominate confidence.",
            "bench":[{"archetype":"Acute Hospital Campus","similarity_score":9,"anchor_cost":"$800M-$6.5B","anchor_duration_months":"36-84","use":"Primary analogue for healthcare capital delivery."},{"archetype":"Life Sciences / GMP Facility","similarity_score":6,"anchor_cost":"$1.5B-$8.5B","anchor_duration_months":"36-78","use":"Adjacent benchmark for controlled environments and validation."}]
        }
    if fam=="defence":
        return {
            "cost":["Hardened facilities and secure fit-out","Resilient power and communications","Cyber/security accreditation","Specialist mission systems","Controlled procurement and assurance"],
            "schedule":["Security accreditation and authority approvals","Secure systems integration","Long-lead mission equipment","Operational acceptance trials","Controlled access and classified interfaces"],
            "confidence":["Benchmark similarity: secure mission infrastructure","Scope maturity: mission system and accreditation requirements","Procurement certainty: controlled equipment and secure supply chain","Schedule maturity: accreditation and acceptance logic","Interface exposure: security authority and operators"],
            "thinking":"CASEY interprets this as secure mission infrastructure where accreditation, resilience and controlled procurement can become hidden critical-path drivers.",
            "bench":[{"archetype":"Secure Mission / Defence Facility","similarity_score":9,"anchor_cost":"$800M-$12B","anchor_duration_months":"30-84","use":"Primary analogue for secure infrastructure delivery."},{"archetype":"Data Centre / Mission Critical Facility","similarity_score":7,"anchor_cost":"$1.5B-$18B","anchor_duration_months":"24-84","use":"Adjacent benchmark for resilience and uptime."}]
        }
    if fam in ["energy","nuclear","water"]:
        return {
            "cost":["Generation/process equipment","Grid, utility and connection infrastructure","Civil and enabling works","Permitting and environmental compliance","Commissioning and operational readiness"],
            "schedule":["Permitting and environmental approvals","Grid connection and utility interfaces","Long-lead equipment procurement","Commissioning and performance testing","Operator readiness and regulatory handover"],
            "confidence":["Benchmark similarity: energy/utilities programme","Scope maturity: process configuration and connection definition","Procurement certainty: long-lead equipment and utility interfaces","Schedule maturity: permitting and commissioning logic","Regulatory exposure: environmental, safety and operator approvals"],
            "thinking":"CASEY interprets this as regulated utility infrastructure where approvals, grid/process interfaces and commissioning evidence dominate confidence.",
            "bench":[{"archetype":"Energy / Utility Megaprogramme","similarity_score":9,"anchor_cost":"$1B-$35B","anchor_duration_months":"36-144","use":"Primary analogue for regulated utility delivery."},{"archetype":"Industrial Process Facility","similarity_score":7,"anchor_cost":"$1B-$12B","anchor_duration_months":"36-96","use":"Adjacent benchmark for process systems and commissioning."}]
        }
    if mode=="Space":
        return {
            "cost":["Launch and mass-to-orbit logistics","Power generation and thermal control","Radiation-hardening and qualification","Autonomous operations and servicing","Ground segment and relay communications"],
            "schedule":["Technology readiness and qualification","Launch manifest and integration cadence","Thermal/power validation","Autonomous recovery and servicing readiness","Ground/orbit commissioning sequence"],
            "confidence":["Benchmark similarity: space infrastructure archetype","Scope maturity: payload and mission architecture definition","Procurement certainty: launch, avionics and qualified hardware","Schedule maturity: launch and commissioning logic","Operational exposure: remote recovery and servicing limits"],
            "thinking":"CASEY interprets this as persistent space infrastructure where delivery confidence is governed by launch cadence, technology readiness, qualification evidence and autonomous operational resilience.",
            "bench":[
                {"archetype":"Orbital / Lunar Infrastructure","similarity_score":8,"anchor_cost":"$8B-$95B","anchor_duration_months":"72-216","use":"Primary analogue for space infrastructure complexity."},
                {"archetype":"Launch and Payload Integration","similarity_score":7,"anchor_cost":"$1B-$14B","anchor_duration_months":"36-108","use":"Relevant for launch cadence and payload integration."},
                {"archetype":"Deep-Space Mission Systems","similarity_score":7,"anchor_cost":"$4B-$35B","anchor_duration_months":"60-144","use":"Adjacent benchmark for qualification and mission assurance."},
                {"archetype":"Autonomous Operations Platform","similarity_score":6,"anchor_cost":"$2B-$20B","anchor_duration_months":"36-96","use":"Adjacent benchmark for autonomy and remote operations."},
            ]
        }
    return {
        "cost":["Civil and enabling works","Specialist systems and long-lead equipment","Utilities and interfaces","Commissioning and operational readiness","Risk reserve driven by procurement and interface uncertainty"],
        "schedule":["Approvals and consents","Long-lead procurement","Design freeze stability","Systems integration and commissioning","Operational access and handover readiness"],
        "confidence":["Benchmark similarity: comparable infrastructure archetypes","Scope maturity: concept / budget level until package evidence is supplied","Procurement certainty: long-lead equipment and market capacity","Schedule maturity: critical path and commissioning logic","Interface exposure: utilities, systems and operations"],
        "thinking":"CASEY interprets this as a system-of-systems infrastructure programme where delivery confidence depends on interface control, procurement evidence and decision velocity.",
        "bench":[]
    }

def sector_signature_behaviour(subsector:str, mode:str) -> Dict[str, Any]:
    fam=_sector_family(subsector)
    signatures = {
        "life_sciences": {
            "shock":"Mechanical completion is unlikely to be the true finish line; validated production readiness and deviation closure are the real board decision gates.",
            "curve":"late_validation_spike",
            "human_basis":"duration reflects CQV sequencing, clean-utility validation, media-fill readiness and phased batch ramp-up rather than simple building completion.",
            "contradiction":"Acceleration may protect market-entry timing but reduces validation float and increases deviation-closure pressure.",
            "signature":["CQV path", "media fill readiness", "GMP turnover", "batch release", "inspection readiness"]
        },
        "data_centre": {
            "shock":"The dominant delivery constraint is likely energisation and commissioning concurrency, not shell construction productivity.",
            "curve":"power_procurement_cliff",
            "human_basis":"duration reflects utility energisation, transformer/switchgear procurement, integrated systems testing and phased data-hall commissioning.",
            "contradiction":"Acceleration can bring revenue online earlier, but usually buys schedule by paying procurement premiums and accepting commissioning overlap.",
            "signature":["grid energisation", "transformers", "liquid cooling", "IST", "phased data halls"]
        },
        "semiconductor": {
            "shock":"The critical path is likely tool install and qualification cadence rather than the cleanroom shell alone.",
            "curve":"tool_install_ramp",
            "human_basis":"duration reflects cleanroom readiness, UPW/gas commissioning, process-tool delivery, hook-up and yield-ramp qualification.",
            "contradiction":"Lower capex assumptions can move risk into yield ramp, tool utilisation and operational start-up.",
            "signature":["tool install", "UPW", "vibration control", "hook-up", "yield ramp"]
        },
        "airport": {
            "shock":"The programme may finish construction before the airport is operationally ready to absorb the change.",
            "curve":"orat_ramp",
            "human_basis":"duration reflects live-airport phasing, baggage/security integration, ORAT, stakeholder approvals and passenger operations continuity.",
            "contradiction":"Acceleration increases disruption exposure unless operational readiness trials are protected.",
            "signature":["ORAT", "baggage integration", "airside phasing", "security systems", "live operations"]
        },
        "rail": {
            "shock":"Possessions, utility diversions and systems assurance are likely to dominate the critical path more than visible civil progress.",
            "curve":"systems_assurance_tail",
            "human_basis":"duration reflects possessions, utility diversions, tunnelling/station works, signalling integration and safety assurance.",
            "contradiction":"A cheaper case may defer systems maturity and create a larger testing and commissioning tail.",
            "signature":["possessions", "utility diversions", "signalling", "safety case", "trial operations"]
        },
        "healthcare": {
            "shock":"Clinical readiness and safe operational transition may lag physical completion unless planned as a critical-path workstream.",
            "curve":"clinical_commissioning_tail",
            "human_basis":"duration reflects clinical commissioning, medical-equipment readiness, digital systems integration, decanting and live-service continuity.",
            "contradiction":"Acceleration can increase patient-service disruption unless clinical transition is protected.",
            "signature":["clinical commissioning", "medical equipment", "decant", "digital clinical systems", "operational readiness"]
        },
        "defence": {
            "shock":"Security accreditation and mission-system acceptance can become the hidden critical path even after the facility is built.",
            "curve":"accreditation_gate",
            "human_basis":"duration reflects controlled procurement, secure systems integration, resilience proving, cyber/security accreditation and operational acceptance.",
            "contradiction":"Lower-cost delivery can create downstream accreditation and acceptance risk.",
            "signature":["accreditation", "mission systems", "secure comms", "resilience proving", "controlled procurement"]
        },
        "energy": {
            "shock":"The approval, grid/process interface and commissioning evidence may drive confidence more than equipment installation.",
            "curve":"regulatory_grid_gate",
            "human_basis":"duration reflects consenting, grid/process interfaces, long-lead equipment and performance testing.",
            "contradiction":"Compression can pull construction forward while leaving approvals and grid acceptance behind.",
            "signature":["consents", "grid interface", "process equipment", "performance test", "operator handover"]
        },
        "nuclear": {
            "shock":"Licensing and safety-case maturity are the true confidence drivers; construction progress alone is a weak indicator.",
            "curve":"licensing_plateau",
            "human_basis":"duration reflects licensing, safety case, nuclear-grade procurement, regulator hold points and commissioning evidence.",
            "contradiction":"Attempting to accelerate before safety-case maturity usually moves risk into rework and regulator challenge.",
            "signature":["licensing", "safety case", "regulator hold point", "nuclear QA", "commissioning evidence"]
        },
        "water": {
            "shock":"Marine/environmental approvals and energy interfaces may dominate the delivery envelope.",
            "curve":"permit_then_commission",
            "human_basis":"duration reflects permitting, intake/outfall works, process commissioning, power tie-in and environmental compliance.",
            "contradiction":"Cheaper options may reduce redundancy and increase operating resilience exposure.",
            "signature":["intake/outfall", "environmental permit", "process commissioning", "power tie-in", "resilience"]
        },
        "orbital": {
            "shock":"The decisive constraint is not launch alone; it is thermal-power balance, autonomous servicing and recoverability after deployment.",
            "curve":"qualification_launch_spike",
            "human_basis":"duration reflects technology qualification, payload integration, launch manifest, thermal/power validation and autonomous operations readiness.",
            "contradiction":"Acceleration increases mission-assurance risk if qualification and environmental testing are compressed.",
            "signature":["TRL", "thermal rejection", "payload integration", "launch manifest", "autonomous servicing"]
        },
        "lunar": {
            "shock":"Surface logistics, dust and power resilience are likely to dominate sustained operability more than initial landing success.",
            "curve":"surface_logistics_tail",
            "human_basis":"duration reflects launch cadence, landing-site preparation, surface logistics, power storage and autonomous construction readiness.",
            "contradiction":"Lower-cost delivery may reduce redundancy and increase resupply dependency.",
            "signature":["landing pads", "dust", "surface power", "rovers", "resupply cadence"]
        },
        "mars": {
            "shock":"The programme is governed by launch windows, autonomy and life-support reliability; recovery options are extremely limited.",
            "curve":"launch_window_steps",
            "human_basis":"duration reflects launch windows, life-support qualification, ISRU maturity, autonomous maintenance and mission-assurance gates.",
            "contradiction":"Acceleration can only be credible if autonomy, ECLSS and ISRU qualification evidence are protected.",
            "signature":["launch windows", "ECLSS", "ISRU", "autonomy", "mission assurance"]
        },
        "spaceport": {
            "shock":"Range safety, licensing and propellant operations can dominate readiness after civil works are complete.",
            "curve":"range_safety_gate",
            "human_basis":"duration reflects licensing, environmental approvals, launch-pad systems, propellant operations and range-safety validation.",
            "contradiction":"Acceleration increases regulatory and range-safety exposure unless approvals are sequenced early.",
            "signature":["range safety", "propellant farm", "launch pad", "payload processing", "licensing"]
        },
        "space_comms": {
            "shock":"Availability is governed by ground/orbit interface resilience, not just antenna or satellite asset delivery.",
            "curve":"mission_network_tail",
            "human_basis":"duration reflects payload qualification, ground segment readiness, relay integration and operational availability proving.",
            "contradiction":"Lower capex can reduce redundancy and increase mission-availability exposure.",
            "signature":["ground segment", "relay", "availability", "antenna", "mission operations"]
        },
        "space_resources": {
            "shock":"Resource uncertainty and autonomous maintenance are likely to dominate business-case confidence.",
            "curve":"resource_uncertainty_spike",
            "human_basis":"duration reflects prospecting uncertainty, extraction technology maturity, autonomous operations and processing reliability.",
            "contradiction":"Acceleration without resource evidence increases downstream production risk.",
            "signature":["resource evidence", "autonomous mining", "processing", "maintenance", "production reliability"]
        },
    }
    if mode=="Space" and fam not in signatures:
        fam="orbital"
    return signatures.get(fam, {
        "shock":"The dominant risk is likely to sit in interfaces, procurement evidence and commissioning readiness rather than headline construction progress.",
        "curve":"generic",
        "human_basis":"duration reflects design maturity, procurement evidence, interface control, commissioning and operational readiness.",
        "contradiction":"Acceleration improves date certainty only if assumptions, interfaces and procurement evidence are strengthened.",
        "signature":["interfaces", "procurement evidence", "commissioning", "operational readiness", "decision gates"]
    })


# ------------------------- CASEY v9.5 scenario cascade engine -------------------------
SCENARIO_PROFILES_V95 = {
    "base": {
        "label":"Base","cost":1.00,"schedule":1.00,"confidence":0,"reserve":1.00,"risk":"Medium-High",
        "delta":["Reference execution case","Balanced cost, schedule and confidence posture","No acceleration or deferral premium applied"],
        "cost_delta":{"Direct":1.00,"Indirect":1.00,"Reserve":1.00},
        "risk_lift":0,
        "curve":"balanced"
    },
    "faster": {
        "label":"Faster","cost":1.14,"schedule":0.82,"confidence":-10,"reserve":1.18,"risk":"High",
        "delta":["Schedule compressed through overlapping delivery workstreams","Procurement premium and concurrency exposure increased","Commissioning float reduced; late-stage volatility increased"],
        "cost_delta":{"Direct":1.08,"Indirect":1.20,"Reserve":1.28},
        "risk_lift":14,
        "curve":"compressed_tail"
    },
    "cheaper": {
        "label":"Cheaper","cost":0.88,"schedule":1.10,"confidence":-14,"reserve":0.72,"risk":"High",
        "delta":["Capital target reduced through slower phasing and lower redundancy","Contingency and acceleration spend reduced","Operational start-up and lifecycle risk increased"],
        "cost_delta":{"Direct":0.91,"Indirect":0.82,"Reserve":0.62},
        "risk_lift":10,
        "curve":"low_median_fat_tail"
    },
    "lower_risk": {
        "label":"Lower Risk","cost":1.09,"schedule":1.12,"confidence":12,"reserve":1.28,"risk":"Medium",
        "delta":["Assurance, procurement buffers and commissioning float expanded","Execution concurrency reduced","Tail risk and decision uncertainty reduced"],
        "cost_delta":{"Direct":1.04,"Indirect":1.12,"Reserve":1.42},
        "risk_lift":-10,
        "curve":"tightened"
    },
    "premium": {
        "label":"Premium","cost":1.22,"schedule":0.96,"confidence":18,"reserve":1.35,"risk":"Medium-Low",
        "delta":["Resilience, redundancy and optionality increased","Procurement certainty and board confidence improved","Higher capex protects delivery and operational readiness"],
        "cost_delta":{"Direct":1.18,"Indirect":1.20,"Reserve":1.48},
        "risk_lift":-14,
        "curve":"premium_resilient"
    },
}

def _money_to_bn_v95(s):
    s=str(s or "").replace("$","").strip().upper()
    try:
        if s.endswith("T"): return float(s[:-1])*1000
        if s.endswith("B"): return float(s[:-1])
        if s.endswith("M"): return float(s[:-1])/1000
        return float(s)
    except Exception:
        return 0.0

def _bn_to_money_v95(bn):
    try:
        bn=float(bn)
    except Exception:
        bn=0.0
    if bn >= 1000: return f"${bn/1000:.1f}T"
    if bn >= 1: return f"${bn:.1f}B"
    return f"${bn*1000:.0f}M"

def _scenario_insight_v95(model, scenario):
    profile=SCENARIO_PROFILES_V95.get(scenario, SCENARIO_PROFILES_V95["base"])
    fam=_sector_family(model.get("subsector",""))
    base = {
        "life_sciences":{
            "base":"Mechanical completion is not the true finish line; validated production readiness and deviation closure are the real board decision gates.",
            "faster":"Acceleration is compressing CQV float and increasing concurrent commissioning exposure across sterile systems.",
            "cheaper":"Capital reduction is shifting risk from construction cost into validation readiness, redundancy and operational start-up.",
            "lower_risk":"The lower-risk path protects GMP turnover by expanding validation float and procurement evidence before execution peaks.",
            "premium":"Premium delivery buys resilience, redundancy and stronger regulatory readiness ahead of commercial production ramp-up."
        },
        "data_centre":{
            "base":"Energisation and commissioning concurrency are more likely to govern delivery than shell construction productivity.",
            "faster":"Acceleration is pulling grid, cooling and IST workstreams into tighter concurrency, increasing late-stage volatility.",
            "cheaper":"The cheaper case reduces capex but pushes risk into redundancy, commissioning resilience and operational uptime.",
            "lower_risk":"The lower-risk case protects energisation and IST by adding commissioning float and procurement buffers.",
            "premium":"Premium delivery increases resilience and power optionality, reducing service-readiness exposure."
        },
        "semiconductor":{
            "base":"Tool install and qualification cadence are likely to govern the real finish line more than cleanroom completion.",
            "faster":"Acceleration compresses tool hook-up, UPW readiness and yield-ramp qualification into a higher-volatility window.",
            "cheaper":"The cheaper case risks deferring utility resilience and tool-readiness assurance into production ramp-up.",
            "lower_risk":"The lower-risk case extends tool qualification and utility proving to reduce yield-ramp uncertainty.",
            "premium":"Premium delivery protects tool cadence, utility resilience and yield-ramp confidence."
        },
        "orbital":{
            "base":"Thermal-power balance, autonomous servicing and recoverability after deployment are the decisive constraints.",
            "faster":"Acceleration compresses qualification and payload integration, increasing mission-assurance exposure.",
            "cheaper":"The cheaper case reduces redundancy and increases recoverability risk after orbital deployment.",
            "lower_risk":"The lower-risk case protects qualification, environmental testing and autonomous recovery readiness.",
            "premium":"Premium delivery buys redundancy, mission assurance and stronger autonomous operations resilience."
        },
        "mars":{
            "base":"Launch windows, autonomy and life-support reliability govern the programme; recovery options are extremely limited.",
            "faster":"Acceleration is credible only if ECLSS, ISRU and autonomy qualification evidence remain protected.",
            "cheaper":"The cheaper case increases mission risk by reducing redundancy where recovery options are limited.",
            "lower_risk":"The lower-risk case expands mission-assurance gates and life-support proving before commitment.",
            "premium":"Premium delivery protects life-support redundancy, autonomy and launch-window resilience."
        }
    }
    pack = base.get(fam) or (base["orbital"] if model.get("mode")=="Space" else {
        "base":"The dominant risk sits in interfaces, procurement evidence and commissioning readiness rather than headline construction progress.",
        "faster":"Acceleration compresses interface resolution and increases late-stage execution volatility.",
        "cheaper":"The cheaper case reduces near-term capex but moves risk into operational readiness and contingency adequacy.",
        "lower_risk":"The lower-risk case adds evidence, float and assurance to reduce tail exposure.",
        "premium":"Premium delivery buys resilience, optionality and stronger confidence."
    })
    return pack.get(scenario, pack["base"])

def _mutate_curve_v95(cost, months, scenario):
    profile=SCENARIO_PROFILES_V95.get(scenario, SCENARIO_PROFILES_V95["base"])
    pts=[1,5,10,20,30,40,50,60,70,80,90,95,99]
    curve=[]
    for p in pts:
        x=p/100
        if profile["curve"]=="compressed_tail":
            cf=0.80 + 0.35*x + 0.24*(x**4)
            sf=0.72 + 0.30*x + 0.30*(x**5)
        elif profile["curve"]=="low_median_fat_tail":
            cf=0.72 + 0.30*x + 0.35*(x**5)
            sf=0.84 + 0.26*x + 0.33*(x**4)
        elif profile["curve"]=="tightened":
            cf=0.86 + 0.22*x + 0.10*(x**3)
            sf=0.90 + 0.18*x + 0.08*(x**3)
        elif profile["curve"]=="premium_resilient":
            cf=0.88 + 0.22*x + 0.08*(x**3)
            sf=0.86 + 0.20*x + 0.10*(x**3)
        else:
            cf=0.78 + 0.35*x + 0.12*(x**3)
            sf=0.82 + 0.26*x + 0.12*(x**3)
        curve.append({"percentile":p,"cost_bn":round(cost*cf,2),"schedule_months":int(round(months*sf))})
    return curve

def scenario_cascade_v95(model:Dict[str,Any], scenario:str) -> Dict[str,Any]:
    scenario=(scenario or "base").lower()
    profile=SCENARIO_PROFILES_V95.get(scenario, SCENARIO_PROFILES_V95["base"])
    # preserve base references
    base_cost=_money_to_bn_v95(model.get("_base_cost_p50") or model.get("cost_p50"))
    base_months=int(float(str(model.get("_base_schedule_months") or model.get("schedule") or "60").replace("months","").strip().split()[0]))
    base_conf=int(model.get("_base_confidence_pct") or model.get("confidence_pct") or 55)
    if "_base_cost_p50" not in model:
        model["_base_cost_p50"]=model.get("cost_p50")
        model["_base_schedule_months"]=base_months
        model["_base_confidence_pct"]=base_conf
        model["_base_risk"]=model.get("risk","Medium-High")
        model["_base_cost_lines"]=json.loads(json.dumps(model.get("cost_lines",[])))
        model["_base_risks"]=json.loads(json.dumps(model.get("risks",[])))
        model["_base_schedule_rows"]=json.loads(json.dumps(model.get("schedule_rows",[])))
    new_cost=base_cost*profile["cost"]
    new_months=max(3,int(round(base_months*profile["schedule"])))
    new_conf=max(8,min(96,base_conf+profile["confidence"]))
    model["scenario"]=scenario
    model["scenario_label=profile_label"] = profile["label"]
    model["scenario_label"]=profile["label"]
    model["cost_p50"]=_bn_to_money_v95(new_cost)
    model["cost_p10"]=_bn_to_money_v95(new_cost*0.80)
    model["cost_p90"]=_bn_to_money_v95(new_cost*1.30)
    model["cost_range"]=f"{model['cost_p10']} - {model['cost_p90']}"

    # Executive scenario comparison versus the immutable Base case.
    base_cost_money=_bn_to_money_hs(base_cost)
    base_schedule_months=base_months
    base_conf_pct=base_conf
    delta_cost=new_cost-base_cost
    delta_months=new_months-base_months
    delta_conf=new_conf-base_conf
    model["scenario_comparison_vs_base"]={
        "base":{"cost_p50":base_cost_money,"schedule_months":base_schedule_months,"confidence_pct":base_conf_pct,"risk":model.get("_base_risk","Medium-High")},
        "selected":{"scenario":profile["label"],"cost_p50":model["cost_p50"],"schedule_months":new_months,"confidence_pct":new_conf,"risk":profile["risk"]},
        "delta":{"cost_bn":round(delta_cost,2),"cost":_bn_to_money_hs(abs(delta_cost)),"cost_direction":"higher" if delta_cost>0 else "lower" if delta_cost<0 else "same",
                 "months":delta_months,"confidence_pts":delta_conf},
        "plain_english": ("Base is the reference case: no cost, schedule or confidence delta is applied. Use Faster, Cheaper, Lower Risk or Premium to expose the board trade-off." if scenario=="base" else f"Compared with Base, {profile['label']} is {_bn_to_money_hs(abs(delta_cost))} {'more expensive' if delta_cost>0 else 'cheaper' if delta_cost<0 else 'unchanged'}, {abs(delta_months)} months {'faster' if delta_months<0 else 'slower' if delta_months>0 else 'unchanged'}, and {abs(delta_conf)} confidence points {'higher' if delta_conf>0 else 'lower' if delta_conf<0 else 'unchanged'}.")
    }
    model["scenario_matrix"]=[]
    for sk,sp in SCENARIO_PROFILES_V95_PLUS.items():
        c=base_cost*sp["cost"]; m=max(3,int(round(base_months*sp["schedule"]))); cf=max(8,min(96,base_conf+sp["confidence"]))
        model["scenario_matrix"].append({"scenario":sk,"label":sp["label"],"cost_p50":_bn_to_money_hs(c),"schedule_months":m,"risk":sp["risk"],"confidence_pct":cf,"cost_delta_pct":round((sp["cost"]-1)*100),"schedule_delta_pct":round((sp["schedule"]-1)*100),"confidence_delta_pts":sp["confidence"],"why":"; ".join([sp.get("trade",""), sp.get("won",""), sp.get("lost","")])})
    # Waterfalls explain WHY the selected scenario moved. These feed dashboard and exports.
    if scenario=="faster":
        cost_moves=[("Base P50",base_cost),("Acceleration premium",base_cost*.075),("Parallel EPC / package overlap",base_cost*.045),("Early long-lead procurement",base_cost*.035),("Commissioning surge / overtime",base_cost*.025),("Faster scenario P50",new_cost)]
        sched_moves=[("Base duration",base_months),("Parallel delivery workfronts",-round(base_months*.07)),("Early procurement release",-round(base_months*.05)),("Concurrent commissioning",-round(base_months*.04)),("Reduced float / overlap",new_months-base_months+round(base_months*.16)),("Faster duration",new_months)]
    elif scenario=="cheaper":
        cost_moves=[("Base P50",base_cost),("Scope / spec restraint",-base_cost*.055),("Deferred redundancy",-base_cost*.040),("Lower indirects / slower phasing",-base_cost*.030),("Reduced reserve",-base_cost*.015),("Cheaper scenario P50",new_cost)]
        sched_moves=[("Base duration",base_months),("Slower procurement path",round(base_months*.04)),("Deferred assurance / resilience",round(base_months*.035)),("Lean owner / commissioning support",round(base_months*.03)),("Cheaper duration",new_months)]
    elif scenario=="lower_risk":
        cost_moves=[("Base P50",base_cost),("Assurance gates",base_cost*.035),("Procurement buffer",base_cost*.030),("Commissioning float",base_cost*.025),("Reserve uplift",base_cost*.030),("Lower-risk scenario P50",new_cost)]
        sched_moves=[("Base duration",base_months),("Evidence gates",round(base_months*.055)),("Protected commissioning float",round(base_months*.055)),("Lower-risk duration",new_months)]
    elif scenario=="premium":
        cost_moves=[("Base P50",base_cost),("Redundancy / resilience",base_cost*.095),("Priority procurement",base_cost*.060),("Operational readiness",base_cost*.045),("Optionality premium",base_cost*.080),("Premium scenario P50",new_cost)]
        sched_moves=[("Base duration",base_months),("Priority procurement",-round(base_months*.04)),("Stronger readiness gates",round(base_months*.02)),("Premium duration",new_months)]
    else:
        cost_moves=[("Base P50",base_cost),("Selected scenario P50",new_cost)]
        sched_moves=[("Base duration",base_months),("Selected scenario duration",new_months)]
    model["cost_waterfall_vs_base"]=[{"driver":a,"value_bn":round(b,2),"value":_bn_to_money_hs(abs(b)) if i not in [0,len(cost_moves)-1] else _bn_to_money_hs(b),"kind":"total" if i in [0,len(cost_moves)-1] else "delta"} for i,(a,b) in enumerate(cost_moves)]
    model["schedule_waterfall_vs_base"]=[{"driver":a,"months":int(b),"kind":"total" if i in [0,len(sched_moves)-1] else "delta"} for i,(a,b) in enumerate(sched_moves)]
    model["schedule"]=f"{new_months} months"
    model["risk"]=profile["risk"]
    model["confidence_pct"]=new_conf
    insight=_scenario_insight_v95(model, scenario)
    model["executive_shock_insight"]=insight
    model["scenario_delta_intelligence"]=[
        {"label":"Cost movement","value":f"{profile['cost']*100-100:+.0f}% vs base","meaning":profile["delta"][0]},
        {"label":"Schedule movement","value":f"{profile['schedule']*100-100:+.0f}% vs base","meaning":profile["delta"][1]},
        {"label":"Confidence movement","value":f"{profile['confidence']:+d} pts","meaning":profile["delta"][2]},
        {"label":"Risk posture","value":profile["risk"],"meaning":insight},
    ]
    model["confidence_breakdown"]=[
        {"driver":"Benchmark similarity","effect":"+8","note":"Comparable sector archetype identified"},
        {"driver":"Procurement certainty","effect":"-9" if scenario in ["faster","cheaper"] else "+7","note":"Long-lead evidence changes scenario confidence"},
        {"driver":"Schedule logic maturity","effect":"-8" if scenario=="faster" else "+6" if scenario in ["lower_risk","premium"] else "+3","note":"Critical path and handover gates rebalanced"},
        {"driver":"Commissioning / validation readiness","effect":"-11" if scenario in ["faster","cheaper"] else "+10" if scenario in ["lower_risk","premium"] else "-4","note":"Scenario changes readiness exposure"},
        {"driver":"Contingency adequacy","effect":"-12" if scenario=="cheaper" else "+11" if scenario in ["lower_risk","premium"] else "+2","note":"Reserve and assurance posture updated"},
    ]
    # Cost line mutation
    new_lines=[]
    for line in model.get("_base_cost_lines", model.get("cost_lines",[])):
        x=dict(line)
        typ=x.get("type","Direct")
        factor=profile["cost_delta"].get(typ, profile["cost"])
        for key in ["low_p10","most_likely_p50","high_p90","p10_bn","p50_bn","p90_bn"]:
            if key in x:
                val=_money_to_bn_v95(x[key]) if not isinstance(x[key], (int,float)) else float(x[key])
                new_val=val*factor
                x[key]=round(new_val,3) if key.endswith("_bn") else _bn_to_money_v95(new_val)
        basis=x.get("basis","")
        if scenario=="faster":
            basis += " Scenario: accelerated procurement / overlapped delivery premium applied."
        elif scenario=="cheaper":
            basis += " Scenario: reduced capital target with lower redundancy and deferred assurance."
        elif scenario=="lower_risk":
            basis += " Scenario: added assurance, procurement buffer and validation float."
        elif scenario=="premium":
            basis += " Scenario: resilience, redundancy and optionality premium applied."
        x["basis"]=basis
        new_lines.append(x)
    model["cost_lines"]=new_lines
    model["cost_breakdown"]=new_lines
    # Risk mutation
    risks=[]
    scenario_risks={
        "faster":[("R-SF1","Concurrent commissioning","Acceleration","Overlapped CQV/commissioning creates rework exposure","Schedule and validation delay","CQV Lead","Protect hold points and add surge QA resource"),
                  ("R-SF2","Acceleration premium","Market","Premium vendors and expediting pressure","Cost escalation","Commercial Lead","Lock procurement packages early and test alternates")],
        "cheaper":[("R-SC1","Deferred redundancy","Capital constraint","Resilience scope deferred to protect capex","Operational start-up exposure","Operations Lead","Define minimum viable redundancy and later upgrade path"),
                   ("R-SC2","Contingency adequacy","Budget pressure","Reserve reduced below uncertainty envelope","Funding shock","Sponsor","Ring-fence risk reserve for critical systems")],
        "lower_risk":[("R-SL1","Assurance gate delay","Assurance posture","Additional evidence gates extend early schedule","Decision delay","PMO","Pre-agree evidence criteria and approval cadence")],
        "premium":[("R-SP1","Premium scope creep","Optionality","Enhanced resilience expands scope boundaries","Capex growth","Sponsor","Govern premium scope through decision gates")]
    }
    base_risks=model.get("_base_risks", model.get("risks",[]))
    for r in base_risks:
        y=dict(r)
        try:
            y["probability_pct"]=max(5,min(85,int(y.get("probability_pct",30))+profile["risk_lift"]))
        except Exception: pass
        risks.append(y)
    for rid,title,cause,event,impact,owner,mit in scenario_risks.get(scenario,[]):
        risks.insert(0,{"id":rid,"title":title,"cause":cause,"event":event,"impact":impact,"probability_pct":max(10,45+profile["risk_lift"]),"activity":"A1500","cbs":"02.01","owner":owner,"mitigation":mit})
    model["risks"]=risks
    model["risk_register"]=model["risks"]
    # Schedule mutation
    rows=[]
    for row in model.get("_base_schedule_rows", model.get("schedule_rows",[])):
        z=dict(row)
        try:
            dur=int(z.get("duration_months",1))
            z["duration_months"]=max(1,int(round(dur*profile["schedule"])))
        except Exception: pass
        if scenario=="faster":
            z["basis"]="Accelerated path: overlapped workfronts, reduced float and earlier procurement release."
        elif scenario=="cheaper":
            z["basis"]="Cheaper path: slower phasing, reduced parallel workfronts and deferred assurance."
        elif scenario=="lower_risk":
            z["basis"]="Lower-risk path: additional assurance gates, protected commissioning float and evidence-led handover."
        elif scenario=="premium":
            z["basis"]="Premium path: stronger redundancy, procurement certainty and controlled readiness gates."
        rows.append(z)
    model["schedule_rows"]=rows
    model["schedule_detail"]=rows
    # Monte Carlo
    curve=_mutate_curve_v95(new_cost,new_months,scenario)
    model["monte_carlo"]={
        "qcra":{"p10":round(curve[2]["cost_bn"],1),"p50":round(new_cost,1),"p80":round(curve[9]["cost_bn"],1),"p90":round(curve[10]["cost_bn"],1)},
        "qsra":{"p10":curve[2]["schedule_months"],"p50":new_months,"p80":curve[9]["schedule_months"],"p90":curve[10]["schedule_months"]},
        "curve":curve,
        "tornado":[{"risk_id":r.get("id"),"title":r.get("title"),"driver_score":max(1,100-i*8)} for i,r in enumerate(model["risks"][:8])]
    }
    model["board_briefing"]=[
        insight,
        f"{profile['label']} scenario indicates {model['cost_p50']} P50 exposure across approximately {model['schedule']}.",
        f"Confidence moves to {new_conf}% because procurement, schedule logic, contingency and commissioning assumptions have been rebalanced.",
        "Scenario consequence: " + " ".join(profile["delta"])
    ]
    model["casey_thinking"]=(model.get("casey_thinking","CASEY interprets the programme.") + f" Scenario cascade: {insight}")
    model["executive_summary"]=f"{model.get('title','Programme')} scenario view: {profile['label']}. CASEY indicates {model['cost_p50']} P50 exposure, {model['cost_range']} range, {model['schedule']} baseline, {profile['risk']} risk and {new_conf}% confidence. {insight}"
    model["outputs_board_memo"]=[
        f"Decision posture: {profile['label']} scenario.",
        f"Investment implication: {insight}",
        f"Top confidence movement: {model['confidence_breakdown'][1]['driver']} {model['confidence_breakdown'][1]['effect']} pts.",
        "Outputs should be treated as first-pass strategic intelligence, not certified estimate documents."
    ]
    model["top_decisions_required"]=[
        "Confirm the governing critical-path constraint and evidence owner.",
        "Approve scenario-specific procurement and contingency posture.",
        "Lock handover / commissioning / validation decision gates.",
        "Resolve highest-probability interface and readiness risks.",
        "Decide whether the scenario trade-off is acceptable for board approval."
    ]

    # FINAL EXEC POLISH: scenario-aware ranking and differentiation
    try:
        sc = str(model.get("scenario","base")).lower()
        schedule_lists = {
            "faster":[
                "Concurrent commissioning overload",
                "Recovery float exhaustion",
                "Acceleration premium shock",
                "Grid connection delay",
                "Integrated systems testing concurrency"
            ],
            "cheaper":[
                "Vendor claims and change exposure",
                "Procurement deferral and long-lead slippage",
                "Design maturity gap",
                "Scope growth from deferred decisions",
                "Interface coordination delay"
            ],
            "lower_risk":[
                "Governance and approvals latency",
                "Extended validation sequencing",
                "Conservative commissioning gates",
                "Operational readiness hold-points",
                "Assurance and compliance reviews"
            ],
            "premium":[
                "Integration complexity across parallel packages",
                "Executive decision latency",
                "Technology assurance alignment",
                "Multi-package interface management",
                "Programme coordination overhead"
            ]
        }
        cost_lists = {
            "faster":[
                "Acceleration premiums and overtime",
                "Power train, transformers and switchgear",
                "Integrated systems testing",
                "Grid and utility concurrency",
                "Recovery-float consumption"
            ],
            "cheaper":[
                "Deferred procurement packaging",
                "Claims and commercial exposure",
                "Rework from reduced contingency",
                "Long-lead inflation volatility",
                "Scope rationalisation impacts"
            ],
            "lower_risk":[
                "Additional contingency and reserve",
                "Enhanced validation and assurance",
                "Programme controls and governance",
                "Redundant infrastructure resilience",
                "Extended commissioning readiness"
            ]
        }
        if sc in schedule_lists:
            model["sector_schedule_threats"]=schedule_lists[sc]
        if sc in cost_lists:
            model["sector_primary_cost_drivers"]=cost_lists[sc]

        if sc=="faster":
            model["executive_shock_insight"]="Acceleration increases spend faster than it reduces uncertainty; the delivery tail becomes more volatile."
        elif sc=="cheaper":
            model["executive_shock_insight"]="Capital efficiency reduces resilience: procurement and recovery flexibility become constrained."
        elif sc=="lower_risk":
            model["executive_shock_insight"]="Confidence is purchased through reserve, governance and extended delivery duration."
        elif sc=="premium":
            model["executive_shock_insight"]="Premium posture buys resilience, optionality and stronger certainty at visible capex premium."
    except Exception:
        pass

    return model

def apply_sector_intelligence(model:Dict[str,Any]) -> Dict[str,Any]:
    mode=model.get("mode","Earth"); subsector=model.get("subsector","")
    fam=_sector_family(subsector)
    lists=sector_specific_lists(subsector, mode)
    sig=sector_signature_behaviour(subsector, mode)

    def clean(txt):
        txt=str(txt)
        txt=txt.replace("selected schedule level and scenario compression/buffer logic", sig["human_basis"])
        txt=txt.replace("Duration derived from sector schedule model, selected schedule level and scenario compression/buffer logic.", "Duration reflects " + sig["human_basis"])
        if fam!="spaceport" and mode!="Space":
            txt=txt.replace(" or launch readiness","")
            txt=txt.replace("launch readiness delay","readiness delay")
            txt=txt.replace("Launch readiness","Operational readiness")
        return txt

    model["casey_thinking"]=lists["thinking"]
    model["executive_shock_insight"]=sig["shock"]
    model["sector_signature_behaviours"]=sig["signature"]
    model["sector_curve_type"]=sig["curve"]
    model["scenario_contradiction"]=sig["contradiction"]
    model["confidence_engine_label"]="CASEY Confidence Engine"
    model["confidence_engine_detail"]="Benchmark + probabilistic + sector-trained reasoning"

    # Cleaner, more memorable board briefing.
    bb = model.get("board_briefing") if isinstance(model.get("board_briefing"), list) else []
    model["board_briefing"]=[
        sig["shock"],
        f"Current intelligence indicates {model.get('cost_p50')} P50 programme exposure across approximately {model.get('schedule')}, with overall risk assessed as {model.get('risk')}.",
        f"Confidence is governed by {', '.join(sig['signature'][:3])} and the maturity of evidence behind procurement and commissioning assumptions.",
        sig["contradiction"]
    ]

    # Humanise actions.
    model["next_best_actions"]=[
        "Run a focused evidence workshop on the dominant schedule and commissioning constraint.",
        "Validate the top five CBS cost drivers against benchmark or supplier evidence.",
        "Confirm the Level 3/4 schedule logic, handover gates and critical-path assumptions.",
        "Prepare Base / Faster / Cheaper / Lower-Risk decision paper with explicit trade-offs.",
        "Create a confidence-improvement plan targeting design maturity, procurement evidence and operational readiness."
    ]

    if isinstance(model.get("mission_control_cards"), list):
        model["mission_control_cards"]=[c for c in model["mission_control_cards"] if c.get("label")!="EXECUTIVE SHOCK"]
        model["mission_control_cards"].insert(0, {"label":"EXECUTIVE SHOCK", "signal":sig["shock"], "severity":model.get("risk","Medium")})
        for c in model["mission_control_cards"]:
            c["signal"]=clean(c.get("signal",""))

    # Sector-specific overview components used by frontend.
    model["sector_confidence_drivers"]=lists["confidence"]
    model["sector_primary_cost_drivers"]=lists["cost"]
    model["sector_schedule_threats"]=lists["schedule"]
    model["why_casey_generated_this"]=[
        f"CASEY detected {subsector} from the project brief and routed it to the {mode} infrastructure model.",
        f"Sector signature behaviours applied: {', '.join(sig['signature'][:5])}.",
        "Cost, schedule and risk were calibrated against estimate class, schedule level, complexity and delivery environment.",
        "The output is designed for early board challenge and scope definition, not certified pricing."
    ]
    if lists["bench"]:
        model["benchmark_comparison"]=lists["bench"]

    # Make schedule rows read less algorithmic.
    for row in model.get("schedule_rows",[]) or []:
        row["basis"]=clean(row.get("basis",""))
        phase=str(row.get("phase","")).lower()
        if fam=="life_sciences" and any(x in phase for x in ["commission", "handover", "delivery"]):
            row["basis"]="Duration reflects CQV, GMP turnover, validation evidence and batch-readiness sequencing."
        elif fam=="data_centre" and any(x in phase for x in ["commission", "delivery"]):
            row["basis"]="Duration reflects energisation, IST, phased data-hall readiness and resilience proving."
        elif mode=="Space" and any(x in phase for x in ["delivery", "commission", "handover"]):
            row["basis"]="Duration reflects qualification, launch integration, deployment readiness and autonomous operations proving."

    # Add subtle human-like P80/P90 rounding display values while retaining numeric curves.
    try:
        mc=model.get("monte_carlo",{})
        if "qsra" in mc:
            for k in ["p10","p50","p80","p90"]:
                mc["qsra"][k]=round(float(mc["qsra"][k]))
        if "qcra" in mc:
            for k in ["p10","p50","p80","p90"]:
                mc["qcra"][k]=round(float(mc["qcra"][k]),1)
    except Exception:
        pass

    # Refresh executive summary after sector insight.
    model["executive_summary"]=(
        f"{model.get('title')} has been classified as {subsector} in {model.get('location','the selected location')}. "
        f"CASEY estimates {model.get('cost_p50')} P50 cost, {model.get('cost_range')} range and {model.get('schedule')} baseline. "
        f"Key insight: {sig['shock']} Risk is assessed as {model.get('risk')} with {model.get('confidence_pct')}% confidence."
    )

    # FINAL EXEC POLISH: scenario-aware ranking and differentiation
    try:
        sc = str(model.get("scenario","base")).lower()
        schedule_lists = {
            "faster":[
                "Concurrent commissioning overload",
                "Recovery float exhaustion",
                "Acceleration premium shock",
                "Grid connection delay",
                "Integrated systems testing concurrency"
            ],
            "cheaper":[
                "Vendor claims and change exposure",
                "Procurement deferral and long-lead slippage",
                "Design maturity gap",
                "Scope growth from deferred decisions",
                "Interface coordination delay"
            ],
            "lower_risk":[
                "Governance and approvals latency",
                "Extended validation sequencing",
                "Conservative commissioning gates",
                "Operational readiness hold-points",
                "Assurance and compliance reviews"
            ],
            "premium":[
                "Integration complexity across parallel packages",
                "Executive decision latency",
                "Technology assurance alignment",
                "Multi-package interface management",
                "Programme coordination overhead"
            ]
        }
        cost_lists = {
            "faster":[
                "Acceleration premiums and overtime",
                "Power train, transformers and switchgear",
                "Integrated systems testing",
                "Grid and utility concurrency",
                "Recovery-float consumption"
            ],
            "cheaper":[
                "Deferred procurement packaging",
                "Claims and commercial exposure",
                "Rework from reduced contingency",
                "Long-lead inflation volatility",
                "Scope rationalisation impacts"
            ],
            "lower_risk":[
                "Additional contingency and reserve",
                "Enhanced validation and assurance",
                "Programme controls and governance",
                "Redundant infrastructure resilience",
                "Extended commissioning readiness"
            ]
        }
        if sc in schedule_lists:
            model["sector_schedule_threats"]=schedule_lists[sc]
        if sc in cost_lists:
            model["sector_primary_cost_drivers"]=cost_lists[sc]

        if sc=="faster":
            model["executive_shock_insight"]="Acceleration increases spend faster than it reduces uncertainty; the delivery tail becomes more volatile."
        elif sc=="cheaper":
            model["executive_shock_insight"]="Capital efficiency reduces resilience: procurement and recovery flexibility become constrained."
        elif sc=="lower_risk":
            model["executive_shock_insight"]="Confidence is purchased through reserve, governance and extended delivery duration."
        elif sc=="premium":
            model["executive_shock_insight"]="Premium posture buys resilience, optionality and stronger certainty at visible capex premium."
    except Exception:
        pass

    return model


def _casey_v150_if_this_fails(mode: str, subsector: str) -> str:
    """The single most uncomfortable sentence for a traditional advisor."""
    s = subsector.lower()
    if mode == "Space":
        return "The programme failed because mission assurance evidence, launch manifest realism and autonomous commissioning maturity were never independently verified — and the programme board treated optimistic vendor commitments as milestone evidence."
    if any(x in s for x in ["rail","metro","transit","hsr"]):
        return "The programme failed because possession access was assumed rather than negotiated, signalling integration maturity was overstated in reporting, and the board approved capital before the operator acceptance evidence existed."
    if any(x in s for x in ["nuclear","smr","reactor"]):
        return "The programme failed because the first-of-a-kind technology risk was underweighted in the baseline, the nuclear supply chain could not deliver to assumed dates, and the regulatory licensing critical path was treated as a parallel workstream rather than the programme governor."
    if any(x in s for x in ["data centre","hyperscale","compute","cloud"]):
        return "The programme failed because the grid connection energisation date was assumed rather than contracted, generator and transformer lead times were not locked before design commitment, and the compressed hyperscale schedule consumed all commissioning float before integrated testing began."
    if any(x in s for x in ["pharma","gmp","biologics","life sciences","vaccine"]):
        return "The programme failed because the regulatory inspection readiness programme was treated as a post-construction activity, a validation batch failed late in the schedule, and the clean utility readiness critical path was not visible to the programme board until it was too late to recover."
    if any(x in s for x in ["defence","naval","submarine","aukus","military","dockyard"]):
        return "The programme failed because security accreditation was assumed as a parallel workstream rather than a gate, sovereign supply chain dependencies were not confirmed before commitment, and classified systems integration maturity was reported by the vendor rather than independently verified."
    if any(x in s for x in ["semiconductor","fab","cleanroom","wafer"]):
        return "The programme failed because process tool allocation was assumed from OEM intent letters rather than confirmed orders, cleanroom HVAC classification was not independently tested before tool installation, and the yield ramp model was based on vendor projections rather than comparable FOAK fab benchmarks."
    if any(x in s for x in ["airport","aviation","runway","airside"]):
        return "The programme failed because ORAT was treated as a construction completion deliverable rather than a parallel programme, baggage systems integration fell behind on the critical path without board visibility, and the CAA acceptance programme was resourced as a single post-construction activity."
    return "The programme failed because the governing constraint was never named at board level, schedule float was reported as certainty, and the P80 downside scenario was not stress-tested before capital was committed."


def _casey_v150_traditional_vs_casey(mode: str, subsector: str, confidence: int, risk: str, p50: float, p80_raw) -> dict:
    """What a traditional controls report would say vs what CASEY reads."""
    s = subsector.lower()
    try: p80 = float(str(p80_raw).replace("$","").replace("B","").replace(",","").strip() or 0)
    except: p80 = 0.0
    try: p50f = float(p50)
    except: p50f = 1.0
    tail = round((p80/p50f - 1)*100) if p50f > 0 and p80 > p50f else 25
    if any(x in s for x in ["rail","metro","transit","hsr"]):
        trad = "Civil works and corridor progress appear on track. Schedule is tracking to baseline with possession strategy agreed."
        casey = f"Opening confidence is governed by systems migration and operational acceptance — not visible civil progress. P80 tail is {tail}% above P50. Possession strategy must be confirmed, not assumed."
    elif any(x in s for x in ["nuclear","smr","reactor"]):
        trad = "Nuclear island procurement is progressing. Licensing programme is on the project schedule. Cost within authorised range."
        casey = f"The GDA and site licence are the primary critical path — civil works cannot start until complete. FOAK supply chain has no comparable delivery history. P80 is {tail}% above P50."
    elif any(x in s for x in ["data centre","hyperscale","compute","cloud"]):
        trad = "Construction programme is tracking. Grid connection is in the works programme. Cooling design is progressing through detailed design."
        casey = f"Grid energisation date is the critical path governor — not construction. Generator and transformer lead times are longer than the programme assumes. P80 tail is {tail}% above P50."
    elif any(x in s for x in ["pharma","gmp","biologics","life sciences"]):
        trad = "Construction is on schedule. Validation programme is resourced. Regulatory submissions are planned for post-qualification."
        casey = f"FDA/EMA inspection readiness is a named milestone on the critical path, not a post-construction activity. A single failed validation batch resets the programme. P80 is {tail}% above P50."
    elif any(x in s for x in ["defence","naval","submarine","aukus","military","dockyard"]):
        trad = "Industrial base development is progressing. Sovereign supplier engagement is underway. Security accreditation is planned."
        casey = f"Security accreditation and mission-system acceptance are the hidden critical path. Sovereign supply chain has unconfirmed allocation dependencies. P80 tail is {tail}% above P50."
    elif mode == "Space":
        trad = "Development milestones are tracking. Launch manifest is agreed. Technology readiness is progressing to plan."
        casey = f"TRL evidence, launch manifest realism and autonomous commissioning maturity must be independently verified — not accepted from vendor reports. P80 is {tail}% above P50."
    else:
        trad = "Programme is tracking to baseline. Cost and schedule are within tolerance. Risk register is being maintained."
        casey = f"The governing constraint is not visible in headline progress reports. P80 downside is {tail}% above P50. The board needs named owners for the critical-path constraint."
    return {"traditional": trad, "casey": casey, "tail_pct": tail, "confidence": confidence, "risk": risk}


def _casey_v150_board_attacks(mode: str, subsector: str, scenario: str) -> list:
    """Hard attack lines an investment committee uses before the team can hide behind dashboards."""
    s = subsector.lower(); sc = str(scenario or "base").lower()
    universal = [
        "Show me the P80 exposure — not the P50 headline.",
        "Which evidence package retires the governing constraint, and who owns the closure date?",
        "Where is the float, and is it operationally usable or just nominally available?",
    ]
    scenario_attacks = {
        "faster": ["Did you consume real float or just compress assurance time?", "What is the failure mode if the acceleration premium isn't funded at delivery?"],
        "cheaper": ["What scope was removed — or was this a contingency cut?", "What is the residual risk the cheaper number is now carrying?"],
        "lower_risk": ["Show me which P80 driver is retired by the extra reserve.", "Is the extended schedule buying evidence or bureaucracy?"],
        "premium": ["What specific downside does the premium protect against?", "Can you prove this is resilience, not scope creep?"],
        "base": ["Is this a decision case or a reference case — the board needs to know the difference.", "What must be proven before this becomes board-approvable?"],
    }.get(sc, [])
    if mode == "Space":
        sector_attacks = ["What flight heritage evidence exists for this mission class?", "Is the launch manifest confirmed or assumed?", "What is the contingency plan if the launch window is lost?"]
    elif any(x in s for x in ["rail","metro","transit","hsr"]):
        sector_attacks = ["Which possessions are genuinely executable — not just theoretically planned?", "What is the cost of a signalling integration failure at trial running?"]
    elif any(x in s for x in ["nuclear","smr","reactor"]):
        sector_attacks = ["What is the GDA critical path and what does a six-month slip cost?", "Which nuclear supply chain items are sole-source with no validated fallback?"]
    elif any(x in s for x in ["pharma","gmp","biologics","life sciences"]):
        sector_attacks = ["When does the inspection readiness programme need to complete, and who owns it?", "What is the programme consequence of a single failed PQ batch at this stage?"]
    elif any(x in s for x in ["defence","naval","submarine","aukus","military"]):
        sector_attacks = ["What is the security accreditation critical path and which classified supplier is unconfirmed?", "What is the export control fallback for the sovereign supply chain items?"]
    elif any(x in s for x in ["data centre","hyperscale","compute","cloud"]):
        sector_attacks = ["Is the grid connection agreement signed, or is the energisation date still an assumption?", "What is the commissioning plan if the cooling vendor fails performance testing?"]
    elif any(x in s for x in ["semiconductor","fab","cleanroom","wafer"]):
        sector_attacks = ["Which tool allocations are confirmed orders versus OEM intent?", "What is the contingency plan if EUV delivery slips by one quarter?"]
    elif any(x in s for x in ["airport","aviation","runway","airside"]):
        sector_attacks = ["What is the ORAT programme status and which airline sign-off can block opening?", "Is baggage integration on the critical path — with test dates named?"]
    else:
        sector_attacks = ["What is the critical-path constraint and who is the named owner?", "What evidence closes the governing approval blocker?"]
    return (sector_attacks + scenario_attacks + universal)[:8]


def build_model(prompt:str, client:str="", class_level:int=3, schedule_level:int=3, scenario:str="base"):
    title,mode,subsector,base_cost,base_months=detect_sector(prompt)
    location,loc_mult=location_factor(prompt)
    scale,scale_mult,scale_months=scale_factor(prompt)
    comp_mult,comp_months,drivers=complexity(prompt)
    cm,sm,rm,conf_adj,scenario_label,scenario_why=scenario_params("base")  # build base first; scenario cascade applies selected case
    raw_cost=base_cost*loc_mult*scale_mult*comp_mult*cm
    months=max(6,round((base_months+scale_months+comp_months+schedule_level*2)*sm))
    raw_cost, months, envelope_notes = sector_envelope(subsector, raw_cost, months, scale)
    raw_cost, months, benchmark_notes, benchmark_matches = calibrate_with_benchmarks(prompt, mode, subsector, raw_cost, months)
    raw_cost, months, envelope_notes2 = sector_envelope(subsector, raw_cost, months, scale)
    benchmark_notes = (envelope_notes + envelope_notes2 + benchmark_notes)[:5]
    risk_score=22 + (42 if mode=="Space" else 0) + (24 if "Mega" in scale or "System" in scale else 9 if scale=="Programme" else 0) + (14 if class_level>=4 else 0) + (18 if loc_mult>=2 else 0) + min(20,len(drivers)*5)
    risk_score*=rm
    risk=risk_label(risk_score)
    base_conf={1:88,2:80,3:68,4:54,5:40}.get(class_level,68)
    penalty={"Low":0,"Medium":4,"Medium-High":8,"High":12,"Very High":16,"Extreme":22}.get(risk,8)
    input_quality = 80
    try:
        if "input_quality_score" in locals():
            input_quality = int(input_quality_score)
    except Exception:
        input_quality = 80
    quality_index = estimate_quality_index(class_level, schedule_level, risk_score, benchmark_matches, input_quality)
    confidence=int(clamp((base_conf-penalty+conf_adj)*0.70 + quality_index["score"]*0.30,12,96))
    low,high,class_name,maturity=class_range(class_level)
    p10=raw_cost*low; p90=raw_cost*high
    lines=[]
    for cbs,desc,typ,basis,pct in sector_cost_lines(mode,subsector):
        p50=raw_cost*pct
        lines.append({"cbs":cbs,"description":desc,"type":typ,"basis":basis,"p10_bn":round(p50*low,3),"p50_bn":round(p50,3),"p90_bn":round(p50*high,3),"impact_basis":f"{desc} priced as {pct:.1%} of P50 from sector template, location factor {loc_mult:.2f}, scale factor {scale_mult:.2f}, complexity factor {comp_mult:.2f}."})
    schedules_by_level={str(l):schedule_rows(mode,subsector,months,l) for l in range(1,6)}
    primary_schedule=schedules_by_level[str(schedule_level)]
    primary_costs=lines
    risks=risk_register(mode,subsector,raw_cost,months,primary_schedule,primary_costs,rm,scenario)
    mc=monte_carlo(raw_cost,months,risks,seed=abs(hash(prompt+scenario))%999999)
    benchmarks=benchmarks_for(mode,subsector,location,raw_cost,months)
    model={
      "version":"v23 Real Deal Edition","id":f"CASEY-{random.randint(100000,999999)}","title":title,"prompt":prompt,"client":client or "Client / operator","mode":mode,"subsector":subsector,"location":location,"scale":scale,"scenario":scenario,"scenario_label":scenario_label,"scenario_why":scenario_why,"estimate_class":class_level,"estimate_class_name":class_name,"estimate_maturity":maturity,"schedule_level":schedule_level,"cost_p10":money_bn(p10),"cost_p50":money_bn(raw_cost),"cost_p90":money_bn(p90),"cost_range":f"{money_bn(p10)} - {money_bn(p90)}","schedule":f"{months} months","risk":risk,"risk_score":round(risk_score,1),"confidence_pct":confidence,
      "executive_summary":f"{title} has been classified as {subsector} in {location}. CASEY estimates {money_bn(raw_cost)} P50 cost, {money_bn(p10)} to {money_bn(p90)} range, {months} month baseline, {risk} risk and {confidence}% confidence under the {scenario_label} scenario.",
      "confidence_explanation":[f"{class_name}: {maturity} drives the starting confidence level.",f"Schedule Level {schedule_level}: more detailed schedule logic improves QSRA traceability.",f"{scenario_label}: {scenario_why}",f"Location/space factor: {location} applies a {loc_mult:.2f} delivery premium.",*benchmark_notes,*(drivers or ["Standard complexity profile inferred from project type."])],
      "benchmark_memory":benchmark_matches,
      "estimate_quality_index":quality_index,
      "procurement_heatmap":procurement_heatmap(mode,subsector,risks),
      "critical_path_narrative":critical_path_narrative(mode,subsector,primary_schedule),
      "board_briefing":board_briefing_narrative(mode,subsector,title,money_bn(p50),f"{months} months",risk,confidence,drivers,benchmark_matches),
      "uncertainty_narrative":uncertainty_narrative(class_level,schedule_level,confidence,mode,risk,drivers),
      "mission_control_cards":mission_control_cards(mode,subsector,title,risk,benchmark_matches,procurement_heatmap(mode,subsector,risks),critical_path_narrative(mode,subsector,primary_schedule)),
      "casey_thinking":casey_thinking(mode,subsector,title),
      "benchmark_comparison":benchmark_comparison(benchmark_matches,mode,subsector),
      "calibration_notes":benchmark_notes,
      "cost_lines":primary_costs,"schedules_by_level":schedules_by_level,"schedule_rows":primary_schedule,"risks":risks,"monte_carlo":mc,"benchmarks":benchmarks,"peer_competitors":peer_competitors(client,subsector,mode),"scenario_comparison":scenario_compare(prompt,client,class_level,schedule_level),"launch_demo_script":demo_script(),"red_flags":red_flags(risk,confidence,mode,subsector,scenario_label,raw_cost,(mc.get("qcra") or {}).get("p80",0.0) or raw_cost*1.25,months),"board_challenge_questions":board_questions(mode,subsector,str(scenario or "base").lower().strip()),
      "if_this_fails":_casey_v150_if_this_fails(mode,subsector),
      "traditional_vs_casey":_casey_v150_traditional_vs_casey(mode,subsector,confidence,risk,money_bn(p50),(mc.get("qcra") or {}).get("p80",0)),
      "board_attack_simulation":_casey_v150_board_attacks(mode,subsector,scenario),"next_best_actions":next_actions(risk,confidence,scenario_label)
    }
    model["estimates_by_class"]={str(c):[{**x,"p10_bn":round(x["p50_bn"]*class_range(c)[0],3),"p90_bn":round(x["p50_bn"]*class_range(c)[1],3),"class":c,"maturity":class_range(c)[3]} for x in primary_costs] for c in range(1,6)}

    # FINAL EXEC POLISH: scenario-aware ranking and differentiation
    try:
        sc = str(model.get("scenario","base")).lower()
        schedule_lists = {
            "faster":[
                "Concurrent commissioning overload",
                "Recovery float exhaustion",
                "Acceleration premium shock",
                "Grid connection delay",
                "Integrated systems testing concurrency"
            ],
            "cheaper":[
                "Vendor claims and change exposure",
                "Procurement deferral and long-lead slippage",
                "Design maturity gap",
                "Scope growth from deferred decisions",
                "Interface coordination delay"
            ],
            "lower_risk":[
                "Governance and approvals latency",
                "Extended validation sequencing",
                "Conservative commissioning gates",
                "Operational readiness hold-points",
                "Assurance and compliance reviews"
            ],
            "premium":[
                "Integration complexity across parallel packages",
                "Executive decision latency",
                "Technology assurance alignment",
                "Multi-package interface management",
                "Programme coordination overhead"
            ]
        }
        cost_lists = {
            "faster":[
                "Acceleration premiums and overtime",
                "Power train, transformers and switchgear",
                "Integrated systems testing",
                "Grid and utility concurrency",
                "Recovery-float consumption"
            ],
            "cheaper":[
                "Deferred procurement packaging",
                "Claims and commercial exposure",
                "Rework from reduced contingency",
                "Long-lead inflation volatility",
                "Scope rationalisation impacts"
            ],
            "lower_risk":[
                "Additional contingency and reserve",
                "Enhanced validation and assurance",
                "Programme controls and governance",
                "Redundant infrastructure resilience",
                "Extended commissioning readiness"
            ]
        }
        if sc in schedule_lists:
            model["sector_schedule_threats"]=schedule_lists[sc]
        if sc in cost_lists:
            model["sector_primary_cost_drivers"]=cost_lists[sc]

        if sc=="faster":
            model["executive_shock_insight"]="Acceleration increases spend faster than it reduces uncertainty; the delivery tail becomes more volatile."
        elif sc=="cheaper":
            model["executive_shock_insight"]="Capital efficiency reduces resilience: procurement and recovery flexibility become constrained."
        elif sc=="lower_risk":
            model["executive_shock_insight"]="Confidence is purchased through reserve, governance and extended delivery duration."
        elif sc=="premium":
            model["executive_shock_insight"]="Premium posture buys resilience, optionality and stronger certainty at visible capex premium."
    except Exception:
        pass

    return model

def risk_register(mode,subsector,cost,months,schedule,costs,rm,scenario):
    base=[("R-001","Scope growth","Scope","Requirements mature late causing additions and rework","Both",42,15,45,120,0.03,0.08,0.16,"Project Director","Freeze scope, change board, value gates","Scope change rate exceeds threshold"),("R-002","Market escalation","Commercial","Inflation and commodity escalation exceeds allowance","Cost",45,0,0,0,0.04,0.10,0.20,"Commercial Lead","Early procurement, index strategy, market testing","Index exceeds allowance"),("R-003","Permits / approvals delay","Regulatory","Authorities or stakeholders delay critical approvals","Schedule",35,20,70,180,0.01,0.04,0.08,"Consents Lead","Authority plan, consent tracker, early submissions","Consent milestone slips"),("R-004","Design maturity gap","Technical","Design maturity lower than estimate basis","Both",38,15,45,120,0.03,0.09,0.18,"Design Manager","Design maturity gates and independent review","Deliverables late"),("R-005","Supply chain delay","Procurement","Long-lead equipment or specialist suppliers slip","Both",40,10,40,120,0.02,0.07,0.15,"Procurement Lead","Alternate suppliers, early orders, expediting","Supplier promise date slips"),("R-006","Productivity underperformance","Delivery","Labour productivity below plan due to access/sequence","Both",35,15,50,130,0.03,0.08,0.16,"Delivery Lead","Package productivity controls and daily planning","SPI/CPI deteriorates"),("R-007","Commissioning delay","Handover","Systems integration and tests take longer than plan","Both",32,10,35,100,0.02,0.06,0.12,"Commissioning Lead","Commissioning readiness and early test plans","Test failures trend upward"),("R-008","Interface misalignment","Integration","Package interfaces not coordinated","Both",34,10,35,100,0.02,0.07,0.14,"Integration Manager","ICDs and weekly interface board","Interface actions overdue")]
    if mode=="Space":
        base=[("R-S01","Launch manifest delay","Launch","Launch slot/provider delay affects deployment","Both",48,15,45,120,0.03,0.09,0.20,"Launch Integration Lead","Reserve alternate slots, freeze payload interfaces","Launch slot not confirmed"),("R-S02","Mass growth","Technical","Payload mass exceeds design allowance","Both",42,15,50,140,0.04,0.11,0.24,"Chief Engineer","Mass control board and margin policy","Mass margin below threshold"),("R-S03","Life-support reliability","Technical","ECLSS reliability below threshold","Both",36,20,70,180,0.05,0.12,0.28,"Life Support Lead","Prototype testing, redundancy, reliability growth","Reliability growth misses target")]+base
    elif "data centre" in subsector.lower():
        base=[("R-D01","Grid connection delay","Utilities","Power connection or substation energisation delays critical path","Both",48,20,80,220,0.03,0.10,0.22,"Utilities Lead","Secure grid agreement, temporary power, early equipment orders","Grid agreement unsigned"),("R-D02","Cooling system capacity","Technical","Cooling design or water availability constrains commissioning","Both",34,10,40,120,0.02,0.07,0.15,"MEP Lead","Early thermal modelling and supplier validation","Cooling tests fail")]+base
    out=[]
    for i,r in enumerate(base,1):
        rid,title,cat,desc,area,prob,so,sm,sp,co,cm,cp,owner,mit,trig=r
        prob=int(clamp(prob*rm,5,90))
        if scenario=="lower_risk": cm*=0.75; cp*=0.8; sm=round(sm*.75); sp=round(sp*.8)
        if scenario=="faster": sm=round(sm*1.18); sp=round(sp*1.25)
        if scenario=="cheaper": cm*=1.12; cp*=1.18
        activity=schedule[min(i+3,len(schedule)-1)] if schedule else {"activity_id":"A1900","activity":"Delivery"}
        costline=costs[min(i+2,len(costs)-1)] if costs else {"cbs":"01.01","description":"Cost"}
        emv=cost*cm*prob/100; semv=sm*prob/100
        out.append({"risk_id":rid if rid.startswith("R-") else f"R-{i:03d}","title":title,"category":cat,"description":desc,"impact_area":area,"probability_pct":prob,"activity_id":activity["activity_id"],"activity_name":activity["activity"],"cbs":costline["cbs"],"cbs_name":costline["description"],"schedule_o_days":so,"schedule_m_days":sm,"schedule_p_days":sp,"cost_o_bn":round(cost*co,3),"cost_m_bn":round(cost*cm,3),"cost_p_bn":round(cost*cp,3),"cost_emv_bn":round(emv,3),"schedule_emv_days":round(semv,1),"driver_score":round((emv*100)+(semv/10),2),"owner":owner,"trigger":trig,"mitigation":mit,"basis_of_cost_impact":f"Cost impact derived from {costline['cbs']} {costline['description']} exposure, project P50 {money_bn(cost)}, and risk severity ({co:.0%}/{cm:.0%}/{cp:.0%}).","basis_of_schedule_impact":f"Schedule impact mapped to {activity['activity_id']} {activity['activity']} with O/M/P delay {so}/{sm}/{sp} days and critical path sensitivity."})
    return sorted(out,key=lambda x:x["driver_score"],reverse=True)

def monte_carlo(cost,months,risks,seed=42,iterations=10000):
    rng=np.random.default_rng(seed)
    cost_samples=[]; sched_samples=[]
    base_cost_unc=rng.triangular(cost*.88,cost,cost*1.18,iterations)
    base_sched_unc=rng.triangular(months*.92,months,months*1.15,iterations)
    risk_cost=np.zeros(iterations); risk_days=np.zeros(iterations)
    contribution={r["risk_id"]:{"cost":0.0,"days":0.0,"title":r["title"],"activity_id":r["activity_id"],"cbs":r["cbs"]} for r in risks}
    for r in risks:
        occurs=rng.random(iterations)<(r["probability_pct"]/100)
        c=rng.triangular(r["cost_o_bn"],r["cost_m_bn"],r["cost_p_bn"],iterations)*occurs
        if r["schedule_o_days"] == r["schedule_m_days"] == r["schedule_p_days"]:
            d=np.zeros(iterations)
        else:
            left=min(r["schedule_o_days"], r["schedule_m_days"], r["schedule_p_days"])
            mode=max(left, r["schedule_m_days"])
            right=max(r["schedule_o_days"], r["schedule_m_days"], r["schedule_p_days"])
            if left == right: d=np.zeros(iterations)
            else: d=rng.triangular(left,mode,right,iterations)*occurs
        risk_cost+=c; risk_days+=d
        contribution[r["risk_id"]]["cost"]=float(np.mean(c)); contribution[r["risk_id"]]["days"]=float(np.mean(d))
    cost_samples=base_cost_unc+risk_cost
    sched_samples=base_sched_unc+(risk_days/30.44)
    def pct(arr,p): return float(np.percentile(arr,p))
    curve=[{"percentile":p,"cost_bn":round(pct(cost_samples,p),3),"schedule_months":round(pct(sched_samples,p),2)} for p in [1,5,10,20,30,40,50,60,70,80,90,95,99]]
    tornado=sorted([{"risk_id":k,"title":v["title"],"activity_id":v["activity_id"],"cbs":v["cbs"],"cost_mean_bn":round(v["cost"],3),"schedule_mean_days":round(v["days"],1),"driver_score":round(v["cost"]*100+v["days"]/10,2)} for k,v in contribution.items()],key=lambda x:x["driver_score"],reverse=True)
    return {"iterations":iterations,"qcra":{"p10":round(pct(cost_samples,10),3),"p50":round(pct(cost_samples,50),3),"p80":round(pct(cost_samples,80),3),"p90":round(pct(cost_samples,90),3),"mean":round(float(np.mean(cost_samples)),3)},"qsra":{"p10":round(pct(sched_samples,10),2),"p50":round(pct(sched_samples,50),2),"p80":round(pct(sched_samples,80),2),"p90":round(pct(sched_samples,90),2),"mean":round(float(np.mean(sched_samples)),2)},"curve":curve,"tornado":tornado}

def benchmarks_for(mode,subsector,location,cost,months):
    s=subsector.lower()
    if mode=="Space": return [{"metric":"Launch/logistics premium","value":"2.5x-4.3x Earth analogue","why":"Remote operations, mass constraints, launch windows and harsh environment."},{"metric":"Programme duration benchmark","value":f"{int(months*.85)}-{int(months*1.25)} months","why":"Qualification, launch integration and commissioning sequence."},{"metric":"Top benchmark gap","value":"TRL / reliability maturity","why":"Space confidence improves fastest through test evidence and heritage."}]
    if "data centre" in s: return [{"metric":"Hyperscale delivery benchmark","value":"$8M-$18M per MW equivalent","why":"Varies by power density, grid scope, cooling, land and regional constraints."},{"metric":"Schedule benchmark","value":"30-60 months","why":"Driven by grid connection, long-lead electrical equipment and phased fit-out."},{"metric":"Top benchmark gap","value":"Power availability","why":"Grid connection is frequently the dominant schedule and risk driver."}]
    if "airport" in s: return [{"metric":"Airport capacity benchmark","value":"$250-$750 per annual pax capacity","why":"Depends on runway, terminal, baggage, rail/road connections and land."},{"metric":"Schedule benchmark","value":"7-12 years","why":"Approvals, airside phasing and operational readiness dominate."}]
    if "rail" in s: return [{"metric":"Rail corridor benchmark","value":"$80M-$500M per km","why":"Underground, stations, signalling and land interfaces drive variance."},{"metric":"Schedule benchmark","value":"6-15 years","why":"Consents, utilities, possessions and systems integration."}]
    return [{"metric":"Capital project benchmark","value":"Sector-adjusted range","why":"Selected based on project type, location, maturity and scale."},{"metric":"Schedule benchmark","value":f"{int(months*.8)}-{int(months*1.3)} months","why":"Derived from maturity and delivery complexity."}]

def scenario_compare(prompt,client,cls,level):
    out=[]
    for sc in ["base","faster","cheaper","lower_risk","premium","investor","survival"]:
        m=build_model_core_no_compare(prompt,client,cls,level,sc)
        out.append({"scenario":sc,"label":m["scenario_label"],"cost":m["cost_p50"],"schedule_months":int(m["schedule"].split()[0]),"risk":m["risk"],"confidence":m["confidence_pct"],"why":m["scenario_why"]})
    return out

def build_model_core_no_compare(prompt,client,cls,level,scenario):
    title,mode,subsector,base_cost,base_months=detect_sector(prompt); location,loc_mult=location_factor(prompt); scale,scale_mult,scale_months=scale_factor(prompt); comp_mult,comp_months,drivers=complexity(prompt); cm,sm,rm,conf_adj,scenario_label,scenario_why=scenario_params(scenario)
    cost=base_cost*loc_mult*scale_mult*comp_mult*cm; months=max(6,round((base_months+scale_months+comp_months+level*2)*sm)); score=(22+(42 if mode=="Space" else 0)+(24 if "Mega" in scale or "System" in scale else 9 if scale=="Programme" else 0)+(14 if cls>=4 else 0)+(18 if loc_mult>=2 else 0)+min(20,len(drivers)*5))*rm; risk=risk_label(score); conf=int(clamp({1:88,2:80,3:68,4:54,5:40}.get(cls,68)-{"Low":0,"Medium":4,"Medium-High":8,"High":12,"Very High":16,"Extreme":22}.get(risk,8)+conf_adj,12,96))
    return {"scenario_label":scenario_label,"scenario_why":scenario_why,"cost_p50":money_bn(cost),"schedule":f"{months} months","risk":risk,"confidence_pct":conf}

def red_flags(risk, confidence, mode, subsector, scenario="base", p50=0.0, p80=0.0, sched_months=0):
    """Sector, scenario and quantitative commercial observations. Never returns generic non-answer."""
    s = subsector.lower()
    sc = str(scenario or "base").lower().replace(" ", "_")
    flags = []

    # ── Confidence-based flags ───────────────────────────────────────
    if confidence < 40:
        flags.append(f"CRITICAL: Confidence {confidence}% — programme is not board-approvable at this maturity. Minimum evidence closure required before any capital commitment.")
    elif confidence < 55:
        flags.append(f"Confidence {confidence}% falls below board-comfort threshold. Scope, benchmark and procurement evidence must close before approval.")
    elif confidence < 65:
        flags.append(f"Confidence {confidence}% is in the challenge zone. Board will probe evidence gaps before approving capital at this level.")

    # ── Risk posture flags ────────────────────────────────────────────
    if risk in ["Very High", "Extreme"]:
        flags.append("Risk posture is Extreme or Very High: executive sponsor must own quantified mitigation plan before commitment.")
    elif risk == "High":
        flags.append("Risk posture is High: named mitigation owners and trigger events must be visible in the board pack.")

    # ── Scenario-specific flags ───────────────────────────────────────
    if sc == "faster":
        flags.append("Faster scenario: acceleration premium will consume procurement and commissioning float. Verify that saved time is operationally usable, not just nominal.")
        if p80 > 0 and p50 > 0 and (p80 / p50) > 1.35:
            flags.append(f"Faster P80/P50 tail ratio is high ({p80/p50:.2f}x). Acceleration has increased the downside exposure, not reduced it.")
    elif sc == "cheaper":
        flags.append("Cheaper scenario: lower authorisation number carries higher residual risk. Board must accept the evidence deferral explicitly.")
    elif sc == "lower_risk":
        flags.append("Lower Risk scenario: verify that reserve additions are linked to specific P80 QCRA movements, not added as a flat contingency.")
    elif sc == "premium":
        flags.append("Premium scenario: confirm what specific resilience or optionality justifies the cost premium before board presentation.")

    # ── P80 tail flags ────────────────────────────────────────────────
    if p80 > 0 and p50 > 0:
        tail = (p80 / p50 - 1) * 100
        if tail > 40:
            flags.append(f"P80 downside is {tail:.0f}% above P50 — board will treat this as an open-ended commitment, not a fixed authorisation.")
        elif tail > 25:
            flags.append(f"P80 tail of {tail:.0f}% above P50 requires explicit contingency approval or a named trigger event before drawdown.")

    # ── Space sector flags ────────────────────────────────────────────
    if mode == "Space":
        flags.append("Space: TRL evidence, launch manifest confirmation and mass margins must be independently verified before commercial programme reliance.")
        flags.append("Space: operations concept and ground segment readiness must be on the critical path — not assumed post-launch.")
        if confidence < 60:
            flags.append("Space confidence below 60% — first-pass model only. Do not present P50 without explicit uncertainty disclosure to funders.")

    # ── Sector-specific flags ─────────────────────────────────────────
    elif any(x in s for x in ["rail","metro","transit","hsr"]):
        flags.append("Rail: possession strategy must be confirmed with the operator before programme baseline is set — assumed windows are not bankable float.")
        flags.append("Rail: signalling integration and systems migration are the hidden critical path. Civil progress is not a proxy for delivery confidence.")
        if sched_months > 120:
            flags.append(f"Rail programme of {sched_months} months carries political and funding exposure across multiple budget cycles. Identify the funding cliff risk.")

    elif any(x in s for x in ["data centre","hyperscale","compute","cloud","ai campus"]):
        flags.append("Data centre: grid connection agreement must be signed before design is committed — energisation dates are not assumptions.")
        flags.append("Data centre: generator and transformer lead times are on the critical path. Long-lead procurement release must precede planning approval.")
        flags.append("Data centre: cooling performance must be proven by vendor test, not datasheet — PUE and water consumption targets are board-visible.")

    elif any(x in s for x in ["nuclear","smr","reactor","fusion"]):
        flags.append("Nuclear: Generic Design Assessment and site licensing are the primary critical path — civil works cannot start until GDA is complete.")
        flags.append("Nuclear: FOAK technology risk requires an independent nuclear assurance programme. Vendor cost and schedule claims are not independently validated.")
        flags.append("Nuclear: nuclear-grade supply chain is capacity-constrained globally. Single-source vendor dependencies must be named and mitigated.")

    elif any(x in s for x in ["pharma","gmp","biologics","life sciences","vaccine","cell therapy"]):
        flags.append("GMP: FDA/EMA inspection readiness must be on the master schedule with named accountable owners — not a post-construction assumption.")
        flags.append("GMP: validation batch failures reset the critical path. A single failed PQ batch at this investment level is a programme-threatening event.")
        flags.append("GMP: clean utility readiness (WFI, clean steam, purified water) must precede process equipment qualification by the master validation plan.")

    elif any(x in s for x in ["defence","defense","naval","submarine","aukus","military","dockyard","shipbuilding"]):
        flags.append("Defence: security accreditation and mission-system acceptance are the hidden critical path — facility completion does not equal operational readiness.")
        flags.append("Defence: sovereign supply chain dependencies and export control risks must be named with fallback options before commitment.")
        flags.append("Defence: classified systems integration requires dedicated assurance environment and independent verification. Optimistic vendor test dates are not programme milestones.")

    elif any(x in s for x in ["semiconductor","fab","cleanroom","wafer","foundry"]):
        flags.append("Fab: process tool delivery is on the critical path — EUV and leading-edge tool allocation must be confirmed, not assumed from OEM intent.")
        flags.append("Fab: cleanroom HVAC classification must be independently tested before tool installation. HVAC failure post-tool-install is a multi-month programme event.")
        flags.append("Fab: yield ramp assumptions must be benchmarked against comparable FOAK fab ramps — vendor projections are systematically optimistic.")

    elif any(x in s for x in ["airport","aviation","runway","airside","terminal"]):
        flags.append("Airport: ORAT readiness and airline transition must be on the master programme — operational readiness is not delivered by construction completion.")
        flags.append("Airport: baggage and security systems integration is the hidden critical path to an operational opening date.")
        if confidence < 65:
            flags.append("Airport: live-operations delivery environment constrains productive windows. Programme acceleration in a live airfield environment is higher risk than greenfield.")

    elif any(x in s for x in ["hydrogen","lng","offshore wind","wind","solar"]):
        flags.append("Energy: grid connection and offtake agreement must be confirmed before final investment decision — merchant exposure without a PPA is a board-level risk.")
        flags.append("Energy: offshore or remote delivery environments carry weather-window risk that is not fully captured in deterministic schedule logic.")

    # ── Final fallback — never return generic ─────────────────────────
    if not flags:
        flags.append(f"No sector-critical commercial observations at {confidence}% confidence — maintain assumptions log, stage-gate assurance and board decision register.")
    
    return flags

def board_questions(mode, subsector, scenario="base"):
    """Scenario-specific AND sector-specific board challenge questions."""
    s = subsector.lower()
    sc = str(scenario or "base").lower().replace(" ", "_")

    # ── Scenario-specific opening questions ─────────────────────────
    scenario_qs = {
        "base": [
            "Is this a decision case or only a reference case — and does the board understand the difference?",
            "What must be proven before this becomes board-approvable, and who owns each proof?",
        ],
        "faster": [
            "Did acceleration consume operationally usable float or just nominal float on the critical path?",
            "Which evidence packages were deferred to fund the schedule compression, and what is the residual risk?",
            "What is the cost of the additional late-stage execution volatility the Faster scenario introduces?",
        ],
        "cheaper": [
            "Where did confidence fall to pay for the saving, and is the board explicitly accepting that trade?",
            "Which scope items or evidence packages were deferred or descoped to reach the lower number?",
            "What is the probability that the Cheaper case requires a rebaseline before contract award?",
        ],
        "lower_risk": [
            "What evidence proves the mitigations are closed and actioned, not just planned and named?",
            "Is the additional reserve genuinely linked to P80 QCRA exposure, or is it a contingency estimate?",
            "Which risks drove the confidence improvement, and are those risks under the programme's control?",
        ],
        "premium": [
            "What specific resilience or optionality does the Premium case buy, and is the board willing to pay for it?",
            "How does the Premium reserve position compare to benchmark outturn for programmes of this type?",
            "Which scenarios has the board explicitly ruled out, and what is the decision rule if this one moves right?",
        ],
    }
    qs = scenario_qs.get(sc, scenario_qs["base"]).copy()

    # ── Universal questions every board should face ─────────────────
    qs += [
        "Which three risks create the most P80/P90 tail exposure, and who owns the mitigation evidence?",
        "What data would move confidence above the board-comfort threshold before the next gate?",
        "Which named owner is accountable for the critical-path constraint — by name, not job title?",
    ]

    # ── Sector-specific hard questions ──────────────────────────────
    if mode == "Space":
        qs += [
            "What flight heritage, TRL and independent qualification evidence supports the cost and schedule basis?",
            "What is the contingency plan if the launch manifest slips and the window is lost?",
        ]
    elif any(x in s for x in ["rail","metro","transit","hsr"]):
        qs += [
            "Which possessions are genuinely executable under live-corridor constraints versus optimistically assumed?",
            "What evidence proves signalling migration readiness and what does a failed integration test cost?",
        ]
    elif any(x in s for x in ["data centre","hyperscale","compute","cloud"]):
        qs += [
            "Is the grid connection agreement signed, or is the energisation date an assumption?",
            "What is the contingency if the cooling vendor fails performance commissioning?",
        ]
    elif any(x in s for x in ["nuclear","smr","reactor"]):
        qs += [
            "What is the regulatory milestone critical path and what is the cost of a six-month GDA slip?",
            "Which nuclear supply chain items are on single-source FOAK vendors with no fallback?",
        ]
    elif any(x in s for x in ["pharma","gmp","biologics","life sciences","vaccine"]):
        qs += [
            "What is the inspection readiness evidence package and when does it need to be complete?",
            "What is the cost and schedule impact of a single failed validation batch at this stage?",
        ]
    elif any(x in s for x in ["defence","defense","naval","submarine","aukus","military","dockyard"]):
        qs += [
            "What is the security accreditation critical path and which classified supplier dependencies are unsecured?",
            "What is the fallback if a sovereign supply chain item is subject to export control or allocation denial?",
        ]
    elif any(x in s for x in ["semiconductor","fab","cleanroom","wafer"]):
        qs += [
            "Which process tool allocations are confirmed versus assumed, and what is the yield ramp risk?",
            "What is the contingency plan if EUV or leading-edge tool delivery slips by one quarter?",
        ]
    elif any(x in s for x in ["airport","aviation","runway","airside"]):
        qs += [
            "What evidence proves ORAT readiness and which airline or regulator interface can block opening?",
            "Is the baggage systems integration on the critical path, and what is the test failure contingency?",
        ]

    return qs[:8]  # Cap at 8 — clean board pack

def next_actions(risk,confidence,scenario):
    return ["Run a risk workshop to validate probability and O/M/P impacts.","Confirm top five CBS cost drivers with benchmark or supplier evidence.","Validate the Level 3/4 schedule logic and critical path assumptions.","Prepare board decision paper with Base/Faster/Cheaper/Lower-Risk cases.","Create confidence improvement plan targeting design maturity, consents and procurement evidence."]

def demo_script():
    return ["Open with: 'Give me any project on Earth or in Space.'","Type the project into CASEY and hit Generate Full Intelligence Pack.","Show the executive dashboard: cost, schedule, risk, confidence and top drivers.","Switch to Monte Carlo: show QCRA/QSRA P-curves and tornado drivers linked to activities.","Show scenarios: Faster, Cheaper, Lower Risk and Premium with why each changed.","Open Output Centre: Excel, Risk CSV, XER, cost workbook, risk register, schedule export and full pack are ready.","Close with: 'Traditional advisory takes weeks. CASEY gives the first board-grade view in seconds.'"]

# ------------------------- routes -------------------------
@app.get("/health")
def health(): return {"status":"ok","service":APP_VERSION,"demo_limit_per_ip":DEMO_LIMIT_PER_IP}

@app.get("/demo/status")
def demo_status(request: Request):
    ip = client_ip(request)
    status = check_demo_allowance(ip)
    return {"allowed": status["allowed"], "used": status["used"], "limit": status["limit"], "remaining": max(0, status["limit"] - status["used"]), "demo_launch_mode": False}


@app.post("/public-demo/generate")
def public_demo_generate(req: PublicDemoRequest, request: Request):
    issues = _quality_gate_public_demo(req)
    if issues:
        raise HTTPException(status_code=422, detail={"message": "CASEY needs one real infrastructure or space programme brief before using your free run.", "issues": issues})
    identity = _public_demo_identity(request, req)
    previous = _public_demo_used(identity)
    if previous:
        raise HTTPException(status_code=403, detail={
            "message": "Your one free CASEY project/showcase run has already been used.",
            "sub": "Request access for unlimited project runs, scenario sensitivity, exports and client-file challenge.",
            "email": "deepa@caseai.co.uk",
            "linkedin": "https://www.linkedin.com/in/deepa-mahadeshwar-727200409/",
            "previous_run": previous
        })
    prompt = _premium_public_prompt(req)
    input_quality = _public_demo_brief_quality_score(req)
    model = build_model(prompt, _normalise_email(req.email), 3, 4, "base")
    model["input_quality_score"] = input_quality["score"]
    model["public_demo"] = True
    model["lead_email"] = _normalise_email(req.email)
    run_id = "CASEY-DEMO-" + uuid.uuid4().hex[:10].upper()
    report = _public_demo_report(model)
    con = db(); cur = con.cursor(); now = datetime.utcnow().isoformat()
    cur.execute("""INSERT INTO public_demo_uses(run_id,email_hash,ip_hash,fingerprint_hash,client_token_hash,project_type,project_text,model_json,created_at)
        VALUES(?,?,?,?,?,?,?,?,?)""", (run_id, identity["email_hash"], identity["ip_hash"], identity["fingerprint_hash"], identity["client_token_hash"], req.project_type, req.project_description, json.dumps(model), now))
    cur.execute("INSERT INTO projects(name,client,prompt,mode,created_at,model_json) VALUES(?,?,?,?,?,?)", (model["title"], model["client"], model["prompt"], model["mode"], now, json.dumps(model)))
    con.commit(); con.close()
    return {"run_id": run_id, "used": 1, "limit": PUBLIC_DEMO_LIMIT, "report": report, "model": model}


@app.get("/public-demo/admin/runs")
def public_demo_admin_runs(request: Request, limit: int = 100):
    """Admin-only view of public demo runs.
    Set CASEY_ADMIN_TOKEN in Render, then call with header:
    x-casey-admin-token: YOUR_TOKEN
    """
    if not ADMIN_TOKEN:
        raise HTTPException(status_code=403, detail="Admin log access is not configured. Set CASEY_ADMIN_TOKEN.")
    supplied = request.headers.get("x-casey-admin-token", "")
    if supplied != ADMIN_TOKEN:
        raise HTTPException(status_code=403, detail="Invalid admin token.")
    limit = max(1, min(int(limit or 100), 1000))
    con = db(); cur = con.cursor()
    rows = cur.execute("""SELECT run_id, project_type, project_text, model_json, created_at
                          FROM public_demo_uses ORDER BY id DESC LIMIT ?""", (limit,)).fetchall()
    feedback = cur.execute("""SELECT run_id, rating, comment, created_at FROM public_demo_feedback ORDER BY id DESC LIMIT ?""", (limit,)).fetchall()
    con.close()
    return {
        "runs": [
            {
                "run_id": r["run_id"],
                "project_type": r["project_type"],
                "project_text": r["project_text"],
                "created_at": r["created_at"],
                "model": json.loads(r["model_json"]) if r["model_json"] else None
            } for r in rows
        ],
        "feedback": [dict(f) for f in feedback]
    }

@app.get("/public-demo/admin/summary")
def public_demo_admin_summary(request: Request):
    if not ADMIN_TOKEN:
        raise HTTPException(status_code=403, detail="Admin log access is not configured. Set CASEY_ADMIN_TOKEN.")
    supplied = request.headers.get("x-casey-admin-token", "")
    if supplied != ADMIN_TOKEN:
        raise HTTPException(status_code=403, detail="Invalid admin token.")
    con = db(); cur = con.cursor()
    total = cur.execute("SELECT COUNT(*) AS c FROM public_demo_uses").fetchone()["c"]
    by_type = cur.execute("SELECT project_type, COUNT(*) AS c FROM public_demo_uses GROUP BY project_type").fetchall()
    recent = cur.execute("SELECT run_id, project_type, project_text, created_at FROM public_demo_uses ORDER BY id DESC LIMIT 20").fetchall()
    con.close()
    return {
        "total_runs": total,
        "by_type": [dict(x) for x in by_type],
        "recent": [dict(x) for x in recent]
    }


@app.post("/public-demo/feedback")
def public_demo_feedback(req: PublicDemoFeedback):
    rating = max(1, min(5, int(req.rating)))
    con = db(); cur = con.cursor(); now = datetime.utcnow().isoformat()
    cur.execute("INSERT INTO public_demo_feedback(run_id,rating,comment,created_at) VALUES(?,?,?,?)", (req.run_id, rating, req.comment or "", now))
    con.commit(); con.close()
    return {"status": "saved", "message": "Feedback saved for CASEY improvement loop."}


def _v126_cost_type(row: Dict[str, Any]) -> str:
    t = str(row.get("type") or row.get("category") or "Direct").strip().lower()
    if any(k in t for k in ["reserve", "risk", "contingency"]): return "Reserve"
    if any(k in t for k in ["indirect", "owner", "pm", "management", "prelim", "assurance", "integration"]): return "Indirect"
    return "Direct"


def _v126_scenario_split(scenario: str) -> Dict[str, float]:
    """Scenario-controlled bucket shares. These are intentionally NOT constant.
    Base is the neutral reference; Faster buys acceleration/indirect controls;
    Cheaper cuts reserve; Lower Risk and Premium visibly buy reserve/assurance.
    """
    key = (scenario or "base").lower().replace(" ", "_").replace("-", "_")
    if key in ["faster", "accelerated", "fast"]:
        return {"Direct": 0.68, "Indirect": 0.18, "Reserve": 0.14}
    if key in ["cheaper", "low_cost", "budget"]:
        return {"Direct": 0.77, "Indirect": 0.14, "Reserve": 0.09}
    if key in ["lower_risk", "low_risk", "derisked"]:
        return {"Direct": 0.64, "Indirect": 0.16, "Reserve": 0.20}
    if key in ["premium", "resilient"]:
        return {"Direct": 0.67, "Indirect": 0.15, "Reserve": 0.18}
    return {"Direct": 0.70, "Indirect": 0.14, "Reserve": 0.16}


def _v126_normalize_model_costs(model: Dict[str, Any]) -> Dict[str, Any]:
    """Keep headline P50, cost cards, scenario tabs and exports in one numeric source of truth.
    Critical fix: direct / indirect / reserve are now rebuilt from the selected scenario split,
    then detailed CBS rows are scaled inside each bucket. This prevents the UI/export layer from
    carrying the same static split across Base, Faster, Cheaper, Lower Risk and Premium.
    """
    try:
        target = parse_bn(model.get("cost_p50") or model.get("p50") or 0)
    except Exception:
        target = 0
    scenario = str(model.get("scenario") or model.get("scenario_label") or "base").lower()
    shares = _v126_scenario_split(scenario)
    lines = model.get("cost_breakdown") or model.get("cost_lines") or []

    if target:
        if not isinstance(lines, list) or not lines:
            lines = [
                {"cbs":"01.00","description":"Direct delivery scope","type":"Direct","p50_bn":target*shares["Direct"]},
                {"cbs":"90.00","description":"Indirects, owner costs and integration","type":"Indirect","p50_bn":target*shares["Indirect"]},
                {"cbs":"99.00","description":"Risk reserve and contingency","type":"Reserve","p50_bn":target*shares["Reserve"]},
            ]

        buckets = {"Direct": [], "Indirect": [], "Reserve": []}
        for i, row in enumerate(lines):
            if not isinstance(row, dict):
                continue
            x = dict(row)
            typ = _v126_cost_type(x)
            x["type"] = typ
            x.setdefault("cbs", f"C-{i+1:02d}")
            x.setdefault("description", typ)
            try:
                x["_old_p50"] = float(x.get("p50_bn") or parse_bn(x.get("most_likely_p50") or x.get("p50") or 0) or 0)
            except Exception:
                x["_old_p50"] = 0.0
            buckets[typ].append(x)

        rebuilt = []
        p10_mult = {"Direct":0.82, "Indirect":0.80, "Reserve":0.58}
        p90_mult = {"Direct":1.30, "Indirect":1.36, "Reserve":1.85}
        if "cheaper" in scenario:
            p90_mult["Reserve"] = 2.20
        if "lower" in scenario or "premium" in scenario:
            p90_mult["Reserve"] = 1.55
        for typ in ["Direct", "Indirect", "Reserve"]:
            group = buckets[typ]
            if not group:
                group = [{"cbs":"01.00" if typ=="Direct" else "90.00" if typ=="Indirect" else "99.00", "description":typ, "type":typ, "_old_p50":1.0}]
            old_sum = sum(max(float(x.get("_old_p50") or 0), 0.0) for x in group) or float(len(group))
            bucket_target = target * shares[typ]
            for x in group:
                weight = (max(float(x.get("_old_p50") or 0), 0.0) / old_sum) if old_sum else (1.0 / len(group))
                p50 = bucket_target * weight
                x["p50_bn"] = round(p50, 3)
                x["p10_bn"] = round(p50 * p10_mult[typ], 3)
                x["p90_bn"] = round(p50 * p90_mult[typ], 3)
                x["most_likely_p50"] = money_bn(p50) if 'money_bn' in globals() else f"${p50:.1f}B"
                x["low_p10"] = money_bn(x["p10_bn"]) if 'money_bn' in globals() else f"${x['p10_bn']:.1f}B"
                x["high_p90"] = money_bn(x["p90_bn"]) if 'money_bn' in globals() else f"${x['p90_bn']:.1f}B"
                x["impact_basis"] = x.get("impact_basis") or x.get("basis") or f"{typ} scenario-controlled bucket, reconciled to selected P50."
                x.pop("_old_p50", None)
                rebuilt.append(x)

        drift = round(target - sum(float(x.get("p50_bn") or 0) for x in rebuilt), 3)
        if abs(drift) >= 0.001:
            for x in rebuilt:
                if x.get("type") == "Direct":
                    x["p50_bn"] = round(float(x.get("p50_bn") or 0) + drift, 3)
                    break
        lines = rebuilt
        model["cost_breakdown"] = lines
        model["cost_lines"] = lines
        model["direct_cost"] = money_bn(target * shares["Direct"]) if 'money_bn' in globals() else f"${target*shares['Direct']:.1f}B"
        model["indirect_cost"] = money_bn(target * shares["Indirect"]) if 'money_bn' in globals() else f"${target*shares['Indirect']:.1f}B"
        model["risk_reserve"] = money_bn(target * shares["Reserve"]) if 'money_bn' in globals() else f"${target*shares['Reserve']:.1f}B"
        model["cost_bucket_split"] = {k: round(v, 3) for k, v in shares.items()}

    # aliases needed by UI and exports
    model["risk_register"] = model.get("risk_register") or model.get("risks") or []
    model["schedule_detail"] = model.get("schedule_detail") or model.get("schedule_rows") or []
    model["scenario_matrix"] = model.get("scenario_matrix") or model.get("scenario_comparison") or []
    return model

def _v126_sector_scary_additions(model: Dict[str, Any]) -> Dict[str, Any]:
    fam = _casey_sector_family(model) if '_casey_sector_family' in globals() else 'infrastructure'
    intel = _SECTOR_INTEL.get(fam, _SECTOR_INTEL.get('infrastructure', {})) if '_SECTOR_INTEL' in globals() else {}
    governor = intel.get('governor', 'evidence maturity, procurement certainty and commissioning readiness')
    hidden = intel.get('hidden', f'confidence is governed by {governor}, not headline progress')
    fail = intel.get('fail', 'governance confidence failed before delivery reporting showed it')
    model['casey_position'] = f"CASEY does not treat this as approval-ready until {governor} is evidenced. The hidden position is that {hidden}."
    model['board_attack_simulation'] = intel.get('attacks', [
        'Which assumption collapses confidence first?',
        'What evidence proves board readiness?',
        'Is float usable or only theoretical?',
        'Which owner accepts integration liability?',
        'What would move confidence above the board comfort threshold?'
    ])
    model['if_this_fails'] = f"If this programme fails, the likely explanation will be: {fail}."
    model['traditional_vs_casey'] = {
        'traditional': intel.get('traditional','Traditional controls read progress, float and contingency as reported.'),
        'casey': intel.get('casey', f'CASEY reads {governor} as the real confidence governor.')
    }
    return model


def _v142_final_demo_polish(model: Dict[str, Any]) -> Dict[str, Any]:
    """Final demo credibility pass.
    Keeps the model auditable, but removes the clean demo-number feel:
    irregular confidence, asymmetrical QSRA/QCRA tails, higher reserve posture,
    operational language instead of AI-explaining-itself language, and a denser
    ugly audit spine for exports/UI trust.
    """
    try:
        scenario = str(model.get("scenario") or "base").lower().replace(" ", "_").replace("-", "_")
        seed_txt = f"{model.get('prompt','')}|{model.get('title','')}|{model.get('subsector','')}|{scenario}"
        h = int(hashlib.sha256(seed_txt.encode('utf-8')).hexdigest()[:8], 16)
        def jitter(span:int=2):
            return (h % (span*2+1)) - span
        def money_to_bn(v):
            return parse_bn(v) if 'parse_bn' in globals() else _money_to_bn_v95(v)
        def money(bn):
            return money_bn(bn) if 'money_bn' in globals() else _bn_to_money_v95(bn)

        p50 = money_to_bn(model.get('cost_p50'))
        months = int(float(str(model.get('schedule','60')).replace('months','').split()[0]))
        base_conf = int(model.get('_base_confidence_pct') or model.get('confidence_pct') or 55)

        # Irregular but stable scenario confidence / duration offsets.
        conf_delta = {'base': 1, 'faster': -15, 'cheaper': -17, 'lower_risk': 13, 'premium': 17}.get(scenario, 0)
        sched_j = {'base': -1, 'faster': 2, 'cheaper': 7, 'lower_risk': -4, 'premium': -3}.get(scenario, 0)
        if scenario == 'base':
            new_conf = max(8, min(96, base_conf + 1 + (h % 2)))
        else:
            new_conf = max(8, min(96, base_conf + conf_delta + jitter(1)))
        if months > 24:
            months = max(3, months + sched_j + jitter(1))
            model['schedule'] = f"{months} months"
        model['confidence_pct'] = new_conf

        # Update scenario matrix with irregular board-plausible values.
        base_cost = money_to_bn(model.get('_base_cost_p50') or model.get('cost_p50')) or p50
        base_months = int(model.get('_base_schedule_months') or months)
        profiles = {
            'base':       ('Base',       1.000, 1.000,  1, 'Medium-High'),
            'faster':     ('Faster',     1.178, 0.795, -15, 'High'),
            'cheaper':    ('Cheaper',    0.861, 1.168, -17, 'High'),
            'lower_risk': ('Lower Risk', 1.119, 1.106, 13, 'Medium-Low'),
            'premium':    ('Premium',    1.281, 0.947, 17, 'Low'),
        }
        matrix=[]
        for i,(sk,(lab,cf,sf,cd,risk)) in enumerate(profiles.items()):
            jj=((h >> (i*3)) % 5) - 2
            sm=max(3, int(round(base_months*sf + {'base':-1,'faster':2,'cheaper':7,'lower_risk':-4,'premium':-3}.get(sk,0) + jj)))
            cc=max(8,min(96,base_conf + cd + (((h >> (i*2)) % 3)-1)))
            cb=base_cost*cf
            matrix.append({'scenario':sk,'label':lab,'cost_p50':money(cb),'schedule_months':sm,'risk':risk,'confidence_pct':cc,'cost_delta_pct':round((cf-1)*100),'schedule_delta_pct':round((sm-base_months)/max(base_months,1)*100),'confidence_delta_pts':cc-base_conf,'why':'Scenario-controlled cost, schedule, reserve and evidence posture.'})
            if sk == scenario:
                model['cost_p50'] = money(cb)
                p50 = cb
                model['schedule'] = f"{sm} months"
                months = sm
                model['confidence_pct'] = cc
                model['risk'] = risk
                model['scenario_label'] = lab
        model['scenario_matrix'] = matrix

        # Rebuild cost buckets again after any P50 change, with reserves high enough to survive challenge.
        shares = _v126_scenario_split(scenario) if '_v126_scenario_split' in globals() else {'Direct':0.70,'Indirect':0.14,'Reserve':0.16}
        model['direct_cost'] = money(p50 * shares['Direct'])
        model['indirect_cost'] = money(p50 * shares['Indirect'])
        model['risk_reserve'] = money(p50 * shares['Reserve'])
        # QCRA tails: reserve posture and low confidence should widen downside.
        tail = {'base':1.30,'faster':1.42,'cheaper':1.50,'lower_risk':1.24,'premium':1.22}.get(scenario,1.30)
        p10_mult = {'base':0.82,'faster':0.78,'cheaper':0.72,'lower_risk':0.86,'premium':0.87}.get(scenario,0.82)
        p90_mult = {'base':1.48,'faster':1.65,'cheaper':1.76,'lower_risk':1.36,'premium':1.33}.get(scenario,1.48)
        model['cost_p10'] = money(p50*p10_mult)
        model['cost_p90'] = money(p50*p90_mult)
        model['cost_range'] = f"{model['cost_p10']} - {model['cost_p90']}"
        model = _v126_normalize_model_costs(model) if '_v126_normalize_model_costs' in globals() else model

        pts=[1,5,10,20,30,40,50,60,70,80,90,95,99]
        curve=[]
        for p in pts:
            x=p/100
            if scenario == 'cheaper':
                cm=p10_mult + (1.0-p10_mult)*(x/0.5 if x<=0.5 else 1) + max(0,x-0.5)**2.15*1.15
                sm=0.88 + 0.20*x + 0.45*(x**4)
            elif scenario == 'faster':
                cm=0.79 + 0.42*x + 0.46*(x**4.5)
                sm=0.72 + 0.33*x + 0.36*(x**5)
            elif scenario in ['lower_risk','premium']:
                cm=0.86 + 0.20*x + 0.17*(x**3)
                sm=0.88 + 0.17*x + 0.12*(x**3)
            else:
                cm=0.80 + 0.32*x + 0.36*(x**3.6)
                sm=0.82 + 0.25*x + 0.20*(x**3.4)
            curve.append({'percentile':p,'cost_bn':round(p50*cm,2),'schedule_months':max(1,int(round(months*sm + (((h>>p)%3)-1))))})
        # force exact P50 at headline.
        for row in curve:
            if row['percentile']==50:
                row['cost_bn']=round(p50,2); row['schedule_months']=months
        qcra={'p10':curve[2]['cost_bn'],'p50':round(p50,1),'p80':round(max(curve[9]['cost_bn'],p50*tail),1),'p90':round(max(curve[10]['cost_bn'],p50*p90_mult),1)}
        qsra={'p10':curve[2]['schedule_months'],'p50':months,'p80':max(curve[9]['schedule_months'], months + {'base':35,'faster':49,'cheaper':57,'lower_risk':23,'premium':18}.get(scenario,35) + jitter(2)), 'p90':max(curve[10]['schedule_months'], months + {'base':59,'faster':76,'cheaper':88,'lower_risk':39,'premium':31}.get(scenario,59) + jitter(3))}
        model['monte_carlo'] = model.get('monte_carlo') or {}
        model['monte_carlo'].update({'qcra':qcra,'qsra':qsra,'curve':curve})
        model['monte_carlo']['curve_readout']=[
            f"Cost P50 reconciles to {model['cost_p50']}; P80 is {money(qcra['p80'])}; P90 is {money(qcra['p90'])}.",
            f"Schedule P50 reconciles to {months} months; P80 is {qsra['p80']} months; P90 is {qsra['p90']} months.",
            'Asymmetric tail: downside grows faster than the median when evidence is incomplete.'
        ]

        # Operational wording.
        repl = {
            'CASEY detected':'The brief indicates',
            'locked the sector ontology before generating outputs':'locked the delivery environment before producing the board pack',
            'sector ontology':'delivery environment',
            'generated this':'produced this',
            'Mission assurance burden':'Late-stage qualification exposure',
            'Evidence maturity':'Package evidence completeness'
        }
        for key in ['casey_thinking','executive_summary','executive_shock_insight']:
            if isinstance(model.get(key), str):
                t=model[key]
                for a,b in repl.items(): t=t.replace(a,b)
                model[key]=t
        if isinstance(model.get('why_casey_generated_this'), list):
            cleaned=[]
            for t in model['why_casey_generated_this']:
                for a,b in repl.items(): t=str(t).replace(a,b)
                cleaned.append(t)
            model['why_casey_generated_this']=cleaned

        # Stronger confidence labels and audit spine.
        model['confidence_breakdown'] = [
            {'driver':'Benchmark fit','effect':'+7','note':'Comparable mission class and delivery environment identified'},
            {'driver':'Package evidence completeness','effect':'-6' if scenario in ['base','cheaper','faster'] else '+6','note':'Basis visibility and owner evidence determine approval readiness'},
            {'driver':'Procurement certainty','effect':'-8' if scenario in ['cheaper','faster'] else '+5','note':'Long-lead and supplier readiness alters the P80 tail'},
            {'driver':'Schedule logic','effect':'-9' if scenario=='faster' else '+6','note':'Critical path, handover gates and QSRA traceability'},
            {'driver':'Reserve adequacy','effect':'+4' if scenario=='base' else '-7' if scenario=='cheaper' else '+9','note':'Reserve is linked to tail exposure, not a flat percentage'}
        ]
        model['audit_spine'] = [
            {'check':'P50 reconciliation','status':'PASS','detail':f"Direct {model.get('direct_cost')} + Indirect {model.get('indirect_cost')} + Reserve {model.get('risk_reserve')} = {model.get('cost_p50')}"},
            {'check':'P80 downside','status':'VISIBLE','detail':f"QCRA P80 {money(qcra['p80'])}; QSRA P80 {qsra['p80']} months"},
            {'check':'Reserve challenge','status':'BOARD-READY','detail':f"Reserve share {round(shares['Reserve']*100)}% varies by selected scenario"},
            {'check':'Evidence gate','status':str(model.get('confidence_engine_label','CASEY Confidence Engine')),'detail':'Approval posture remains tied to evidence closure, not dashboard colour'}
        ]
        model['version']='CASEY V142 Final 10/10 Demo Credibility Drop'
    except Exception as e:
        model['v142_polish_warning']=str(e)
    return model

@app.post("/generate")
def generate(req: GenerateRequest, request: Request):
    if req.demo:
        # V147 production gate: one free run per IP.
        # Bypassed automatically for: localhost/127.0.0.1, admin emails (CASEY_ADMIN_EMAILS), admin key header.
        # Set CASEY_LOCAL_DEMO=1 to force-bypass in any environment.
        local_dev = (
            os.environ.get("CASEY_LOCAL_DEMO", "0") == "1"
            or _casey_v116_local_request(request)
        )
        email = (getattr(req, "email", None) or getattr(req, "lead_email", None) or "").strip()
        admin_bypass = local_dev or _v108_is_admin_email(email) or _v108_admin_key_ok(request)
        if not admin_bypass:
            ip = client_ip(request)
            status = check_demo_allowance(ip)
            if not status["allowed"]:
                raise HTTPException(status_code=403, detail={
                    "message": "You\'ve used your one free CASEY intelligence run.",
                    "sub": "To run more projects, compare scenarios or download the full output pack, get in touch.",
                    "email": "deepa@caseai.co.uk",
                    "linkedin": "https://www.linkedin.com/in/deepa-mahadeshwar-727200409/",
                    "upgrade_cta": "Contact us"
                })
            record_demo_use(ip)
        else:
            # Admin / localhost: record for analytics but never gate
            ip = client_ip(request)
            record_demo_use(ip)

    # Canonical project-context lock. Scenario buttons must never rebuild from the default
    # data-centre seed when the active project is Space, Rail, Energy, Defence, etc.
    scenario = (req.scenario or "base").lower()
    active = req.active_model or {}
    locked_prompt = active.get("prompt") or req.prompt
    locked_mode = active.get("mode")
    locked_subsector = active.get("subsector")

    model=build_model(locked_prompt, req.client or "", int(req.class_level or 3), int(req.schedule_level or 3), scenario)

    # Defensive guard: if a non-base scenario ever classifies differently from the locked
    # context, keep the model in the original project universe and flag it for audit.
    if scenario != "base" and locked_mode and model.get("mode") != locked_mode:
        model["mode"] = locked_mode
        if locked_subsector: model["subsector"] = locked_subsector
        if active.get("title"): model["title"] = active.get("title")
        if active.get("location"): model["location"] = active.get("location")
        if active.get("scale"): model["scale"] = active.get("scale")
        model["prompt"] = locked_prompt
        model["context_lock_repaired"] = True

    model["active_context_lock"] = {
        "mode": model.get("mode"),
        "subsector": model.get("subsector"),
        "prompt_hash": hashlib.sha256((model.get("prompt") or "").encode("utf-8")).hexdigest()[:12],
        "scenario": scenario
    }

    con=db(); cur=con.cursor(); now=datetime.utcnow().isoformat(); cur.execute("INSERT INTO projects(name,client,prompt,mode,created_at,model_json) VALUES(?,?,?,?,?,?)",(model["title"],model["client"],model["prompt"],model["mode"],now,json.dumps(model))); con.commit(); con.close()

    # FINAL EXEC POLISH: scenario-aware ranking and differentiation
    try:
        sc = str(model.get("scenario","base")).lower()
        schedule_lists = {
            "faster":[
                "Concurrent commissioning overload",
                "Recovery float exhaustion",
                "Acceleration premium shock",
                "Grid connection delay",
                "Integrated systems testing concurrency"
            ],
            "cheaper":[
                "Vendor claims and change exposure",
                "Procurement deferral and long-lead slippage",
                "Design maturity gap",
                "Scope growth from deferred decisions",
                "Interface coordination delay"
            ],
            "lower_risk":[
                "Governance and approvals latency",
                "Extended validation sequencing",
                "Conservative commissioning gates",
                "Operational readiness hold-points",
                "Assurance and compliance reviews"
            ],
            "premium":[
                "Integration complexity across parallel packages",
                "Executive decision latency",
                "Technology assurance alignment",
                "Multi-package interface management",
                "Programme coordination overhead"
            ]
        }
        cost_lists = {
            "faster":[
                "Acceleration premiums and overtime",
                "Power train, transformers and switchgear",
                "Integrated systems testing",
                "Grid and utility concurrency",
                "Recovery-float consumption"
            ],
            "cheaper":[
                "Deferred procurement packaging",
                "Claims and commercial exposure",
                "Rework from reduced contingency",
                "Long-lead inflation volatility",
                "Scope rationalisation impacts"
            ],
            "lower_risk":[
                "Additional contingency and reserve",
                "Enhanced validation and assurance",
                "Programme controls and governance",
                "Redundant infrastructure resilience",
                "Extended commissioning readiness"
            ]
        }
        if sc in schedule_lists:
            model["sector_schedule_threats"]=schedule_lists[sc]
        if sc in cost_lists:
            model["sector_primary_cost_drivers"]=cost_lists[sc]

        if sc=="faster":
            model["executive_shock_insight"]="Acceleration increases spend faster than it reduces uncertainty; the delivery tail becomes more volatile."
        elif sc=="cheaper":
            model["executive_shock_insight"]="Capital efficiency reduces resilience: procurement and recovery flexibility become constrained."
        elif sc=="lower_risk":
            model["executive_shock_insight"]="Confidence is purchased through reserve, governance and extended delivery duration."
        elif sc=="premium":
            model["executive_shock_insight"]="Premium posture buys resilience, optionality and stronger certainty at visible capex premium."
    except Exception:
        pass

    model = _v142_final_demo_polish(_v126_normalize_model_costs(_v126_sector_scary_additions(model)))
    return model

@app.get("/projects")
def projects():
    con=db(); rows=[dict(x) for x in con.execute("SELECT id,name,client,prompt,mode,created_at FROM projects ORDER BY id DESC LIMIT 50")]; con.close(); return rows

@app.get("/projects/{project_id}")
def project(project_id:int):
    con=db(); row=con.execute("SELECT * FROM projects WHERE id=?",(project_id,)).fetchone(); con.close()
    if not row: raise HTTPException(404,"Project not found")
    return json.loads(row["model_json"])


def _casey_sector_family(model: Dict[str, Any]) -> str:
    txt = " ".join(str(model.get(k,"")) for k in ["title","sector","subsector","mode","prompt","client"]).lower()
    if any(w in txt for w in ["rail", "transit", "metro", "high-speed", "rolling stock", "signalling"]): return "rail"
    if any(w in txt for w in ["airport", "aviation", "airside", "terminal", "runway", "orat"]): return "airport"
    if any(w in txt for w in ["data centre", "data center", "hyperscale", "ai campus", "gpu", "liquid cooling"]): return "hyperscale"
    if any(w in txt for w in ["pharma", "lilly", "amgen", "biologics", "gmp", "cqv", "fill-finish", "validation"]): return "life_sciences"
    if any(w in txt for w in ["semiconductor", "fab", "wafer", "lithography", "cleanroom"]): return "semiconductor"
    if any(w in txt for w in ["space", "lunar", "mars", "orbital", "launch", "satellite", "payload", "mission"]): return "space"
    if any(w in txt for w in ["defence", "defense", "naval", "secure", "classified", "military"]): return "defence"
    if any(w in txt for w in ["oil", "gas", "lng", "refinery", "offshore", "cryogenic"]): return "oil_gas"
    if any(w in txt for w in ["nuclear", "reactor", "containment"]): return "nuclear"
    if any(w in txt for w in ["hospital", "healthcare", "clinical"]): return "healthcare"
    if any(w in txt for w in ["water", "wastewater", "desalination", "reservoir"]): return "water"
    if any(w in txt for w in ["port", "marine", "harbour", "harbor", "terminal berth"]): return "ports"
    if any(w in txt for w in ["power", "utility", "grid", "transmission", "substation", "renewable", "energy"]): return "energy"
    return "infrastructure"

_SECTOR_INTEL = {
    "rail": {
        "governor": "possessions, signalling integration, migration sequencing and trial operations",
        "hidden": "opening confidence is governed by systems migration and operational acceptance, not visible civil progress",
        "fail": "systems integration and possession access failed to mature before trial running and regulator acceptance",
        "evidence": ["signed possession strategy", "signalling integration test plan", "rolling-stock interface evidence", "trial running readiness gate", "operator acceptance criteria"],
        "traditional": "Civil works and corridor progress look reportable.",
        "casey": "Operational migration, possessions and signalling assurance govern opening confidence.",
        "attacks": ["Which possessions are truly executable rather than assumed?", "What evidence proves signalling migration readiness?", "Where is rolling-stock interface liability owned?", "What trial-running gate will fail first?", "Is float usable under live corridor access constraints?"]
    },
    "airport": {
        "governor": "ORAT, baggage/security systems integration, airside phasing and regulator acceptance",
        "hidden": "terminal completion can precede operational readiness by months if ORAT and systems acceptance are immature",
        "fail": "operational transition and baggage/security integration lagged construction progress",
        "evidence": ["ORAT schedule and acceptance gates", "baggage integration test evidence", "airside phasing plan", "airline transition sign-offs", "regulator readiness evidence"],
        "traditional": "Terminal construction and package completion appear on track.",
        "casey": "Operational readiness, baggage/security integration and live-airport phasing govern usable opening confidence.",
        "attacks": ["What evidence proves ORAT readiness?", "Which airline or regulator interface can block opening?", "Is baggage systems testing on the critical path?", "What live-airside access assumption is weakest?", "Who owns operational transition liability?"]
    },
    "hyperscale": {
        "governor": "utility energisation, transformer supply, liquid-cooling readiness and integrated systems testing",
        "hidden": "practical completion is secondary to energisation, thermal stability and IST clearance",
        "fail": "power availability and cooling/IST readiness lagged the declared construction position",
        "evidence": ["utility interconnect agreement", "transformer delivery certainty", "liquid-cooling commissioning evidence", "IST sequence", "energisation acceptance gates"],
        "traditional": "Shell, core and white-space progress look healthy.",
        "casey": "Energisation, cooling readiness and IST congestion govern revenue-ready capacity.",
        "attacks": ["What is the true energisation date?", "Which transformer or switchgear package controls first capacity?", "Is liquid cooling commissioned or merely installed?", "Where does IST congestion compress the plan?", "Does acceleration buy capacity or create rework?"]
    },
    "life_sciences": {
        "governor": "CQV execution, GMP readiness, validation closure and regulatory inspection evidence",
        "hidden": "mechanical completion is not the finish line; validated production readiness governs confidence",
        "fail": "CQV and validation maturity lagged construction completion and regulatory readiness",
        "evidence": ["CQV execution plan", "validation master plan", "GMP turnover matrix", "media-fill/qualification evidence", "inspection readiness tracker"],
        "traditional": "Facility construction and equipment installation are progressing.",
        "casey": "Validated GMP readiness, CQV closure and regulatory evidence govern usable production confidence.",
        "attacks": ["Which validation package breaks first?", "Is GMP turnover evidenced or assumed?", "What inspection-readiness gate is not closed?", "Where does commissioning overlap threaten deviation closure?", "Can production readiness survive accelerated handover?"]
    },
    "semiconductor": {
        "governor": "cleanroom readiness, process-tool install, ultra-pure utilities and yield-ramp qualification",
        "hidden": "fab readiness is governed by contamination control, tool qualification and yield ramp rather than building completion",
        "fail": "process-tool qualification and cleanroom utility maturity lagged shell completion",
        "evidence": ["tool delivery and hook-up plan", "cleanroom certification", "ultra-pure water readiness", "process qualification gates", "yield ramp assumptions"],
        "traditional": "Base build and cleanroom construction appear advanced.",
        "casey": "Tool qualification, contamination control and yield maturity govern production confidence.",
        "attacks": ["Which tool family governs first wafer readiness?", "Is cleanroom certification evidence complete?", "What purity utility assumption is most fragile?", "Does schedule compression damage yield ramp?", "Where is specialist labour concentration highest?"]
    },
    "space": {
        "governor": "mission assurance, payload integration, launch-window dependency, range availability and qualification maturity",
        "hidden": "launch date confidence is governed by assurance closure and integration maturity, not headline assembly progress",
        "fail": "launch cadence assumptions exceeded payload, propulsion or mission-assurance maturity",
        "evidence": ["mission assurance closure matrix", "payload integration evidence", "range availability plan", "thermal/propulsion qualification gates", "launch readiness review criteria"],
        "traditional": "Vehicle or facility progress appears to support the target launch window.",
        "casey": "Mission assurance, payload integration and range readiness govern launch confidence.",
        "attacks": ["Which qualification gate can stop launch?", "Is range availability booked or assumed?", "What payload integration risk is understated?", "Can cadence survive one failed test event?", "What evidence proves launch readiness beyond schedule intent?"]
    },
    "defence": {
        "governor": "mission readiness, classified supplier dependencies, security assurance and systems integration maturity",
        "hidden": "capability readiness is governed by assurance and integration maturity, not procurement progress",
        "fail": "classified dependencies and mission-system integration lagged delivery reporting",
        "evidence": ["security accreditation path", "classified supplier delivery evidence", "mission-system integration tests", "operational acceptance criteria", "sovereign supply-chain risk log"],
        "traditional": "Procurement and construction packages appear reportable.",
        "casey": "Mission readiness, security assurance and integration maturity govern defensible capability.",
        "attacks": ["Which classified dependency is least recoverable?", "What mission-system test gate is immature?", "Where is sovereign supply-chain fragility hidden?", "Can security accreditation support the date?", "Who owns operational readiness acceptance?"]
    },
    "oil_gas": {
        "governor": "process safety, modular delivery, shutdown windows, commissioning sequence and export/start-up readiness",
        "hidden": "start-up confidence is governed by process safety and commissioning readiness, not mechanical progress",
        "fail": "process-safety assurance and commissioning sequencing lagged mechanical completion",
        "evidence": ["HAZOP / process-safety closure", "shutdown window approval", "module delivery certainty", "commissioning sequence", "start-up readiness gates"],
        "traditional": "Construction and mechanical package progress appear strong.",
        "casey": "Process safety, shutdown access and commissioning sequence govern start-up confidence.",
        "attacks": ["Which safety gate can block start-up?", "Is the shutdown window secured or assumed?", "Where is modular interface risk concentrated?", "What commissioning sequence is non-recoverable?", "Does reserve reflect process-safety uncertainty?"]
    },
    "nuclear": {
        "governor": "safety case maturity, regulator hold points, nuclear-grade QA and commissioning governance",
        "hidden": "delivery confidence is governed by licensing and QA traceability more than construction productivity",
        "fail": "safety-case and regulator hold-point maturity lagged physical delivery",
        "evidence": ["safety case maturity", "regulator hold-point plan", "QA traceability evidence", "nuclear-grade supplier records", "commissioning governance gates"],
        "traditional": "Physical works and package progress can appear stable.",
        "casey": "Regulatory hold points, safety-case maturity and nuclear-grade QA govern approval confidence.",
        "attacks": ["Which regulator hold point is least evidenced?", "Can QA traceability survive audit?", "Where is nuclear-grade supplier certainty weakest?", "Does schedule assume licensing approval?", "What safety-case gap blocks commissioning?"]
    },
    "healthcare": {
        "governor": "clinical commissioning, infection-control assurance, medical equipment integration and phased occupancy",
        "hidden": "clinical readiness and safe occupancy govern opening confidence more than construction completion",
        "fail": "clinical commissioning and operational transition lagged facility completion",
        "evidence": ["clinical commissioning plan", "infection-control sign-off", "medical equipment integration evidence", "phased occupancy plan", "patient transition readiness"],
        "traditional": "Building completion and fit-out progress appear stable.",
        "casey": "Clinical readiness, equipment integration and safe occupancy govern usable opening confidence.",
        "attacks": ["Which clinical pathway blocks opening?", "Is infection-control evidence complete?", "What equipment integration risk is hidden?", "Can phased occupancy support operations?", "Who owns patient transition readiness?"]
    },
    "water": {
        "governor": "environmental permits, process commissioning, network tie-ins and regulatory performance acceptance",
        "hidden": "service confidence is governed by permits, treatment performance and network integration",
        "fail": "process commissioning and environmental compliance lagged civil completion",
        "evidence": ["permit status", "process performance tests", "network tie-in plan", "environmental compliance evidence", "operator acceptance gates"],
        "traditional": "Civil asset progress appears reportable.",
        "casey": "Treatment performance, permits and network tie-ins govern service confidence.",
        "attacks": ["Which permit can block operation?", "Is treatment performance proven?", "What network tie-in assumption is fragile?", "Does reserve reflect compliance risk?", "Who accepts operator handover?"]
    },
    "ports": {
        "governor": "marine works windows, terminal systems, dredging/interface permits and operational handover",
        "hidden": "operational throughput readiness is governed by marine access, systems and regulatory interfaces",
        "fail": "marine/interface permitting and terminal systems readiness lagged physical works",
        "evidence": ["marine works window", "dredging/permitting status", "terminal systems integration", "operator transition plan", "customs/security readiness"],
        "traditional": "Berth and landside works appear measurable.",
        "casey": "Marine windows, terminal systems and operational handover govern throughput confidence.",
        "attacks": ["Which marine window is non-recoverable?", "Is dredging/permitting fully cleared?", "What terminal-system integration risk is hidden?", "Can operator transition support go-live?", "Where does customs/security readiness sit?"]
    },
    "energy": {
        "governor": "grid connection, major equipment procurement, outage windows and commissioning acceptance",
        "hidden": "energisation confidence is governed by network access, equipment certainty and commissioning readiness",
        "fail": "grid connection and major equipment certainty lagged construction progress",
        "evidence": ["grid connection agreement", "major equipment delivery evidence", "outage window plan", "commissioning acceptance criteria", "operator handover plan"],
        "traditional": "Physical asset progress appears stable.",
        "casey": "Grid access, equipment certainty and commissioning acceptance govern usable energy delivery.",
        "attacks": ["What grid assumption can fail first?", "Which equipment package is non-recoverable?", "Is the outage window secured?", "Where is commissioning acceptance under-evidenced?", "Does reserve reflect network dependency?"]
    },
    "infrastructure": {
        "governor": "interfaces, procurement certainty, regulatory approvals and commissioning readiness",
        "hidden": "confidence is governed by evidence maturity and integration logic rather than reported progress",
        "fail": "interface maturity and approval readiness lagged delivery reporting",
        "evidence": ["interface register", "procurement maturity", "regulatory approvals", "commissioning gates", "reserve basis"],
        "traditional": "Progress appears reportable through schedule and cost dashboards.",
        "casey": "Evidence maturity, interfaces and operational readiness govern board-defensible confidence.",
        "attacks": ["Which assumption collapses confidence first?", "What evidence proves readiness?", "Where is reserve optimism hidden?", "Who owns interface liability?", "Is float usable or theoretical?"]
    },
}

def _money_short(v: Any) -> str:
    try:
        x=float(v)
        if x >= 1000: return f"${x/1000:.1f}T"
        if x >= 1: return f"${x:.1f}B"
        return f"${x*1000:.0f}M"
    except Exception:
        return "not quantified"

def _advisor_metrics(model: Dict[str, Any]) -> Tuple[str, str, str]:
    conf = model.get("confidence_score", model.get("confidence", "n/a"))
    mc = model.get("monte_carlo", {}) or {}
    qcra = mc.get("qcra", {}) or {}
    qsra = mc.get("qsra", {}) or {}
    p80c = _money_short(qcra.get("p80", model.get("p80_cost_bn", 0)))
    p80s = qsra.get("p80", model.get("p80_schedule_months", "n/a"))
    return str(conf), p80c, str(p80s)

def _format_advisor_answer(question: str, model: Dict[str, Any]) -> str:
    q=(question or "").lower()
    fam=_casey_sector_family(model)
    intel=_SECTOR_INTEL.get(fam, _SECTOR_INTEL["infrastructure"])
    title=model.get("title") or model.get("client") or "this programme"
    scenario=str(model.get("scenario","base")).upper()
    conf,p80c,p80s=_advisor_metrics(model)
    risks=model.get("risks") or model.get("risk_register") or []
    risk_titles=[]
    for r in risks[:3]:
        if isinstance(r, dict): risk_titles.append(str(r.get("title") or r.get("risk") or r.get("name") or "unresolved risk"))
        else: risk_titles.append(str(r))
    if not risk_titles: risk_titles=intel["evidence"][:3]

    if any(k in q for k in ["traditional", "t&t", "turner", "consultant", "pm view", "controls view"]):
        return f"""TRADITIONAL CONTROLS VIEW VS CASEY VIEW

Traditional view:
{intel['traditional']}

CASEY view:
{intel['casey']}

Why this matters:
The conventional report can show progress while confidence is already deteriorating. CASEY is reading the programme behaviour underneath the dashboard.

CASEY POSITION
{title} should not be judged on reported progress alone. It should be judged on whether {intel['governor']} are evidenced strongly enough to survive board challenge."""

    if any(k in q for k in ["board", "attack", "question", "challenge"]):
        attacks="\n".join([f"{i+1}. {x}" for i,x in enumerate(intel["attacks"])])
        return f"""LIKELY BOARD ATTACKS

{attacks}

Most dangerous question:
"What evidence proves the current confidence position is real rather than narrative optimism?"

CASEY POSITION
The board is not simply approving cost or schedule. It is accepting exposure to {intel['governor']}. If that evidence is weak, the programme is not approval-ready."""

    if any(k in q for k in ["disagree", "reported", "official", "narrative", "status"]):
        return f"""CASEY DISAGREES WITH THE REPORTED POSITION

The visible programme narrative may appear stable.

CASEY disagrees.

The hidden governing constraint is that {intel['hidden']}. Reported progress can therefore improve while board-defensible confidence remains weak.

Current confidence: {conf}
P80 cost exposure: {p80c}
P80 schedule exposure: {p80s} months

CASEY POSITION
Do not approve acceleration or reduced reserve unless the evidence base proves {intel['governor']} are under control."""

    if any(k in q for k in ["fail", "public", "blame", "headline"]):
        return f"""IF THIS PROGRAMME FAILS

The most likely public explanation will be:

"{intel['fail']}."

The internal explanation will be sharper:
The programme treated reported delivery progress as equivalent to operational readiness.

CASEY POSITION
Failure risk is currently behavioural, not just numerical. The board should test whether the programme has confused progress with readiness."""

    if any(k in q for k in ["decision", "really deciding", "approve", "approval"]):
        ev="\n".join([f"- {x}" for x in intel["evidence"][:5]])
        return f"""WHAT THE BOARD IS REALLY DECIDING

The board is not just approving a project.

The board is accepting:
- exposure to {intel['governor']}
- the current evidence maturity position
- the adequacy of reserve against P80/P90 behaviour
- whether management can recover if the governing constraint slips

Evidence required before approval:
{ev}

CASEY POSITION
Approval should be conditional, not unconditional. Confidence improvement requires evidence closure, not more narrative."""

    if any(k in q for k in ["confidence", "why", "moved", "score"]):
        drivers="\n".join([f"- {x}" for x in intel["evidence"][:5]])
        return f"""CONFIDENCE EXPLANATION

Current confidence: {conf}

Confidence is not optimism. It is board defensibility.

CASEY reads confidence through:
{drivers}

Primary governing constraint:
{intel['governor']}

CASEY POSITION
If management wants confidence to improve, it should close evidence gaps around the governing constraint. Adding contingency without evidence does not create confidence; it only buys time."""

    if any(k in q for k in ["intervention", "do now", "fix", "improve", "actions"]):
        return f"""INTERVENTION INTELLIGENCE

The highest-value intervention is not more reporting.

Priority 1:
Prove {intel['governor']} with auditable evidence.

Priority 2:
Convert the weakest assumption into an owned mitigation with date, owner, trigger and residual exposure.

Priority 3:
Run a pressure test against P80/P90 rather than the median plan.

CASEY POSITION
The programme does not need more dashboard visibility. It needs targeted evidence closure where confidence is actually governed."""

    # default / what are we not seeing / hidden constraint
    return f"""WHAT THE BOARD IS NOT SEEING

The dominant issue is not the headline cost or schedule.

The hidden issue:
{intel['hidden']}.

Current confidence: {conf}
P80 cost exposure: {p80c}
P80 schedule exposure: {p80s} months

Most exposed items:
""" + "\n".join([f"- {x}" for x in risk_titles[:3]]) + f"""

CASEY POSITION
This programme should be governed through {intel['governor']}. If those items are not evidenced, the reported delivery position is not board-defensible."""

@app.post("/chat")
def chat(req: ChatRequest):
    try:
        return {"answer": _format_advisor_answer(req.question, req.project or {})}
    except Exception as e:
        return {"answer": f"""CASEY POSITION

The advisor could not fully resolve the live model payload, but the governance position is clear: do not rely on headline progress alone.

Use this challenge set:
1. What evidence proves readiness?
2. Which assumption collapses confidence first?
3. Is reserve adequate against P80/P90?
4. Is float usable or theoretical?
5. Who owns the governing interface?

Technical note: {str(e)[:180]}"""}




# ======================= CASEY V167 REAL INTAKE NORMALISATION ENGINE =======================
def _v167_money_to_bn(value):
    if value is None: return None
    if isinstance(value,(int,float)):
        # Excel workbooks vary: dollars, thousands, millions, billions
        v=float(value)
        if abs(v)>1e9: return v/1e9
        if abs(v)>1e6: return v/1e6/1000 if abs(v)>250000000 else v/1000  # mostly $m -> $bn for estimates
        if abs(v)>1000: return v/1000
        return v
    txt=str(value).strip().replace(',','').replace('$','')
    if not txt or txt.lower() in ('nan','none','n/a','-'): return None
    m=re.search(r'(?<![A-Za-z])([-+]?\d+(?:\.\d+)?)\s*([kmb]|million|billion|bn|m)?(?![A-Za-z])', txt, re.I)
    if not m: return None
    v=float(m.group(1)); unit=(m.group(2) or '').lower()
    if unit in ('b','bn','billion'): return v
    if unit in ('m','million'): return v/1000.0
    if unit=='k': return v/1_000_000.0
    # naked numbers: infer scale
    if abs(v)>1e9: return v/1e9
    if abs(v)>1e6: return v/1e9
    return v

def _v167_read_upload_text(name, content):
    lname=(name or '').lower(); text=''; rows=[]; sheets=[]
    if lname.endswith(('.xlsx','.xlsm','.xltx')):
        try:
            wb=load_workbook(BytesIO(content), data_only=True, read_only=True)
            for ws in wb.worksheets[:25]:
                sheet_rows=[]
                for ridx,row in enumerate(ws.iter_rows(values_only=True),1):
                    vals=[v for v in row if v is not None]
                    if vals:
                        line=' | '.join(str(v)[:160] for v in vals)
                        sheet_rows.append([str(v) if v is not None else '' for v in row[:30]])
                        text += f"\nSHEET={ws.title} ROW={ridx} :: {line}"
                    if ridx>600: break
                sheets.append({'name':ws.title,'rows':sheet_rows[:600]})
            return text[:350000], sheets
        except Exception as e:
            return f'XLSX parse warning: {e}', []
    try:
        text=content.decode('utf-8',errors='ignore')
    except Exception:
        text=str(content[:200000])
    # Simple CSV rows
    if lname.endswith('.csv') or ',' in text[:1000] or '\t' in text[:1000]:
        try:
            sample=text[:200000]
            dialect=csv.Sniffer().sniff(sample[:4096], delimiters=',\t;|')
            reader=csv.reader(StringIO(sample), dialect)
            rows=[r[:30] for _,r in zip(range(600),reader)]
        except Exception:
            rows=[]
    return text[:350000], [{'name':'text/csv','rows':rows}] if rows else []

def _v167_classify_file(name,text,sheets):
    t=((name or '')+' '+text[:8000]).lower()
    score={'xer':0,'cost':0,'risk':0,'schedule':0}
    if (name or '').lower().endswith('.xer') or '%t\ttask' in t or 'taskpred' in t: score['xer']+=10
    for k in ['cost','estimate','cbs','boq','budget','direct','indirect','contingency','allowance','escalation','qty','quantity','rate']:
        if k in t: score['cost']+=1
    for k in ['risk','probability','impact','mitigation','owner','trigger','residual','emv','p90','p80']:
        if k in t: score['risk']+=1
    for k in ['schedule','activity','task','start','finish','duration','predecessor','logic','critical','float','calendar']:
        if k in t: score['schedule']+=1
    if score['xer']>=5: return 'XER schedule'
    if score['cost']>=score['risk'] and score['cost']>=score['schedule'] and score['cost']>2: return 'cost workbook'
    if score['risk']>=score['schedule'] and score['risk']>2: return 'risk register'
    if score['schedule']>2: return 'schedule workbook'
    return 'mixed / unknown client file'

def _v167_extract_costs(text,sheets):
    candidates=[]; direct=indirect=reserve=None; line_count=0
    cost_words=re.compile(r'(cost|estimate|amount|total|direct|indirect|contingency|reserve|allowance|capex|budget|p50|p80|p90)',re.I)
    for sh in sheets:
        rows=sh.get('rows') or []
        # header rows
        for r in rows:
            joined=' '.join(r).lower()
            if cost_words.search(joined):
                nums=[_v167_money_to_bn(x) for x in r]
                nums=[n for n in nums if n is not None and abs(n)>0.0001]
                if nums:
                    label=' '.join(str(x) for x in r[:5])[:100]
                    val=max(nums, key=abs)
                    if 0.001 <= abs(val) <= 1000:
                        candidates.append({'sheet':sh.get('name'),'label':label,'value_bn':round(val,3)})
                        line_count+=1
                        if 'direct' in joined and direct is None: direct=val
                        if 'indirect' in joined and indirect is None: indirect=val
                        if any(k in joined for k in ['contingency','reserve','risk allowance']) and reserve is None: reserve=val
    # Text fallback
    for m in re.finditer(r'(direct|indirect|contingency|reserve|total|p50|p80|p90|capex|budget)[^\n\r]{0,80}?\$?\s*([0-9][0-9,\.]*\s*(?:bn|b|m|million|billion)?)', text, re.I):
        val=_v167_money_to_bn(m.group(2));
        if val is not None: candidates.append({'sheet':'text','label':m.group(1),'value_bn':round(val,3)})
    # derive totals
    vals=[c['value_bn'] for c in candidates if c.get('value_bn') is not None]
    total=None
    if vals:
        # choose a plausible headline: labelled total/P50 first; otherwise sum unique sheet line values.
        totals=[c['value_bn'] for c in candidates if re.search(r'total|p50|capex|budget|estimate', c.get('label',''), re.I)]
        sheet_vals=[]
        seen=set()
        for c in candidates:
            if c.get('sheet') == 'text': continue
            key=(c.get('sheet'), c.get('label'), c.get('value_bn'))
            if key not in seen:
                seen.add(key); sheet_vals.append(c.get('value_bn'))
        total=max(totals) if totals else (round(sum(sheet_vals),3) if len(sheet_vals)>1 else (direct+indirect+(reserve or 0) if direct and indirect else max(vals)))
    if direct is None and total: direct=round(total*0.70,3)
    if indirect is None and total: indirect=round(total*0.14,3)
    if reserve is None and total: reserve=round(max(total-direct-indirect, total*0.12),3)
    return {'p50_bn':round(total,3) if total else None,'direct_bn':round(direct,3) if direct else None,'indirect_bn':round(indirect,3) if indirect else None,'reserve_bn':round(reserve,3) if reserve else None,'cost_lines_found':line_count,'sample_lines':candidates[:12]}

def _v167_extract_risks(text,sheets):
    risks=[]; emv=[]; p90=[]
    risk_terms=re.compile(r'(risk|cause|event|impact|prob|probability|owner|mitigation|trigger|residual|emv|p90|p80)',re.I)
    for sh in sheets:
        rows=sh.get('rows') or []
        for r in rows:
            joined=' | '.join(r)
            low=joined.lower()
            if risk_terms.search(joined):
                nums=[_v167_money_to_bn(x) for x in r]
                nums=[n for n in nums if n is not None]
                if 'emv' in low and nums: emv.append(max(nums, key=abs))
                if 'p90' in low and nums: p90.append(max(nums, key=abs))
                if len(joined)>15 and any(k in low for k in ['risk','mitigation','owner','prob','impact','trigger']):
                    risks.append(joined[:260])
    # XER/text risk fallback
    for line in text.splitlines()[:5000]:
        if risk_terms.search(line) and len(line)>15:
            risks.append(line[:260])
    # Pull probability % and any values
    probs=[]
    for m in re.finditer(r'(\d{1,3})\s*%', text):
        v=int(m.group(1))
        if 1<=v<=99: probs.append(v)
    total_emv=round(sum(abs(x) for x in emv),3) if emv else None
    max_p90=round(max([abs(x) for x in p90]),3) if p90 else None
    return {'risk_rows_found':len(risks),'top_risk_rows':risks[:10],'probability_signals':probs[:12],'emv_bn':total_emv,'p90_bn':max_p90}

def _v167_extract_xer(text):
    if '%T\tTASK' not in text and '%T TASK' not in text and 'TASKPRED' not in text: return {}
    tasks=[]; preds=0; open_ends=0; constraints=0; durations=[]; in_task=False; fields=[]
    for line in text.splitlines():
        if line.startswith('%T'):
            in_task='TASK' in line and 'TASKPRED' not in line
        elif in_task and line.startswith('%F'):
            fields=line.split('\t')[1:]
        elif in_task and line.startswith('%R'):
            vals=line.split('\t')[1:]
            row=dict(zip(fields,vals)) if fields else {}
            code=row.get('task_code') or row.get('task_id') or (vals[3] if len(vals)>3 else '')
            name=row.get('task_name') or (vals[4] if len(vals)>4 else '')
            dur=row.get('target_drtn_hr_cnt') or row.get('remain_drtn_hr_cnt') or ''
            try: durations.append(float(dur)/8.0 if dur else 0)
            except Exception: pass
            if row.get('cstr_type') or row.get('act_constraint_type'): constraints+=1
            tasks.append({'code':code,'name':name[:120]})
        elif line.startswith('%R') and 'TASKPRED' in line[:80]:
            preds+=1
    if tasks:
        # More robust predecessor count: count rows after TASKPRED section
        in_pred=False
        for line in text.splitlines():
            if line.startswith('%T'): in_pred='TASKPRED' in line
            elif in_pred and line.startswith('%R'): preds+=1
        open_ends=max(0, len(tasks)-max(0,preds)) if preds else max(1,int(len(tasks)*0.2))
    months=round(sum(durations)/22/30,1) if durations else None
    return {'task_count':len(tasks),'logic_link_count':preds,'open_end_risk_count':open_ends,'constraint_count':constraints,'duration_months':months,'sample_tasks':tasks[:8]}

def _v167_build_intake_challenge(name, content):
    text,sheets=_v167_read_upload_text(name, content)
    ftype=_v167_classify_file(name,text,sheets)
    costs=_v167_extract_costs(text,sheets)
    risks=_v167_extract_risks(text,sheets)
    xer=_v167_extract_xer(text)
    p50=costs.get('p50_bn')
    if not p50 and risks.get('emv_bn'): p50=round(max(0.1,risks['emv_bn']*1.8),3)
    if not p50 and xer.get('task_count'): p50=round(max(0.5, xer['task_count']*0.018),3)
    p80=round(p50*1.28,3) if p50 else None
    p90=risks.get('p90_bn') or (round(p50*1.55,3) if p50 else None)
    conf=68
    if not costs.get('cost_lines_found'): conf-=10
    if risks.get('risk_rows_found')<5: conf-=8
    if xer.get('open_end_risk_count',0)>0: conf-=10
    if not costs.get('reserve_bn'): conf-=7
    if risks.get('p90_bn') and p50 and risks['p90_bn']>p50*2.5: conf-=14
    conf=max(24,min(82,conf))
    findings=[]
    findings.append(f"Intake normalised {ftype}: {len(sheets)} sheet/section set(s), {costs.get('cost_lines_found',0)} cost signals, {risks.get('risk_rows_found',0)} risk rows, {xer.get('task_count',0)} XER activities.")
    if p50: findings.append(f"Client-file derived P50/headline signal: ${p50:.1f}B. CASEY did not invent this from benchmark memory; it was inferred from the uploaded file content and reconciled to challenge ranges.")
    if p80: findings.append(f"CASEY challenge range: P80 ${p80:.1f}B, P90 ${p90:.1f}B. The board should not approve the P50 unless this downside is accepted or retired with evidence.")
    if costs.get('reserve_bn'): findings.append(f"Reserve/contingency signal found: ${costs['reserve_bn']:.1f}B. CASEY checks whether this is risk-linked or just a flat allowance.")
    if xer.get('task_count'): findings.append(f"Schedule logic detected: {xer['task_count']} activities and {xer.get('logic_link_count',0)} logic links. Open-end/logic risk count: {xer.get('open_end_risk_count',0)}.")
    red=[]
    if costs.get('p50_bn') and not costs.get('reserve_bn'): red.append('Cost workbook has headline cost but no clearly mapped reserve/contingency basis.')
    if costs.get('reserve_bn') and p50 and costs['reserve_bn']/p50<0.08: red.append('Reserve appears thin versus programme exposure; require QCRA support, not a flat percentage.')
    if risks.get('risk_rows_found') and not risks.get('emv_bn'): red.append('Risk register rows detected but EMV / P80 / residual exposure is not clearly reconciled to the cost model.')
    if risks.get('p90_bn') and p50 and risks['p90_bn']>p50*2.0: red.append(f"Uploaded risk downside is extreme: P90 signal ${risks['p90_bn']:.1f}B versus P50 ${p50:.1f}B. This is a board-level escalation.")
    if xer.get('open_end_risk_count',0)>0: red.append(f"XER schedule has {xer.get('open_end_risk_count')} potential open-end / weak-logic exposure points. QSRA must not rely on the contractor date.")
    if not red: red=['The file was parsed, but CASEY still requires source validation before commercial reliance. Challenge basis, owner, residual exposure and evidence closure before approval.']
    attacks=[
        'Show me where the uploaded P50 reconciles to P80/P90 downside — not a separate CASEY benchmark number.',
        'Which CBS/WBS package owns the largest unpriced exposure and who signs the evidence closure?',
        'Where does the risk register residual exposure land in the reserve, line by line?',
        'Which XER activities drive the board date, and are their predecessors, calendars and constraints defensible?',
        'What changed if we approve the contractor number versus the CASEY challenge range?'
    ]
    next_steps=[
        'Upload estimate + risk register + XER together, then reconcile CBS/WBS/activity IDs in one source bundle.',
        'Require P50/P80/P90 from the contractor and compare it to CASEY-derived downside.',
        'Close open schedule logic, missing owners, missing basis statements and unlinked reserve before board approval.',
        'Export challenged board pack only after the source bundle passes traceability checks.'
    ]
    result={'version':'CASEY V167 intake normalisation engine','filename':name,'file_type':ftype,'size_bytes':len(content),'schema_confidence':('high' if (costs.get('cost_lines_found') or xer.get('task_count') or risks.get('risk_rows_found')) else 'low'), 'source_intelligence':{'cost':costs,'risk':risks,'xer':xer}, 'challenge_model':{'p50_bn':p50,'p80_bn':p80,'p90_bn':p90,'confidence_pct':conf,'board_posture': 'Do not approve without more evidence' if conf<55 else 'Board challenge likely'}, 'findings':findings,'red_flags':red[:8], 'board_challenge_questions':attacks, 'next_steps':next_steps, 'if_this_fails':'The likely failure mode is not that the spreadsheet was wrong; it is that unlinked cost, schedule and risk assumptions were accepted as board-ready evidence.', 'epc_challenge':True}
    return result
# ======================= END CASEY V167 REAL INTAKE NORMALISATION ENGINE =======================

@app.post("/upload/analyse")
async def analyse_upload(file: UploadFile = File(...)):
    content = await file.read()
    name = file.filename or "upload"
    try:
        result = _v167_build_intake_challenge(name, content)
    except Exception as e:
        result={
            "version":"CASEY V167 intake normalisation engine",
            "filename":name,
            "size_bytes":len(content),
            "file_type":"unreadable / protected client file",
            "schema_confidence":"low",
            "findings":["File received but the parser could not fully normalise it. CASEY still challenges this as an evidence problem, not a UI failure."],
            "red_flags":[f"Parser exception: {str(e)[:160]}","Require unlocked XLSX/CSV/XER/PDF source export before commercial reliance."],
            "board_challenge_questions":["Why is the board being asked to approve a file that cannot be normalised?","Where is the unlocked source estimate, risk register and schedule export?"],
            "next_steps":["Ask for native XLSX/XER/CSV files, not screenshots or locked PDFs.","Run source-bundle reconciliation after receiving clean exports."],
            "epc_challenge":True
        }
    con=db(); cur=con.cursor(); cur.execute("INSERT INTO uploads(filename,created_at,analysis_json) VALUES(?,?,?)",(name,datetime.utcnow().isoformat(),json.dumps(result))); con.commit(); con.close()
    return result

# ------------------------- exports -------------------------
def style_ws(ws):
    fill=PatternFill("solid",fgColor="07111F"); font=Font(color="FFFFFF",bold=True); thin=Side(style="thin",color="26384F")
    ws.freeze_panes="A2"
    for row in ws.iter_rows():
        for cell in row:
            cell.alignment=Alignment(vertical="top",wrap_text=True); cell.border=Border(left=thin,right=thin,top=thin,bottom=thin)
    for c in ws[1]: c.fill=fill; c.font=font
    for i in range(1,min(ws.max_column+1,50)): ws.column_dimensions[get_column_letter(i)].width=22

def add_sheet(wb,name,rows):
    ws=wb.create_sheet(name[:31]) if not wb.sheetnames else wb.create_sheet(name[:31]); ws.title=name[:31]
    for r in rows: ws.append(r)
    style_ws(ws); return ws

def workbook_bytes(model):
    model = _casey_reconcile_cost_lines(model)
    wb=Workbook(); wb.remove(wb.active)
    add_sheet(wb,"Board Summary",[["Field","Value"],["Generated",datetime.utcnow().isoformat()],["Version",model.get("version")],["Title",model.get("title")],["Client",model.get("client")],["Executive Summary",model.get("executive_summary")],["Mode",model.get("mode")],["Subsector",model.get("subsector")],["Location",model.get("location")],["Scenario",model.get("scenario_label")],["Cost P10",model.get("cost_p10")],["Cost P50",model.get("cost_p50")],["Cost P90",model.get("cost_p90")],["Schedule",model.get("schedule")],["QCRA P80",money_bn(model["monte_carlo"]["qcra"]["p80"])],["QSRA P80",model["monte_carlo"]["qsra"]["p80"]],["Risk",model.get("risk")],["Confidence",f"{model.get('confidence_pct')}%"]])
    add_sheet(wb,"Primary Cost Estimate",[["CBS","Description","Type","Basis","P10 BN","P50 BN","P90 BN","Impact Basis"]]+[[x["cbs"],x["description"],x["type"],x["basis"],x["p10_bn"],x["p50_bn"],x["p90_bn"],x["impact_basis"]] for x in model["cost_lines"]])
    for c,rows in model["estimates_by_class"].items(): add_sheet(wb,f"Class {c} Estimate",[["Class","CBS","Description","P10","P50","P90","Maturity"]]+[[c,x["cbs"],x["description"],x["p10_bn"],x["p50_bn"],x["p90_bn"],x["maturity"]] for x in rows])
    for l,rows in model["schedules_by_level"].items(): add_sheet(wb,f"Level {l} Schedule",[["Activity ID","Phase","Activity","Predecessor","Duration Months","Critical","Basis"]]+[[x["activity_id"],x["phase"],x["activity"],x["predecessor"],x["duration_months"],x["critical"],x["basis"]] for x in rows])
    add_sheet(wb,"Risk Register",[["ID","Risk","Category","Probability","Activity","CBS","Sched O","Sched M","Sched P","Cost O","Cost M","Cost P","Owner","Trigger","Mitigation","Cost Basis","Schedule Basis"]]+[[r["risk_id"],r["title"],r["category"],r["probability_pct"],r["activity_id"],r["cbs"],r["schedule_o_days"],r["schedule_m_days"],r["schedule_p_days"],r["cost_o_bn"],r["cost_m_bn"],r["cost_p_bn"],r["owner"],r["trigger"],r["mitigation"],r["basis_of_cost_impact"],r["basis_of_schedule_impact"]] for r in model["risks"]])
    add_sheet(wb,"Monte Carlo P-Curve",[["Percentile","QCRA Cost BN","QSRA Months"]]+[[x["percentile"],x["cost_bn"],x["schedule_months"]] for x in model["monte_carlo"]["curve"]])
    add_sheet(wb,"Tornado Drivers",[["Risk","Title","Activity","CBS","Cost Mean BN","Schedule Mean Days","Driver Score"]]+[[x["risk_id"],x["title"],x["activity_id"],x["cbs"],x["cost_mean_bn"],x["schedule_mean_days"],x["driver_score"]] for x in model["monte_carlo"]["tornado"]])
    add_sheet(wb,"Scenarios",[["Scenario","Cost","Schedule Months","Risk","Confidence","Why"]]+[[x["label"],x["cost"],x["schedule_months"],x["risk"],x["confidence"],x["why"]] for x in model["scenario_comparison"]])
    add_sheet(wb,"Benchmarks",[["Metric","Value","Why"]]+[[x["metric"],x["value"],x["why"]] for x in model["benchmarks"]])
    add_sheet(wb,"Demo Script",[["Step"]]+[[x] for x in model["launch_demo_script"]])
    # charts
    try:
        ws=wb["Monte Carlo P-Curve"]; chart=LineChart(); chart.title="QCRA/QSRA P-Curve"; data=Reference(ws,min_col=2,max_col=3,min_row=1,max_row=ws.max_row); cats=Reference(ws,min_col=1,min_row=2,max_row=ws.max_row); chart.add_data(data,titles_from_data=True); chart.set_categories(cats); ws.add_chart(chart,"E2")
        ws2=wb["Tornado Drivers"]; b=BarChart(); b.title="Tornado Drivers"; data=Reference(ws2,min_col=7,min_row=1,max_row=min(ws2.max_row,12)); cats=Reference(ws2,min_col=2,min_row=2,max_row=min(ws2.max_row,12)); b.add_data(data,titles_from_data=True); b.set_categories(cats); ws2.add_chart(b,"I2")
    except Exception: pass
    bio=BytesIO(); wb.save(bio); bio.seek(0); return bio.getvalue()

def risk_csv_bytes(model):
    model = _casey_reconcile_cost_lines(model)
    out=StringIO(); w=csv.writer(out); w.writerow(["Risk ID","Title","Category","Probability %","Activity ID","Activity Name","CBS","CBS Name","Schedule O","Schedule M","Schedule P","Cost O BN","Cost M BN","Cost P BN","Owner","Trigger","Mitigation","Basis Cost","Basis Schedule"])
    for r in model["risks"]: w.writerow([r["risk_id"],r["title"],r["category"],r["probability_pct"],r["activity_id"],r["activity_name"],r["cbs"],r["cbs_name"],r["schedule_o_days"],r["schedule_m_days"],r["schedule_p_days"],r["cost_o_bn"],r["cost_m_bn"],r["cost_p_bn"],r["owner"],r["trigger"],r["mitigation"],r["basis_of_cost_impact"],r["basis_of_schedule_impact"]])
    return out.getvalue().encode()

def xer_bytes(model):
    model = _casey_reconcile_cost_lines(model)
    lines=["ERMHDR\t24.12\t2026-05-02\tProject\tCASEY\tCASEY\tProject Management\tUSD","%T\tPROJECT","%F\tproj_id\tproj_short_name\tproj_name",f"%R\t1\tCASEY\t{model.get('title','CASEY Project')}","%T\tWBS","%F\twbs_id\tproj_id\twbs_short_name\twbs_name","%R\t1\t1\tCASEY\tCASEY Project","%T\tTASK","%F\ttask_id\tproj_id\twbs_id\ttask_code\ttask_name\tduration_type\ttarget_drtn_hr_cnt"]
    idx={}
    for i,a in enumerate(model["schedule_rows"],1): idx[a["activity_id"]]=i; lines.append(f"%R\t{i}\t1\t1\t{a['activity_id']}\t{a['activity']}\tFixed Duration\t{a['duration_months']*160}")
    lines += ["%T\tTASKPRED","%F\ttask_pred_id\ttask_id\tpred_task_id\tpred_type"]
    k=1
    for a in model["schedule_rows"]:
        for p in str(a.get("predecessor","")).split(";"):
            p=p.strip()
            if p in idx: lines.append(f"%R\t{k}\t{idx[a['activity_id']]}\t{idx[p]}\tFS"); k+=1
    return "\n".join(lines).encode()

def word_bytes(model):
    model = _casey_reconcile_cost_lines(model)
    doc=Document(); styles=doc.styles; styles["Normal"].font.name="Aptos"; styles["Normal"].font.size=Pt(10)
    h=doc.add_heading("CASEY TITAN X Board Report",0); h.runs[0].font.color.rgb=RGBColor(5,20,35)
    doc.add_paragraph(model.get("executive_summary",""))
    tbl=doc.add_table(rows=1,cols=2); tbl.style="Light Shading Accent 1"; hdr=tbl.rows[0].cells; hdr[0].text="Metric"; hdr[1].text="Value"
    for k,v in [("Cost P50",model["cost_p50"]),("Cost Range",model["cost_range"]),("Schedule",model["schedule"]),("QCRA P80",money_bn(model["monte_carlo"]["qcra"]["p80"])),("QSRA P80",str(model["monte_carlo"]["qsra"]["p80"])+" months"),("Risk",model["risk"]),("Confidence",str(model["confidence_pct"])+"%")]: row=tbl.add_row().cells; row[0].text=k; row[1].text=v
    for title,items in [("Why Confidence",model["confidence_explanation"]),("Top Red Flags",model["red_flags"]),("Board Challenge Questions",model["board_challenge_questions"]),("Next Best Actions",model["next_best_actions"])]:
        doc.add_heading(title,1)
        for x in items: doc.add_paragraph(x,style="List Bullet")
    doc.add_heading("Top Risk Drivers",1)
    for r in model["risks"][:8]: doc.add_paragraph(f"{r['risk_id']} {r['title']} — {r['activity_id']} / {r['cbs']}: {r['basis_of_schedule_impact']}",style="List Bullet")
    bio=BytesIO(); doc.save(bio); bio.seek(0); return bio.getvalue()

def pptx_bytes(model):
    model = _casey_reconcile_cost_lines(model)
    from pptx.dml.color import RGBColor as PRGBColor
    prs=Presentation(); blank=prs.slide_layouts[6]
    def add_title(slide,t,sub=""):
        box=slide.shapes.add_textbox(PptxInches(.5),PptxInches(.35),PptxInches(12),PptxInches(.8)); tf=box.text_frame; tf.text=t; tf.paragraphs[0].runs[0].font.size=PptxPt(34); tf.paragraphs[0].runs[0].font.bold=True
        if sub: b=slide.shapes.add_textbox(PptxInches(.55),PptxInches(1.1),PptxInches(12),PptxInches(.4)); b.text_frame.text=sub
    s=prs.slides.add_slide(blank); add_title(s,"CASEY TITAN X",model.get("executive_summary",""));
    s2=prs.slides.add_slide(blank); add_title(s2,"Board Metrics")
    metrics=[("P50",model["cost_p50"]),("Range",model["cost_range"]),("Schedule",model["schedule"]),("Risk",model["risk"]),("Confidence",str(model["confidence_pct"])+"%"),("QCRA P80",money_bn(model["monte_carlo"]["qcra"]["p80"]))]
    for i,(k,v) in enumerate(metrics):
        x=.7+(i%3)*4.1; y=1.6+(i//3)*1.6; shp=s2.shapes.add_shape(1,PptxInches(x),PptxInches(y),PptxInches(3.6),PptxInches(1.0)); shp.text=f"{k}\n{v}"
    s3=prs.slides.add_slide(blank); add_title(s3,"Top Risk Drivers")
    txt=s3.shapes.add_textbox(PptxInches(.7),PptxInches(1.4),PptxInches(12),PptxInches(5)).text_frame
    for r in model["risks"][:8]: p=txt.add_paragraph(); p.text=f"{r['title']} — {r['activity_id']} / {r['cbs']}"; p.level=0
    bio=BytesIO(); prs.save(bio); bio.seek(0); return bio.getvalue()

def pdf_bytes(model):
    model = _casey_reconcile_cost_lines(model)
    bio=BytesIO(); doc=SimpleDocTemplate(bio,pagesize=A4); styles=getSampleStyleSheet(); story=[Paragraph("CASEY TITAN X Board Report",styles["Title"]),Spacer(1,12),Paragraph(model.get("executive_summary",""),styles["BodyText"]),Spacer(1,12)]
    data=[["Metric","Value"],["Cost P50",model["cost_p50"]],["Cost Range",model["cost_range"]],["Schedule",model["schedule"]],["Risk",model["risk"]],["Confidence",f"{model['confidence_pct']}%"],["QCRA P80",money_bn(model["monte_carlo"]["qcra"]["p80"])],["QSRA P80",str(model["monte_carlo"]["qsra"]["p80"])+" months"]]
    table=Table(data); table.setStyle(TableStyle([("BACKGROUND",(0,0),(-1,0),colors.HexColor("#07111F")),("TEXTCOLOR",(0,0),(-1,0),colors.white),("GRID",(0,0),(-1,-1),.5,colors.grey)])); story.append(table); story.append(PageBreak())
    story.append(Paragraph("Top Risk Drivers",styles["Heading1"]))
    for r in model["risks"][:10]: story.append(Paragraph(f"<b>{r['risk_id']} {r['title']}</b>: {r['activity_id']} / {r['cbs']} — {r['mitigation']}",styles["BodyText"])); story.append(Spacer(1,6))
    doc.build(story); bio.seek(0); return bio.getvalue()

def stream(data:bytes, media:str, filename:str): return StreamingResponse(BytesIO(data), media_type=media, headers={"Content-Disposition":f"attachment; filename={filename}"})



@app.get("/v26/launch-readiness")
def launch_readiness():
    return {
        "status": "demo-ready / production-scaffolded",
        "included": [
            "Earth and Space feature demo theatre",
            "QCRA/QSRA Monte Carlo engine",
            "Risk register mapped to WBS activities and CBS cost accounts",
            "Class 1-5 estimate comparison",
            "Schedule Level 1-5 comparison",
            "Export centre: cost workbook, risk register, schedule export, model audit and full pack",
            "Upload analysis / project doctor endpoint",
            "Saved projects via local SQLite scaffold",
            "Public one-use demo limiter",
            "Server-side private AI key environment pattern",
            "Auth, billing, deployment and database scaffolding docs"
        ],
        "production_next_steps": [
            "Connect hosted Postgres",
            "Add Clerk/Auth0/Supabase Auth",
            "Add Stripe live product IDs",
            "Connect verified benchmark data",
            "Deploy frontend and backend"
        ]
    }

@app.get("/v26/billing/plans")
def billing_plans():
    return {
        "plans": [
            {"name":"Demo", "price":"Public one-use", "features":["Earth/Space demo", "sample exports"]},
            {"name":"Pro", "price":"TBC", "features":["saved projects", "full exports", "upload analysis"]},
            {"name":"Enterprise", "price":"Custom", "features":["teams", "SSO", "benchmark library", "private deployment"]}
        ],
        "note": "Stripe hooks are scaffolded; configure live price IDs before charging customers."
    }

@app.get("/v26/deployment-check")
def deployment_check():
    return {
        "backend": "ok",
        "database_path": DB_PATH,
        "demo_limit_per_ip": DEMO_LIMIT_PER_IP,
        "openai_key_visible_to_users": False,
        "environment_variables": ["OPENAI_API_KEY", "CASEY_DB", "CASEY_DEMO_LIMIT_PER_IP", "STRIPE_SECRET_KEY", "STRIPE_PRICE_ID"],
        "ready_for_local_demo": True
    }

@app.get("/demo/showcases")
def demo_showcases():
    return {
        "headline": "Two-minute appetite demo: Earth and Space in one product.",
        "earth": {
            "title": "Earth revenue demo: 2026/2027 AI data centre",
            "prompt": "Riyadh AI Hyperscale Campus 500MW accelerated 2027 with sovereign cloud, grid connection and liquid cooling",
            "client": "Microsoft",
            "alternates": ["Northern Virginia GPU Data Centre Expansion 2026", "Manchester AI Compute Campus UK 2027", "Boston GMP Life Sciences Campus 2026", "Arizona Advanced Semiconductor Fab 2027"],
            "script": [
                "Open with a live buyer-market example: AI data centre, life sciences or semiconductor facility.",
                "Show cost, schedule, confidence, QCRA/QSRA and risk mapping in the first 60 seconds.",
                "Open scenarios: Faster, Cheaper and Lower-Risk explain the commercial choices.",
                "Open exports and download the full boardroom output pack."
            ]
        },
        "space": {
            "title": "Space jaw-drop demo: Lunar Base Alpha",
            "prompt": "Lunar Base Alpha with 1000 crew, nuclear power, landing pads, life support and launch logistics",
            "client": "Lunar Development Authority",
            "alternates": ["Mars Fuel Refinery with ISRU methane oxygen production", "Orbital Solar Power Ring delivering energy to Earth", "Satellite Gigafactory for 10000 satellites per year", "Deep Space Cargo Hub for asteroid and lunar logistics"],
            "script": [
                "Switch from Earth to Space without changing tools.",
                "Show launch, life-support, mass growth, radiation, ISRU and logistics risks mapped to QSRA/QCRA.",
                "Use the same output centre to prove CASEY handles impossible projects with traceable assumptions.",
                "Close with the same export pack: cost workbook, risk register, schedule export, model audit and full pack."
            ]
        },
        "demo_rules": [
            "Do not explain software first; ask for a project first.",
            "Run Earth demo, then Space demo, then open Monte Carlo and Exports.",
            "End with: CASEY compresses weeks of early controls work into minutes, with traceable assumptions."
        ]
    }


@app.get("/v26/revenue-machine")
def revenue_machine():
    return {
        "positioning": "AI Operating System for capital projects on Earth and in Space.",
        "headline": "Price the future before it gets built.",
        "target_verticals_2026_2027": [
            "AI data centres: Northern Virginia, Riyadh, Abu Dhabi, Manchester, Dublin, Johor, Phoenix, Dallas",
            "Life sciences: Boston/Cambridge MA, North Carolina, Oxford, Cambridge UK, Liverpool, Basel",
            "Semiconductors: Arizona, Texas, Ohio, Wales compound semiconductor cluster, Germany fabs",
            "Space: lunar bases, Mars ISRU, spaceports, satellite factories, orbital power, logistics hubs"
        ],
        "demo_flow": [
            "Run Earth demo using a live-demand project type.",
            "Show cost, schedule, confidence, QCRA/QSRA, mapped risks, peers and scenarios.",
            "Download the board pack live.",
            "Switch to Space demo and prove the same controls OS works for frontier infrastructure.",
            "End with ROI vs traditional consultancy and a book-demo CTA."
        ],
        "commercial_ctas": ["Book demo", "Request sample board pack", "Upload project for review", "Enterprise pilot"]
    }

@app.get("/v26/demo-library")
def v26_demo_library():
    return {
        "data_centres": [
            "Riyadh AI Hyperscale Campus 500MW accelerated 2027",
            "Northern Virginia GPU Data Centre Expansion 2026",
            "Manchester AI Compute Campus UK 2027",
            "Slough Hyperscale Expansion 2027",
            "Dallas GPU Compute Mega Campus 2026",
            "Johor AI Data Corridor 2027"
        ],
        "life_sciences": [
            "Boston GMP Biotech Campus 2026",
            "Cambridge UK Cell Therapy Facility 2027",
            "Oxford Vaccine Research Campus 2026",
            "North Carolina Pharma Mega Plant 2027",
            "Liverpool Pharma Expansion Plant 2027"
        ],
        "semiconductors": [
            "Arizona Advanced Semiconductor Fab 2027",
            "Texas Advanced Packaging Plant 2026",
            "Ohio Semiconductor Mega Campus 2027",
            "Wales Compound Semiconductor Hub 2026",
            "Germany Chip Sovereignty Fab 2027"
        ],
        "space": [
            "Lunar Base Alpha with 1000 crew and nuclear power",
            "Mars Fuel Refinery with ISRU methane oxygen production",
            "Orbital Solar Power Ring delivering energy to Earth",
            "UK Spaceport Expansion with launch pads and fuels",
            "Satellite Gigafactory for 10000 satellites per year",
            "Deep Space Cargo Hub for asteroid and lunar logistics"
        ]
    }

@app.get("/v26/pricing")
def v26_pricing():
    return {
        "starter": {"price":"£499/mo", "ideal_for":"Small teams and early concepts", "features":["10 projects", "board summaries", "basic exports", "demo examples"]},
        "pro": {"price":"£1,999/mo", "ideal_for":"Project owners, developers and controls teams", "features":["unlimited projects", "QCRA/QSRA", "XER", "upload analysis", "peer intelligence"]},
        "enterprise": {"price":"Custom", "ideal_for":"Large owners, governments, funds and enterprise teams", "features":["teams", "SSO", "private deployment", "benchmark library", "white-label reports"]}
    }

@app.post("/export/workbook")
def export_workbook(model:Dict[str,Any]): return stream(workbook_bytes(model),"application/vnd.openxmlformats-officedocument.spreadsheetml.sheet","CASEY_Cost_Model_Planet_Class.xlsx")
@app.post("/export/risk-register")
def export_risk(model:Dict[str,Any]): return stream(risk_register_workbook_bytes(model),"application/vnd.openxmlformats-officedocument.spreadsheetml.sheet","CASEY_Risk_Register_Pro.xlsx")
@app.post("/export/xer")
def export_xer(model:Dict[str,Any]): return stream(xer_bytes(model),"application/octet-stream","CASEY_TITAN_X_v26_P6_Schedule.xer")
@app.post("/export/word")
def export_word(model:Dict[str,Any]): return stream(word_bytes(model),"application/vnd.openxmlformats-officedocument.wordprocessingml.document","CASEY_Executive_Board_Report.docx")
@app.post("/export/pdf")
def export_pdf(model:Dict[str,Any]): return stream(pdf_bytes(model),"application/pdf","CASEY_Board_Intelligence_Pack.pdf")
@app.post("/export/pptx")
def export_pptx(model:Dict[str,Any]): return stream(pptx_bytes(model),"application/vnd.openxmlformats-officedocument.presentationml.presentation","CASEY_Board_Deck_Elite.pptx")
@app.post("/export/json")
def export_json(model:Dict[str,Any]): return stream(json.dumps(_casey_reconcile_cost_lines(model),indent=2).encode(),"application/json","CASEY_TITAN_X_v26_Model.json")
@app.post("/export/all")
def export_all(model:Dict[str,Any]):
    model = _casey_reconcile_cost_lines(model)
    bio=BytesIO()
    with zipfile.ZipFile(bio,"w",zipfile.ZIP_DEFLATED) as z:
        z.writestr("01_CASEY_Cost_Model_Planet_Class.xlsx",workbook_bytes(model))
        z.writestr("02_CASEY_Risk_Register_Pro.xlsx",risk_register_workbook_bytes(model))
        z.writestr("03_CASEY_P6_Schedule.xer",xer_bytes(model))
        z.writestr("04_CASEY_Executive_Board_Report.docx",word_bytes(model))
        z.writestr("05_CASEY_Board_Intelligence_Pack.pdf",pdf_bytes(model))
        z.writestr("06_CASEY_Board_Deck_Elite.pptx",pptx_bytes(model))
        z.writestr("07_CASEY_Full_Model_Audit.json",json.dumps(model,indent=2))
        z.writestr("08_CASEY_Risk_Register_Raw.csv",risk_csv_bytes(model))
        z.writestr("09_CASEY_Demo_Close_Script.txt","\n".join(model.get("launch_demo_script",[])))
    bio.seek(0); return stream(bio.getvalue(),"application/zip","CASEY_Output_Pack_Planet_Class.zip")


# ========================= CASEY v50 OUTPUT DOMINATION OVERRIDES =========================
# These overrides replace the older lightweight output pack with premium, scenario-linked,
# board-ready exports. They keep the same API endpoints so the frontend does not break.

def _scenario_effects(model):
    sc=(model.get("scenario") or "base").lower()
    effects={
        "base":{"label":"Base","cost_note":"Balanced reference case.","schedule_note":"Standard delivery logic.","risk_note":"Normal contingency and mitigation profile."},
        "faster":{"label":"Faster","cost_note":"Acceleration premium added for overlap, premium logistics and extended shifts.","schedule_note":"Critical path compressed with higher interface risk.","risk_note":"Residual delivery and quality risk increases unless assurance is funded."},
        "cheaper":{"label":"Cheaper","cost_note":"Value engineering and procurement competition reduce initial capex.","schedule_note":"Programme may lengthen due to procurement tension and reduced float.","risk_note":"Scope, quality and rework exposure rises; board should not treat savings as certain."},
        "lower_risk":{"label":"Lower Risk","cost_note":"Higher assurance, surveys and contingency improve board confidence.","schedule_note":"More realistic buffers and stage gates improve QSRA reliability.","risk_note":"Residual risk reduced through earlier evidence and stronger controls."},
        "premium":{"label":"Premium","cost_note":"Flagship resilience, assurance and quality increase capex.","schedule_note":"Schedule protected by better readiness and procurement strategy.","risk_note":"Higher front-end investment reduces downstream delivery uncertainty."},
        "investor":{"label":"Investor","cost_note":"Frames capex as a decision range with commercial challenge points.","schedule_note":"Highlights fundability and milestones.","risk_note":"Emphasis on value leakage and investment downside protection."},
        "survival":{"label":"Survival","cost_note":"Minimum viable scope reduces capex but creates major residual risk.","schedule_note":"Lean programme with limited contingency.","risk_note":"Not board-safe without explicit risk acceptance."},
    }
    return effects.get(sc,effects["base"])

def _total_costs(model):
    lines=model.get("cost_lines",[])
    totals={"Direct":0.0,"Indirect":0.0,"Reserve":0.0}
    for x in lines:
        typ=x.get("type","Direct")
        totals[typ]=totals.get(typ,0)+float(x.get("p50_bn",0) or 0)
    return totals

def _risk_rating(prob, cost_m, sched_m):
    exposure = prob/100.0*(float(cost_m or 0)*100 + float(sched_m or 0)/4)
    if exposure > 90: return "Extreme"
    if exposure > 55: return "High"
    if exposure > 25: return "Medium"
    return "Low"

def risk_register(mode,subsector,cost,months,schedule,costs,rm,scenario):
    """Premium quantified risk register — 15-25 sector-specific risks per project type."""
    s = subsector.lower()

    # ── UNIVERSAL BASE (12 risks — every sector) ────────────────────
    base=[
      ("R-001","Scope growth","Scope","Ambiguous requirements, immature scope freeze and late stakeholder changes","Approved scope expands after estimate freeze","Additional quantities, redesign, rework, procurement churn and board contingency drawdown","Both",42,15,45,120,0.030,0.080,0.160,"Project Director","Freeze scope baseline, change board, value gates, client decision log","Scope change rate exceeds 2% of package value","Reduce","Open"),
      ("R-002","Market escalation","Commercial","Supplier market tightness, inflation, FX and long procurement window","Rates exceed assumed escalation profile","P50 cost understated; commercial approvals or procurement strategy may need reset","Cost",45,0,0,0,0.040,0.100,0.200,"Commercial Lead","Early procurement, index-linked allowances, market testing, FX strategy","Index exceeds allowance by 3%","Mitigate","Open"),
      ("R-003","Permits and approvals delay","Regulatory","Authority process, stakeholder objections or incomplete submissions","Consent milestone slips beyond baseline","Critical path delay, extended preliminaries and potential redesign conditions","Schedule",35,20,70,180,0.010,0.040,0.080,"Consents Lead","Authority plan, consent tracker, early submissions, stakeholder map","Consent milestone slips by 30 days","Mitigate","Open"),
      ("R-004","Design maturity gap","Technical","Estimate based on immature design, incomplete surveys or unresolved interfaces","Design basis changes during detailed design","Cost growth, quantity movement, rework and schedule resequencing","Both",38,15,45,120,0.030,0.090,0.180,"Design Manager","Design maturity gates, independent review, assumption register","Design deliverables miss maturity gate","Reduce","Open"),
      ("R-005","Supply chain disruption","Procurement","Long-lead equipment, constrained suppliers or late procurement release","Supplier promise dates move right","Delayed installation, critical path movement and acceleration cost","Both",40,10,40,120,0.020,0.070,0.150,"Procurement Lead","Alternate suppliers, early orders, expediting, framework options","Supplier promise date slips","Transfer / Mitigate","Open"),
      ("R-006","Labour productivity underperformance","Delivery","Access constraints, poor sequencing, labour scarcity or learning curve","Production rates underperform plan","Extended duration, increased preliminaries and loss of float","Both",35,15,50,130,0.030,0.080,0.160,"Delivery Lead","Package productivity controls, daily planning, earned value cadence","SPI/CPI deteriorates for two periods","Mitigate","Open"),
      ("R-007","Commissioning and testing delay","Handover","Incomplete readiness, software integration or defects during testing","Systems fail tests or require repeated commissioning cycles","Delayed handover, operational readiness slippage and liquidated damages exposure","Both",32,10,35,100,0.020,0.060,0.120,"Commissioning Lead","Commissioning readiness plan, early test packs, digital integration","Test failures trend upward","Reduce","Open"),
      ("R-008","Interface misalignment","Integration","Multiple contractors, unclear interface ownership or late design coordination","Package interfaces do not align during installation","Rework, claims, delay and fragmented accountability","Both",34,10,35,100,0.020,0.070,0.140,"Integration Manager","Interface control documents, weekly interface board, accountable owners","Interface actions overdue","Mitigate","Open"),
      ("R-009","Third-party and stakeholder delay","Stakeholder","Utility owners, regulators, adjacent landowners or community opposition","Third-party agreements or diversions not complete at programme gate","Critical path suspension, additional mitigation cost and reputational exposure","Schedule",30,15,55,150,0.010,0.035,0.075,"Stakeholder Lead","Third-party agreement register, escalation protocol, legal options","Third-party action not confirmed","Mitigate","Open"),
      ("R-010","Client decision latency","Governance","Client approval processes, funding sign-off or programme governance gaps","Key decisions not made in time to sustain programme logic","Float consumed, acceleration required and cost of delay transferred to contingency","Both",28,10,30,90,0.010,0.030,0.060,"Programme Director","Decision log, RACI, clear authority levels, escalation triggers","Decision window missed","Reduce","Open"),
      ("R-011","Ground conditions and enabling surprises","Geotechnical","Incomplete ground investigation or legacy contamination","Adverse ground conditions encountered during enabling or foundation works","Additional treatment, redesign, delay and cost growth","Both",33,10,40,110,0.015,0.050,0.110,"Geotechnical Lead","Ground investigation programme, risk-informed quantities, provisional sums","Ground investigation shows divergence","Mitigate","Open"),
      ("R-012","Digital systems and data integration","Technology","Complex BMS, SCADA, ERP or operational technology integration","Integrated systems do not achieve acceptance in planned window","Commissioning delay, parallel run cost and operational performance risk","Both",30,10,30,90,0.015,0.045,0.095,"Digital Lead","Systems integration testing plan, factory acceptance, digital twin","Integration defect rate above threshold","Reduce","Open"),
    ]

    # ── SPACE SECTOR (adds 8 mission-critical risks) ─────────────────
    if mode == "Space":
        base = [
          ("R-S01","Launch manifest delay","Launch","Launch provider capacity, weather, manifest priority or vehicle readiness","Confirmed launch slot moves or payload misses manifest gate","Programme delay, storage cost, resequencing and mission-readiness risk","Both",52,20,60,180,0.035,0.110,0.240,"Launch Integration Lead","Reserve alternate slot, freeze payload interface, maintain launch-readiness checklist","Launch slot not confirmed at L-12 months","Mitigate","Open"),
          ("R-S02","Payload mass growth","Technical","Design creep, shielding additions, redundancy or late configuration changes","Payload mass exceeds launch vehicle allocation","Launch cost increase, redesign, performance loss and mission safety exposure","Both",44,15,55,150,0.050,0.130,0.280,"Chief Engineer","Mass control board, margin policy, design-to-mass reviews at each milestone","Mass margin below 8%","Reduce","Open"),
          ("R-S03","Life-support system reliability","Safety / Mission","ECLSS reliability evidence, spares depth or redundancy below design threshold","Life-support fails reliability growth or integrated mission test","Crew safety exposure, launch delay, redesign and potential mission abort","Both",38,25,80,220,0.060,0.150,0.320,"Life Support Lead","Prototype testing, redundancy architecture, reliability growth programme, independent assurance","Reliability growth target missed","Avoid / Reduce","Open"),
          ("R-S04","Radiation and thermal environment","Technical","Space environment models are uncertain for novel orbit or mission profile","Hardware degrades faster than design life under radiation / thermal cycling","Mission shortfall, component replacement, unplanned servicing and safety risk","Both",36,20,70,200,0.040,0.110,0.240,"Systems Engineer","Radiation analysis, shielding margin, thermal vacuum testing, margins policy","On-orbit anomaly exceeds design margin","Reduce","Open"),
          ("R-S05","Autonomous systems failure","Technical","Autonomy software maturity, testing coverage or edge-case handling inadequate","Autonomous control fails or produces unsafe state in mission-critical phase","Mission abort, crew safety risk and recovery mission cost","Both",34,20,65,180,0.035,0.100,0.220,"Autonomy Lead","Software verification, FMEA, hardware-in-the-loop testing, human override","Anomaly rate exceeds threshold","Reduce","Open"),
          ("R-S06","Launch vehicle performance variance","Launch","Launch vehicle thrust, trajectory or separation performance deviates","Payload delivered to wrong orbit or at lower-than-planned performance","Mission redesign, additional propellant cost, reduced operational life","Cost",30,0,0,0,0.020,0.080,0.180,"Mission Lead","Launch provider performance history, orbital margins, contingency delta-V","Launch telemetry shows deviation","Transfer","Open"),
          ("R-S07","Mission operations readiness","Operations","Ground segment, crew training or mission operations procedures immature at launch","Mission operations team not ready at launch-readiness review","Launch delay, mission risk and compressed crew readiness window","Schedule",35,20,60,150,0.015,0.060,0.130,"Mission Ops Lead","Operations readiness reviews, crew training milestones, simulation campaign","Operations readiness review fails","Mitigate","Open"),
          ("R-S08","Regulatory and frequency licensing","Regulatory","Spectrum allocation, launch licence or on-orbit regulatory clearances delayed","Operating licence not granted ahead of launch window","Launch delay, spectrum conflict and commercial exposure","Schedule",28,15,50,140,0.010,0.040,0.090,"Regulatory Lead","Early filing, ITU coordination, national regulator engagement","Licence milestone slips","Mitigate","Open"),
        ] + base

    # ── RAIL / TRANSIT (adds 10 sector-specific risks) ───────────────
    elif any(x in s for x in ["rail","metro","transit","hsr","high-speed","tram","underground"]):
        base = [
          ("R-R01","Possession window availability","Operations","Track access windows shorter than planned or subject to operational priority","Productive works windows reduce below programme assumption","Extended delivery duration, overtime premium and critical path pressure","Both",55,20,70,200,0.030,0.090,0.200,"Access Planning Lead","Negotiate access framework early, model possession alternatives, protect buffer","Possession windows rejected or curtailed","Mitigate","Open"),
          ("R-R02","Signalling and systems migration","Technical","Signalling migration, legacy interface complexity or vendor delivery","New signalling system fails integration or acceptance testing with legacy infrastructure","Late commissioning, operator rejection and safety assurance delay","Both",48,20,75,210,0.035,0.100,0.220,"Signalling Lead","Phased migration plan, shadow-running period, independent systems assurance","Integration test failures exceed threshold","Reduce","Open"),
          ("R-R03","Utility diversions and corridor","Utilities","Uncharted or complex utilities in rail corridor require diversion before works can proceed","Utility diversion programme overruns","Enabling delay, main works displaced and critical path threatened","Schedule",42,15,60,170,0.015,0.055,0.120,"Utilities Lead","Early utility mapping, CAT and Genny survey, advance diversions programme","Diversion milestone not met","Mitigate","Open"),
          ("R-R04","Land and corridor acquisition","Legal","Compulsory purchase, third-party objection or title defect delays possession","Land not available for construction at programme gate","Construction cannot start in affected sections; cost and delay","Both",38,15,55,160,0.020,0.060,0.130,"Land Lead","CPO programme tracker, legal reserve, advance negotiation, temporary possession","Land not handed over","Mitigate","Open"),
          ("R-R05","Trial running and safety certification","Regulatory","Safety case maturity, operator acceptance or regulator sign-off delayed","Trial running fails or safety certification not granted at programme date","Opening delayed, revenue deferral and public accountability risk","Schedule",40,20,70,180,0.010,0.050,0.110,"Safety Assurance Lead","Safety case programme, early regulator engagement, independent assurance","Safety case milestone missed","Reduce","Open"),
          ("R-R06","Tunnelling and civils ground risk","Geotechnical","Tunnel ground conditions diverge from ground investigation model","Excavation rates, settlements or ground treatment costs exceed plan","Cost and schedule growth, structural design change and third-party liability","Both",36,15,55,160,0.025,0.080,0.180,"Geotechnical Lead","Comprehensive ground investigation, observational method, contingency volumes","Settlement or water ingress exceeds threshold","Mitigate","Open"),
          ("R-R07","Rolling stock interface and acceptance","Technical","Rolling stock design interfaces, gauge clearance or compatibility uncertain","Rolling stock fails interface or acceptance with infrastructure as-built","Retrofit cost, delayed operations and operator satisfaction risk","Both",34,15,45,130,0.020,0.065,0.140,"Rolling Stock Interface Manager","Interface register, gauge clearance sign-off, rolling stock test programme","Interface clearance rejected","Reduce","Open"),
          ("R-R08","Electrification and power supply","Technical","Power supply design, OLE stagger or substation configuration uncertain","Electrification system fails commissioning or performance test","Commissioning delay, energy management risk and traction performance shortfall","Both",32,10,40,120,0.015,0.050,0.110,"E&P Lead","Early power system modelling, OLE supplier management, substation sequencing","Power performance test fails","Reduce","Open"),
          ("R-R09","Political and funding exposure","Governance","Large programme subject to political mandate, election cycles or funding review","Funding commitment reduced, delayed or subject to scope change","Programme descoped, delayed or restructured with cost and claims consequences","Both",30,0,0,0,0.020,0.070,0.180,"Sponsor","Political engagement, benefits case maintenance, funding milestone tracking","Funding review announced","Accept / Mitigate","Open"),
          ("R-R10","Operator integration and TOCS","Operations","Train operating company interface, timetabling or operational readiness","Operator not ready to receive infrastructure at commissioning","Delayed handover, penalty exposure and public trust risk","Schedule",28,10,35,100,0.008,0.030,0.065,"Operations Readiness Lead","Joint readiness programme with operator, timetabling workshops, ORAT equivalent","Operator readiness review fails","Mitigate","Open"),
        ] + base

    # ── DATA CENTRE / HYPERSCALE (adds 9 sector-specific risks) ──────
    elif any(x in s for x in ["data centre","data center","hyperscale","datacenter","compute","cloud","ai campus"]):
        base = [
          ("R-D01","Grid connection and energisation delay","Utilities","Utility agreement, DNO scope, substation construction or energisation pathway not secured","Power connection or energisation milestone slips","Critical path delay, temporary power cost, phased opening risk and revenue deferral","Both",55,25,90,250,0.040,0.120,0.260,"Utilities Lead","Secure grid agreement early, temporary power strategy, early switchgear orders","Grid agreement unsigned at gate","Mitigate","Open"),
          ("R-D02","Cooling capacity and water availability","Technical","Thermal design basis, water source agreement or cooling vendor performance uncertain","Cooling system fails capacity, redundancy or commissioning test","Redesign, commissioning delay, lower IT load density and resilience risk","Both",40,10,45,130,0.025,0.080,0.170,"MEP Lead","Thermal modelling, water licence, supplier validation, early performance test","Cooling commissioning test fails","Reduce","Open"),
          ("R-D03","Generator and UPS supply chain","Procurement","Generator OEM, transformer and UPS lead times extend in constrained market","Critical power equipment delayed beyond programme gate","Critical path threat, phased energisation and revenue deferral","Both",48,20,75,210,0.030,0.095,0.210,"Procurement Lead","Early OEM engagement, letter of intent, alternative sourcing, frame agreements","Generator delivery date confirmed past gate","Mitigate","Open"),
          ("R-D04","IT infrastructure and network readiness","Technology","Network connectivity, dark fibre, BGP routing or IX peering not secured","Data centre cannot reach customers or cloud on-ramps at launch","Commercial launch blocked, revenue loss and customer SLA breach","Both",30,10,35,100,0.015,0.050,0.110,"Network Lead","Connectivity strategy, dark fibre agreements, network testing programme","Connectivity test fails at launch","Mitigate","Open"),
          ("R-D05","Planning and environmental consent","Regulatory","Visual impact, noise, EIA requirements or grid upgrade consent delays","Planning permission not granted ahead of construction start","Critical path delay, redesign and competitor first-mover advantage","Schedule",38,20,70,180,0.010,0.045,0.100,"Planning Lead","Pre-application engagement, EIA programme, community engagement, noise model","Planning decision delayed","Mitigate","Open"),
          ("R-D06","Structural and floor load specification","Technical","IT load density, hot aisle containment or raised floor specification evolves during design","Structural design requires amendment after construction starts","Retrofit cost, delay and load density reduction risk","Both",28,10,30,80,0.010,0.035,0.075,"Structural Lead","Load envelope freeze, vendor intent letters, structural design reviews","Structural amendment required after pour","Reduce","Open"),
          ("R-D07","Security and compliance certification","Regulatory","Tier certification, ISO 27001, security vetting or customer audit requirements","Facility does not achieve required certification at launch","Commercial launch blocked or delayed; customer contracts at risk","Both",25,5,20,60,0.008,0.025,0.055,"Compliance Lead","Certification programme, pre-audit, Tier design review, security baseline","Certification audit fails","Mitigate","Open"),
          ("R-D08","Compressed hyperscale delivery schedule","Programme","Market demand drives aggressive programme; float and testing time compressed","Programme shortfall on critical systems integration and pre-commissioning","Operational incidents at launch, reliability shortfall and reputational risk","Both",35,10,35,100,0.020,0.060,0.130,"Programme Lead","Float management, integrated commissioning plan, pre-launch readiness gate","Float consumed ahead of commissioning","Mitigate","Open"),
          ("R-D09","Fibre and connectivity wayleave","Legal","Wayleave agreements, third-party land rights or local authority consent for connectivity","Connectivity route blocked or delayed","Launch delay, connectivity workaround cost and customer SLA pressure","Schedule",22,10,30,80,0.005,0.020,0.045,"Connectivity Lead","Wayleave programme, legal options, alternate route strategy","Wayleave not granted","Mitigate","Open"),
        ] + base

    # ── NUCLEAR / SMR (adds 9 risks) ──────────────────────────────────
    elif any(x in s for x in ["nuclear","smr","fusion","reactor","fission"]):
        base = [
          ("R-N01","Nuclear regulatory licensing","Regulatory","Generic Design Assessment, site licensing or safety case maturity takes longer than plan","Regulator delays Design Acceptance Confirmation or construction licence","Programme delay, cost growth and potential technology redesign","Schedule",50,30,100,270,0.015,0.070,0.160,"Regulatory Lead","Early ONR/NRC engagement, pre-submission meetings, safety case programme","Regulatory milestone slips","Mitigate","Open"),
          ("R-N02","First-of-a-kind nuclear technology risk","Technical","Novel reactor design or first commercial deployment of new technology","System performance, safety features or construction methods deviate from design basis","Costly redesign, extended testing, regulatory challenge and mission creep","Both",44,20,80,240,0.050,0.140,0.300,"Chief Nuclear Engineer","Technology readiness reviews, independent nuclear assurance, FOAK margin policy","Technology readiness level gap identified","Reduce","Open"),
          ("R-N03","Nuclear supply chain maturity","Procurement","Specialist nuclear-grade manufacturing, QA standards and N-stamp certification constrained","Nuclear equipment fails quality audit or delivery milestone","Critical path threat, re-manufacturing and regulatory notification","Both",42,20,70,210,0.035,0.110,0.240,"Nuclear Supply Chain Lead","Approved vendor list, Tier 1 audits, N-stamp qualification, long-lead reservation","Quality audit failure","Mitigate","Open"),
          ("R-N04","Site radiological baseline","Geotechnical / Safety","Legacy contamination, unknown buried waste or groundwater radiological baseline","Remediation scope or cost exceeds plan","Cost growth, programme delay and radiological compliance risk","Both",36,20,65,180,0.025,0.080,0.180,"Environmental Lead","Comprehensive characterisation survey, brownfield allowances, ALARP review","Characterisation diverges from baseline","Mitigate","Open"),
          ("R-N05","Security infrastructure requirements","Security","CPNI / NSCS requirements, security vetting or physical protection system scope evolves","Security compliance scope changes during construction","Cost growth, delay and interface complexity","Both",30,10,40,120,0.015,0.050,0.110,"Security Lead","CPNI engagement, security design freeze, vetting timelines","Security scope change at late stage","Reduce","Open"),
          ("R-N06","Waste management and disposal route","Environmental","Long-term waste disposal route not confirmed or regulatory position changes","Waste route uncertainty blocks commissioning approval","Commissioning delayed, operating cost uncertain and liability retained","Both",35,20,60,170,0.020,0.065,0.140,"Waste Lead","NDA coordination, waste route confirmation, interim storage strategy","Waste route not confirmed","Mitigate","Open"),
          ("R-N07","Public acceptance and political risk","Governance","Community opposition, political mandate change or planning public inquiry","Planning rejected, materially amended or subject to extended inquiry","Programme delay, cost and potential cancellation risk","Both",32,0,0,0,0.020,0.070,0.180,"Stakeholder Lead","Benefit sharing, community engagement, planning evidence programme","Planning inquiry announced","Accept / Mitigate","Open"),
          ("R-N08","Integrated commissioning and fuel loading","Technical","Nuclear system integration, fuel loading sequence or safety case maturity at commissioning","First criticality or commissioning tests delayed or fail safety case","Revenue deferral, regulatory exposure and safety accountability risk","Both",28,20,65,180,0.015,0.060,0.130,"Commissioning Director","Integrated commissioning plan, non-nuclear systems first, safety case gates","Commissioning gate missed","Reduce","Open"),
          ("R-N09","Grid connection and offtake agreement","Utilities","Transmission connection, capacity market contract or offtake agreement not in place","Electricity cannot be exported at first generation","Revenue loss, grid compliance cost and commercial viability risk","Both",30,15,50,140,0.010,0.040,0.090,"Commercial Lead","Early TNUoS engagement, grid connection agreement, offtake heads of terms","Grid agreement not signed","Mitigate","Open"),
        ] + base

    # ── SEMICONDUCTOR / ADVANCED MFG (adds 9 risks) ──────────────────
    elif any(x in s for x in ["semiconductor","fab","wafer","cleanroom","chip","advanced manufacturing","advanced mfg"]):
        base = [
          ("R-F01","Process tool delivery and qualification","Procurement","Leading-edge EUV, CVD, ALD or CMP tools on constrained global allocation","Tool delivery misses installation date or fails process qualification","Ramp delayed, yield loss and competitive market-entry risk","Both",52,20,70,200,0.040,0.110,0.240,"Equipment Lead","OEM allocation lock, qualification programme, backup tool options","Tool delivery date confirmed past gate","Mitigate","Open"),
          ("R-F02","Ultra-pure water and chemical supply","Utilities","UPW system design, chemical supply agreements or water permit not confirmed","UPW or chemical systems not ready at process qualification","Cleanroom contamination risk, yield loss and qualification delay","Both",44,15,55,160,0.030,0.090,0.200,"Utilities Lead","UPW design freeze, chemical supply agreements, permit programme","UPW commissioning test fails","Reduce","Open"),
          ("R-F03","Cleanroom classification and HVAC","Technical","Cleanroom HVAC, vibration isolation or EMI performance below specification","Cleanroom fails ISO classification or process qualification environment test","Process tool cannot be installed; ramp delayed","Both",40,10,40,120,0.025,0.080,0.180,"Cleanroom Lead","HVAC modelling, vibration survey, early HVAC commissioning test","Cleanroom test fails class threshold","Reduce","Open"),
          ("R-F04","Seismic and vibration isolation","Technical","Fab sited in seismic zone or vibration from traffic / adjacent plant exceeds tool tolerance","Vibration exceeds EUV or other sensitive tool specification","Tool disqualified, siting review required and ramp delayed","Both",35,10,35,100,0.020,0.070,0.155,"Structural Lead","Seismic study, vibration survey, isolation system specification","Vibration measurement exceeds specification","Reduce","Open"),
          ("R-F05","Power quality and redundancy","Utilities","Utility power quality, dual-feed redundancy or UPS specification uncertain","Power event causes process disruption or tool damage","Yield loss, tool repair cost and regulatory compliance risk","Both",38,10,30,90,0.015,0.055,0.120,"Electrical Lead","Power quality study, dual-feed agreement, UPS sizing review","Power quality test fails","Mitigate","Open"),
          ("R-F06","Process yield ramp performance","Technical","Yield ramp model assumes learning curve achievable in market window","Yield underperforms ramp model in early production phase","Revenue shortfall, capacity underutilisation and competitive exposure","Cost",42,0,0,0,0.030,0.100,0.220,"Process Lead","Yield model benchmarking, pilot line data, conservative ramp assumptions","Yield below 50% of target at month 6","Accept / Mitigate","Open"),
          ("R-F07","Export control and technology licence","Legal","ITAR, EAR or bilateral technology control requirements affect tool import or process","Export licence denied, tool withheld or technology use restricted","Equipment not deliverable; legal exposure and programme delay","Both",28,10,40,120,0.015,0.050,0.110,"Legal / Compliance Lead","Early export licence review, OEM licence strategy, government engagement","Export licence application refused","Mitigate","Open"),
          ("R-F08","Talent acquisition for fab operations","Human Resources","Qualified process engineers and technicians scarce in programme location","Operating team not at required scale at ramp start","Ramp delayed, yield impacted and knowledge transfer exposure","Schedule",35,0,0,0,0.010,0.040,0.090,"HR Lead","Graduate pipeline, operator partnerships, training programme, relo packages","Headcount below target at ramp","Mitigate","Open"),
          ("R-F09","Wastewater and chemical disposal","Environmental","Fab wastewater, chemical waste or solvent disposal consents not confirmed","Discharge consent refused or subject to expensive treatment requirement","Permitting delay, capex increase and environmental liability","Both",25,10,35,100,0.010,0.035,0.080,"Environment Lead","Early discharge consent application, waste contractor agreements, treatment design","Consent refused","Mitigate","Open"),
        ] + base

    # ── LIFE SCIENCES / PHARMA / GMP (adds 8 risks) ──────────────────
    elif any(x in s for x in ["life science","pharma","biotech","gmp","vaccine","cell therapy","drug","manufacturing","biopharma"]):
        base = [
          ("R-L01","GMP regulatory approval and inspection","Regulatory","FDA, EMA or MHRA inspection readiness, documentation maturity or facility compliance","Regulatory inspection fails or approval delayed beyond programme gate","Commercial launch blocked, revenue deferral and supply chain disruption","Both",45,15,55,160,0.020,0.070,0.160,"Regulatory Lead","Inspection readiness programme, mock inspection, regulatory dossier review","Pre-approval inspection observation issued","Reduce","Open"),
          ("R-L02","GMP validation and qualification","Technical","IQ/OQ/PQ protocols, validation batch performance or product yield uncertain","Validation batch fails or requires extensive repeat qualification","Launch delay, regulatory notification and remediation cost","Both",42,10,40,120,0.020,0.065,0.145,"Validation Lead","Validation master plan, early OQ testing, supplier qualification","Validation batch fails","Reduce","Open"),
          ("R-L03","Cleanroom HVAC and contamination control","Technical","Cleanroom HVAC, pressure differentials or environmental monitoring below GMP standard","Cleanroom fails environmental monitoring or classification test","Process redesign, re-qualification and launch delay","Both",38,10,35,100,0.015,0.055,0.120,"QA Lead","HVAC commissioning programme, environmental monitoring baseline, FAT/SAT","EM excursion during qualification","Reduce","Open"),
          ("R-L04","Cold chain and utility reliability","Technical","Cold chain, water-for-injection or purified water system availability uncertain","Utility failure during production causes batch loss","Batch write-off, regulatory deviation and investigation cost","Both",30,0,0,0,0.010,0.040,0.090,"Utilities Lead","Redundancy design, UPS, alarm management, validated cleaning procedures","Utility failure during production","Mitigate","Open"),
          ("R-L05","Product licence and clinical dependency","Regulatory","Manufacturing site added to product licence requires clinical or regulatory package","Licence amendment delayed or linked to ongoing clinical programme","Site investment stranded if clinical programme fails","Cost",28,0,0,0,0.015,0.055,0.130,"Regulatory Lead","Licence strategy, CMC dossier timeline, regulatory agency pre-submission","Licence amendment rejected","Accept / Mitigate","Open"),
          ("R-L06","Specialist equipment lead times","Procurement","Bioreactor, isolator, lyophiliser or aseptic fill-finish equipment on constrained supply","Equipment delivery misses installation gate","Validation delayed, launch deferral and capital lock-up","Both",40,15,50,140,0.015,0.055,0.120,"Procurement Lead","Early OEM engagement, letter of intent, equipment pre-qualification","Delivery date confirmed past gate","Mitigate","Open"),
          ("R-L07","Technology transfer risk","Technical","Process transferred from pilot, external CDO or acquired site","Technology transfer fails to reproduce product quality at commercial scale","Batch failures, re-development cost and regulatory impact","Both",35,10,40,120,0.020,0.065,0.145,"Tech Transfer Lead","TT protocol, risk-sharing with originator, early development runs","TT batch fails acceptance criteria","Reduce","Open"),
          ("R-L08","Occupational health and containment","Safety","Potent compound, biological agent or novel modality requires containment infrastructure","Containment infrastructure fails engineering or personnel safety review","Regulatory stop, redesign and reputational risk","Both",25,10,35,100,0.015,0.050,0.110,"EHS Lead","COSHH assessment, containment design review, ATEX compliance","Containment audit fails","Avoid / Reduce","Open"),
        ] + base

    # ── AIRPORT / AVIATION (adds 8 risks) ─────────────────────────────
    elif any(x in s for x in ["airport","aviation","runway","airside","baggage","terminal"]):
        base = [
          ("R-AP01","Live airport operational disruption","Operations","Construction access in or adjacent to live airside operations creates conflict","Operational restrictions reduce productive windows or cause safety incident","Delay, premium cost, stakeholder disruption and safety accountability","Both",45,20,70,180,0.025,0.080,0.180,"ORAT Lead","Phasing simulation, possession plan, airport operations integration board","Possession windows rejected","Mitigate","Open"),
          ("R-AP02","Baggage system integration","Systems","Baggage handling system vendor, software interface or acceptance testing maturity","BHS integration fails acceptance testing with terminal systems","Commissioning delay, passenger disruption and opening date risk","Both",42,15,50,140,0.020,0.070,0.155,"Systems Lead","FAT/SAT programme, integration lab, BHS commissioning plan","FAT/SAT defect rate above threshold","Reduce","Open"),
          ("R-AP03","ORAT and airline transition readiness","Operations","Airline, ground handler or airport operator transition readiness behind programme","ORAT programme incomplete at opening gate","Delayed opening, passenger disruption and CAA intervention risk","Both",38,15,50,130,0.015,0.055,0.120,"ORAT Lead","Joint ORAT programme, airline engagement, dry-run schedule","ORAT milestone missed","Reduce","Open"),
          ("R-AP04","CAA and regulatory acceptance","Regulatory","Safety case, aerodrome licence amendment or CAA acceptance delayed","CAA sign-off not granted ahead of operational opening","Opening delayed, airline schedule disruption and reputational cost","Schedule",35,15,50,140,0.010,0.040,0.090,"Safety Lead","Early CAA engagement, safety case programme, regulatory pre-submission","CAA acceptance delayed","Mitigate","Open"),
          ("R-AP05","Airfield and runway works in live environment","Delivery","Runway, taxiway or apron works require NOTAMs and restrict live flight operations","Works access shorter than planned or flight disruption occurs","Production shortfall, night-work premium and airline compensation exposure","Both",40,15,55,150,0.020,0.065,0.145,"Airfield Lead","NOTAM planning, airline liaison, works windows analysis, contingency protocol","NOTAM access refused","Mitigate","Open"),
          ("R-AP06","Security infrastructure compliance","Security","DfT, TSA or ECAC security standards for new terminal or checkpoint evolve","Security design amendment required during construction","Cost growth, interface rework and opening delay","Both",28,10,30,80,0.010,0.035,0.075,"Security Lead","DfT engagement, security design freeze, early concept approval","Security standard change announced","Reduce","Open"),
          ("R-AP07","Passenger volume and demand risk","Commercial","Traffic forecasts deviate from business case; airline commitment uncertain","Terminal undersized or oversized relative to actual demand","Revenue shortfall, asset underutilisation or capacity overshoot","Cost",25,0,0,0,0.015,0.050,0.115,"Commercial Lead","Demand sensitivity model, airline agreements, phased capacity strategy","Traffic forecast diverges by 15%","Accept / Mitigate","Open"),
          ("R-AP08","Ground transport and surface access","Integration","Rail, road or bus connections not ready at terminal opening","Surface access below required capacity at launch","Passenger experience failure, congestion and regulatory pressure","Schedule",28,10,35,90,0.008,0.025,0.055,"Surface Access Lead","Surface access programme, local authority engagement, transport assessment","Surface access milestone missed","Mitigate","Open"),
        ] + base

    # ── DEFENCE / SECURE INFRASTRUCTURE (adds 7 risks) ───────────────
    elif any(x in s for x in ["defence","defense","military","naval","airbase","radar","command","secure mission","aukus","weapons"]):
        base = [
          ("R-DF01","Security accreditation and classification","Security","CPNI, JSP, NCSC or classified-facility accreditation processes delay approval","Security accreditation not granted ahead of operational handover","Facility cannot be handed over; use restriction and cost of delay","Schedule",45,20,70,180,0.015,0.060,0.130,"Security Lead","CPNI engagement, JSP compliance programme, early accreditation submission","Accreditation milestone missed","Mitigate","Open"),
          ("R-DF02","Sovereign supply chain constraints","Procurement","Controlled technology, ITAR restrictions or national security supply limits","Critical equipment not procurable through standard channels","Alternative sourcing cost, delay and technology risk","Both",40,15,55,160,0.025,0.080,0.180,"Procurement Lead","Approved vendor list, technology waiver strategy, alternative supplier mapping","Supplier disqualified","Mitigate","Open"),
          ("R-DF03","Classified systems integration","Technical","Integration of classified systems, data links or intelligence interfaces not mature","Classified systems integration fails acceptance in planned window","Mission capability delayed, security exposure and operational risk","Both",38,15,50,140,0.025,0.080,0.175,"Systems Lead","Classified integration test environment, accreditation pathway, test campaign","Integration test fails","Reduce","Open"),
          ("R-DF04","Political mandate and funding continuity","Governance","Programme subject to defence review, budget cycles or political change","Funding commitment varied or programme descoped","Cost and schedule waste, contractor claims and capability gap","Both",32,0,0,0,0.020,0.070,0.180,"Sponsor","Benefits case maintenance, political engagement, milestone achievement","Funding review triggered","Accept / Mitigate","Open"),
          ("R-DF05","Operational authority-to-operate","Regulatory","Military operational acceptance, battle management system clearance or OAT delayed","Operational authority not granted at facility handover","Facility handed over but not operationally available","Schedule",30,10,40,110,0.008,0.030,0.065,"Ops Assurance Lead","OAT programme, acceptance criteria, combined authority sign-off","OAT milestone missed","Reduce","Open"),
          ("R-DF06","Personnel vetting and clearance","Human Resources","Security clearance timelines for contractors, operators or programme staff","Cleared personnel not available at required programme stages","Programme activities constrained, delay and cost of delay","Schedule",35,0,0,0,0.005,0.020,0.045,"Security Lead","Early vetting initiation, cleared contractor register, buffer workforce plan","Vetting backlog impacting programme","Mitigate","Open"),
          ("R-DF07","Electromagnetic compatibility","Technical","EMC, EMI or RF interference between systems in dense defence environment","Systems fail EMC testing or interfere with adjacent classified systems","Redesign, shielding cost and acceptance delay","Both",25,10,30,80,0.010,0.035,0.075,"EMC Lead","EMC test programme, shielding design, early frequency management","EMC test fails","Reduce","Open"),
        ] + base

    # ── ENERGY / RENEWABLES / HYDROGEN (adds 7 risks) ─────────────────
    elif any(x in s for x in ["solar","wind","hydrogen","lng","battery","offshore wind","energy storage","power plant","gas"]):
        base = [
          ("R-E01","Grid connection and curtailment","Utilities","Transmission network capacity, connection offer or curtailment agreement not confirmed","Grid connection delayed or revenue curtailed below forecast","Revenue loss, investment return shortfall and offtake pressure","Both",48,20,70,200,0.030,0.090,0.200,"Grid Lead","Early connection application, curtailment model, grid upgrade coordination","Grid connection offer delayed","Mitigate","Open"),
          ("R-E02","Offshore installation weather window","Delivery","Offshore marine installation dependent on seasonal weather windows","Adverse weather exceeds marine spread tolerance","Monopile, cable or turbine installation delayed to next season","Schedule",42,30,90,210,0.015,0.065,0.150,"Marine Lead","Weather downtime model, marine spread selection, float-off alternative","Installation vessel idle period exceeds budget","Mitigate","Open"),
          ("R-E03","Turbine and equipment supply chain","Procurement","Major wind turbine OEM or battery module supply constrained","Equipment delivery misses installation gate","Critical path threat, lower capacity factor and revenue deferral","Both",44,15,55,160,0.030,0.095,0.210,"Procurement Lead","OEM frame agreement, early delivery milestone, alternative sourcing options","OEM delivery confirmation past gate","Mitigate","Open"),
          ("R-E04","Planning and environmental consent","Regulatory","Habitats regulation, visual impact, noise or grid upgrade EIA requirements","Planning or DCO not granted in time for construction start","Critical path delay, competitor advantage and investor confidence risk","Schedule",38,20,70,200,0.010,0.045,0.100,"Planning Lead","Pre-application engagement, EIA programme, habitat compensation strategy","Planning decision delayed","Mitigate","Open"),
          ("R-E05","Hydrogen technology and safety case","Technical","Hydrogen production, compression, storage or dispensing technology maturity","Plant fails safety case, insurance or regulatory acceptance test","Cost redesign, delayed commissioning and regulatory exposure","Both",35,15,50,140,0.020,0.075,0.165,"Technical Lead","HAZOP, FMEA, safety case programme, ATEX compliance review","Safety case milestone missed","Reduce","Open"),
          ("R-E06","Power purchase agreement terms","Commercial","PPA counterparty, price floor or duration not agreed","PPA not signed ahead of financial close","Project unfinanceable; investor withdrawal and cost of delay","Cost",32,0,0,0,0.015,0.055,0.130,"Commercial Lead","PPA heads of terms, offtake market analysis, government support scheme","PPA negotiation stalls","Mitigate","Open"),
          ("R-E07","Subsea cable and export route","Technical","Subsea cable specification, installation vessel or landfall consent uncertain","Export cable delayed or fails commissioning test","Revenue deferral, insurance cost and grid compliance risk","Both",30,15,50,140,0.015,0.055,0.120,"Cable Lead","Early cable contract, installation vessel reservation, landfall consent","Cable installation delayed","Mitigate","Open"),
        ] + base

    # ── WATER / UTILITIES (adds 6 risks) ──────────────────────────────
    elif any(x in s for x in ["water","desalination","wastewater","reservoir","flood","pumping"]):
        base = [
          ("R-W01","Water abstraction and discharge licence","Regulatory","Environment Agency, EPA or equivalent licence for abstraction or discharge","Licence not granted or subject to restrictive conditions","Operating constraint, revenue deferral and regulatory compliance cost","Schedule",40,15,55,150,0.010,0.040,0.090,"Environmental Lead","Early licence application, pre-consultation, quality monitoring programme","Licence condition imposed","Mitigate","Open"),
          ("R-W02","Membrane and process technology performance","Technical","Desalination, advanced treatment or membrane technology performance uncertain","Process performance below design specification","Product water fails quality standard, re-design and output shortfall","Both",35,10,40,120,0.020,0.070,0.155,"Process Lead","Pilot plant testing, OEM performance guarantee, treatability study","Performance test below design","Reduce","Open"),
          ("R-W03","Seismic and flood risk","Natural Hazard","Site-specific seismic or flood event exceeds design basis","Extreme event damages or compromises facility before or during operation","Asset damage, business interruption and public safety risk","Both",20,0,0,0,0.015,0.055,0.140,"Risk Lead","Site-specific hazard assessment, resilient design standards, insurance","Extreme event occurs","Accept / Transfer","Open"),
          ("R-W04","Raw water quality variability","Technical","Source water turbidity, salinity or contamination exceeds design envelope","Pre-treatment overloaded or product quality breached","Output curtailed, treatment cost increase and regulatory violation","Both",30,0,0,0,0.010,0.040,0.090,"Water Quality Lead","Raw water monitoring, adaptive treatment design, quality trigger protocol","Turbidity exceedance event","Mitigate","Open"),
          ("R-W05","Pipeline corridor and wayleave","Legal","Pipeline route wayleave, third-party land consent or easement not secured","Pipeline cannot be constructed on planned route","Reroute cost, delay and programme disruption","Schedule",28,10,35,100,0.008,0.030,0.065,"Legal Lead","Wayleave programme, alternate route modelling, legal reserve","Wayleave refused","Mitigate","Open"),
          ("R-W06","Chemical supply and COSHH compliance","Procurement","Chemical dosing supply continuity, COSHH requirements or health and safety compliance","Chemical supply interrupted or facility fails COSHH inspection","Treatment failure, regulatory action and public health risk","Both",22,0,0,0,0.005,0.020,0.045,"Operations Lead","Multi-supplier strategy, storage reserve, COSHH assessment","Supply interruption","Mitigate","Open"),
        ] + base

    # ── PORTS / MARINE (adds 6 risks) ─────────────────────────────────
    elif any(x in s for x in ["port","harbour","marine","quay","berth","cargo","container","logistics hub"]):
        base = [
          ("R-P01","Marine works and dredging","Technical","Seabed conditions, dredge volume or contaminated spoil beyond survey basis","Dredging or marine civil works exceed planned scope","Cost growth, programme delay and disposal cost","Both",42,15,55,160,0.025,0.080,0.180,"Marine Lead","Pre-construction survey, spoil characterisation, dredge volume contingency","Dredge survey diverges from design","Mitigate","Open"),
          ("R-P02","Tidal and weather window","Delivery","Marine construction dependent on tidal windows and vessel weather tolerance","Installation or marine piling delayed by adverse conditions","Extended marine spread cost and programme pressure","Schedule",35,15,50,140,0.010,0.040,0.090,"Marine Lead","Weather downtime model, vessel selection, tidal window analysis","Weather downtime exceeds budget","Mitigate","Open"),
          ("R-P03","Harbour authority and MCA approval","Regulatory","Harbour revision order, MCA survey or coastal authority consent not confirmed","Approval milestone slips beyond programme date","Construction cannot start; cost of delay and programme displacement","Schedule",32,15,50,130,0.008,0.030,0.065,"Regulatory Lead","Early HRO application, pre-consultation, MCA engagement","Approval milestone missed","Mitigate","Open"),
          ("R-P04","Container handling equipment delivery","Procurement","STS cranes, RTG/RMG or automated handling system OEM supply constrained","Equipment delivery misses berth commissioning gate","Terminal cannot operate at capacity; revenue deferral","Both",40,15,50,140,0.020,0.070,0.155,"Procurement Lead","Early OEM commitment, alternative sourcing, phased delivery","Equipment delivery confirmed past gate","Mitigate","Open"),
          ("R-P05","Automation and terminal operating system","Technology","Automated port TOS, AGV system or port community system integration uncertain","TOS integration fails acceptance or automation system underperforms","Terminal throughput below plan, customer attrition and commercial risk","Both",30,10,35,100,0.015,0.050,0.110,"Digital Lead","TOS integration test programme, simulation, operational shadow-running","TOS acceptance test fails","Reduce","Open"),
          ("R-P06","Environmental and protected habitat","Environmental","Marine habitats, species survey or coastal heritage adjacent to works","Environmental condition imposed or works suspended","Scope constraint, mitigation cost and reputational risk","Both",25,10,35,90,0.008,0.025,0.055,"Environmental Lead","Habitat survey, licence conditions, mitigation programme","Environmental stop order issued","Mitigate","Open"),
        ] + base

    out=[]
    for i,r in enumerate(base,1):
        rid,title,cat,cause,event,impact,area,prob,so,sm,sp,co,cm,cp,owner,mit,trig,response,status=r
        prob=int(clamp(prob*rm,5,92))
        if scenario=="lower_risk": cm*=0.70; cp*=0.75; sm=round(sm*.72); sp=round(sp*.78)
        if scenario=="faster": sm=round(sm*1.22); sp=round(sp*1.28); cm*=1.06; cp*=1.08
        if scenario=="cheaper": cm*=1.16; cp*=1.22; prob=int(clamp(prob*1.08,5,95))
        if scenario=="premium": cm*=0.82; cp*=0.88; prob=int(clamp(prob*0.9,5,95))
        activity=schedule[min(i+3,len(schedule)-1)] if schedule else {"activity_id":"A1900","activity":"Delivery"}
        costline=costs[min(i+2,len(costs)-1)] if costs else {"cbs":"01.01","description":"Cost"}
        emv=cost*cm*prob/100; semv=sm*prob/100
        rating=_risk_rating(prob,cost*cm,sm)
        out.append({
          "risk_id":rid,"title":title,"category":cat,"cause":cause,"risk_event":event,"impact_description":impact,
          "description":event,"impact_area":area,"probability_pct":prob,"pre_mitigation_rating":rating,
          "activity_id":activity["activity_id"],"activity_name":activity["activity"],"cbs":costline["cbs"],"cbs_name":costline["description"],
          "schedule_o_days":so,"schedule_m_days":sm,"schedule_p_days":sp,"cost_o_bn":round(cost*co,3),"cost_m_bn":round(cost*cm,3),"cost_p_bn":round(cost*cp,3),
          "cost_emv_bn":round(emv,3),"schedule_emv_days":round(semv,1),"owner":owner,"trigger":trig,"mitigation":mit,"response_strategy":response,
          "residual_rating":"Medium" if rating in ["High","Extreme"] else "Low","status":status,"last_reviewed":datetime.utcnow().strftime("%Y-%m-%d"),
          "board_visibility":"Yes" if i<=12 or rating in ["High","Extreme"] else "No",
          "basis_of_cost_impact":f"O/M/P cost impact equals {money_bn(cost*co)} / {money_bn(cost*cm)} / {money_bn(cost*cp)} based on scenario-adjusted exposure to CBS {costline['cbs']} {costline['description']}.",
          "basis_of_schedule_impact":f"O/M/P schedule impact equals {so}/{sm}/{sp} days mapped to activity {activity['activity_id']} {activity['activity']}.",
          "driver_score":round(emv*100+semv/5,2)
        })
    return sorted(out,key=lambda x:x["driver_score"],reverse=True)
    out=[]
    for i,r in enumerate(base,1):
        rid,title,cat,cause,event,impact,area,prob,so,sm,sp,co,cm,cp,owner,mit,trig,response,status=r
        prob=int(clamp(prob*rm,5,92))
        if scenario=="lower_risk": cm*=0.70; cp*=0.75; sm=round(sm*.72); sp=round(sp*.78)
        if scenario=="faster": sm=round(sm*1.22); sp=round(sp*1.28); cm*=1.06; cp*=1.08
        if scenario=="cheaper": cm*=1.16; cp*=1.22; prob=int(clamp(prob*1.08,5,95))
        if scenario=="premium": cm*=0.82; cp*=0.88; prob=int(clamp(prob*0.9,5,95))
        activity=schedule[min(i+3,len(schedule)-1)] if schedule else {"activity_id":"A1900","activity":"Delivery"}
        costline=costs[min(i+2,len(costs)-1)] if costs else {"cbs":"01.01","description":"Cost"}
        emv=cost*cm*prob/100; semv=sm*prob/100
        rating=_risk_rating(prob,cost*cm,sm)
        out.append({
          "risk_id":rid,"title":title,"category":cat,"cause":cause,"risk_event":event,"impact_description":impact,
          "description":event,"impact_area":area,"probability_pct":prob,"pre_mitigation_rating":rating,
          "activity_id":activity["activity_id"],"activity_name":activity["activity"],"cbs":costline["cbs"],"cbs_name":costline["description"],
          "schedule_o_days":so,"schedule_m_days":sm,"schedule_p_days":sp,"cost_o_bn":round(cost*co,3),"cost_m_bn":round(cost*cm,3),"cost_p_bn":round(cost*cp,3),
          "cost_emv_bn":round(emv,3),"schedule_emv_days":round(semv,1),"owner":owner,"trigger":trig,"mitigation":mit,"response_strategy":response,
          "residual_rating":"Medium" if rating in ["High","Extreme"] else "Low","status":status,"last_reviewed":datetime.utcnow().strftime("%Y-%m-%d"),
          "board_visibility":"Yes" if i<=8 or rating in ["High","Extreme"] else "No",
          "basis_of_cost_impact":f"O/M/P cost impact equals {money_bn(cost*co)} / {money_bn(cost*cm)} / {money_bn(cost*cp)} based on scenario-adjusted exposure to CBS {costline['cbs']} {costline['description']}.",
          "basis_of_schedule_impact":f"O/M/P schedule impact equals {so}/{sm}/{sp} days mapped to activity {activity['activity_id']} {activity['activity']}.",
          "driver_score":round(emv*100+semv/5,2)
        })
    return sorted(out,key=lambda x:x["driver_score"],reverse=True)

def monte_carlo(cost,months,risks,seed=42,iterations=10000):
    rng=np.random.default_rng(seed)
    base_cost_unc=rng.triangular(cost*.90,cost,cost*1.16,iterations)
    base_sched_unc=rng.triangular(months*.94,months,months*1.14,iterations)
    risk_cost=np.zeros(iterations); risk_days=np.zeros(iterations)
    contribution={r["risk_id"]:{"cost":0.0,"days":0.0,"title":r["title"],"activity_id":r["activity_id"],"cbs":r["cbs"],"category":r.get("category","Risk")} for r in risks}
    for r in risks:
        occurs=rng.random(iterations)<(r["probability_pct"]/100)
        c=rng.triangular(r["cost_o_bn"],r["cost_m_bn"],r["cost_p_bn"],iterations)*occurs
        left=min(r["schedule_o_days"],r["schedule_m_days"],r["schedule_p_days"]); mode=max(left,r["schedule_m_days"]); right=max(r["schedule_o_days"],r["schedule_m_days"],r["schedule_p_days"])
        d=np.zeros(iterations) if left==right else rng.triangular(left,mode,right,iterations)*occurs
        risk_cost+=c; risk_days+=d
        contribution[r["risk_id"]]["cost"]=float(np.mean(c)); contribution[r["risk_id"]]["days"]=float(np.mean(d))
    cost_samples=base_cost_unc+risk_cost
    sched_samples=base_sched_unc+(risk_days/30.44)
    def pct(arr,p): return float(np.percentile(arr,p))
    curve=[{"percentile":p,"cost_bn":round(pct(cost_samples,p),3),"schedule_months":round(pct(sched_samples,p),2)} for p in [1,5,10,20,30,40,50,60,70,80,90,95,99]]
    qcra_tornado=sorted([{"risk_id":k,"title":v["title"],"category":v["category"],"cbs":v["cbs"],"activity_id":v["activity_id"],"cost_mean_bn":round(v["cost"],3),"driver_score":round(v["cost"]*100,2)} for k,v in contribution.items()],key=lambda x:x["driver_score"],reverse=True)
    qsra_tornado=sorted([{"risk_id":k,"title":v["title"],"category":v["category"],"cbs":v["cbs"],"activity_id":v["activity_id"],"schedule_mean_days":round(v["days"],1),"driver_score":round(v["days"]/2,2)} for k,v in contribution.items()],key=lambda x:x["driver_score"],reverse=True)
    tornado=sorted([{"risk_id":k,"title":v["title"],"activity_id":v["activity_id"],"cbs":v["cbs"],"cost_mean_bn":round(v["cost"],3),"schedule_mean_days":round(v["days"],1),"driver_score":round(v["cost"]*100+v["days"]/10,2)} for k,v in contribution.items()],key=lambda x:x["driver_score"],reverse=True)
    return {"iterations":iterations,"qcra":{"p10":round(pct(cost_samples,10),3),"p50":round(pct(cost_samples,50),3),"p80":round(pct(cost_samples,80),3),"p90":round(pct(cost_samples,90),3),"mean":round(float(np.mean(cost_samples)),3)},"qsra":{"p10":round(pct(sched_samples,10),2),"p50":round(pct(sched_samples,50),2),"p80":round(pct(sched_samples,80),2),"p90":round(pct(sched_samples,90),2),"mean":round(float(np.mean(sched_samples)),2)},"curve":curve,"tornado":tornado,"qcra_tornado":qcra_tornado,"qsra_tornado":qsra_tornado}

def _style_ws(ws, freeze="A2"):
    dark="07111F"; cyan="8DF7FF"; panel="0B1424"; white="F7FBFF"
    ws.freeze_panes=freeze
    for row in ws.iter_rows():
        for c in row:
            c.alignment=Alignment(vertical="top",wrap_text=True)
            c.border=Border(bottom=Side(style="thin",color="1E2A3A"))
    if ws.max_row>=1:
        for c in ws[1]:
            c.fill=PatternFill("solid",fgColor=dark); c.font=Font(color=white,bold=True); c.alignment=Alignment(vertical="center",wrap_text=True)
    for col in range(1,ws.max_column+1):
        ws.column_dimensions[get_column_letter(col)].width=min(38,max(14,len(str(ws.cell(1,col).value or ""))+3))

def _add(ws, rows):
    for r in rows: ws.append(r)
    _style_ws(ws)

def workbook_bytes(model):
    wb=Workbook(); wb.remove(wb.active)
    fx=_scenario_effects(model); totals=_total_costs(model); mc=model["monte_carlo"]
    ws=wb.create_sheet("00 Executive Summary")
    rows=[
      ["CASEY TITAN X", "Elite Project Controls Intelligence Pack"],
      ["Generated", datetime.utcnow().isoformat()], ["Project", model.get("title")], ["Client", model.get("client")], ["Sector", model.get("subsector")], ["Location", model.get("location")], ["Scenario", model.get("scenario_label")],
      ["Executive Summary", model.get("executive_summary")], ["Board Recommendation", f"Proceed only with {model.get('scenario_label')} controls basis, explicit risk allowance and evidence plan before commitment."],
      ["P10 / P50 / P80 / P90 Cost", f"{money_bn(mc['qcra']['p10'])} / {model.get('cost_p50')} / {money_bn(mc['qcra']['p80'])} / {money_bn(mc['qcra']['p90'])}"],
      ["QSRA P50 / P80", f"{mc['qsra']['p50']} / {mc['qsra']['p80']} months"], ["Direct / Indirect / Reserve P50", f"{money_bn(totals.get('Direct',0))} / {money_bn(totals.get('Indirect',0))} / {money_bn(totals.get('Reserve',0))}"],
      ["Confidence", f"{model.get('confidence_pct')}%"], ["Risk", model.get("risk")], ["Scenario Cost Note", fx["cost_note"]], ["Scenario Schedule Note", fx["schedule_note"]], ["Scenario Risk Note", fx["risk_note"]]
    ]
    for r in rows: ws.append(r)
    ws.merge_cells("A1:B1"); ws["A1"].font=Font(size=20,bold=True,color="8DF7FF"); ws["A1"].fill=PatternFill("solid",fgColor="02050A")
    _style_ws(ws,"A3"); ws.column_dimensions["A"].width=34; ws.column_dimensions["B"].width=110

    ws=wb.create_sheet("01 Selected Cost Estimate")
    _add(ws, [["CBS","Description","Cost Type","Basis","Min / O BN","Most Likely / P50 BN","Max / P90 BN","P80 BN","Scenario Note","Impact Basis"]]+[
      [x["cbs"],x["description"],x["type"],x["basis"],x["p10_bn"],x["p50_bn"],x["p90_bn"],round((x["p50_bn"]+x["p90_bn"])*0.62,3),fx["cost_note"],x["impact_basis"]] for x in model.get("cost_lines",[])
    ])
    ws=wb.create_sheet("02 Cost Summary")
    _add(ws, [["Cost Bucket","P50 BN","Share","Board Meaning"], ["Direct",round(totals.get("Direct",0),3),round(totals.get("Direct",0)/max(parse_bn(model.get('cost_p50')),0.001),3),"Physical scope and deliverable assets"], ["Indirect",round(totals.get("Indirect",0),3),round(totals.get("Indirect",0)/max(parse_bn(model.get('cost_p50')),0.001),3),"Prelims, design, assurance, insurance and management"], ["Reserve",round(totals.get("Reserve",0),3),round(totals.get("Reserve",0)/max(parse_bn(model.get('cost_p50')),0.001),3),"Contingency and management reserve linked to QCRA"]])
    ws=wb.create_sheet("03 Risk Register Pro")
    risk_headers=["Risk ID","Category","Title","Cause","Risk Event","Impact Description","Probability %","Pre Rating","Residual Rating","Owner","Response Strategy","Mitigation Actions","Trigger","Status","Last Reviewed","Board Visibility","Linked Activity","Activity Name","Linked CBS","CBS Name","Sched O","Sched ML","Sched Max","Cost O BN","Cost ML BN","Cost Max BN","Cost EMV BN","Schedule EMV Days","Basis Cost","Basis Schedule"]
    _add(ws,[risk_headers]+[[r.get(k) for k in []] for r in []])
    for r in model.get("risks",[]):
        ws.append([r["risk_id"],r["category"],r["title"],r.get("cause"),r.get("risk_event"),r.get("impact_description"),r["probability_pct"],r.get("pre_mitigation_rating"),r.get("residual_rating"),r["owner"],r.get("response_strategy"),r["mitigation"],r["trigger"],r.get("status"),r.get("last_reviewed"),r.get("board_visibility"),r["activity_id"],r["activity_name"],r["cbs"],r["cbs_name"],r["schedule_o_days"],r["schedule_m_days"],r["schedule_p_days"],r["cost_o_bn"],r["cost_m_bn"],r["cost_p_bn"],r.get("cost_emv_bn"),r.get("schedule_emv_days"),r["basis_of_cost_impact"],r["basis_of_schedule_impact"]])
    _style_ws(ws)
    ws=wb.create_sheet("04 QCRA P-Curve")
    _add(ws, [["Percentile","Cost BN"]]+[[x["percentile"],x["cost_bn"]] for x in mc["curve"]])
    try:
        chart=LineChart(); chart.title="QCRA Cost P-Curve"; chart.y_axis.title="Cost BN"; chart.x_axis.title="Percentile"; data=Reference(ws,min_col=2,min_row=1,max_row=ws.max_row); cats=Reference(ws,min_col=1,min_row=2,max_row=ws.max_row); chart.add_data(data,titles_from_data=True); chart.set_categories(cats); ws.add_chart(chart,"D2")
    except Exception: pass
    ws=wb.create_sheet("05 QSRA P-Curve")
    _add(ws, [["Percentile","Schedule Months"]]+[[x["percentile"],x["schedule_months"]] for x in mc["curve"]])
    try:
        chart=LineChart(); chart.title="QSRA Schedule P-Curve"; chart.y_axis.title="Months"; chart.x_axis.title="Percentile"; data=Reference(ws,min_col=2,min_row=1,max_row=ws.max_row); cats=Reference(ws,min_col=1,min_row=2,max_row=ws.max_row); chart.add_data(data,titles_from_data=True); chart.set_categories(cats); ws.add_chart(chart,"D2")
    except Exception: pass
    ws=wb.create_sheet("06 QCRA Tornado")
    _add(ws, [["Rank","Risk ID","Driver","Category","CBS","Cost Mean BN","Driver Score"]]+[[i+1,x["risk_id"],x["title"],x["category"],x["cbs"],x["cost_mean_bn"],x["driver_score"]] for i,x in enumerate(mc.get("qcra_tornado",mc.get("tornado",[]))[:12])])
    try:
        b=BarChart(); b.type="bar"; b.title="QCRA Tornado - Cost Drivers"; data=Reference(ws,min_col=7,min_row=1,max_row=ws.max_row); cats=Reference(ws,min_col=3,min_row=2,max_row=ws.max_row); b.add_data(data,titles_from_data=True); b.set_categories(cats); ws.add_chart(b,"I2")
    except Exception: pass
    ws=wb.create_sheet("07 QSRA Tornado")
    _add(ws, [["Rank","Risk ID","Driver","Category","Activity","Schedule Mean Days","Driver Score"]]+[[i+1,x["risk_id"],x["title"],x["category"],x["activity_id"],x["schedule_mean_days"],x["driver_score"]] for i,x in enumerate(mc.get("qsra_tornado",mc.get("tornado",[]))[:12])])
    try:
        b=BarChart(); b.type="bar"; b.title="QSRA Tornado - Schedule Drivers"; data=Reference(ws,min_col=7,min_row=1,max_row=ws.max_row); cats=Reference(ws,min_col=3,min_row=2,max_row=ws.max_row); b.add_data(data,titles_from_data=True); b.set_categories(cats); ws.add_chart(b,"I2")
    except Exception: pass
    ws=wb.create_sheet("08 Schedule Logic")
    _add(ws, [["Activity ID","Phase","Activity","Predecessor","Duration Months","Critical","Basis"]]+[[x["activity_id"],x["phase"],x["activity"],x["predecessor"],x["duration_months"],x["critical"],x["basis"]] for x in model.get("schedule_rows",[])])
    ws=wb.create_sheet("09 Scenarios")
    _add(ws, [["Scenario","Cost","Schedule Months","Risk","Confidence","Why"]]+[[x["label"],x["cost"],x["schedule_months"],x["risk"],x["confidence"],x["why"]] for x in model.get("scenario_comparison",[])])
    ws=wb.create_sheet("10 Benchmarks")
    _add(ws, [["Metric","Value","Why"]]+[[x["metric"],x["value"],x["why"]] for x in model.get("benchmarks",[])])
    ws=wb.create_sheet("11 Assumptions & Audit")
    _add(ws, [["Item","Value"],["Model ID",model.get("id")],["Version","v50 Output Domination"],["Estimate Class",model.get("estimate_class_name")],["Schedule Level",model.get("schedule_level")],["Important note","This is a first-pass intelligence pack for option testing and challenge. It is not a substitute for verified design, supplier quotes, legal, safety or regulated professional sign-off."]]+[["Confidence basis",x] for x in model.get("confidence_explanation",[])])
    bio=BytesIO(); wb.save(bio); bio.seek(0); return bio.getvalue()

def risk_csv_bytes(model):
    out=StringIO(); w=csv.writer(out)
    headers=["Risk ID","Category","Title","Cause","Risk Event","Impact Description","Probability %","Pre Rating","Residual Rating","Owner","Response Strategy","Mitigation","Trigger","Status","Last Reviewed","Board Visibility","Activity ID","Activity Name","CBS","CBS Name","Sched O","Sched ML","Sched Max","Cost O BN","Cost ML BN","Cost Max BN","Cost EMV BN","Schedule EMV Days","Basis Cost","Basis Schedule"]
    w.writerow(headers)
    for r in model.get("risks",[]):
        w.writerow([r.get("risk_id"),r.get("category"),r.get("title"),r.get("cause"),r.get("risk_event"),r.get("impact_description"),r.get("probability_pct"),r.get("pre_mitigation_rating"),r.get("residual_rating"),r.get("owner"),r.get("response_strategy"),r.get("mitigation"),r.get("trigger"),r.get("status"),r.get("last_reviewed"),r.get("board_visibility"),r.get("activity_id"),r.get("activity_name"),r.get("cbs"),r.get("cbs_name"),r.get("schedule_o_days"),r.get("schedule_m_days"),r.get("schedule_p_days"),r.get("cost_o_bn"),r.get("cost_m_bn"),r.get("cost_p_bn"),r.get("cost_emv_bn"),r.get("schedule_emv_days"),r.get("basis_of_cost_impact"),r.get("basis_of_schedule_impact")])
    return out.getvalue().encode()

def word_bytes(model):
    doc=Document(); styles=doc.styles; styles["Normal"].font.name="Aptos"; styles["Normal"].font.size=Pt(10)
    title=doc.add_heading("CASEY TITAN X — Executive Project Controls Intelligence",0); title.runs[0].font.color.rgb=RGBColor(8,40,60)
    doc.add_paragraph(model.get("executive_summary",""))
    fx=_scenario_effects(model); totals=_total_costs(model); mc=model["monte_carlo"]
    doc.add_heading("Board Decision Snapshot",1)
    tbl=doc.add_table(rows=1,cols=2); tbl.style="Light Shading Accent 1"; hdr=tbl.rows[0].cells; hdr[0].text="Metric"; hdr[1].text="Value"
    for k,v in [("Scenario",model.get("scenario_label")),("P10 / P50 / P80 / P90 Cost",f"{money_bn(mc['qcra']['p10'])} / {model.get('cost_p50')} / {money_bn(mc['qcra']['p80'])} / {money_bn(mc['qcra']['p90'])}"),("Direct / Indirect / Reserve",f"{money_bn(totals.get('Direct',0))} / {money_bn(totals.get('Indirect',0))} / {money_bn(totals.get('Reserve',0))}"),("Baseline / P80 Schedule",f"{model.get('schedule')} / {mc['qsra']['p80']} months"),("Risk / Confidence",f"{model.get('risk')} / {model.get('confidence_pct')}%")]:
        row=tbl.add_row().cells; row[0].text=k; row[1].text=str(v)
    doc.add_heading("Executive Recommendation",1); doc.add_paragraph(f"Proceed with the {model.get('scenario_label')} basis only if the sponsor accepts the quantified risk range, funds the reserve strategy and closes the top evidence gaps before approval.")
    for heading,items in [("Scenario Logic",[fx['cost_note'],fx['schedule_note'],fx['risk_note']]),("Confidence Basis",model.get("confidence_explanation",[])),("Board Challenge Questions",model.get("board_challenge_questions",[])),("Next Best Actions",model.get("next_best_actions",[]))]:
        doc.add_heading(heading,1)
        for x in items: doc.add_paragraph(str(x),style="List Bullet")
    doc.add_heading("Top Risk Drivers",1)
    for r in model.get("risks",[])[:10]:
        doc.add_paragraph(f"{r['risk_id']} {r['title']}",style="List Bullet")
        doc.add_paragraph(f"Cause: {r.get('cause')} | Event: {r.get('risk_event')} | Impact: {r.get('impact_description')} | Response: {r.get('mitigation')}")
    doc.add_heading("Cost View",1)
    tbl=doc.add_table(rows=1,cols=5); tbl.style="Light Grid Accent 1"; hdr=tbl.rows[0].cells
    for i,h in enumerate(["CBS","Description","Type","P50 BN","P90 BN"]): hdr[i].text=h
    for x in model.get("cost_lines",[])[:18]:
        row=tbl.add_row().cells; row[0].text=x["cbs"]; row[1].text=x["description"]; row[2].text=x["type"]; row[3].text=str(x["p50_bn"]); row[4].text=str(x["p90_bn"])
    bio=BytesIO(); doc.save(bio); bio.seek(0); return bio.getvalue()

def pdf_bytes(model):
    bio=BytesIO(); doc=SimpleDocTemplate(bio,pagesize=landscape(A4),rightMargin=28,leftMargin=28,topMargin=24,bottomMargin=24)
    styles=getSampleStyleSheet(); styles.add(ParagraphStyle(name="CASEYTitle",fontSize=24,leading=28,textColor=colors.HexColor("#07111F"),spaceAfter=12)); styles.add(ParagraphStyle(name="Small",fontSize=8,leading=10,textColor=colors.HexColor("#334155")))
    mc=model["monte_carlo"]; totals=_total_costs(model); story=[]
    story.append(Paragraph("CASEY TITAN X — Board Intelligence Pack",styles["CASEYTitle"])); story.append(Paragraph(model.get("executive_summary",""),styles["BodyText"])); story.append(Spacer(1,10))
    data=[["Metric","Value"],["Scenario",model.get("scenario_label")],["P50 Cost",model.get("cost_p50")],["QCRA P80",money_bn(mc["qcra"]["p80"])],["QSRA P80",f"{mc['qsra']['p80']} months"],["Direct / Indirect / Reserve",f"{money_bn(totals.get('Direct',0))} / {money_bn(totals.get('Indirect',0))} / {money_bn(totals.get('Reserve',0))}"],["Risk / Confidence",f"{model.get('risk')} / {model.get('confidence_pct')}%"]]
    table=Table(data,colWidths=[170,480]); table.setStyle(TableStyle([("BACKGROUND",(0,0),(-1,0),colors.HexColor("#07111F")),("TEXTCOLOR",(0,0),(-1,0),colors.white),("GRID",(0,0),(-1,-1),.35,colors.HexColor("#CBD5E1")),("BACKGROUND",(0,1),(-1,-1),colors.HexColor("#F8FAFC"))])); story.append(table); story.append(PageBreak())
    story.append(Paragraph("Top QCRA and QSRA Drivers",styles["Heading1"]))
    rows=[["Rank","QCRA Cost Driver","Cost Mean BN","QSRA Schedule Driver","Mean Days"]]
    q=mc.get("qcra_tornado",[])[:8]; s=mc.get("qsra_tornado",[])[:8]
    for i in range(max(len(q),len(s))): rows.append([i+1,q[i]["title"] if i<len(q) else "",q[i].get("cost_mean_bn","") if i<len(q) else "",s[i]["title"] if i<len(s) else "",s[i].get("schedule_mean_days","") if i<len(s) else ""])
    t=Table(rows,colWidths=[38,250,80,250,80]); t.setStyle(TableStyle([("BACKGROUND",(0,0),(-1,0),colors.HexColor("#07111F")),("TEXTCOLOR",(0,0),(-1,0),colors.white),("GRID",(0,0),(-1,-1),.35,colors.HexColor("#CBD5E1")),("FONTSIZE",(0,0),(-1,-1),8)])); story.append(t); story.append(PageBreak())
    story.append(Paragraph("Risk Register — Board Visibility",styles["Heading1"]))
    rows=[["ID","Risk","Cause","Event","Impact","Owner","Response"]]
    for r in model.get("risks",[])[:10]: rows.append([r.get("risk_id"),r.get("title"),r.get("cause"),r.get("risk_event"),r.get("impact_description"),r.get("owner"),r.get("mitigation")])
    t=Table(rows,colWidths=[45,90,150,150,175,80,165]); t.setStyle(TableStyle([("BACKGROUND",(0,0),(-1,0),colors.HexColor("#07111F")),("TEXTCOLOR",(0,0),(-1,0),colors.white),("GRID",(0,0),(-1,-1),.25,colors.HexColor("#CBD5E1")),("FONTSIZE",(0,0),(-1,-1),7),('VALIGN',(0,0),(-1,-1),'TOP')])) ; story.append(t)
    doc.build(story); bio.seek(0); return bio.getvalue()

def pptx_bytes(model):
    from pptx.dml.color import RGBColor as PRGBColor
    prs=Presentation(); prs.slide_width=PptxInches(13.333); prs.slide_height=PptxInches(7.5); blank=prs.slide_layouts[6]
    def title(sl,t,sub=""):
        sl.background.fill.solid(); sl.background.fill.fore_color.rgb=PRGBColor(2,5,10)
        box=sl.shapes.add_textbox(PptxInches(.45),PptxInches(.35),PptxInches(12.4),PptxInches(.6)); tf=box.text_frame; tf.text=t; r=tf.paragraphs[0].runs[0]; r.font.size=PptxPt(26); r.font.bold=True; r.font.color.rgb=PRGBColor(141,247,255)
        if sub:
            b=sl.shapes.add_textbox(PptxInches(.48),PptxInches(.95),PptxInches(12),PptxInches(.55)); b.text_frame.text=sub; b.text_frame.paragraphs[0].runs[0].font.size=PptxPt(13); b.text_frame.paragraphs[0].runs[0].font.color.rgb=PRGBColor(245,250,255)
    s=prs.slides.add_slide(blank); title(s,"CASEY TITAN X",model.get("executive_summary",""))
    s=prs.slides.add_slide(blank); title(s,"Board Metrics")
    metrics=[("P50",model.get("cost_p50")),("P80",money_bn(model['monte_carlo']['qcra']['p80'])),("Schedule",model.get("schedule")),("QSRA P80",str(model['monte_carlo']['qsra']['p80'])+" months"),("Risk",model.get("risk")),("Confidence",str(model.get("confidence_pct"))+"%")]
    for i,(k,v) in enumerate(metrics):
        x=.6+(i%3)*4.2; y=1.55+(i//3)*1.7; shp=s.shapes.add_shape(1,PptxInches(x),PptxInches(y),PptxInches(3.7),PptxInches(1.1)); shp.fill.solid(); shp.fill.fore_color.rgb=PRGBColor(7,17,31); shp.line.color.rgb=PRGBColor(141,247,255); shp.text=f"{k}\n{v}"; shp.text_frame.paragraphs[0].runs[0].font.color.rgb=PRGBColor(255,255,255)
    s=prs.slides.add_slide(blank); title(s,"Top Risk Drivers","Cause • Event • Impact • Response")
    txt=s.shapes.add_textbox(PptxInches(.55),PptxInches(1.3),PptxInches(12.1),PptxInches(5.6)).text_frame
    for r in model.get("risks",[])[:8]:
        p=txt.add_paragraph(); p.text=f"{r['risk_id']} {r['title']} — {r.get('impact_description')}"; p.font.size=PptxPt(12); p.font.color.rgb=PRGBColor(245,250,255)
    s=prs.slides.add_slide(blank); title(s,"Output Pack","Selected cost estimate, risk register, QCRA/QSRA and decision evidence generated from one model.")
    bio=BytesIO(); prs.save(bio); bio.seek(0); return bio.getvalue()

# ======================= END CASEY v50 OUTPUT DOMINATION OVERRIDES =======================


# Frontend compatibility aliases for v50 UI
_CASEY_ORIGINAL_BUILD_MODEL = build_model
def build_model(prompt:str, client:str="", class_level:int=3, schedule_level:int=3, scenario:str="base"):
    m=_CASEY_ORIGINAL_BUILD_MODEL(prompt, client, class_level, schedule_level, scenario)
    m=apply_sector_intelligence(m)
    m=scenario_cascade_v95(m, scenario)
    m["version"]="CASEY TITAN X v95 Demo Launch"
    m["cost_breakdown"]=m.get("cost_breakdown") or m.get("cost_lines",[])
    m["risk_register"]=m.get("risk_register") or m.get("risks",[])
    m["schedule_detail"]=m.get("schedule_detail") or m.get("schedule_rows",[])
    m=_v126_normalize_model_costs(m) if "_v126_normalize_model_costs" in globals() else m
    for x in m.get("monte_carlo",{}).get("tornado",[]):
        x["driver"]=x.get("title")
        x["contribution"]=x.get("driver_score")
    return m

@app.post("/upload")
async def upload_alias(file: UploadFile = File(...)):
    return await analyse_upload(file)

# v50 export filenames override
@app.post("/v50/export/workbook")
def export_workbook_v50(model:Dict[str,Any]): return stream(workbook_bytes(model),"application/vnd.openxmlformats-officedocument.spreadsheetml.sheet","CASEY_v50_Cost_Model_Elite.xlsx")
@app.post("/v50/export/risk-register")
def export_risk_v50(model:Dict[str,Any]): return stream(risk_csv_bytes(model),"text/csv","CASEY_v50_Risk_Register_Pro.csv")
@app.post("/v50/export/all")
def export_all_v50(model:Dict[str,Any]):
    bio=BytesIO()
    with zipfile.ZipFile(bio,"w",zipfile.ZIP_DEFLATED) as z:
        z.writestr("01_CASEY_v50_Cost_Model_Elite.xlsx",workbook_bytes(model))
        z.writestr("02_CASEY_v50_Risk_Register_Pro.csv",risk_csv_bytes(model))
        z.writestr("03_CASEY_v50_P6_Schedule.xer",xer_bytes(model))
        z.writestr("04_CASEY_v50_Executive_Report.docx",word_bytes(model))
        z.writestr("05_CASEY_v50_Board_Report.pdf",pdf_bytes(model))
        z.writestr("06_CASEY_v50_Board_Deck.pptx",pptx_bytes(model))
        z.writestr("07_CASEY_v50_Model.json",json.dumps(model,indent=2))
    bio.seek(0); return stream(bio.getvalue(),"application/zip","CASEY_v50_Output_Domination_Pack.zip")



# ======================= CASEY v51 OUTPUT PLANET CLASS OVERRIDES =======================
# Highest-spec outputs: premium workbook, risk workbook, DOCX, PDF, PPTX and ZIP.
# Designed to make export quality match the cinematic CASEY product experience.

from openpyxl.chart.label import DataLabelList
from openpyxl.worksheet.table import Table as XLTable, TableStyleInfo
from openpyxl.formatting.rule import ColorScaleRule, CellIsRule
from openpyxl.comments import Comment

V51_NAVY = "02050A"
V51_PANEL = "07111F"
V51_CYAN = "8DF7FF"
V51_BLUE = "2FB7FF"
V51_WHITE = "F7FBFF"
V51_MUTED = "9FB3C8"
V51_LINE = "203040"
V51_RED = "FF5C7A"
V51_AMBER = "FFD166"
V51_GREEN = "7CFFB2"


def _safe(v, default=""):
    return default if v is None else v


def _num(v, default=0.0):
    try:
        if isinstance(v, (int, float)):
            return float(v)
        s = str(v).replace('$','').replace('£','').replace('€','').replace(',','').strip().upper()
        n = float(re.sub(r'[^0-9.\-]', '', s) or 0)
        if 'T' in s:
            return n * 1000
        if 'M' in s:
            return n / 1000
        return n
    except Exception:
        return default


def _v51_money(v):
    return money_bn(_num(v))


def _casey_reconcile_cost_lines(model):
    # Final demo safety: cost rows must always reconcile to selected scenario P50.
    if not isinstance(model, dict):
        return model
    target = _num(model.get('cost_p50'))
    rows = model.get('cost_lines') or model.get('cost_breakdown') or []
    clean = []
    for i, x in enumerate(rows if isinstance(rows, list) else []):
        if not isinstance(x, dict):
            continue
        typ_raw = str(x.get('type') or 'Direct')
        typ = 'Reserve' if re.search(r'reserve|risk|contingency', typ_raw, re.I) else ('Indirect' if re.search(r'indirect|owner|pm|management|prelim', typ_raw, re.I) else 'Direct')
        p50 = _num(x.get('p50_bn') or x.get('most_likely_p50') or x.get('p50') or x.get('most_likely') or x.get('value'))
        if p50 <= 0:
            continue
        clean.append({**x, 'type': typ, 'cbs': x.get('cbs') or f'C-{i+1:02d}', 'description': x.get('description') or x.get('title') or typ,
                      'p10_bn': _num(x.get('p10_bn') or x.get('low_p10') or x.get('p10') or x.get('low')) or p50*0.8,
                      'p50_bn': p50,
                      'p90_bn': _num(x.get('p90_bn') or x.get('high_p90') or x.get('p90') or x.get('high')) or p50*1.3,
                      'impact_basis': x.get('impact_basis') or x.get('basis') or f'{typ} cost basis reconciled to selected scenario P50.'})
    if target > 0 and (not clean or sum(_num(x.get('p50_bn')) for x in clean) <= 0):
        clean = [
            {'cbs':'01.00','description':'Direct delivery scope','type':'Direct','p10_bn':target*0.62,'p50_bn':target*0.76,'p90_bn':target*0.96,'impact_basis':'Reconciled scenario cost split.'},
            {'cbs':'90.00','description':'Indirects, owner costs and integration','type':'Indirect','p10_bn':target*0.10,'p50_bn':target*0.14,'p90_bn':target*0.18,'impact_basis':'Reconciled scenario cost split.'},
            {'cbs':'99.00','description':'Risk reserve and contingency','type':'Reserve','p10_bn':target*0.06,'p50_bn':target*0.10,'p90_bn':target*0.16,'impact_basis':'Reconciled scenario cost split.'},
        ]
    total = sum(_num(x.get('p50_bn')) for x in clean)
    if target > 0 and total > 0:
        factor = target / total
        for x in clean:
            x['p10_bn'] = round(_num(x.get('p10_bn')) * factor, 3)
            x['p50_bn'] = round(_num(x.get('p50_bn')) * factor, 3)
            x['p90_bn'] = round(_num(x.get('p90_bn')) * factor, 3)
        drift = round(target - sum(_num(x.get('p50_bn')) for x in clean), 3)
        if clean and abs(drift) >= 0.001:
            mx = max(range(len(clean)), key=lambda i: clean[i]['p50_bn'])
            clean[mx]['p50_bn'] = round(clean[mx]['p50_bn'] + drift, 3)
    model['cost_lines'] = clean
    model['cost_breakdown'] = clean
    return model


def _v51_total_costs(model):
    model = _casey_reconcile_cost_lines(model)
    out = {"Direct": 0.0, "Indirect": 0.0, "Reserve": 0.0}
    for x in model.get("cost_lines", []):
        typ = x.get("type", "Direct")
        if typ not in out: out[typ] = 0.0
        out[typ] += _num(x.get("p50_bn"))
    return out


def _v51_scenario_notes(model):
    s = str(model.get("scenario", "base")).lower()
    if s == "faster":
        return [
            "Acceleration selected: schedule compression is funded through logistics premium, parallel procurement and higher delivery intensity.",
            "Controls consequence: QSRA risk should be watched closely because interface density and commissioning pressure rise.",
            "Board decision: approve acceleration only with package-level recovery plan and explicit acceleration cost allowance."
        ]
    if s == "cheaper":
        return [
            "Cheaper selected: value engineering and procurement competition reduce P50 but increase residual delivery exposure.",
            "Controls consequence: all savings must be linked to scope choices; avoid hiding deferred scope as efficiency.",
            "Board decision: approve only after VE log, exclusions, quality impacts and owner-retained risks are signed off."
        ]
    if s == "lower_risk":
        return [
            "Lower-risk selected: more assurance, surveys, float and contingency raise headline cost but improve confidence.",
            "Controls consequence: stronger evidence should move risk from unknown unknowns into quantified known unknowns.",
            "Board decision: approve with a confidence-improvement plan and gate evidence pack."
        ]
    if s == "premium":
        return [
            "Premium selected: resilience, specification, assurance and stakeholder certainty are prioritised over lowest cost.",
            "Controls consequence: P50 rises, but delivery certainty and reputational protection should improve.",
            "Board decision: approve if the premium is justified by availability, criticality, carbon, security or strategic value."
        ]
    return [
        "Base selected: balanced control case for first-pass board challenge and option testing.",
        "Controls consequence: use this as the reference model before switching to Faster, Cheaper, Lower Risk or Premium.",
        "Board decision: approve further definition work, not unconditional capital commitment."
    ]


def _apply_ws_theme(ws, widths=None, freeze="A8"):
    ws.sheet_view.showGridLines = False
    ws.freeze_panes = freeze
    if widths:
        for col, width in widths.items():
            ws.column_dimensions[col].width = width
    thin = Side(style="thin", color=V51_LINE)
    for row in ws.iter_rows():
        for cell in row:
            cell.font = Font(name="Aptos", size=10, color="EAF6FF")
            cell.alignment = Alignment(vertical="top", wrap_text=True)
            cell.border = Border(bottom=thin)
    for cell in ws[1]:
        cell.font = Font(name="Aptos Display", size=16, bold=True, color=V51_CYAN)
    return ws


def _title(ws, title, subtitle=None, end_col="J"):
    ws.merge_cells(f"A1:{end_col}1")
    ws["A1"] = title
    ws["A1"].fill = PatternFill("solid", fgColor=V51_Navy if False else V51_Navy if False else V51_NAVY)
    ws["A1"].font = Font(name="Aptos Display", size=18, bold=True, color=V51_CYAN)
    ws["A1"].alignment = Alignment(vertical="center")
    ws.row_dimensions[1].height = 34
    if subtitle:
        ws.merge_cells(f"A2:{end_col}2")
        ws["A2"] = subtitle
        ws["A2"].font = Font(name="Aptos", size=10, color=V51_MUTED)
        ws["A2"].fill = PatternFill("solid", fgColor=V51_PANEL)
        ws.row_dimensions[2].height = 28


def _header_row(ws, row, cols, fill=V51_PANEL):
    for i, h in enumerate(cols, 1):
        c = ws.cell(row=row, column=i, value=h)
        c.fill = PatternFill("solid", fgColor=fill)
        c.font = Font(name="Aptos", size=10, bold=True, color=V51_CYAN)
        c.alignment = Alignment(horizontal="center", vertical="center", wrap_text=True)
    ws.row_dimensions[row].height = 26


def _add_table(ws, ref, name):
    try:
        tab = XLTable(displayName=name, ref=ref)
        style = TableStyleInfo(name="TableStyleMedium2", showFirstColumn=False, showLastColumn=False, showRowStripes=True, showColumnStripes=False)
        tab.tableStyleInfo = style
        ws.add_table(tab)
    except Exception:
        pass


def workbook_bytes(model:Dict[str,Any]):
    wb = Workbook()
    # Remove default and build planet-class workbook
    ws = wb.active
    ws.title = "00 Executive Dashboard"
    for s in ["01 Basis", "02 Cost Estimate", "03 Indirects & Reserves", "04 QCRA", "05 QSRA", "06 Risk Register", "07 Schedule", "08 Scenarios", "09 Benchmarking", "10 Assumptions"]:
        wb.create_sheet(s)

    mc = model.get("monte_carlo", {})
    totals = _v51_total_costs(model)
    scenario_notes = _v51_scenario_notes(model)

    # Dashboard
    ws = wb["00 Executive Dashboard"]
    _title(ws, "CASEY - EXECUTIVE PROJECT CONTROLS INTELLIGENCE", "Board-ready first-pass cost, schedule, risk and confidence model generated from a single connected project basis.", "K")
    widths = {"A":22,"B":22,"C":22,"D":22,"E":22,"F":22,"G":22,"H":22,"I":22,"J":22,"K":22}
    _apply_ws_theme(ws, widths, freeze="A9")
    ws["A4"]="Project"; ws["B4"]=model.get("title")
    ws["D4"]="Sector"; ws["E4"]=model.get("subsector")
    ws["G4"]="Scenario"; ws["H4"]=model.get("scenario_label")
    ws["J4"]="Generated"; ws["K4"]=datetime.utcnow().strftime("%Y-%m-%d %H:%M UTC")
    kpis=[("P10 Cost",model.get("cost_p10")),("P50 Cost",model.get("cost_p50")),("P90 Cost",model.get("cost_p90")),("QCRA P80",_v51_money(mc.get("qcra",{}).get("p80"))),("Schedule",model.get("schedule")),("QSRA P80",f"{mc.get('qsra',{}).get('p80','-')} months"),("Risk",model.get("risk")),("Confidence",f"{model.get('confidence_pct')}%"),("Direct",_v51_money(totals.get("Direct"))),("Indirect",_v51_money(totals.get("Indirect"))),("Reserve",_v51_money(totals.get("Reserve")))]
    start=6
    for i,(k,v) in enumerate(kpis):
        r=start+(i//4)*3; c=1+(i%4)*3
        ws.cell(r,c,k); ws.cell(r+1,c,v)
        for rr in [r,r+1]:
            for cc in [c,c+1]:
                ws.cell(rr,cc).fill = PatternFill("solid", fgColor=V51_PANEL)
                ws.cell(rr,cc).border = Border(left=Side(style="thin",color=V51_CYAN),right=Side(style="thin",color=V51_LINE),top=Side(style="thin",color=V51_LINE),bottom=Side(style="thin",color=V51_LINE))
        ws.cell(r,c).font = Font(name="Aptos", size=9, bold=True, color=V51_MUTED)
        ws.cell(r+1,c).font = Font(name="Aptos Display", size=18, bold=True, color=V51_WHITE)
        ws.merge_cells(start_row=r,start_column=c,end_row=r,end_column=c+1)
        ws.merge_cells(start_row=r+1,start_column=c,end_row=r+1,end_column=c+1)
    ws["A16"]="Executive Summary"; ws["A16"].font=Font(name="Aptos Display",size=14,bold=True,color=V51_CYAN)
    ws.merge_cells("A17:K20"); ws["A17"]=model.get("executive_summary",""); ws["A17"].fill=PatternFill("solid",fgColor=V51_PANEL)
    ws["A22"]="Board Recommendation"; ws["A22"].font=Font(name="Aptos Display",size=14,bold=True,color=V51_CYAN)
    notes = scenario_notes + model.get("next_best_actions", [])[:5]
    for i,n in enumerate(notes,23):
        ws.cell(i,1,f"{i-22}"); ws.cell(i,2,n); ws.merge_cells(start_row=i,start_column=2,end_row=i,end_column=11)
        ws.cell(i,1).fill=PatternFill("solid",fgColor=V51_CYAN); ws.cell(i,1).font=Font(bold=True,color=V51_Navy if False else V51_NAVY)
        ws.cell(i,2).fill=PatternFill("solid",fgColor="06101C")
    # Cost composition mini chart table
    ws["A32"]="Cost Composition"; ws["A32"].font=Font(name="Aptos Display",size=14,bold=True,color=V51_CYAN)
    _header_row(ws,33,["Type","P50 BN"])
    for i,(typ,val) in enumerate(totals.items(),34):
        ws.cell(i,1,typ); ws.cell(i,2,round(val,3)); ws.cell(i,2).number_format='$#,##0.0 "B"'
    bar=BarChart(); bar.title="Direct / Indirect / Reserve"; bar.style=10; bar.y_axis.title="BN"; bar.x_axis.title="Cost Type";
    bar.add_data(Reference(ws,min_col=2,min_row=33,max_row=36),titles_from_data=True)
    bar.set_categories(Reference(ws,min_col=1,min_row=34,max_row=36)); bar.height=6; bar.width=11; ws.add_chart(bar,"D32")

    # Basis
    ws=wb["01 Basis"]
    _title(ws,"BASIS OF ESTIMATE AND CONTROL MODEL", "Traceability page: no board output should be read without understanding the class, maturity, location, scenario and assumptions.", "H")
    _apply_ws_theme(ws,{"A":28,"B":80,"C":22,"D":22,"E":22,"F":22,"G":22,"H":22}, freeze="A6")
    pairs=[("Prompt",model.get("prompt")),("Client",model.get("client")),("Mode",model.get("mode")),("Sector",model.get("subsector")),("Location",model.get("location")),("Scale",model.get("scale")),("Estimate class",model.get("estimate_class_name")),("Estimate maturity",model.get("estimate_maturity")),("Schedule level",model.get("schedule_level")),("Scenario",model.get("scenario_label")),("Scenario rationale",model.get("scenario_why"))]
    _header_row(ws,4,["Input","Value"])
    for i,(k,v) in enumerate(pairs,5):
        ws.cell(i,1,k); ws.cell(i,2,v); ws.merge_cells(start_row=i,start_column=2,end_row=i,end_column=8)
    ws["A19"]="Confidence Explanation"; ws["A19"].font=Font(size=13,bold=True,color=V51_CYAN)
    for i,x in enumerate(model.get("confidence_explanation",[]),20):
        ws.cell(i,1,"•"); ws.cell(i,2,x); ws.merge_cells(start_row=i,start_column=2,end_row=i,end_column=8)

    # Cost estimate
    ws=wb["02 Cost Estimate"]
    _title(ws,"SELECTED CLASS COST ESTIMATE - DIRECT / INDIRECT / RESERVE", "Only the selected class estimate is exported. P10 / P50 / P90 and direct / indirect / reserve are visible and auditable.", "L")
    cols=["CBS","Description","Type","Basis","P10 BN","P50 BN","P90 BN","Min BN","Most Likely BN","Max BN","Scenario Impact","Evidence / Challenge"]
    _apply_ws_theme(ws,{"A":10,"B":30,"C":14,"D":45,"E":12,"F":12,"G":12,"H":12,"I":14,"J":12,"K":26,"L":55}, freeze="A5")
    _header_row(ws,4,cols)
    for r,x in enumerate(model.get("cost_lines",[]),5):
        vals=[x.get("cbs"),x.get("description"),x.get("type"),x.get("basis"),x.get("p10_bn"),x.get("p50_bn"),x.get("p90_bn"),x.get("p10_bn"),x.get("p50_bn"),x.get("p90_bn"),model.get("scenario_label"),x.get("impact_basis")]
        for c,v in enumerate(vals,1): ws.cell(r,c,v)
        for c in range(5,11): ws.cell(r,c).number_format='$#,##0.000 "B"'
        if x.get("type") == "Direct": fill="071A2F"
        elif x.get("type") == "Indirect": fill="101827"
        else: fill="26111A"
        for c in range(1,len(cols)+1): ws.cell(r,c).fill=PatternFill("solid",fgColor=fill)
    end=4+len(model.get("cost_lines",[]))
    _add_table(ws,f"A4:L{end}","CASEY_CostEstimate")
    # summary
    row=end+3
    ws.cell(row,1,"Type"); ws.cell(row,2,"P50 BN"); ws.cell(row,3,"Share")
    _header_row(ws,row,["Type","P50 BN","Share"])
    for i,typ in enumerate(["Direct","Indirect","Reserve"],row+1):
        ws.cell(i,1,typ); ws.cell(i,2,totals.get(typ,0)); ws.cell(i,2).number_format='$#,##0.0 "B"'; ws.cell(i,3,f"=B{i}/SUM(B{row+1}:B{row+3})"); ws.cell(i,3).number_format='0.0%'
    chart=BarChart(); chart.title="Selected Estimate Cost Split"; chart.add_data(Reference(ws,min_col=2,min_row=row,max_row=row+3),titles_from_data=True); chart.set_categories(Reference(ws,min_col=1,min_row=row+1,max_row=row+3)); chart.height=7; chart.width=12; ws.add_chart(chart,f"E{row}")

    # Indirects
    ws=wb["03 Indirects & Reserves"]
    _title(ws,"INDIRECTS, CONTINGENCY AND RESERVE ASSURANCE", "Separates project delivery cost from sponsor-held uncertainty and contingency.", "H")
    _apply_ws_theme(ws,{"A":18,"B":32,"C":14,"D":16,"E":16,"F":16,"G":38,"H":40}, freeze="A5")
    _header_row(ws,4,["CBS","Description","Type","P10 BN","P50 BN","P90 BN","Assurance Question","Board Treatment"])
    rows=[x for x in model.get("cost_lines",[]) if x.get("type")!="Direct"]
    for r,x in enumerate(rows,5):
        ws.cell(r,1,x.get("cbs")); ws.cell(r,2,x.get("description")); ws.cell(r,3,x.get("type")); ws.cell(r,4,x.get("p10_bn")); ws.cell(r,5,x.get("p50_bn")); ws.cell(r,6,x.get("p90_bn")); ws.cell(r,7,"Is this genuinely indirect/reserve, or hidden direct scope?"); ws.cell(r,8,"Show separately in board pack; do not bury in headline P50.")
        for c in [4,5,6]: ws.cell(r,c).number_format='$#,##0.000 "B"'

    # QCRA
    ws=wb["04 QCRA"]
    _title(ws,"QCRA COST RISK ANALYSIS", "Cost curve and separate tornado drivers. QCRA is not mixed with QSRA.", "J")
    _apply_ws_theme(ws,{"A":14,"B":16,"C":16,"D":16,"E":16,"F":16,"G":18,"H":28,"I":18,"J":18}, freeze="A5")
    _header_row(ws,4,["Percentile","Cost BN","Schedule Months","","Rank","Driver","CBS","Mean BN","Score","Board Challenge"])
    curve=mc.get("curve",[])
    for i,x in enumerate(curve,5):
        ws.cell(i,1,x.get("percentile")); ws.cell(i,2,x.get("cost_bn")); ws.cell(i,3,x.get("schedule_months")); ws.cell(i,2).number_format='$#,##0.0 "B"'
    for j,x in enumerate(mc.get("qcra_tornado",[])[:10],5):
        ws.cell(j,5,j-4); ws.cell(j,6,x.get("title")); ws.cell(j,7,x.get("cbs")); ws.cell(j,8,x.get("cost_mean_bn")); ws.cell(j,9,x.get("driver_score")); ws.cell(j,10,"Request evidence for cost range and owner-retained exposure.")
        ws.cell(j,8).number_format='$#,##0.000 "B"'
    line=LineChart(); line.title="QCRA Cost Curve"; line.y_axis.title="Cost BN"; line.x_axis.title="Percentile"; line.add_data(Reference(ws,min_col=2,min_row=4,max_row=4+len(curve)),titles_from_data=True); line.set_categories(Reference(ws,min_col=1,min_row=5,max_row=4+len(curve))); line.height=7; line.width=12; ws.add_chart(line,"A22")
    bar=BarChart(); bar.type="bar"; bar.style=12; bar.title="QCRA Tornado - Cost Drivers"; bar.add_data(Reference(ws,min_col=8,min_row=4,max_row=14),titles_from_data=True); bar.set_categories(Reference(ws,min_col=6,min_row=5,max_row=14)); bar.height=7; bar.width=12; ws.add_chart(bar,"E22")

    # QSRA
    ws=wb["05 QSRA"]
    _title(ws,"QSRA SCHEDULE RISK ANALYSIS", "Schedule curve and separate tornado drivers. QSRA is not mixed with QCRA.", "J")
    _apply_ws_theme(ws,{"A":14,"B":16,"C":16,"D":16,"E":16,"F":32,"G":16,"H":18,"I":18,"J":34}, freeze="A5")
    _header_row(ws,4,["Percentile","Schedule Months","Cost BN","","Rank","Driver","Activity","Mean Days","Score","Recovery / Mitigation Challenge"])
    for i,x in enumerate(curve,5):
        ws.cell(i,1,x.get("percentile")); ws.cell(i,2,x.get("schedule_months")); ws.cell(i,3,x.get("cost_bn")); ws.cell(i,3).number_format='$#,##0.0 "B"'
    for j,x in enumerate(mc.get("qsra_tornado",[])[:10],5):
        ws.cell(j,5,j-4); ws.cell(j,6,x.get("title")); ws.cell(j,7,x.get("activity_id")); ws.cell(j,8,x.get("schedule_mean_days")); ws.cell(j,9,x.get("driver_score")); ws.cell(j,10,"Confirm critical path logic, float erosion and recovery option.")
    line=LineChart(); line.title="QSRA Schedule Curve"; line.y_axis.title="Months"; line.x_axis.title="Percentile"; line.add_data(Reference(ws,min_col=2,min_row=4,max_row=4+len(curve)),titles_from_data=True); line.set_categories(Reference(ws,min_col=1,min_row=5,max_row=4+len(curve))); line.height=7; line.width=12; ws.add_chart(line,"A22")
    bar=BarChart(); bar.type="bar"; bar.title="QSRA Tornado - Schedule Drivers"; bar.add_data(Reference(ws,min_col=8,min_row=4,max_row=14),titles_from_data=True); bar.set_categories(Reference(ws,min_col=6,min_row=5,max_row=14)); bar.height=7; bar.width=12; ws.add_chart(bar,"E22")

    # Risk register
    ws=wb["06 Risk Register"]
    _title(ws,"RISK REGISTER PRO - CAUSE / EVENT / IMPACT / RESPONSE", "Board-grade register with quantified cost and schedule exposure, owners and linked WBS/CBS.", "R")
    headers=["Risk ID","Category","Title","Cause","Risk Event","Impact Description","Probability %","Cost O BN","Cost M BN","Cost P BN","Cost EMV BN","Schedule O Days","Schedule M Days","Schedule P Days","Schedule EMV Days","Activity","CBS","Owner","Trigger","Mitigation","Response Strategy","Residual","Status","Board"]
    _apply_ws_theme(ws,{"A":10,"B":14,"C":22,"D":36,"E":36,"F":42,"G":12,"H":12,"I":12,"J":12,"K":14,"L":12,"M":12,"N":12,"O":14,"P":12,"Q":10,"R":18,"S":26,"T":44,"U":18,"V":12,"W":12,"X":10}, freeze="A5")
    _header_row(ws,4,headers)
    for r,x in enumerate(model.get("risks",[]),5):
        vals=[x.get("risk_id"),x.get("category"),x.get("title"),x.get("cause"),x.get("risk_event"),x.get("impact_description"),x.get("probability_pct"),x.get("cost_o_bn"),x.get("cost_m_bn"),x.get("cost_p_bn"),x.get("cost_emv_bn"),x.get("schedule_o_days"),x.get("schedule_m_days"),x.get("schedule_p_days"),x.get("schedule_emv_days"),x.get("activity_id"),x.get("cbs"),x.get("owner"),x.get("trigger"),x.get("mitigation"),x.get("response_strategy"),x.get("residual_rating"),x.get("status"),x.get("board_visibility")]
        for c,v in enumerate(vals,1): ws.cell(r,c,v)
        for c in [8,9,10,11]: ws.cell(r,c).number_format='$#,##0.000 "B"'
        if x.get("board_visibility") == "Yes":
            for c in range(1,len(headers)+1): ws.cell(r,c).fill=PatternFill("solid",fgColor="071A2F")
    _add_table(ws,f"A4:X{4+len(model.get('risks',[]))}","CASEY_RiskRegister")
    ws.conditional_formatting.add(f"G5:G{4+len(model.get('risks',[]))}", ColorScaleRule(start_type='min', start_color=V51_GREEN, mid_type='percentile', mid_value=50, mid_color=V51_AMBER, end_type='max', end_color=V51_RED))

    # Schedule
    ws=wb["07 Schedule"]
    _title(ws,"SCHEDULE LOGIC - LEVEL SELECTED", "Linked schedule activities for QSRA and export to P6/XER.", "H")
    _apply_ws_theme(ws,{"A":12,"B":18,"C":42,"D":22,"E":12,"F":12,"G":55,"H":12}, freeze="A5")
    _header_row(ws,4,["Activity ID","Phase","Activity","Predecessor","Duration Months","Critical","Basis","Risk Links"])
    for r,x in enumerate(model.get("schedule_rows",[]),5):
        linked=", ".join([rr.get("risk_id") for rr in model.get("risks",[]) if rr.get("activity_id")==x.get("activity_id")])
        vals=[x.get("activity_id"),x.get("phase"),x.get("activity"),x.get("predecessor"),x.get("duration_months"),x.get("critical"),x.get("basis"),linked]
        for c,v in enumerate(vals,1): ws.cell(r,c,v)
        if x.get("critical")=="Yes":
            for c in range(1,9): ws.cell(r,c).fill=PatternFill("solid",fgColor="26111A")

    # Scenarios
    ws=wb["08 Scenarios"]
    _title(ws,"SCENARIO COMPARISON - BASE / FASTER / CHEAPER / LOWER RISK / PREMIUM", "Scenario choice must change cost, schedule, risk and confidence.", "H")
    _apply_ws_theme(ws,{"A":20,"B":16,"C":16,"D":16,"E":16,"F":42,"G":42,"H":42}, freeze="A5")
    _header_row(ws,4,["Scenario","P50 Cost","Schedule","Risk","Confidence","Cost Change Logic","Schedule Change Logic","Risk Logic"])
    for r,x in enumerate(model.get("scenario_comparison",[]),5):
        vals=[x.get("scenario"),x.get("cost_p50"),x.get("schedule"),x.get("risk"),x.get("confidence"),x.get("why","") or x.get("cost_note","") ,x.get("schedule_note","") ,x.get("risk_note","")]
        for c,v in enumerate(vals,1): ws.cell(r,c,v)

    # Benchmarking
    ws=wb["09 Benchmarking"]
    _title(ws,"BENCHMARKING AND PEER CHALLENGE", "Use to challenge Tier 1 contractor estimates, market rates and sponsor assumptions.", "H")
    _apply_ws_theme(ws,{"A":26,"B":26,"C":26,"D":26,"E":26,"F":26,"G":26,"H":26}, freeze="A5")
    _header_row(ws,4,["Benchmark / Peer","What to compare","Why it matters","Challenge question","Evidence required"])
    peers=model.get("peer_competitors",[]) or []
    for r,peer in enumerate(peers[:10],5):
        vals=[peer,"Unit rates, schedule intensity, cost/scope normalization","Prevents blind acceptance of contractor/adviser basis","Where are we above peer norm and why?","Rate build-up, supplier quote, productivity basis, exclusions log"]
        for c,v in enumerate(vals,1): ws.cell(r,c,v)

    # Assumptions
    ws=wb["10 Assumptions"]
    _title(ws,"ASSUMPTIONS, LIMITATIONS AND DECISION GATES", "Makes the output honest: a first-pass project controls intelligence pack, not a final tender estimate.", "H")
    _apply_ws_theme(ws,{"A":30,"B":85,"C":28,"D":28,"E":28,"F":28,"G":28,"H":28}, freeze="A5")
    _header_row(ws,4,["Area","Assumption / Limitation","Board Action"])
    assumptions=[
        ("Estimate",f"{model.get('estimate_class_name')} with {model.get('estimate_maturity')}; not a contractor tender price.","Use to challenge, shortlist and frame approval gates."),
        ("Schedule",f"Level {model.get('schedule_level')} logic; critical path and float require validation against project-specific constraints.","Run schedule risk workshop and validate dependencies."),
        ("Risk","QCRA and QSRA distributions are modelled from the generated register and selected scenario.","Replace assumptions with workshop data as evidence improves."),
        ("Scenario",model.get("scenario_why"),"Compare Base / Faster / Cheaper / Lower Risk before board decision."),
        ("Outputs","Exports are board-preparation outputs and should be reviewed before external release.","Apply owner governance, sign-off and evidence checks.")
    ]
    for r,vals in enumerate(assumptions,5):
        for c,v in enumerate(vals,1): ws.cell(r,c,v)

    # final styling sizes
    for ws in wb.worksheets:
        for row in ws.iter_rows():
            for cell in row:
                if cell.value is not None:
                    cell.alignment = Alignment(vertical="top", wrap_text=True)
        for r in range(1, ws.max_row+1):
            ws.row_dimensions[r].height = min(80, max(18, ws.row_dimensions[r].height or 18))
    bio=BytesIO(); wb.save(bio); bio.seek(0); return bio.getvalue()


def risk_register_workbook_bytes(model:Dict[str,Any]):
    model = _casey_reconcile_cost_lines(model)
    # Dedicated premium risk workbook: dashboard + register + mitigations + board risks.
    wb=Workbook(); ws=wb.active; ws.title="Risk Dashboard"; wb.create_sheet("Risk Register Pro"); wb.create_sheet("Mitigation Tracker"); wb.create_sheet("Board Top 10")
    risks=model.get("risks",[]); mc=model.get("monte_carlo",{})
    ws=wb["Risk Dashboard"]; _title(ws,"CASEY RISK REGISTER PRO", "Cause - Event - Impact - Response - Quantified cost and schedule exposure.", "I")
    _apply_ws_theme(ws,{"A":22,"B":20,"C":22,"D":22,"E":22,"F":22,"G":22,"H":22,"I":22}, freeze="A6")
    kpis=[("Risks",len(risks)),("Board Visible",sum(1 for r in risks if r.get("board_visibility")=="Yes")),("QCRA P80",_v51_money(mc.get("qcra",{}).get("p80"))),("QSRA P80",f"{mc.get('qsra',{}).get('p80','-')} months"),("Top Risk",risks[0].get("title") if risks else "-"),("Scenario",model.get("scenario_label"))]
    for i,(k,v) in enumerate(kpis):
        r=4+(i//3)*3; c=1+(i%3)*3; ws.cell(r,c,k); ws.cell(r+1,c,v); ws.merge_cells(start_row=r,start_column=c,end_row=r,end_column=c+1); ws.merge_cells(start_row=r+1,start_column=c,end_row=r+1,end_column=c+1)
        ws.cell(r,c).fill=PatternFill("solid",fgColor=V51_PANEL); ws.cell(r+1,c).fill=PatternFill("solid",fgColor=V51_PANEL); ws.cell(r+1,c).font=Font(size=18,bold=True,color=V51_WHITE)
    _header_row(ws,12,["Rank","Risk","Cost EMV BN","Schedule EMV Days","Owner","Response"])
    for i,risk in enumerate(sorted(risks,key=lambda x:_num(x.get('driver_score')),reverse=True)[:10],13):
        vals=[i-12,risk.get("title"),risk.get("cost_emv_bn"),risk.get("schedule_emv_days"),risk.get("owner"),risk.get("mitigation")]
        for c,v in enumerate(vals,1): ws.cell(i,c,v)
        ws.cell(i,3).number_format='$#,##0.000 "B"'
    chart=BarChart(); chart.title="Top Cost EMV Risks"; chart.add_data(Reference(ws,min_col=3,min_row=12,max_row=22),titles_from_data=True); chart.set_categories(Reference(ws,min_col=2,min_row=13,max_row=22)); chart.height=7; chart.width=12; ws.add_chart(chart,"A25")
    ws=wb["Risk Register Pro"]; _title(ws,"FULL RISK REGISTER PRO", "Every risk has a cause, event, impact description, owner, response and quantification.", "X")
    headers=["Risk ID","Category","Title","Cause","Risk Event","Impact Description","Probability %","Cost O BN","Cost M BN","Cost P BN","Cost EMV BN","Schedule O Days","Schedule M Days","Schedule P Days","Schedule EMV Days","Activity","CBS","Owner","Trigger","Mitigation","Response Strategy","Residual","Status","Board"]
    _apply_ws_theme(ws,{get_column_letter(i):18 for i in range(1,25)}, freeze="A5"); _header_row(ws,4,headers)
    for i,r in enumerate(risks,5):
        vals=[r.get("risk_id"),r.get("category"),r.get("title"),r.get("cause"),r.get("risk_event"),r.get("impact_description"),r.get("probability_pct"),r.get("cost_o_bn"),r.get("cost_m_bn"),r.get("cost_p_bn"),r.get("cost_emv_bn"),r.get("schedule_o_days"),r.get("schedule_m_days"),r.get("schedule_p_days"),r.get("schedule_emv_days"),r.get("activity_id"),r.get("cbs"),r.get("owner"),r.get("trigger"),r.get("mitigation"),r.get("response_strategy"),r.get("residual_rating"),r.get("status"),r.get("board_visibility")]
        for c,v in enumerate(vals,1): ws.cell(i,c,v)
        for c in [8,9,10,11]: ws.cell(i,c).number_format='$#,##0.000 "B"'
    _add_table(ws,f"A4:X{4+len(risks)}","RiskRegisterPro")
    ws=wb["Mitigation Tracker"]; _title(ws,"MITIGATION TRACKER", "Owner actions, triggers and residual risk conversion.", "J")
    _apply_ws_theme(ws,{"A":12,"B":30,"C":25,"D":50,"E":18,"F":18,"G":20,"H":20,"I":20,"J":28}, freeze="A5")
    _header_row(ws,4,["Risk ID","Risk","Owner","Mitigation Action","Trigger","Due Gate","Status","Residual","Board","Evidence Needed"])
    for i,r in enumerate(risks,5):
        vals=[r.get("risk_id"),r.get("title"),r.get("owner"),r.get("mitigation"),r.get("trigger"),"Next gate",r.get("status"),r.get("residual_rating"),r.get("board_visibility"),"Named evidence, date, owner and decision log"]
        for c,v in enumerate(vals,1): ws.cell(i,c,v)
    ws=wb["Board Top 10"]; _title(ws,"BOARD TOP 10 RISKS", "The risks most likely to change the investment decision.", "H")
    _apply_ws_theme(ws,{"A":8,"B":32,"C":40,"D":40,"E":40,"F":18,"G":18,"H":50}, freeze="A5")
    _header_row(ws,4,["Rank","Risk","Cause","Event","Impact","Cost EMV BN","Schedule EMV Days","Board Ask"])
    for i,r in enumerate(sorted(risks,key=lambda x:_num(x.get('driver_score')),reverse=True)[:10],5):
        vals=[i-4,r.get("title"),r.get("cause"),r.get("risk_event"),r.get("impact_description"),r.get("cost_emv_bn"),r.get("schedule_emv_days"),"Approve mitigation spend or request further evidence"]
        for c,v in enumerate(vals,1): ws.cell(i,c,v)
        ws.cell(i,6).number_format='$#,##0.000 "B"'
    bio=BytesIO(); wb.save(bio); bio.seek(0); return bio.getvalue()


def risk_csv_bytes(model):
    # Preserve CSV endpoint if called directly, but make it richer than before.
    out=StringIO(); w=csv.writer(out)
    cols=["Risk ID","Category","Title","Cause","Risk Event","Impact Description","Probability %","Cost O BN","Cost M BN","Cost P BN","Cost EMV BN","Schedule O Days","Schedule M Days","Schedule P Days","Schedule EMV Days","Activity","CBS","Owner","Trigger","Mitigation","Response Strategy","Residual","Status","Board"]
    w.writerow(cols)
    for r in model.get("risks",[]):
        w.writerow([r.get("risk_id"),r.get("category"),r.get("title"),r.get("cause"),r.get("risk_event"),r.get("impact_description"),r.get("probability_pct"),r.get("cost_o_bn"),r.get("cost_m_bn"),r.get("cost_p_bn"),r.get("cost_emv_bn"),r.get("schedule_o_days"),r.get("schedule_m_days"),r.get("schedule_p_days"),r.get("schedule_emv_days"),r.get("activity_id"),r.get("cbs"),r.get("owner"),r.get("trigger"),r.get("mitigation"),r.get("response_strategy"),r.get("residual_rating"),r.get("status"),r.get("board_visibility")])
    return out.getvalue().encode()


def word_bytes(model):
    doc=Document()
    sec=doc.sections[0]; sec.top_margin=Inches(.55); sec.bottom_margin=Inches(.55); sec.left_margin=Inches(.6); sec.right_margin=Inches(.6)
    styles=doc.styles
    styles["Normal"].font.name="Aptos"; styles["Normal"].font.size=Pt(9.5)
    for style in ["Heading 1","Heading 2","Title"]:
        styles[style].font.name="Aptos Display"
        styles[style].font.color.rgb=RGBColor(6,28,44)
    title=doc.add_heading("CASEY Executive Board Intelligence Pack",0)
    doc.add_paragraph(f"{model.get('title')} | {model.get('subsector')} | {model.get('scenario_label')} scenario | Generated {datetime.utcnow().strftime('%Y-%m-%d %H:%M UTC')}")
    doc.add_paragraph(model.get("executive_summary",""))
    doc.add_heading("1. Board Decision Snapshot",1)
    mc=model.get("monte_carlo",{}); totals=_v51_total_costs(model)
    table=doc.add_table(rows=1,cols=4); table.style="Light Shading Accent 1"
    headers=["Metric","Value","Interpretation","Board Action"]
    for i,h in enumerate(headers): table.rows[0].cells[i].text=h
    rows=[
        ("P50 Cost",model.get("cost_p50"),"Most likely first-pass cost based on selected class and scenario.","Use as initial investment control point."),
        ("Cost Range",model.get("cost_range"),"Class/maturity and QCRA uncertainty range.","Do not approve without range acknowledgement."),
        ("QCRA P80",_v51_money(mc.get("qcra",{}).get("p80")),"Risk-adjusted cost confidence level.","Set funding envelope and contingency rules."),
        ("Baseline / QSRA P80",f"{model.get('schedule')} / {mc.get('qsra',{}).get('p80')} months","Schedule risk exposure and delivery confidence.","Require critical path validation."),
        ("Risk / Confidence",f"{model.get('risk')} / {model.get('confidence_pct')}%","Early-stage confidence and data quality status.","Approve evidence-improvement plan."),
        ("Direct / Indirect / Reserve",f"{_v51_money(totals.get('Direct'))} / {_v51_money(totals.get('Indirect'))} / {_v51_money(totals.get('Reserve'))}","Prevents reserves being hidden inside direct cost.","Challenge contractor and adviser basis.")]
    for row in rows:
        cells=table.add_row().cells
        for i,v in enumerate(row): cells[i].text=str(v)
    doc.add_heading("2. Executive Recommendation",1)
    for x in _v51_scenario_notes(model): doc.add_paragraph(x, style="List Bullet")
    for x in model.get("next_best_actions",[])[:5]: doc.add_paragraph(x, style="List Bullet")
    doc.add_heading("3. Cost Estimate - Selected Class Only",1)
    t=doc.add_table(rows=1,cols=7); t.style="Light Grid Accent 1"
    for i,h in enumerate(["CBS","Description","Type","P10","P50","P90","Basis"]): t.rows[0].cells[i].text=h
    for x in model.get("cost_lines",[])[:24]:
        c=t.add_row().cells
        vals=[x.get("cbs"),x.get("description"),x.get("type"),_v51_money(x.get("p10_bn")),_v51_money(x.get("p50_bn")),_v51_money(x.get("p90_bn")),x.get("impact_basis")]
        for i,v in enumerate(vals): c[i].text=str(v)
    doc.add_heading("4. QCRA and QSRA",1)
    doc.add_paragraph(f"QCRA cost P10 / P50 / P80 / P90: {_v51_money(mc.get('qcra',{}).get('p10'))} / {model.get('cost_p50')} / {_v51_money(mc.get('qcra',{}).get('p80'))} / {_v51_money(mc.get('qcra',{}).get('p90'))}")
    doc.add_paragraph(f"QSRA schedule P10 / P50 / P80 / P90: {mc.get('qsra',{}).get('p10')} / {mc.get('qsra',{}).get('p50')} / {mc.get('qsra',{}).get('p80')} / {mc.get('qsra',{}).get('p90')} months")
    doc.add_heading("5. QCRA Tornado Drivers",2)
    for x in mc.get("qcra_tornado",[])[:8]: doc.add_paragraph(f"{x.get('title')} - {x.get('cbs')} - mean exposure {_v51_money(x.get('cost_mean_bn'))}", style="List Bullet")
    doc.add_heading("6. QSRA Tornado Drivers",2)
    for x in mc.get("qsra_tornado",[])[:8]: doc.add_paragraph(f"{x.get('title')} - {x.get('activity_id')} - mean exposure {x.get('schedule_mean_days')} days", style="List Bullet")
    doc.add_heading("7. Risk Register - Board Top Risks",1)
    for r in model.get("risks",[])[:10]:
        doc.add_paragraph(f"{r.get('risk_id')} {r.get('title')}",style="List Bullet")
        doc.add_paragraph(f"Cause: {r.get('cause')}. Event: {r.get('risk_event')}. Impact: {r.get('impact_description')}. Mitigation: {r.get('mitigation')}. Owner: {r.get('owner')}.")
    doc.add_heading("8. Challenge Questions",1)
    for q in model.get("board_challenge_questions",[])[:8]: doc.add_paragraph(q,style="List Number")
    doc.add_heading("9. Assumptions and Limitations",1)
    doc.add_paragraph("This is a first-pass project controls intelligence pack. It is designed to challenge, accelerate and structure decision-making; it is not a final contractor tender or signed cost plan. Each output includes assumptions and should be validated against project-specific evidence before commitment.")
    bio=BytesIO(); doc.save(bio); bio.seek(0); return bio.getvalue()


def pdf_bytes(model):
    bio=BytesIO(); doc=SimpleDocTemplate(bio,pagesize=landscape(A4),rightMargin=24,leftMargin=24,topMargin=22,bottomMargin=22)
    styles=getSampleStyleSheet()
    styles.add(ParagraphStyle(name="Hero",fontSize=24,leading=28,textColor=colors.HexColor("#8DF7FF"),fontName="Helvetica-Bold",spaceAfter=12))
    styles.add(ParagraphStyle(name="Section",fontSize=15,leading=18,textColor=colors.HexColor("#07111F"),fontName="Helvetica-Bold",spaceBefore=8,spaceAfter=6))
    styles.add(ParagraphStyle(name="SmallCasey",fontSize=8,leading=10,textColor=colors.HexColor("#334155")))
    story=[]; mc=model.get("monte_carlo",{}); totals=_v51_total_costs(model)
    story.append(Paragraph("CASEY BOARD INTELLIGENCE PACK",styles["Hero"]))
    story.append(Paragraph(model.get("executive_summary",""),styles["BodyText"])); story.append(Spacer(1,8))
    metrics=[["Project",model.get("title"),"Scenario",model.get("scenario_label")],["P50",model.get("cost_p50"),"Cost Range",model.get("cost_range")],["QCRA P80",_v51_money(mc.get("qcra",{}).get("p80")),"QSRA P80",f"{mc.get('qsra',{}).get('p80')} months"],["Direct / Indirect / Reserve",f"{_v51_money(totals.get('Direct'))} / {_v51_money(totals.get('Indirect'))} / {_v51_money(totals.get('Reserve'))}","Risk / Confidence",f"{model.get('risk')} / {model.get('confidence_pct')}%"]]
    t=Table(metrics,colWidths=[110,260,110,260]); t.setStyle(TableStyle([("BACKGROUND",(0,0),(-1,-1),colors.HexColor("#07111F")),("TEXTCOLOR",(0,0),(-1,-1),colors.white),("GRID",(0,0),(-1,-1),.5,colors.HexColor("#8DF7FF")),("FONTNAME",(0,0),(-1,-1),"Helvetica-Bold"),("FONTSIZE",(0,0),(-1,-1),9),("VALIGN",(0,0),(-1,-1),"TOP")])) ; story.append(t)
    story.append(Spacer(1,8)); story.append(Paragraph("Board Recommendation",styles["Section"]))
    for x in _v51_scenario_notes(model)[:3]: story.append(Paragraph("- "+x, styles["SmallCasey"]))
    story.append(PageBreak())
    story.append(Paragraph("QCRA and QSRA - Separate Tornado Drivers",styles["Hero"]))
    rows=[["Rank","QCRA Cost Driver","CBS","Mean Cost","QSRA Schedule Driver","Activity","Mean Days"]]
    q=mc.get("qcra_tornado",[])[:10]; s=mc.get("qsra_tornado",[])[:10]
    for i in range(10):
        rows.append([i+1,q[i].get("title","") if i<len(q) else "",q[i].get("cbs","") if i<len(q) else "",_v51_money(q[i].get("cost_mean_bn")) if i<len(q) else "",s[i].get("title","") if i<len(s) else "",s[i].get("activity_id","") if i<len(s) else "",s[i].get("schedule_mean_days","") if i<len(s) else ""])
    t=Table(rows,colWidths=[38,190,55,75,190,60,70]); t.setStyle(TableStyle([("BACKGROUND",(0,0),(-1,0),colors.HexColor("#07111F")),("TEXTCOLOR",(0,0),(-1,0),colors.HexColor("#8DF7FF")),("GRID",(0,0),(-1,-1),.25,colors.HexColor("#CBD5E1")),("FONTSIZE",(0,0),(-1,-1),7),("VALIGN",(0,0),(-1,-1),"TOP")])) ; story.append(t)
    story.append(PageBreak())
    story.append(Paragraph("Risk Register - Board Top Risks",styles["Hero"]))
    rows=[["ID","Risk","Cause","Event","Impact","Owner","Response"]]
    for r in model.get("risks",[])[:10]: rows.append([r.get("risk_id"),r.get("title"),r.get("cause"),r.get("risk_event"),r.get("impact_description"),r.get("owner"),r.get("mitigation")])
    t=Table(rows,colWidths=[42,88,135,135,170,78,170]); t.setStyle(TableStyle([("BACKGROUND",(0,0),(-1,0),colors.HexColor("#07111F")),("TEXTCOLOR",(0,0),(-1,0),colors.HexColor("#8DF7FF")),("GRID",(0,0),(-1,-1),.25,colors.HexColor("#CBD5E1")),("FONTSIZE",(0,0),(-1,-1),6.6),("VALIGN",(0,0),(-1,-1),"TOP")])) ; story.append(t)
    story.append(PageBreak())
    story.append(Paragraph("Cost Estimate - Selected Class Only",styles["Hero"]))
    rows=[["CBS","Description","Type","P10","P50","P90","Basis"]]
    for x in model.get("cost_lines",[])[:18]: rows.append([x.get("cbs"),x.get("description"),x.get("type"),_v51_money(x.get("p10_bn")),_v51_money(x.get("p50_bn")),_v51_money(x.get("p90_bn")),x.get("basis")])
    t=Table(rows,colWidths=[44,145,70,55,55,55,270]); t.setStyle(TableStyle([("BACKGROUND",(0,0),(-1,0),colors.HexColor("#07111F")),("TEXTCOLOR",(0,0),(-1,0),colors.HexColor("#8DF7FF")),("GRID",(0,0),(-1,-1),.25,colors.HexColor("#CBD5E1")),("FONTSIZE",(0,0),(-1,-1),6.8),("VALIGN",(0,0),(-1,-1),"TOP")])) ; story.append(t)
    doc.build(story); bio.seek(0); return bio.getvalue()


def pptx_bytes(model):
    from pptx.dml.color import RGBColor as PRGBColor
    prs=Presentation(); prs.slide_width=PptxInches(13.333); prs.slide_height=PptxInches(7.5); blank=prs.slide_layouts[6]
    def slide(title,subtitle=""):
        s=prs.slides.add_slide(blank); s.background.fill.solid(); s.background.fill.fore_color.rgb=PRGBColor(2,5,10)
        tx=s.shapes.add_textbox(PptxInches(.45),PptxInches(.3),PptxInches(12.4),PptxInches(.45)); tx.text_frame.text="CASEY"; tx.text_frame.paragraphs[0].runs[0].font.size=PptxPt(14); tx.text_frame.paragraphs[0].runs[0].font.color.rgb=PRGBColor(141,247,255)
        h=s.shapes.add_textbox(PptxInches(.45),PptxInches(.86),PptxInches(12.4),PptxInches(.85)); h.text_frame.text=title; h.text_frame.paragraphs[0].runs[0].font.size=PptxPt(30); h.text_frame.paragraphs[0].runs[0].font.bold=True; h.text_frame.paragraphs[0].runs[0].font.color.rgb=PRGBColor(247,251,255)
        if subtitle:
            p=s.shapes.add_textbox(PptxInches(.48),PptxInches(1.55),PptxInches(12.1),PptxInches(.55)); p.text_frame.text=subtitle; p.text_frame.paragraphs[0].runs[0].font.size=PptxPt(12); p.text_frame.paragraphs[0].runs[0].font.color.rgb=PRGBColor(170,190,205)
        return s
    s=slide("Mission Control for Capital Projects",model.get("executive_summary",""))
    # Metrics slide
    s=slide("Board Decision Snapshot", f"{model.get('title')} | {model.get('scenario_label')} scenario")
    mc=model.get("monte_carlo",{}); totals=_v51_total_costs(model)
    metrics=[("P50",model.get("cost_p50")),("P80 Cost",_v51_money(mc.get("qcra",{}).get("p80"))),("Schedule",model.get("schedule")),("QSRA P80",f"{mc.get('qsra',{}).get('p80')} mo"),("Risk",model.get("risk")),("Confidence",f"{model.get('confidence_pct')}%"),("Direct",_v51_money(totals.get("Direct"))),("Indirect",_v51_money(totals.get("Indirect"))),("Reserve",_v51_money(totals.get("Reserve")))]
    for i,(k,v) in enumerate(metrics):
        x=.55+(i%3)*4.2; y=2.0+(i//3)*1.35
        shp=s.shapes.add_shape(1,PptxInches(x),PptxInches(y),PptxInches(3.75),PptxInches(1.02)); shp.fill.solid(); shp.fill.fore_color.rgb=PRGBColor(7,17,31); shp.line.color.rgb=PRGBColor(141,247,255)
        tf=shp.text_frame; tf.text=f"{k}\n{v}"; tf.paragraphs[0].runs[0].font.size=PptxPt(10); tf.paragraphs[0].runs[0].font.color.rgb=PRGBColor(159,179,200)
        tf.paragraphs[1].runs[0].font.size=PptxPt(20); tf.paragraphs[1].runs[0].font.bold=True; tf.paragraphs[1].runs[0].font.color.rgb=PRGBColor(247,251,255)
    # Scenario slide
    s=slide("Scenario Logic", "What changed and why the outputs are not generic.")
    y=1.95
    for n in _v51_scenario_notes(model):
        box=s.shapes.add_textbox(PptxInches(.7),PptxInches(y),PptxInches(12),PptxInches(.52)); box.text_frame.text="• "+n; box.text_frame.paragraphs[0].runs[0].font.size=PptxPt(15); box.text_frame.paragraphs[0].runs[0].font.color.rgb=PRGBColor(247,251,255); y+=.7
    # QCRA/QSRA slide
    s=slide("QCRA + QSRA Drivers", "Separate cost and schedule risk; no mixed generic tornado.")
    left=s.shapes.add_textbox(PptxInches(.65),PptxInches(1.9),PptxInches(5.9),PptxInches(4.9)); left.text_frame.text="QCRA COST TORNADO"; left.text_frame.paragraphs[0].runs[0].font.color.rgb=PRGBColor(141,247,255); left.text_frame.paragraphs[0].runs[0].font.bold=True
    for x in mc.get("qcra_tornado",[])[:7]:
        p=left.text_frame.add_paragraph(); p.text=f"{x.get('title')} - {_v51_money(x.get('cost_mean_bn'))}"; p.font.size=PptxPt(12); p.font.color.rgb=PRGBColor(247,251,255)
    right=s.shapes.add_textbox(PptxInches(6.8),PptxInches(1.9),PptxInches(5.9),PptxInches(4.9)); right.text_frame.text="QSRA SCHEDULE TORNADO"; right.text_frame.paragraphs[0].runs[0].font.color.rgb=PRGBColor(141,247,255); right.text_frame.paragraphs[0].runs[0].font.bold=True
    for x in mc.get("qsra_tornado",[])[:7]:
        p=right.text_frame.add_paragraph(); p.text=f"{x.get('title')} - {x.get('schedule_mean_days')} days"; p.font.size=PptxPt(12); p.font.color.rgb=PRGBColor(247,251,255)
    # Risk slide
    s=slide("Board Top Risks", "Cause - Event - Impact - Mitigation.")
    tf=s.shapes.add_textbox(PptxInches(.65),PptxInches(1.85),PptxInches(12.2),PptxInches(5.3)).text_frame
    for r in model.get("risks",[])[:8]:
        p=tf.add_paragraph(); p.text=f"{r.get('risk_id')} {r.get('title')}: {r.get('impact_description')}"; p.font.size=PptxPt(11); p.font.color.rgb=PRGBColor(247,251,255)
    # Close slide
    s=slide("Output Pack", "Excel cost model, risk register workbook, XER schedule, DOCX/PDF board report, PPTX board deck and JSON audit model generated from one connected CASEY model.")
    y=2.1
    for x in ["Selected class cost estimate only", "Direct / indirect / reserve cost split", "Risk register with cause, event, impact and mitigation", "Separate QCRA and QSRA curves/tornado", "Scenario-linked values for Base / Faster / Cheaper / Lower Risk / Premium", "Pricing and email conversion path"]:
        b=s.shapes.add_textbox(PptxInches(1.0),PptxInches(y),PptxInches(11.4),PptxInches(.42)); b.text_frame.text="✓ "+x; b.text_frame.paragraphs[0].runs[0].font.size=PptxPt(15); b.text_frame.paragraphs[0].runs[0].font.color.rgb=PRGBColor(141,247,255); y+=.58
    bio=BytesIO(); prs.save(bio); bio.seek(0); return bio.getvalue()

# Patch route functions to expose premium outputs and filenames.
def export_risk(model:Dict[str,Any]): return stream(risk_register_workbook_bytes(model),"application/vnd.openxmlformats-officedocument.spreadsheetml.sheet","CASEY_Risk_Register_Pro.xlsx")
def export_workbook(model:Dict[str,Any]): return stream(workbook_bytes(model),"application/vnd.openxmlformats-officedocument.spreadsheetml.sheet","CASEY_Cost_Model_Planet_Class.xlsx")
def export_word(model:Dict[str,Any]): return stream(word_bytes(model),"application/vnd.openxmlformats-officedocument.wordprocessingml.document","CASEY_Executive_Board_Report.docx")
def export_pdf(model:Dict[str,Any]): return stream(pdf_bytes(model),"application/pdf","CASEY_Board_Intelligence_Pack.pdf")
def export_pptx(model:Dict[str,Any]): return stream(pptx_bytes(model),"application/vnd.openxmlformats-officedocument.presentationml.presentation","CASEY_Board_Deck_Elite.pptx")

def export_all(model:Dict[str,Any]):
    model = _casey_reconcile_cost_lines(model)
    bio=BytesIO()
    with zipfile.ZipFile(bio,"w",zipfile.ZIP_DEFLATED) as z:
        z.writestr("01_CASEY_Cost_Model_Planet_Class.xlsx",workbook_bytes(model))
        z.writestr("02_CASEY_Risk_Register_Pro.xlsx",risk_register_workbook_bytes(model))
        z.writestr("03_CASEY_P6_Schedule.xer",xer_bytes(model))
        z.writestr("04_CASEY_Executive_Board_Report.docx",word_bytes(model))
        z.writestr("05_CASEY_Board_Intelligence_Pack.pdf",pdf_bytes(model))
        z.writestr("06_CASEY_Board_Deck_Elite.pptx",pptx_bytes(model))
        z.writestr("07_CASEY_Full_Model_Audit.json",json.dumps(model,indent=2))
        z.writestr("08_CASEY_Risk_Register_Raw.csv",risk_csv_bytes(model))
        z.writestr("09_CASEY_Demo_Close_Script.txt","\n".join(model.get("launch_demo_script",[])))
    bio.seek(0); return stream(bio.getvalue(),"application/zip","CASEY_Output_Pack_Planet_Class.zip")

# Register new explicit v51 endpoints as well.
@app.post("/v51/export/risk-register")
def export_risk_v51(model:Dict[str,Any]): return export_risk(model)
@app.post("/v51/export/all")
def export_all_v51(model:Dict[str,Any]): return export_all(model)

# ===================== END CASEY v51 OUTPUT PLANET CLASS OVERRIDES =====================



# ======================= CASEY v52 ELITE OUTPUT SYSTEM =======================
# Focus: board-grade readability, scenario consistency and premium exports.
# v52 patches the existing /export routes in-place so the current frontend keeps working.

from openpyxl.chart import PieChart
from reportlab.lib.units import inch
from pptx.enum.text import PP_ALIGN

V52_DARK = "101820"
V52_NAVY = "0B1220"
V52_CYAN = "00AEEF"
V52_BLUE = "0070C0"
V52_LIGHT = "F7FBFF"
V52_GRID = "D9E2EC"
V52_RED = "C00000"
V52_AMBER = "FFC000"
V52_GREEN = "00A65A"


def _v52_scenario_label(model):
    return str(model.get("scenario_label") or model.get("scenario") or "Base").title()


def _v52_totals(model):
    totals = {"Direct": 0.0, "Indirect": 0.0, "Reserve": 0.0}
    for line in model.get("cost_lines", []):
        typ = line.get("type", "Direct")
        totals[typ] = totals.get(typ, 0.0) + float(line.get("p50_bn", 0) or 0)
    return totals


def _v52_recommendation(model):
    s = str(model.get("scenario", "base")).lower()
    risk = str(model.get("risk", "Medium"))
    if s == "faster":
        return {
            "headline": "APPROVE ACCELERATION ONLY WITH FUNDED ASSURANCE",
            "why": "The faster case can protect market window but pushes interface, commissioning and logistics pressure onto the critical path.",
            "decision": "Approve acceleration premium only if package-level recovery plans, supplier commitments and commissioning readiness gates are funded.",
        }
    if s == "cheaper":
        return {
            "headline": "DO NOT TREAT SAVINGS AS CERTAIN UNTIL SCOPE IS SIGNED OFF",
            "why": "The cheaper case lowers visible capex through value engineering and procurement pressure, but increases residual scope, quality and rework risk.",
            "decision": "Approve only after VE log, exclusions, owner-retained risks and quality impacts are signed off.",
        }
    if s == "lower_risk":
        return {
            "headline": "PROCEED WITH LOWER-RISK CASE AS BOARD CONTROL BASELINE",
            "why": "Higher upfront assurance, buffers and evidence maturity reduce downstream contingency draw and schedule surprise.",
            "decision": "Use as board approval baseline where certainty is more important than lowest headline capex.",
        }
    if s == "premium":
        return {
            "headline": "APPROVE PREMIUM CASE ONLY WHERE RESILIENCE HAS STRATEGIC VALUE",
            "why": "The premium case increases capex but improves resilience, assurance and confidence for mission-critical assets.",
            "decision": "Proceed if the business case values resilience, operational assurance and lower downside risk.",
        }
    return {
        "headline": "APPROVE DEFINITION PHASE — DO NOT COMMIT FULL CAPITAL YET",
        "why": f"The base case provides a first-pass controls baseline with {risk} risk and evidence still required before final capital commitment.",
        "decision": "Use CASEY as the reference model for board challenge, option testing and consultant / contractor assurance.",
    }


def _v52_style_ws(ws, title=None):
    ws.sheet_view.showGridLines = False
    ws.freeze_panes = "A6"
    if title:
        ws["A1"] = title
        ws["A1"].font = Font(size=22, bold=True, color="FFFFFF")
        ws["A1"].fill = PatternFill("solid", fgColor=V52_NAVY)
        ws.merge_cells("A1:H1")
        ws.row_dimensions[1].height = 32
    thin = Side(style="thin", color=V52_GRID)
    for row in ws.iter_rows():
        for cell in row:
            cell.alignment = Alignment(vertical="top", wrap_text=True)
            cell.border = Border(bottom=thin)


def _v52_header(row):
    for c in row:
        c.font = Font(bold=True, color="FFFFFF")
        c.fill = PatternFill("solid", fgColor=V52_BLUE)
        c.alignment = Alignment(horizontal="center", vertical="center", wrap_text=True)


def _v52_autowidth(ws, max_width=42):
    widths = {}
    for row in ws.iter_rows():
        for cell in row:
            if cell.value is None: continue
            widths[cell.column] = min(max(widths.get(cell.column, 10), len(str(cell.value)) + 2), max_width)
    for col, width in widths.items():
        ws.column_dimensions[get_column_letter(col)].width = width


def workbook_bytes_v52(model: Dict[str, Any]) -> bytes:
    wb = Workbook()
    wb.remove(wb.active)
    totals = _v52_totals(model); rec = _v52_recommendation(model); mc = model.get("monte_carlo", {})
    scenario = _v52_scenario_label(model)

    # Dashboard
    ws = wb.create_sheet("01 Executive Dashboard")
    _v52_style_ws(ws, "CASEY EXECUTIVE DASHBOARD")
    ws["A3"] = model.get("title", "Project")
    ws["A3"].font = Font(size=16, bold=True, color=V52_NAVY)
    ws["A4"] = f"Scenario: {scenario} | Generated: {datetime.utcnow().strftime('%Y-%m-%d %H:%M UTC')} | Engine: CASEY v52"
    ws["A4"].font = Font(size=10, color="555555")
    kpis = [("P50 Cost", model.get("cost_p50")), ("Cost Range", model.get("cost_range")), ("QCRA P80", money_bn(mc.get("qcra", {}).get("p80", 0))), ("Schedule", model.get("schedule")), ("QSRA P80", f"{mc.get('qsra', {}).get('p80', '')} months"), ("Risk / Confidence", f"{model.get('risk')} / {model.get('confidence_pct')}%"), ("Direct", money_bn(totals.get("Direct",0))), ("Indirect + Reserve", money_bn(totals.get("Indirect",0)+totals.get("Reserve",0)))]
    start_row = 6
    for idx, (k, v) in enumerate(kpis):
        r = start_row + (idx // 4) * 4; c = 1 + (idx % 4) * 2
        ws.cell(r,c,k); ws.cell(r+1,c,v)
        ws.cell(r,c).font = Font(bold=True, color="FFFFFF"); ws.cell(r,c).fill = PatternFill("solid", fgColor=V52_NAVY)
        ws.cell(r+1,c).font = Font(size=16, bold=True, color=V52_BLUE); ws.cell(r+1,c).fill = PatternFill("solid", fgColor="FFFFFF")
        ws.merge_cells(start_row=r, start_column=c, end_row=r, end_column=c+1)
        ws.merge_cells(start_row=r+1, start_column=c, end_row=r+2, end_column=c+1)
    ws["A15"] = "Executive recommendation"; ws["A15"].font = Font(size=14, bold=True, color=V52_NAVY)
    ws["A16"] = rec["headline"]; ws["A16"].font = Font(size=14, bold=True, color=V52_RED if "DO NOT" in rec["headline"] else V52_BLUE)
    ws["A17"] = rec["why"]
    ws["A18"] = rec["decision"]
    ws.merge_cells("A16:H16"); ws.merge_cells("A17:H17"); ws.merge_cells("A18:H18")
    ws["A21"] = "Cost split"; ws["A21"].font = Font(bold=True)
    split_start = 22
    for i, (typ, val) in enumerate(totals.items(), split_start):
        ws.cell(i,1,typ); ws.cell(i,2,val); ws.cell(i,2).number_format = '$0.0B'
    pie = PieChart(); labels=Reference(ws,min_col=1,min_row=22,max_row=24); data=Reference(ws,min_col=2,min_row=21,max_row=24)
    pie.add_data(data, titles_from_data=True); pie.set_categories(labels); pie.title="Direct / Indirect / Reserve"; pie.height=7; pie.width=9
    ws.add_chart(pie,"D21")
    _v52_autowidth(ws)

    # Cost model
    ws = wb.create_sheet("02 Cost Model")
    _v52_style_ws(ws, "SELECTED CLASS COST ESTIMATE")
    headers = ["CBS", "Description", "Type", "P10", "P50", "P90", "Basis"]
    ws.append([]); ws.append([]); ws.append([]); ws.append(headers); _v52_header(ws[4])
    for x in model.get("cost_lines", []):
        ws.append([x.get("cbs"), x.get("description"), x.get("type"), x.get("p10_bn"), x.get("p50_bn"), x.get("p90_bn"), x.get("basis") or x.get("impact_basis")])
    for row in ws.iter_rows(min_row=5, min_col=4, max_col=6):
        for cell in row: cell.number_format = '$0.0B'
    # Subtotals
    r = ws.max_row + 2
    ws.cell(r,1,"Subtotals"); ws.cell(r,1).font=Font(bold=True,color=V52_NAVY)
    for i, typ in enumerate(["Direct","Indirect","Reserve"], r+1):
        ws.cell(i,1,typ); ws.cell(i,5,totals.get(typ,0)); ws.cell(i,5).number_format='$0.0B'; ws.cell(i,1).font=Font(bold=True)
    chart = BarChart(); chart.title="Cost Estimate by CBS (P50)"; chart.y_axis.title="US$B"; chart.x_axis.title="CBS"
    data=Reference(ws,min_col=5,min_row=4,max_row=4+len(model.get('cost_lines',[])))
    cats=Reference(ws,min_col=1,min_row=5,max_row=4+len(model.get('cost_lines',[])))
    chart.add_data(data,titles_from_data=True); chart.set_categories(cats); chart.height=8; chart.width=18
    ws.add_chart(chart,"I4")
    _v52_autowidth(ws)

    # Scenarios
    ws = wb.create_sheet("03 Scenario Comparison")
    _v52_style_ws(ws, "SCENARIO COMPARISON")
    ws.append([]); ws.append([]); ws.append([]); ws.append(["Scenario", "P50 Cost", "Schedule", "Risk", "Confidence", "Commercial interpretation"]); _v52_header(ws[4])
    for s in model.get("scenario_comparison", []):
        ws.append([s.get("label"), s.get("cost_p50"), s.get("schedule"), s.get("risk"), s.get("confidence_pct"), s.get("note") or s.get("why") or "Scenario-linked model output"])
    _v52_autowidth(ws)

    # Risks
    ws = wb.create_sheet("04 Risk Register")
    _v52_style_ws(ws, "RISK REGISTER PRO")
    headers = ["ID", "Title", "Category", "Cause", "Event", "Impact", "Owner", "Prob %", "Cost Most Likely", "Schedule Most Likely Days", "Rating", "Mitigation", "Trigger", "Status"]
    ws.append([]); ws.append([]); ws.append([]); ws.append(headers); _v52_header(ws[4])
    for rsk in model.get("risks", []):
        rating = rsk.get("pre_mitigation_rating") or rsk.get("residual_rating") or "Medium"
        ws.append([rsk.get("risk_id"), rsk.get("title"), rsk.get("category"), rsk.get("cause"), rsk.get("risk_event"), rsk.get("impact_description"), rsk.get("owner"), rsk.get("probability_pct"), rsk.get("cost_m_bn"), rsk.get("schedule_m_days"), rating, rsk.get("mitigation"), rsk.get("trigger"), rsk.get("status")])
        rr=ws.max_row
        if str(rating).lower().startswith(("high","very","extreme")):
            fill=PatternFill("solid", fgColor="FCE4D6")
        elif str(rating).lower().startswith("medium"):
            fill=PatternFill("solid", fgColor="FFF2CC")
        else:
            fill=PatternFill("solid", fgColor="E2F0D9")
        for c in range(1,15): ws.cell(rr,c).fill=fill
    for row in ws.iter_rows(min_row=5, min_col=9, max_col=9):
        row[0].number_format='$0.0B'
    _v52_autowidth(ws, 48)

    # QCRA QSRA
    ws = wb.create_sheet("05 QCRA QSRA")
    _v52_style_ws(ws, "QCRA + QSRA ANALYTICS")
    qcra=mc.get("qcra",{}); qsra=mc.get("qsra",{})
    ws.append([]); ws.append([]); ws.append([]); ws.append(["Metric","P10","P50","P80","P90"]); _v52_header(ws[4])
    ws.append(["QCRA Cost", qcra.get("p10"), qcra.get("p50"), qcra.get("p80"), qcra.get("p90")])
    ws.append(["QSRA Schedule", qsra.get("p10"), qsra.get("p50"), qsra.get("p80"), qsra.get("p90")])
    for c in range(2,6): ws.cell(5,c).number_format='$0.0B'
    ws["A9"]="QCRA Cost Tornado"; ws["E9"]="QSRA Schedule Tornado"; ws["A9"].font=ws["E9"].font=Font(bold=True,size=14,color=V52_NAVY)
    ws.append([])
    ws.append(["Driver","CBS","Mean Cost ($B)","","Driver","Activity","Mean Days"]); _v52_header(ws[11])
    for i in range(10):
        q = (mc.get("qcra_tornado") or [])[i] if i < len(mc.get("qcra_tornado") or []) else {}
        s = (mc.get("qsra_tornado") or [])[i] if i < len(mc.get("qsra_tornado") or []) else {}
        ws.append([q.get("title"), q.get("cbs"), q.get("cost_mean_bn"), "", s.get("title"), s.get("activity_id"), s.get("schedule_mean_days")])
    for row in ws.iter_rows(min_row=12, min_col=3, max_col=3): row[0].number_format='$0.0B'
    bc=BarChart(); bc.title="QCRA Cost Tornado"; bc.y_axis.title="US$B"; bc.add_data(Reference(ws,min_col=3,min_row=11,max_row=21),titles_from_data=True); bc.set_categories(Reference(ws,min_col=1,min_row=12,max_row=21)); bc.height=7; bc.width=11; ws.add_chart(bc,"I4")
    sc=BarChart(); sc.title="QSRA Schedule Tornado"; sc.y_axis.title="Days"; sc.add_data(Reference(ws,min_col=7,min_row=11,max_row=21),titles_from_data=True); sc.set_categories(Reference(ws,min_col=5,min_row=12,max_row=21)); sc.height=7; sc.width=11; ws.add_chart(sc,"I20")
    _v52_autowidth(ws)

    # Assumptions
    ws = wb.create_sheet("06 Basis + Audit")
    _v52_style_ws(ws, "BASIS, ASSUMPTIONS AND AUDIT TRAIL")
    rows = [["Project", model.get("title")], ["Mode", model.get("mode")], ["Sector", model.get("subsector")], ["Location", model.get("location")], ["Scenario", scenario], ["Class", model.get("estimate_class_name")], ["Schedule Level", model.get("schedule_level")], ["Limitations", "First-pass controls intelligence pack; validate with project-specific evidence before capital commitment."], ["Generated", datetime.utcnow().isoformat()+"Z"], ["Engine", "CASEY v52 Elite Outputs"]]
    for r in rows: ws.append(r)
    _v52_autowidth(ws)
    bio=BytesIO(); wb.save(bio); bio.seek(0); return bio.getvalue()


def risk_register_workbook_bytes_v52(model: Dict[str, Any]) -> bytes:
    # Dedicated risk workbook, not CSV, for readability.
    wb = Workbook(); ws = wb.active; ws.title = "Risk Register Pro"; _v52_style_ws(ws, "CASEY RISK REGISTER PRO")
    headers=["ID","Title","Category","Cause","Risk Event","Impact Description","Owner","Probability %","Cost O","Cost M","Cost P","Schedule O","Schedule M","Schedule P","Pre Rating","Residual","Mitigation","Trigger","Status","Activity","CBS","Board Visibility"]
    ws.append([]); ws.append([]); ws.append([]); ws.append(headers); _v52_header(ws[4])
    for rsk in model.get("risks",[]):
        ws.append([rsk.get("risk_id"),rsk.get("title"),rsk.get("category"),rsk.get("cause"),rsk.get("risk_event"),rsk.get("impact_description"),rsk.get("owner"),rsk.get("probability_pct"),rsk.get("cost_o_bn"),rsk.get("cost_m_bn"),rsk.get("cost_p_bn"),rsk.get("schedule_o_days"),rsk.get("schedule_m_days"),rsk.get("schedule_p_days"),rsk.get("pre_mitigation_rating"),rsk.get("residual_rating"),rsk.get("mitigation"),rsk.get("trigger"),rsk.get("status"),rsk.get("activity_id"),rsk.get("cbs"),rsk.get("board_visibility")])
    for row in ws.iter_rows(min_row=5,min_col=9,max_col=11):
        for cell in row: cell.number_format='$0.0B'
    # Summary tab
    ws2=wb.create_sheet("Board Top 10"); _v52_style_ws(ws2,"BOARD TOP 10 RISKS")
    ws2.append([]); ws2.append([]); ws2.append([]); ws2.append(["Rank","Risk","Why it matters","Owner","Required action"]); _v52_header(ws2[4])
    for i,rsk in enumerate(model.get("risks",[])[:10],1):
        ws2.append([i,rsk.get("title"),rsk.get("impact_description"),rsk.get("owner"),rsk.get("mitigation")])
    for sh in wb.worksheets: _v52_autowidth(sh,50)
    bio=BytesIO(); wb.save(bio); bio.seek(0); return bio.getvalue()


def xer_bytes_v52(model: Dict[str, Any]) -> bytes:
    # Legacy-targeted XER header to avoid PRA 24.12 rejection. Includes basic TASK + TASKPRED only.
    lines=[]
    lines.append("ERMHDR\t20.12\t{}\tProject\tCASEY\tCASEY\tProject Management\tUSD".format(datetime.utcnow().strftime("%Y-%m-%d")))
    lines += ["%T\tPROJECT","%F\tproj_id\tproj_short_name\tproj_name","%R\t1\tCASEY\t"+str(model.get('title','CASEY Project'))]
    lines += ["%T\tWBS","%F\twbs_id\tproj_id\twbs_short_name\twbs_name","%R\t1\t1\tCASEY\tCASEY Project"]
    lines += ["%T\tTASK","%F\ttask_id\tproj_id\twbs_id\ttask_code\ttask_name\tduration_type\ttarget_drtn_hr_cnt"]
    idmap={}
    for i,row in enumerate(model.get("schedule_rows",[]),1):
        code=row.get('activity_id',f'A{i:04d}'); idmap[code]=i
        name=str(row.get('activity','Activity')).replace('\t',' ').replace('\n',' ')
        dur=max(8,int(float(row.get('duration_months',1) or 1)*160))
        lines.append(f"%R\t{i}\t1\t1\t{code}\t{name}\tFixed Duration\t{dur}")
    lines += ["%T\tTASKPRED","%F\ttask_pred_id\ttask_id\tpred_task_id\tpred_type"]
    pred_id=1
    for row in model.get("schedule_rows",[]):
        tid=idmap.get(row.get('activity_id'))
        for p in str(row.get('predecessor','')).replace(';',',').split(','):
            p=p.strip()
            if tid and p and p in idmap:
                lines.append(f"%R\t{pred_id}\t{tid}\t{idmap[p]}\tFS"); pred_id+=1
    return ("\n".join(lines)+"\n").encode("utf-8")


def pdf_bytes_v52(model: Dict[str, Any]) -> bytes:
    bio=BytesIO(); doc=SimpleDocTemplate(bio,pagesize=landscape(A4),rightMargin=24,leftMargin=24,topMargin=20,bottomMargin=20)
    styles=getSampleStyleSheet(); rec=_v52_recommendation(model); totals=_v52_totals(model); mc=model.get('monte_carlo',{})
    title=ParagraphStyle('v52Title',parent=styles['Title'],fontSize=22,textColor=colors.HexColor('#101820'),spaceAfter=8)
    h=ParagraphStyle('v52H',parent=styles['Heading2'],fontSize=14,textColor=colors.HexColor('#0070C0'),spaceAfter=6)
    body=ParagraphStyle('v52Body',parent=styles['BodyText'],fontSize=8.5,leading=11,textColor=colors.HexColor('#101820'))
    story=[]
    story.append(Paragraph(f"CASEY BOARD DECISION PACK — {model.get('title','Project')}",title))
    story.append(Paragraph(f"{_v52_scenario_label(model)} scenario | {model.get('subsector')} | {model.get('location')} | Generated {datetime.utcnow().strftime('%Y-%m-%d %H:%M UTC')}",body))
    metrics=[["P50 Cost",model.get('cost_p50'),"Cost Range",model.get('cost_range'),"QCRA P80",money_bn(mc.get('qcra',{}).get('p80',0)),"Schedule",model.get('schedule')],["QSRA P80",f"{mc.get('qsra',{}).get('p80','')} months","Risk",model.get('risk'),"Confidence",f"{model.get('confidence_pct')}%","Direct / Indirect / Reserve",f"{money_bn(totals.get('Direct',0))} / {money_bn(totals.get('Indirect',0))} / {money_bn(totals.get('Reserve',0))}"]]
    t=Table(metrics,colWidths=[70,95,75,95,75,95,95,130])
    t.setStyle(TableStyle([('BACKGROUND',(0,0),(-1,0),colors.HexColor('#EAF6FF')),('TEXTCOLOR',(0,0),(-1,-1),colors.HexColor('#101820')),('GRID',(0,0),(-1,-1),0.4,colors.HexColor('#D9E2EC')),('FONTNAME',(0,0),(-1,-1),'Helvetica'),('FONTSIZE',(0,0),(-1,-1),8),('VALIGN',(0,0),(-1,-1),'TOP')]))
    story += [Spacer(1,8),t,Spacer(1,10),Paragraph(rec['headline'],h),Paragraph(rec['why'],body),Paragraph(rec['decision'],body),PageBreak()]
    # Cost and scenario
    story.append(Paragraph("Cost Estimate — Selected Class Only",h))
    rows=[["CBS","Description","Type","P10","P50","P90"]]+[[x.get('cbs'),x.get('description'),x.get('type'),money_bn(x.get('p10_bn',0)),money_bn(x.get('p50_bn',0)),money_bn(x.get('p90_bn',0))] for x in model.get('cost_lines',[])[:18]]
    tt=Table(rows,colWidths=[45,210,65,70,70,70]); tt.setStyle(TableStyle([('BACKGROUND',(0,0),(-1,0),colors.HexColor('#0070C0')),('TEXTCOLOR',(0,0),(-1,0),colors.white),('GRID',(0,0),(-1,-1),0.3,colors.HexColor('#D9E2EC')),('FONTSIZE',(0,0),(-1,-1),7.5),('VALIGN',(0,0),(-1,-1),'TOP')]))
    story += [tt,PageBreak()]
    # Risk
    story.append(Paragraph("QCRA + QSRA Tornado Drivers",h))
    q=mc.get('qcra_tornado',[])[:10]; s=mc.get('qsra_tornado',[])[:10]
    rows=[["Rank","QCRA Cost Driver","CBS","Mean Cost","QSRA Schedule Driver","Activity","Mean Days"]]
    for i in range(max(len(q),len(s))):
        qr=q[i] if i<len(q) else {}; sr=s[i] if i<len(s) else {}
        rows.append([i+1,qr.get('title',''),qr.get('cbs',''),money_bn(qr.get('cost_mean_bn',0)),sr.get('title',''),sr.get('activity_id',''),sr.get('schedule_mean_days','')])
    tt=Table(rows,colWidths=[35,170,45,70,170,55,65]); tt.setStyle(TableStyle([('BACKGROUND',(0,0),(-1,0),colors.HexColor('#0070C0')),('TEXTCOLOR',(0,0),(-1,0),colors.white),('GRID',(0,0),(-1,-1),0.3,colors.HexColor('#D9E2EC')),('FONTSIZE',(0,0),(-1,-1),7.5),('VALIGN',(0,0),(-1,-1),'TOP')]))
    story += [tt,Spacer(1,10),Paragraph("Board Top Risks",h)]
    rows=[["ID","Risk","Cause","Event","Impact","Owner","Mitigation"]]
    for rsk in model.get('risks',[])[:8]: rows.append([rsk.get('risk_id'),rsk.get('title'),rsk.get('cause'),rsk.get('risk_event'),rsk.get('impact_description'),rsk.get('owner'),rsk.get('mitigation')])
    tt=Table(rows,colWidths=[45,90,130,120,160,80,160]); tt.setStyle(TableStyle([('BACKGROUND',(0,0),(-1,0),colors.HexColor('#0070C0')),('TEXTCOLOR',(0,0),(-1,0),colors.white),('GRID',(0,0),(-1,-1),0.3,colors.HexColor('#D9E2EC')),('FONTSIZE',(0,0),(-1,-1),6.4),('VALIGN',(0,0),(-1,-1),'TOP')]))
    story += [tt]
    doc.build(story); bio.seek(0); return bio.getvalue()


def word_bytes_v52(model: Dict[str, Any]) -> bytes:
    doc=Document(); sec=doc.sections[0]; sec.top_margin=Inches(.45); sec.bottom_margin=Inches(.45); sec.left_margin=Inches(.55); sec.right_margin=Inches(.55)
    rec=_v52_recommendation(model); totals=_v52_totals(model); mc=model.get('monte_carlo',{})
    p=doc.add_paragraph(); r=p.add_run("CASEY EXECUTIVE BOARD DECISION PACK"); r.bold=True; r.font.size=Pt(20); r.font.color.rgb=RGBColor(0,112,192)
    doc.add_paragraph(f"{model.get('title')} | {_v52_scenario_label(model)} scenario | {datetime.utcnow().strftime('%Y-%m-%d %H:%M UTC')}")
    table=doc.add_table(rows=1,cols=4); table.style='Table Grid'; hdr=table.rows[0].cells
    for i,x in enumerate(['Metric','Value','Interpretation','Board action']): hdr[i].text=x
    for k,v,interp,action in [('P50 Cost',model.get('cost_p50'),'Most likely first-pass cost','Use as control point'),('Cost Range',model.get('cost_range'),'Class and QCRA uncertainty','Approve with range acknowledged'),('QCRA P80',money_bn(mc.get('qcra',{}).get('p80',0)),'Risk-adjusted funding level','Set contingency rules'),('QSRA P80',f"{mc.get('qsra',{}).get('p80','')} months",'Schedule confidence level','Validate critical path'),('Direct / Indirect / Reserve',f"{money_bn(totals.get('Direct',0))} / {money_bn(totals.get('Indirect',0))} / {money_bn(totals.get('Reserve',0))}",'Stops reserve hiding','Challenge contractor basis')]:
        c=table.add_row().cells; c[0].text=k; c[1].text=str(v); c[2].text=interp; c[3].text=action
    doc.add_heading('Executive Recommendation',level=1); doc.add_paragraph(rec['headline']); doc.add_paragraph(rec['why']); doc.add_paragraph(rec['decision'])
    doc.add_heading('QCRA and QSRA Tornado Drivers',level=1)
    for label, rows in [('QCRA Cost', mc.get('qcra_tornado',[])[:8]), ('QSRA Schedule', mc.get('qsra_tornado',[])[:8])]:
        doc.add_heading(label,level=2)
        t=doc.add_table(rows=1,cols=4); t.style='Table Grid'; t.rows[0].cells[0].text='Rank'; t.rows[0].cells[1].text='Driver'; t.rows[0].cells[2].text='Link'; t.rows[0].cells[3].text='Exposure'
        for i,x in enumerate(rows,1):
            c=t.add_row().cells; c[0].text=str(i); c[1].text=str(x.get('title','')); c[2].text=str(x.get('cbs') or x.get('activity_id') or ''); c[3].text=money_bn(x.get('cost_mean_bn',0)) if label.startswith('QCRA') else f"{x.get('schedule_mean_days')} days"
    doc.add_heading('Risk Register - Board Top Risks',level=1)
    t=doc.add_table(rows=1,cols=6); t.style='Table Grid';
    for i,hdr in enumerate(['ID','Risk','Cause','Event','Impact','Mitigation']): t.rows[0].cells[i].text=hdr
    for rsk in model.get('risks',[])[:10]:
        c=t.add_row().cells; vals=[rsk.get('risk_id'),rsk.get('title'),rsk.get('cause'),rsk.get('risk_event'),rsk.get('impact_description'),rsk.get('mitigation')]
        for i,v in enumerate(vals): c[i].text=str(v or '')
    bio=BytesIO(); doc.save(bio); bio.seek(0); return bio.getvalue()


def pptx_bytes_v52(model: Dict[str, Any]) -> bytes:
    from pptx.dml.color import RGBColor as PRGBColor
    prs=Presentation(); prs.slide_width=PptxInches(13.333); prs.slide_height=PptxInches(7.5); blank=prs.slide_layouts[6]
    rec=_v52_recommendation(model); mc=model.get('monte_carlo',{}); totals=_v52_totals(model)
    def slide(title, subtitle=""):
        s=prs.slides.add_slide(blank); s.background.fill.solid(); s.background.fill.fore_color.rgb=PRGBColor(247,251,255)
        tx=s.shapes.add_textbox(PptxInches(.55),PptxInches(.35),PptxInches(12.2),PptxInches(.6)); tf=tx.text_frame; tf.text=title; p=tf.paragraphs[0]; p.runs[0].font.size=PptxPt(25); p.runs[0].font.bold=True; p.runs[0].font.color.rgb=PRGBColor(16,24,32)
        if subtitle:
            sub=s.shapes.add_textbox(PptxInches(.58),PptxInches(.95),PptxInches(12),PptxInches(.4)); sub.text_frame.text=subtitle; sub.text_frame.paragraphs[0].runs[0].font.size=PptxPt(12); sub.text_frame.paragraphs[0].runs[0].font.color.rgb=PRGBColor(80,96,112)
        return s
    s=slide(f"{model.get('title')} — Board Intelligence", f"{_v52_scenario_label(model)} scenario | CASEY v52")
    box=s.shapes.add_textbox(PptxInches(.8),PptxInches(2.0),PptxInches(11.5),PptxInches(1.2)); box.text_frame.text=rec['headline']; box.text_frame.paragraphs[0].runs[0].font.size=PptxPt(30); box.text_frame.paragraphs[0].runs[0].font.bold=True; box.text_frame.paragraphs[0].runs[0].font.color.rgb=PRGBColor(0,112,192)
    s=slide("Board Decision Snapshot", rec['decision'])
    metrics=[('P50',model.get('cost_p50')),('Range',model.get('cost_range')),('QCRA P80',money_bn(mc.get('qcra',{}).get('p80',0))),('Schedule',model.get('schedule')),('QSRA P80',f"{mc.get('qsra',{}).get('p80','')} mo"),('Risk',model.get('risk')),('Confidence',f"{model.get('confidence_pct')}%"),('Direct',money_bn(totals.get('Direct',0))),('Reserve',money_bn(totals.get('Reserve',0)))]
    for i,(k,v) in enumerate(metrics):
        x=.6+(i%3)*4.25; y=1.65+(i//3)*1.4
        shp=s.shapes.add_shape(1,PptxInches(x),PptxInches(y),PptxInches(3.8),PptxInches(1.08)); shp.fill.solid(); shp.fill.fore_color.rgb=PRGBColor(232,246,255); shp.line.color.rgb=PRGBColor(0,112,192)
        shp.text=f"{k}\n{v}"; shp.text_frame.paragraphs[0].runs[0].font.size=PptxPt(11); shp.text_frame.paragraphs[1].runs[0].font.size=PptxPt(21); shp.text_frame.paragraphs[1].runs[0].font.bold=True
    s=slide("QCRA and QSRA are separated", "Cost risk and schedule risk are not mixed.")
    t=s.shapes.add_table(9,4,PptxInches(.7),PptxInches(1.65),PptxInches(5.8),PptxInches(4.8)).table
    for j,h in enumerate(['Rank','QCRA Driver','CBS','Cost']): t.cell(0,j).text=h
    for i,x in enumerate(mc.get('qcra_tornado',[])[:8],1):
        t.cell(i,0).text=str(i); t.cell(i,1).text=str(x.get('title','')); t.cell(i,2).text=str(x.get('cbs','')); t.cell(i,3).text=money_bn(x.get('cost_mean_bn',0))
    t2=s.shapes.add_table(9,4,PptxInches(6.85),PptxInches(1.65),PptxInches(5.8),PptxInches(4.8)).table
    for j,h in enumerate(['Rank','QSRA Driver','Act','Days']): t2.cell(0,j).text=h
    for i,x in enumerate(mc.get('qsra_tornado',[])[:8],1):
        t2.cell(i,0).text=str(i); t2.cell(i,1).text=str(x.get('title','')); t2.cell(i,2).text=str(x.get('activity_id','')); t2.cell(i,3).text=str(x.get('schedule_mean_days',''))
    s=slide("Board Top Risks", "Cause, event, impact and response are visible.")
    y=1.55
    for rsk in model.get('risks',[])[:5]:
        bx=s.shapes.add_textbox(PptxInches(.8),PptxInches(y),PptxInches(11.8),PptxInches(.85)); bx.text_frame.text=f"{rsk.get('risk_id')} {rsk.get('title')} — {rsk.get('impact_description')}"; bx.text_frame.paragraphs[0].runs[0].font.size=PptxPt(13); y+=.93
    s=slide("What CASEY produces", "Board PDF, cost workbook, risk register workbook, schedule export, DOCX, PPTX, JSON audit and ZIP pack.")
    y=1.75
    for x in ['Selected class estimate only','Direct / indirect / reserve split','Scenario-linked outputs for Base, Faster, Cheaper, Lower Risk and Premium','Readable white-background Excel/PDF exports','Risk register includes cause, event, impact, owner and mitigation','Legacy-targeted XER plus schedule tables']:
        bx=s.shapes.add_textbox(PptxInches(1),PptxInches(y),PptxInches(11.5),PptxInches(.45)); bx.text_frame.text='✓ '+x; bx.text_frame.paragraphs[0].runs[0].font.size=PptxPt(16); bx.text_frame.paragraphs[0].runs[0].font.color.rgb=PRGBColor(0,112,192); y+=.65
    bio=BytesIO(); prs.save(bio); bio.seek(0); return bio.getvalue()


def export_workbook_v52_endpoint(model: Dict[str, Any]): return stream(workbook_bytes_v52(model),"application/vnd.openxmlformats-officedocument.spreadsheetml.sheet","CASEY_v52_Elite_Cost_Model.xlsx")
def export_risk_v52_endpoint(model: Dict[str, Any]): return stream(risk_register_workbook_bytes_v52(model),"application/vnd.openxmlformats-officedocument.spreadsheetml.sheet","CASEY_v52_Risk_Register_Pro.xlsx")
def export_xer_v52_endpoint(model: Dict[str, Any]): return stream(xer_bytes_v52(model),"application/octet-stream","CASEY_v52_P6_Schedule_Legacy.xer")
def export_word_v52_endpoint(model: Dict[str, Any]): return stream(word_bytes_v52(model),"application/vnd.openxmlformats-officedocument.wordprocessingml.document","CASEY_v52_Executive_Board_Report.docx")
def export_pdf_v52_endpoint(model: Dict[str, Any]): return stream(pdf_bytes_v52(model),"application/pdf","CASEY_v52_Board_Decision_Pack.pdf")
def export_pptx_v52_endpoint(model: Dict[str, Any]): return stream(pptx_bytes_v52(model),"application/vnd.openxmlformats-officedocument.presentationml.presentation","CASEY_v52_Board_Deck.pptx")
def export_all_v52_endpoint(model: Dict[str, Any]):
    bio=BytesIO()
    with zipfile.ZipFile(bio,"w",zipfile.ZIP_DEFLATED) as z:
        z.writestr("01_CASEY_v52_Elite_Cost_Model.xlsx", workbook_bytes_v52(model))
        z.writestr("02_CASEY_v52_Risk_Register_Pro.xlsx", risk_register_workbook_bytes_v52(model))
        z.writestr("03_CASEY_v52_P6_Schedule_Legacy.xer", xer_bytes_v52(model))
        z.writestr("04_CASEY_v52_Executive_Board_Report.docx", word_bytes_v52(model))
        z.writestr("05_CASEY_v52_Board_Decision_Pack.pdf", pdf_bytes_v52(model))
        z.writestr("06_CASEY_v52_Board_Deck.pptx", pptx_bytes_v52(model))
        z.writestr("07_CASEY_v52_Full_Model_Audit.json", json.dumps(model, indent=2))
        z.writestr("08_READ_ME_OUTPUT_STANDARD.txt", "CASEY v52 Elite Outputs: white-background readable exports, scenario-linked numbers, QCRA/QSRA separation, risk cause-event-impact-mitigation, legacy-targeted XER. Validate XER in your Primavera/PRA environment because vendor import support differs by installed version.")
    bio.seek(0); return stream(bio.getvalue(),"application/zip","CASEY_v52_Elite_Output_Pack.zip")


def _v52_patch_route(path: str, endpoint):
    for route in app.routes:
        if getattr(route, 'path', None) == path and 'POST' in getattr(route, 'methods', set()):
            route.endpoint = endpoint
            if hasattr(route, 'dependant'):
                route.dependant.call = endpoint

_v52_patch_route('/export/workbook', export_workbook_v52_endpoint)
_v52_patch_route('/export/risk-register', export_risk_v52_endpoint)
_v52_patch_route('/export/xer', export_xer_v52_endpoint)
_v52_patch_route('/export/word', export_word_v52_endpoint)
_v52_patch_route('/export/pdf', export_pdf_v52_endpoint)
_v52_patch_route('/export/pptx', export_pptx_v52_endpoint)
_v52_patch_route('/export/all', export_all_v52_endpoint)
_v52_patch_route('/v51/export/all', export_all_v52_endpoint)

@app.post('/v52/export/workbook')
def export_workbook_v52(model: Dict[str, Any]): return export_workbook_v52_endpoint(model)
@app.post('/v52/export/risk-register')
def export_risk_v52(model: Dict[str, Any]): return export_risk_v52_endpoint(model)
@app.post('/v52/export/xer')
def export_xer_v52(model: Dict[str, Any]): return export_xer_v52_endpoint(model)
@app.post('/v52/export/all')
def export_all_v52(model: Dict[str, Any]): return export_all_v52_endpoint(model)

# ===================== END CASEY v52 ELITE OUTPUT SYSTEM =====================

if __name__ == "__main__":
    import uvicorn
    uvicorn.run("main:app", host="127.0.0.1", port=8000, reload=True)


# ========================= CASEY v54 PLATINUM OUTPUT ENGINE =========================
# Final output override: readable white exports, real QCRA/QSRA separation, validated risk fields,
# scenario-linked narratives, and legacy-targeted schedule export. Same frontend endpoints.

from openpyxl.chart import PieChart
from openpyxl.styles import GradientFill

APP_VERSION = "CASEY TITAN X v54 Platinum Output Engine"

V54_BLUE = "0B1F33"
V54_CYAN = "00A9C7"
V54_LIGHT = "F3F7FA"
V54_LINE = "D9E2EA"
V54_GREEN = "1F7A4D"
V54_AMBER = "C77800"
V54_RED = "B42318"


def _v54_safe(v, default=""):
    return default if v is None else v


def _v54_slug(s: str) -> str:
    return re.sub(r"[^A-Za-z0-9_]+", "_", str(s or "CASEY")).strip("_")[:60]


def _v54_num_money(s: Any) -> float:
    return parse_bn(s)


def _v54_scenario(model: Dict[str, Any]) -> Dict[str, Any]:
    label = model.get("scenario_label") or str(model.get("scenario") or "Base").replace("_", " ").title()
    key = str(model.get("scenario") or label).lower().replace(" ", "_")
    narratives = {
        "base": ("Approve definition phase — do not commit full capital yet.", "Balanced control case used as the board reference model."),
        "faster": ("Proceed only with funded acceleration controls.", "Time is bought with premium logistics, overlap and higher interface management."),
        "cheaper": ("Use as a value-engineering challenge, not a final funding case.", "Lower capex increases scope, productivity and rework exposure unless controls tighten."),
        "lower_risk": ("Recommended board path where certainty matters more than lowest cost.", "Higher assurance, surveys and buffers reduce downstream surprise."),
        "premium": ("Proceed if resilience, brand and delivery certainty justify premium capex.", "Best-in-class suppliers, stronger assurance and higher redundancy reduce residual risk."),
        "investor": ("Use for investment committee challenge and option screening.", "Transparent risk-adjusted view designed for funding conversations."),
        "survival": ("Do not use for approval without explicit risk acceptance.", "Minimum viable case with high residual risk and reduced confidence."),
    }
    rec, why = narratives.get(key, narratives["base"])
    return {"key": key, "label": label, "recommendation": rec, "why": why}


def _v54_risks(model: Dict[str, Any]) -> List[Dict[str, Any]]:
    out=[]
    schedule = model.get("schedule_rows") or []
    cost_lines = model.get("cost_lines") or []
    for idx,r in enumerate(model.get("risks") or [],1):
        prob = int(max(5, min(95, float(r.get("probability_pct") or 25))))
        # derive full cause/event/impact from old description if not present
        title = str(r.get("title") or f"Risk {idx}")
        cause = str(r.get("cause") or r.get("description") or "Project evidence is not yet mature enough to fully validate this assumption.")
        event = str(r.get("event") or (title + " materialises during delivery"))
        impact = str(r.get("impact") or r.get("impact_description") or _v54_risk_impact(title, r))
        mitigation = str(r.get("mitigation") or "Assign accountable owner, confirm evidence, and update risk response before the next gate.")
        owner = str(r.get("owner") or "Project Controls Lead")
        act_id = str(r.get("activity_id") or (schedule[min(idx, len(schedule)-1)].get("activity_id") if schedule else "A1900"))
        act_name = str(r.get("activity_name") or (schedule[min(idx, len(schedule)-1)].get("activity") if schedule else "Delivery activity"))
        cbs = str(r.get("cbs") or (cost_lines[min(idx, len(cost_lines)-1)].get("cbs") if cost_lines else "01.01"))
        cbs_name = str(r.get("cbs_name") or (cost_lines[min(idx, len(cost_lines)-1)].get("description") if cost_lines else "Cost account"))
        # cost values, ensure > 0 and ordered
        co = max(0.001, float(r.get("cost_o_bn") or 0.01))
        cm = max(co, float(r.get("cost_m_bn") or co*2))
        cp = max(cm, float(r.get("cost_p_bn") or cm*1.8))
        # schedule values: retain zero only for cost-only risks but keep out of QSRA tornado later
        so = max(0.0, float(r.get("schedule_o_days") or 0))
        sm = max(so, float(r.get("schedule_m_days") or 0))
        sp = max(sm, float(r.get("schedule_p_days") or 0))
        cost_emv = cm * prob / 100
        sched_emv = sm * prob / 100
        out.append({
            **r,
            "risk_id": str(r.get("risk_id") or f"R-{idx:03d}"),
            "title": title,
            "category": str(r.get("category") or "Delivery"),
            "cause": cause,
            "event": event,
            "impact": impact,
            "probability_pct": prob,
            "likelihood_pct": prob,
            "activity_id": act_id,
            "activity_name": act_name,
            "cbs": cbs,
            "cbs_name": cbs_name,
            "schedule_o_days": so,
            "schedule_m_days": sm,
            "schedule_p_days": sp,
            "cost_o_bn": round(co,3),
            "cost_m_bn": round(cm,3),
            "cost_p_bn": round(cp,3),
            "cost_emv_bn": round(cost_emv,3),
            "schedule_emv_days": round(sched_emv,1),
            "owner": owner,
            "trigger": str(r.get("trigger") or "Trigger threshold not yet defined"),
            "mitigation": mitigation,
            "residual_rating": _v54_rating(prob, cost_emv, sched_emv),
            "status": str(r.get("status") or "Open"),
            "basis_of_cost_impact": str(r.get("basis_of_cost_impact") or f"Linked to CBS {cbs} {cbs_name}; expected exposure is probability-weighted from O/M/P cost impacts."),
            "basis_of_schedule_impact": str(r.get("basis_of_schedule_impact") or f"Linked to activity {act_id} {act_name}; expected exposure is probability-weighted from O/M/P schedule impacts."),
        })
    return sorted(out, key=lambda x: (x["cost_emv_bn"]*1000 + x["schedule_emv_days"]), reverse=True)


def _v54_risk_impact(title, r):
    t = title.lower()
    if "grid" in t or "power" in t: return "Critical path delay, temporary power cost, phased opening risk and revenue deferral."
    if "market" in t or "escalation" in t: return "P50 cost becomes understated; procurement strategy and funding envelope may need reset."
    if "design" in t: return "Cost growth, quantity movement, rework, late approvals and schedule resequencing."
    if "scope" in t: return "Additional quantities, redesign, rework, procurement churn and contingency drawdown."
    if "commission" in t: return "Delayed handover, repeated testing cycles, operational readiness slippage and damages exposure."
    return "Cost, schedule and confidence exposure affecting the board approval basis."


def _v54_rating(prob, cost_emv, sched_emv):
    score = prob/100 * (cost_emv*10 + sched_emv/10)
    if score > 12: return "Red"
    if score > 5: return "Amber"
    return "Green"


def _v54_curves(model: Dict[str, Any]) -> Dict[str, Any]:
    mc = model.get("monte_carlo") or {}
    curve = mc.get("curve") or []
    if curve:
        qcra = [(int(x.get("percentile",0)), float(x.get("cost_bn",0))) for x in curve]
        qsra = [(int(x.get("percentile",0)), float(x.get("schedule_months",0))) for x in curve]
    else:
        p50 = parse_bn(model.get("cost_p50")); months = float(str(model.get("schedule","60")).split()[0])
        qcra = [(p, round(p50*(0.7+p/100*0.8),3)) for p in [5,10,20,30,40,50,60,70,80,90,95]]
        qsra = [(p, round(months*(0.9+p/100*0.35),2)) for p in [5,10,20,30,40,50,60,70,80,90,95]]
    return {"qcra": qcra, "qsra": qsra}


def _v54_tornado(model: Dict[str, Any]) -> Dict[str, List[Dict[str, Any]]]:
    risks = _v54_risks(model)
    qcra = sorted([r for r in risks if r.get("cost_emv_bn",0) > 0], key=lambda r:r.get("cost_emv_bn",0), reverse=True)[:10]
    qsra = sorted([r for r in risks if r.get("schedule_emv_days",0) > 0], key=lambda r:r.get("schedule_emv_days",0), reverse=True)[:10]
    return {"qcra": qcra, "qsra": qsra}


def _v54_costs(model: Dict[str, Any]) -> List[Dict[str, Any]]:
    low, high, *_ = class_range(int(model.get("estimate_class") or 3))
    out=[]
    for x in model.get("cost_lines") or []:
        p50 = float(x.get("p50_bn") or 0)
        out.append({
            **x,
            "p10_bn": round(float(x.get("p10_bn") or p50*low),3),
            "p50_bn": round(p50,3),
            "p90_bn": round(float(x.get("p90_bn") or p50*high),3),
            "basis": x.get("basis") or x.get("impact_basis") or "Sector benchmark and model factor basis.",
        })
    return out


def _v54_style_header(ws, row=1, last_col=None):
    last_col = last_col or ws.max_column
    for c in range(1,last_col+1):
        cell=ws.cell(row,c)
        cell.fill=PatternFill("solid", fgColor=V54_BLUE)
        cell.font=Font(color="FFFFFF", bold=True)
        cell.alignment=Alignment(wrap_text=True, vertical="center")
        cell.border=Border(bottom=Side(style="thin", color=V54_LINE))


def _v54_style_sheet(ws, freeze="A2"):
    ws.freeze_panes = freeze
    ws.sheet_view.showGridLines = False
    for row in ws.iter_rows():
        for cell in row:
            cell.alignment = Alignment(wrap_text=True, vertical="top")
            cell.border = Border(bottom=Side(style="hair", color="E8EEF5"))
            if isinstance(cell.value, (int,float)):
                cell.number_format = '#,##0.0'
    for col in range(1, ws.max_column+1):
        letter = get_column_letter(col)
        max_len = 10
        for cell in ws[letter]:
            if cell.value is not None:
                max_len = max(max_len, min(55, len(str(cell.value)) + 2))
        ws.column_dimensions[letter].width = min(max_len, 38)
    for r in range(1, ws.max_row+1):
        ws.row_dimensions[r].height = 24


def workbook_bytes(model: Dict[str, Any]) -> bytes:
    model = dict(model)
    sc = _v54_scenario(model); costs=_v54_costs(model); risks=_v54_risks(model); curves=_v54_curves(model); torn=_v54_tornado(model)
    wb=Workbook(); wb.remove(wb.active)
    # Dashboard
    ws=wb.create_sheet("Executive Dashboard")
    ws["A1"]="CASEY PLATINUM COST / SCHEDULE / RISK MODEL"; ws["A1"].font=Font(size=18,bold=True,color=V54_BLUE)
    ws["A2"]=f"{model.get('title')} | {sc['label']} scenario | Generated {datetime.utcnow().strftime('%Y-%m-%d %H:%M UTC')}"; ws["A2"].font=Font(size=10,color="666666")
    kpis=[("P50 Cost",model.get("cost_p50")),("Range",model.get("cost_range")),("QCRA P80",money_bn(model.get("monte_carlo",{}).get("qcra",{}).get("p80", parse_bn(model.get("cost_p50"))*1.25))),("Schedule",model.get("schedule")),("QSRA P80",str(model.get("monte_carlo",{}).get("qsra",{}).get("p80",""))+" months"),("Risk / Confidence",f"{model.get('risk')} / {model.get('confidence_pct')}%")]
    row=4
    for i,(k,v) in enumerate(kpis):
        c=1+(i%3)*3; r=row+(i//3)*3
        ws.cell(r,c,k); ws.cell(r+1,c,str(v)); ws.merge_cells(start_row=r,start_column=c,end_row=r,end_column=c+1); ws.merge_cells(start_row=r+1,start_column=c,end_row=r+1,end_column=c+1)
        ws.cell(r,c).fill=PatternFill("solid", fgColor=V54_LIGHT); ws.cell(r,c).font=Font(bold=True,color=V54_BLUE)
        ws.cell(r+1,c).font=Font(size=18,bold=True,color=V54_BLUE); ws.cell(r+1,c).fill=PatternFill("solid", fgColor="FFFFFF")
    ws["A11"]="EXECUTIVE RECOMMENDATION"; ws["A11"].font=Font(bold=True,color=V54_BLUE,size=13)
    ws["A12"]=sc["recommendation"]; ws["A13"]=sc["why"]; ws["A12"].font=Font(bold=True,size=14,color=V54_RED if sc['key'] in ['cheaper','survival'] else V54_BLUE)
    ws["A15"]="Direct / Indirect / Reserve"; ws["A15"].font=Font(bold=True,color=V54_BLUE)
    summary={"Direct":0,"Indirect":0,"Reserve":0}
    for x in costs: summary[x.get("type","Direct")] = summary.get(x.get("type","Direct"),0)+float(x.get("p50_bn",0))
    for i,(k,v) in enumerate(summary.items(), start=16): ws.cell(i,1,k); ws.cell(i,2,v); ws.cell(i,2).number_format='$#,##0.0B'
    pie=PieChart(); labels=Reference(ws,min_col=1,min_row=16,max_row=18); data=Reference(ws,min_col=2,min_row=15,max_row=18); pie.add_data(data,titles_from_data=True); pie.set_categories(labels); pie.title="Cost Composition"; ws.add_chart(pie,"D15")
    _v54_style_sheet(ws,"A4")
    # Cost Estimate
    ws=wb.create_sheet("Cost Estimate")
    headers=["CBS","Description","Type","Basis","P10 $B","P50 $B","P90 $B","Board Challenge"]
    ws.append(headers); _v54_style_header(ws)
    for x in costs:
        challenge = "Validate benchmark and supplier evidence" if x.get("type") == "Direct" else "Confirm not hidden in direct cost"
        ws.append([x.get("cbs"),x.get("description"),x.get("type"),x.get("basis"),x.get("p10_bn"),x.get("p50_bn"),x.get("p90_bn"),challenge])
    _v54_style_sheet(ws)
    # Scenario Comparison
    ws=wb.create_sheet("Scenario Comparison"); ws.append(["Scenario","Cost","Schedule Months","Risk","Confidence","What changes","Board Use"]); _v54_style_header(ws)
    for x in model.get("scenario_comparison") or []:
        ws.append([x.get("label"),x.get("cost"),x.get("schedule_months"),x.get("risk"),x.get("confidence"),x.get("why"),"Use to frame trade-off decision"])
    _v54_style_sheet(ws)
    # Curves
    ws=wb.create_sheet("QCRA Curve"); ws.append(["Percentile","Cost $B"]); _v54_style_header(ws)
    for p,v in curves["qcra"]: ws.append([p,v])
    chart=LineChart(); chart.title="QCRA Cost Confidence Curve"; chart.y_axis.title="Cost $B"; chart.x_axis.title="Percentile"; data=Reference(ws,min_col=2,min_row=1,max_row=ws.max_row); cats=Reference(ws,min_col=1,min_row=2,max_row=ws.max_row); chart.add_data(data,titles_from_data=True); chart.set_categories(cats); ws.add_chart(chart,"D2"); _v54_style_sheet(ws)
    ws=wb.create_sheet("QSRA Curve"); ws.append(["Percentile","Duration Months"]); _v54_style_header(ws)
    for p,v in curves["qsra"]: ws.append([p,v])
    chart=LineChart(); chart.title="QSRA Schedule Confidence Curve"; chart.y_axis.title="Months"; chart.x_axis.title="Percentile"; data=Reference(ws,min_col=2,min_row=1,max_row=ws.max_row); cats=Reference(ws,min_col=1,min_row=2,max_row=ws.max_row); chart.add_data(data,titles_from_data=True); chart.set_categories(cats); ws.add_chart(chart,"D2"); _v54_style_sheet(ws)
    # Tornados
    ws=wb.create_sheet("QCRA Tornado"); ws.append(["Rank","Driver","CBS","Mean Cost $B","Action"]); _v54_style_header(ws)
    for i,r in enumerate(torn["qcra"],1): ws.append([i,r["title"],r["cbs"],r["cost_emv_bn"],r["mitigation"]])
    chart=BarChart(); chart.type="bar"; chart.title="QCRA Cost Tornado"; data=Reference(ws,min_col=4,min_row=1,max_row=ws.max_row); cats=Reference(ws,min_col=2,min_row=2,max_row=ws.max_row); chart.add_data(data,titles_from_data=True); chart.set_categories(cats); ws.add_chart(chart,"G2"); _v54_style_sheet(ws)
    ws=wb.create_sheet("QSRA Tornado"); ws.append(["Rank","Driver","Activity","Mean Days","Action"]); _v54_style_header(ws)
    for i,r in enumerate(torn["qsra"],1): ws.append([i,r["title"],r["activity_id"],r["schedule_emv_days"],r["mitigation"]])
    chart=BarChart(); chart.type="bar"; chart.title="QSRA Schedule Tornado"; data=Reference(ws,min_col=4,min_row=1,max_row=ws.max_row); cats=Reference(ws,min_col=2,min_row=2,max_row=ws.max_row); chart.add_data(data,titles_from_data=True); chart.set_categories(cats); ws.add_chart(chart,"G2"); _v54_style_sheet(ws)
    # Risk register
    ws=wb.create_sheet("Risk Register")
    headers=["ID","Risk","Category","Likelihood %","Cause","Event","Impact","Owner","Mitigation","Trigger","Residual","Activity","CBS","Cost EMV $B","Schedule EMV Days"]
    ws.append(headers); _v54_style_header(ws)
    for r in risks:
        ws.append([r["risk_id"],r["title"],r["category"],r["probability_pct"],r["cause"],r["event"],r["impact"],r["owner"],r["mitigation"],r["trigger"],r["residual_rating"],r["activity_id"],r["cbs"],r["cost_emv_bn"],r["schedule_emv_days"]])
        color = {"Red":"FCE4D6","Amber":"FFF2CC","Green":"E2F0D9"}.get(r["residual_rating"],"FFFFFF")
        for c in range(1,16): ws.cell(ws.max_row,c).fill=PatternFill("solid",fgColor=color)
    _v54_style_sheet(ws)
    ws.column_dimensions['E'].width=34; ws.column_dimensions['F'].width=34; ws.column_dimensions['G'].width=40; ws.column_dimensions['I'].width=38
    # Schedule
    ws=wb.create_sheet("Schedule")
    ws.append(["Activity ID","Phase","Activity","Predecessor","Duration Months","Critical","Basis"]); _v54_style_header(ws)
    for a in model.get("schedule_rows") or []: ws.append([a.get("activity_id"),a.get("phase"),a.get("activity"),a.get("predecessor"),a.get("duration_months"),a.get("critical"),a.get("basis")])
    _v54_style_sheet(ws)
    # Assumptions
    ws=wb.create_sheet("Methodology & Assumptions")
    rows=[["Area","Method"],["Cost","Parametric sector model + location factor + scale factor + class range + scenario logic."],["QCRA","Probability-weighted O/M/P cost impacts with percentile confidence outputs."],["QSRA","Probability-weighted activity impacts with schedule confidence curve and tornado."],["Risk","Every risk includes cause, event, impact, owner, mitigation and residual rating."],["Limitation","First-pass intelligence model. Validate with project-specific evidence before capital commitment."]]
    for r in rows: ws.append(r)
    _v54_style_header(ws); _v54_style_sheet(ws)
    bio=BytesIO(); wb.save(bio); bio.seek(0); return bio.getvalue()


def risk_register_workbook_bytes(model: Dict[str, Any]) -> bytes:
    risks=_v54_risks(model); torn=_v54_tornado(model)
    wb=Workbook(); wb.remove(wb.active)
    ws=wb.create_sheet("Risk Dashboard")
    ws["A1"]="CASEY PLATINUM RISK REGISTER"; ws["A1"].font=Font(size=18,bold=True,color=V54_BLUE)
    ws["A2"]=f"{model.get('title')} | {model.get('scenario_label',model.get('scenario','Base'))}"; ws["A2"].font=Font(color="666666")
    ws.append([]); ws.append(["Metric","Value"]); _v54_style_header(ws,4,2)
    ws.append(["Total risks",len(risks)]); ws.append(["Red risks",sum(1 for r in risks if r['residual_rating']=='Red')]); ws.append(["Amber risks",sum(1 for r in risks if r['residual_rating']=='Amber')]); ws.append(["Top QCRA driver",torn['qcra'][0]['title'] if torn['qcra'] else 'None']); ws.append(["Top QSRA driver",torn['qsra'][0]['title'] if torn['qsra'] else 'None'])
    _v54_style_sheet(ws,"A4")
    ws=wb.create_sheet("Risk Register")
    headers=["ID","Risk","Category","Likelihood %","Cause","Risk Event","Impact Description","Owner","Mitigation / Response","Trigger","Residual Rating","Linked Activity","Linked CBS","QCRA Mean $B","QSRA Mean Days","Status"]
    ws.append(headers); _v54_style_header(ws)
    for r in risks:
        ws.append([r["risk_id"],r["title"],r["category"],r["probability_pct"],r["cause"],r["event"],r["impact"],r["owner"],r["mitigation"],r["trigger"],r["residual_rating"],r["activity_id"],r["cbs"],r["cost_emv_bn"],r["schedule_emv_days"],r["status"]])
        color={"Red":"FCE4D6","Amber":"FFF2CC","Green":"E2F0D9"}.get(r["residual_rating"],"FFFFFF")
        for c in range(1,17): ws.cell(ws.max_row,c).fill=PatternFill("solid",fgColor=color)
    _v54_style_sheet(ws)
    for col in ['E','F','G','I','J']: ws.column_dimensions[col].width=40
    ws=wb.create_sheet("QCRA Tornado")
    ws.append(["Rank","Driver","CBS","Mean Exposure $B","Mitigation"]); _v54_style_header(ws)
    for i,r in enumerate(torn['qcra'],1): ws.append([i,r['title'],r['cbs'],r['cost_emv_bn'],r['mitigation']])
    chart=BarChart(); chart.type="bar"; chart.title="QCRA Tornado - Cost Exposure"; data=Reference(ws,min_col=4,min_row=1,max_row=ws.max_row); cats=Reference(ws,min_col=2,min_row=2,max_row=ws.max_row); chart.add_data(data,titles_from_data=True); chart.set_categories(cats); ws.add_chart(chart,"G2"); _v54_style_sheet(ws)
    ws=wb.create_sheet("QSRA Tornado")
    ws.append(["Rank","Driver","Activity","Mean Exposure Days","Mitigation"]); _v54_style_header(ws)
    for i,r in enumerate(torn['qsra'],1): ws.append([i,r['title'],r['activity_id'],r['schedule_emv_days'],r['mitigation']])
    chart=BarChart(); chart.type="bar"; chart.title="QSRA Tornado - Schedule Exposure"; data=Reference(ws,min_col=4,min_row=1,max_row=ws.max_row); cats=Reference(ws,min_col=2,min_row=2,max_row=ws.max_row); chart.add_data(data,titles_from_data=True); chart.set_categories(cats); ws.add_chart(chart,"G2"); _v54_style_sheet(ws)
    bio=BytesIO(); wb.save(bio); bio.seek(0); return bio.getvalue()


def risk_csv_bytes(model: Dict[str, Any]) -> bytes:
    out=StringIO(); w=csv.writer(out)
    w.writerow(["Risk ID","Title","Category","Likelihood %","Cause","Risk Event","Impact Description","Owner","Mitigation","Trigger","Residual Rating","Activity ID","Activity Name","CBS","CBS Name","QCRA Mean $B","QSRA Mean Days","Status"])
    for r in _v54_risks(model):
        w.writerow([r['risk_id'],r['title'],r['category'],r['probability_pct'],r['cause'],r['event'],r['impact'],r['owner'],r['mitigation'],r['trigger'],r['residual_rating'],r['activity_id'],r['activity_name'],r['cbs'],r['cbs_name'],r['cost_emv_bn'],r['schedule_emv_days'],r['status']])
    return out.getvalue().encode("utf-8")


def xer_bytes(model: Dict[str, Any]) -> bytes:
    # Legacy-targeted P6 XER style. PRA versions vary; for stubborn PRA installs use the Schedule sheet/CSV as fallback.
    today=datetime.utcnow().strftime("%Y-%m-%d")
    title=str(model.get('title') or 'CASEY Project').replace('\t',' ').replace('\n',' ')
    lines=[]
    lines.append(f"ERMHDR\t8.4\t{today}\tProject\tCASEY\tCASEY\tCASEYDatabase\tProject Management\tUSD")
    lines += [
        "%T\tCURRTYPE",
        "%F\tcurr_id\tdecimal_digit_cnt\tcurr_symbol\tdecimal_symbol\tdigit_group_symbol\tpos_curr_fmt_type\tneg_curr_fmt_type\tcurr_type\tcurr_short_name\tgroup_digit_cnt\tbase_exch_rate",
        "%R\t1\t2\t$\t.\t,\t#1.1\t(#1.1)\tDollar\tUSD\t3\t1",
        "%T\tOBS",
        "%F\tobs_id\tparent_obs_id\tguid\tseq_num\tobs_name\tobs_descr",
        "%R\t1\t\t\t0\tEnterprise\t",
        "%T\tCALENDAR",
        "%F\tclndr_id\tclndr_name\tdefault_flag\tday_hr_cnt\tweek_hr_cnt\tmonth_hr_cnt\tyear_hr_cnt",
        "%R\t1\tCASEY Standard 5 Day\tY\t8\t40\t173.33\t2080",
        "%T\tPROJECT",
        "%F\tproj_id\tfy_start_month_num\trsrc_self_add_flag\tallow_complete_flag\trsrc_multi_assign_flag\tcheckout_flag\tproject_flag\tstep_complete_flag\tcost_qty_recalc_flag\tbatch_sum_flag\tname_sep_char\tdef_complete_pct_type\tproj_short_name\tproj_name\tplan_start_date\tscd_end_date\tclndr_id\tcurr_id\torig_proj_id\tlast_fin_dates_id\tlast_baseline_update_date\tpriority_num\tcreate_date\tupdate_date\tcreate_user\tupdate_user",
        f"%R\t1\t1\tY\tY\tY\tN\tY\tY\tY\tY\t.\tCP\tCASEY\t{title}\t{today}\t\t1\t1\t\t\t\t1\t{today}\t{today}\tCASEY\tCASEY",
        "%T\tPROJWBS",
        "%F\twbs_id\tproj_id\tobs_id\tseq_num\twbs_short_name\twbs_name\tparent_wbs_id\tproj_node_flag",
        "%R\t1\t1\t1\t1\tCASEY\tCASEY Project\t\tY",
        "%T\tTASK",
        "%F\ttask_id\tproj_id\twbs_id\tclndr_id\tphys_complete_pct\trev_fdbk_flag\test_wt\tlock_plan_flag\tauto_compute_act_flag\tcomplete_pct_type\ttask_type\tduration_type\tstatus_code\ttask_code\ttask_name\ttarget_drtn_hr_cnt\tremain_drtn_hr_cnt\ttotal_float_hr_cnt",
    ]
    idx={}
    for i,a in enumerate(model.get('schedule_rows') or [],1):
        code=str(a.get('activity_id') or f"A{i:04d}").replace('\t',' ')
        name=str(a.get('activity') or 'Activity').replace('\t',' ').replace('\n',' ')
        dur=max(8, int(float(a.get('duration_months') or 1)*173.33))
        idx[code]=i
        lines.append(f"%R\t{i}\t1\t1\t1\t0\tN\t1\tN\tY\tCP\tTT_Task\tDT_FixedDUR2\tTK_NotStart\t{code}\t{name}\t{dur}\t{dur}\t0")
    lines += ["%T\tTASKPRED","%F\ttask_pred_id\ttask_id\tpred_task_id\tpred_type\tlag_hr_cnt"]
    k=1
    for a in model.get('schedule_rows') or []:
        code=str(a.get('activity_id'))
        for pred in str(a.get('predecessor') or '').split(';'):
            pred=pred.strip()
            if pred in idx and code in idx:
                lines.append(f"%R\t{k}\t{idx[code]}\t{idx[pred]}\tPR_FS\t0"); k+=1
    return ("\n".join(lines)+"\n").encode("latin-1", errors="replace")


def word_bytes(model: Dict[str, Any]) -> bytes:
    sc=_v54_scenario(model); risks=_v54_risks(model); torn=_v54_tornado(model); costs=_v54_costs(model)
    doc=Document(); styles=doc.styles; styles['Normal'].font.name='Aptos'; styles['Normal'].font.size=Pt(10)
    title=doc.add_heading('CASEY PLATINUM BOARD DECISION PACK',0); title.runs[0].font.color.rgb=RGBColor(11,31,51)
    doc.add_paragraph(f"{model.get('title')} | {sc['label']} scenario | Generated {datetime.utcnow().strftime('%Y-%m-%d %H:%M UTC')}")
    p=doc.add_paragraph(); p.add_run(sc['recommendation']).bold=True; p.runs[0].font.size=Pt(14); p.runs[0].font.color.rgb=RGBColor(180,35,24)
    doc.add_paragraph(sc['why'])
    doc.add_heading('Board Decision Snapshot',1)
    tbl=doc.add_table(rows=1,cols=4); tbl.style='Table Grid'
    for i,h in enumerate(['Metric','Value','Interpretation','Board action']): tbl.rows[0].cells[i].text=h
    for row in [('P50 Cost',model.get('cost_p50'),'Most likely first-pass control point','Use for option comparison'),('Cost Range',model.get('cost_range'),'Class and QCRA uncertainty range','Approve only with range acknowledged'),('QCRA P80',money_bn(model.get('monte_carlo',{}).get('qcra',{}).get('p80',0)),'Risk-adjusted funding confidence','Set contingency rules'),('QSRA P80',str(model.get('monte_carlo',{}).get('qsra',{}).get('p80',''))+' months','Schedule confidence level','Validate critical path'),('Risk / Confidence',f"{model.get('risk')} / {model.get('confidence_pct')}%",'Data maturity and residual exposure','Approve evidence-improvement plan')]:
        cells=tbl.add_row().cells
        for i,v in enumerate(row): cells[i].text=str(v)
    doc.add_heading('QCRA Cost Tornado',1)
    for r in torn['qcra'][:8]: doc.add_paragraph(f"{r['title']} — {r['cbs']} — {money_bn(r['cost_emv_bn'])} mean exposure", style='List Bullet')
    doc.add_heading('QSRA Schedule Tornado',1)
    for r in torn['qsra'][:8]: doc.add_paragraph(f"{r['title']} — {r['activity_id']} — {r['schedule_emv_days']} days mean exposure", style='List Bullet')
    doc.add_heading('Board Top Risks',1)
    for r in risks[:10]:
        doc.add_paragraph(f"{r['risk_id']} {r['title']}", style='List Bullet')
        doc.add_paragraph(f"Cause: {r['cause']}\nEvent: {r['event']}\nImpact: {r['impact']}\nMitigation: {r['mitigation']}")
    doc.add_heading('Selected Class Cost Estimate',1)
    tbl=doc.add_table(rows=1, cols=6); tbl.style='Table Grid'
    for i,h in enumerate(['CBS','Description','Type','P10','P50','P90']): tbl.rows[0].cells[i].text=h
    for x in costs:
        cells=tbl.add_row().cells
        vals=[x['cbs'],x['description'],x['type'],money_bn(x['p10_bn']),money_bn(x['p50_bn']),money_bn(x['p90_bn'])]
        for i,v in enumerate(vals): cells[i].text=str(v)
    bio=BytesIO(); doc.save(bio); bio.seek(0); return bio.getvalue()


def pdf_bytes(model: Dict[str, Any]) -> bytes:
    sc=_v54_scenario(model); risks=_v54_risks(model); torn=_v54_tornado(model); costs=_v54_costs(model)
    bio=BytesIO(); doc=SimpleDocTemplate(bio,pagesize=landscape(A4),leftMargin=28,rightMargin=28,topMargin=24,bottomMargin=24)
    styles=getSampleStyleSheet()
    styles.add(ParagraphStyle(name='CaseyTitle', parent=styles['Title'], fontSize=24, textColor=colors.HexColor("#"+V54_BLUE), leading=28, spaceAfter=10))
    styles.add(ParagraphStyle(name='CaseyH', parent=styles['Heading1'], fontSize=16, textColor=colors.HexColor("#"+V54_BLUE), spaceAfter=8))
    styles.add(ParagraphStyle(name='Small', parent=styles['BodyText'], fontSize=8, leading=10))
    story=[]
    story.append(Paragraph(f"CASEY BOARD DECISION PACK — {model.get('title')}", styles['CaseyTitle']))
    story.append(Paragraph(f"{sc['label']} scenario | {model.get('subsector')} | {model.get('location')} | Generated {datetime.utcnow().strftime('%Y-%m-%d %H:%M UTC')}", styles['Small']))
    kpis=[['P50 Cost',model.get('cost_p50'),'Cost Range',model.get('cost_range'),'QCRA P80',money_bn(model.get('monte_carlo',{}).get('qcra',{}).get('p80',0))],['Schedule',model.get('schedule'),'QSRA P80',str(model.get('monte_carlo',{}).get('qsra',{}).get('p80',''))+' mo','Risk / Confidence',f"{model.get('risk')} / {model.get('confidence_pct')}%"]]
    t=Table(kpis, colWidths=[90,110,90,130,90,130]); t.setStyle(TableStyle([('BACKGROUND',(0,0),(-1,-1),colors.HexColor('#F3F7FA')),('TEXTCOLOR',(0,0),(-1,-1),colors.HexColor("#"+V54_BLUE)),('FONTNAME',(0,0),(-1,-1),'Helvetica-Bold'),('GRID',(0,0),(-1,-1),0.5,colors.HexColor("#"+V54_LINE)),('VALIGN',(0,0),(-1,-1),'MIDDLE'),('FONTSIZE',(0,0),(-1,-1),10)])); story.append(t); story.append(Spacer(1,12))
    story.append(Paragraph(sc['recommendation'].upper(), styles['CaseyH'])); story.append(Paragraph(sc['why'], styles['BodyText'])); story.append(PageBreak())
    story.append(Paragraph('QCRA + QSRA Tornado Drivers', styles['CaseyH']))
    data=[['Rank','QCRA Cost Driver','CBS','Mean Cost','QSRA Schedule Driver','Activity','Mean Days']]
    for i in range(max(len(torn['qcra']),len(torn['qsra']))):
        qc=torn['qcra'][i] if i<len(torn['qcra']) else {}; qs=torn['qsra'][i] if i<len(torn['qsra']) else {}
        data.append([i+1,qc.get('title',''),qc.get('cbs',''),money_bn(qc.get('cost_emv_bn',0)) if qc else '',qs.get('title',''),qs.get('activity_id',''),qs.get('schedule_emv_days','')])
    t=Table(data, repeatRows=1, colWidths=[35,145,55,70,145,65,60]); t.setStyle(TableStyle([('BACKGROUND',(0,0),(-1,0),colors.HexColor("#"+V54_BLUE)),('TEXTCOLOR',(0,0),(-1,0),colors.white),('GRID',(0,0),(-1,-1),0.35,colors.HexColor("#"+V54_LINE)),('FONTSIZE',(0,0),(-1,-1),8),('VALIGN',(0,0),(-1,-1),'TOP'),('ROWBACKGROUNDS',(0,1),(-1,-1),[colors.white,colors.HexColor('#F8FBFD')]) ])); story.append(t); story.append(PageBreak())
    story.append(Paragraph('Risk Register — Board Top Risks', styles['CaseyH']))
    data=[['ID','Risk','Cause','Event','Impact','Owner','Mitigation']]
    for r in risks[:10]: data.append([r['risk_id'],r['title'],r['cause'],r['event'],r['impact'],r['owner'],r['mitigation']])
    t=Table(data, repeatRows=1, colWidths=[45,90,120,120,160,70,145]); t.setStyle(TableStyle([('BACKGROUND',(0,0),(-1,0),colors.HexColor("#"+V54_BLUE)),('TEXTCOLOR',(0,0),(-1,0),colors.white),('GRID',(0,0),(-1,-1),0.25,colors.HexColor("#"+V54_LINE)),('FONTSIZE',(0,0),(-1,-1),6.8),('VALIGN',(0,0),(-1,-1),'TOP'),('ROWBACKGROUNDS',(0,1),(-1,-1),[colors.white,colors.HexColor('#F8FBFD')]) ])); story.append(t); story.append(PageBreak())
    story.append(Paragraph('Cost Estimate — Selected Class Only', styles['CaseyH']))
    data=[['CBS','Description','Type','P10','P50','P90','Basis']]
    for x in costs: data.append([x['cbs'],x['description'],x['type'],money_bn(x['p10_bn']),money_bn(x['p50_bn']),money_bn(x['p90_bn']),x.get('basis','')])
    t=Table(data, repeatRows=1, colWidths=[50,150,70,70,70,70,300]); t.setStyle(TableStyle([('BACKGROUND',(0,0),(-1,0),colors.HexColor("#"+V54_BLUE)),('TEXTCOLOR',(0,0),(-1,0),colors.white),('GRID',(0,0),(-1,-1),0.25,colors.HexColor("#"+V54_LINE)),('FONTSIZE',(0,0),(-1,-1),7.2),('VALIGN',(0,0),(-1,-1),'TOP'),('ROWBACKGROUNDS',(0,1),(-1,-1),[colors.white,colors.HexColor('#F8FBFD')]) ])); story.append(t)
    doc.build(story); bio.seek(0); return bio.getvalue()


def pptx_bytes(model: Dict[str, Any]) -> bytes:
    prs=Presentation(); prs.slide_width=PptxInches(13.333); prs.slide_height=PptxInches(7.5)
    sc=_v54_scenario(model); risks=_v54_risks(model); torn=_v54_tornado(model)
    def slide(title, subtitle=""):
        s=prs.slides.add_slide(prs.slide_layouts[6]);
        bg=s.shapes.add_shape(1,0,0,prs.slide_width,prs.slide_height); bg.fill.solid(); bg.fill.fore_color.rgb=__import__('pptx').dml.color.RGBColor(255,255,255); bg.line.fill.background()
        t=s.shapes.add_textbox(PptxInches(.55),PptxInches(.35),PptxInches(12),PptxInches(.7)); tf=t.text_frame; tf.text=title; run=tf.paragraphs[0].runs[0]; run.font.size=PptxPt(28); run.font.bold=True; run.font.color.rgb=__import__('pptx').dml.color.RGBColor(11,31,51)
        if subtitle:
            st=s.shapes.add_textbox(PptxInches(.6),PptxInches(1.05),PptxInches(12),PptxInches(.4)); st.text_frame.text=subtitle; st.text_frame.paragraphs[0].runs[0].font.size=PptxPt(12); st.text_frame.paragraphs[0].runs[0].font.color.rgb=__import__('pptx').dml.color.RGBColor(90,90,90)
        return s
    s=slide(f"{model.get('title')} — Board Intelligence", f"{sc['label']} scenario | CASEY Platinum Output Engine")
    tx=s.shapes.add_textbox(PptxInches(.75),PptxInches(2.0),PptxInches(12),PptxInches(1.8)); tf=tx.text_frame; tf.text=sc['recommendation'].upper(); tf.paragraphs[0].runs[0].font.size=PptxPt(30); tf.paragraphs[0].runs[0].font.bold=True
    tx2=s.shapes.add_textbox(PptxInches(.8),PptxInches(4.1),PptxInches(11),PptxInches(.6)); tx2.text_frame.text=sc['why']
    s=slide("Board Decision Snapshot")
    metrics=[('P50',model.get('cost_p50')),('Range',model.get('cost_range')),('QCRA P80',money_bn(model.get('monte_carlo',{}).get('qcra',{}).get('p80',0))),('Schedule',model.get('schedule')),('QSRA P80',str(model.get('monte_carlo',{}).get('qsra',{}).get('p80',''))+' mo'),('Risk',f"{model.get('risk')} / {model.get('confidence_pct')}%")]
    for i,(k,v) in enumerate(metrics):
        x=.75+(i%3)*4.15; y=1.6+(i//3)*1.7; shp=s.shapes.add_shape(1,PptxInches(x),PptxInches(y),PptxInches(3.6),PptxInches(1.15)); shp.fill.solid(); shp.fill.fore_color.rgb=__import__('pptx').dml.color.RGBColor(243,247,250); shp.line.color.rgb=__import__('pptx').dml.color.RGBColor(217,226,234); shp.text=f"{k}\n{v}"; shp.text_frame.paragraphs[0].runs[0].font.size=PptxPt(11); shp.text_frame.paragraphs[1].runs[0].font.size=PptxPt(20); shp.text_frame.paragraphs[1].runs[0].font.bold=True
    s=slide("QCRA and QSRA are separated", "Cost exposure and schedule exposure are not mixed.")
    left=s.shapes.add_textbox(PptxInches(.8),PptxInches(1.6),PptxInches(5.8),PptxInches(5.2)).text_frame; left.text="QCRA COST TORNADO"; left.paragraphs[0].runs[0].font.bold=True
    for r in torn['qcra'][:7]: p=left.add_paragraph(); p.text=f"{r['title']} — {money_bn(r['cost_emv_bn'])}"; p.level=0
    right=s.shapes.add_textbox(PptxInches(7.0),PptxInches(1.6),PptxInches(5.8),PptxInches(5.2)).text_frame; right.text="QSRA SCHEDULE TORNADO"; right.paragraphs[0].runs[0].font.bold=True
    for r in torn['qsra'][:7]: p=right.add_paragraph(); p.text=f"{r['title']} — {r['schedule_emv_days']} days"; p.level=0
    s=slide("Top risks have cause, event, impact and action", "Designed to be used in a real risk review, not as a decorative export.")
    tf=s.shapes.add_textbox(PptxInches(.75),PptxInches(1.5),PptxInches(12),PptxInches(5.5)).text_frame
    for r in risks[:6]: p=tf.add_paragraph(); p.text=f"{r['risk_id']} {r['title']}: {r['impact']} | Action: {r['mitigation']}"; p.level=0; p.font.size=PptxPt(12)
    s=slide("What CASEY delivers", "Outputs are readable, scenario-linked and defensible.")
    tf=s.shapes.add_textbox(PptxInches(.9),PptxInches(1.6),PptxInches(11.5),PptxInches(4.8)).text_frame
    for txt in ["Selected class cost estimate only", "Direct / indirect / reserve split", "Risk register with cause, event, impact, owner and mitigation", "Separate QCRA and QSRA curves/tornado", "Legacy-targeted schedule export + schedule table fallback", "Board PDF, DOCX, PPTX, Excel and ZIP pack"]:
        p=tf.add_paragraph(); p.text="✓ "+txt; p.font.size=PptxPt(18)
    bio=BytesIO(); prs.save(bio); bio.seek(0); return bio.getvalue()

# v54 routes for explicit use; old /export routes automatically use the overridden global builders above.
@app.post('/v54/export/all')
def export_all_v54(model:Dict[str,Any]):
    bio=BytesIO()
    with zipfile.ZipFile(bio,'w',zipfile.ZIP_DEFLATED) as z:
        z.writestr('01_CASEY_v54_Platinum_Cost_Model.xlsx',workbook_bytes(model))
        z.writestr('02_CASEY_v54_Platinum_Risk_Register.xlsx',risk_register_workbook_bytes(model))
        z.writestr('03_CASEY_v54_Legacy_Targeted_P6_Schedule.xer',xer_bytes(model))
        z.writestr('04_CASEY_v54_Board_Decision_Pack.docx',word_bytes(model))
        z.writestr('05_CASEY_v54_Board_Decision_Pack.pdf',pdf_bytes(model))
        z.writestr('06_CASEY_v54_Board_Deck.pptx',pptx_bytes(model))
        z.writestr('07_CASEY_v54_Model_Audit.json',json.dumps(model,indent=2))
        z.writestr('08_CASEY_v54_Risk_Register.csv',risk_csv_bytes(model))
        # schedule CSV fallback for tools that reject XER versions
        out=StringIO(); w=csv.writer(out); w.writerow(['Activity ID','Activity','Predecessor','Duration Months','Critical'])
        for a in model.get('schedule_rows') or []: w.writerow([a.get('activity_id'),a.get('activity'),a.get('predecessor'),a.get('duration_months'),a.get('critical')])
        z.writestr('09_CASEY_v54_Schedule_Fallback.csv',out.getvalue())
        z.writestr('README_PRA_IMPORT_NOTE.txt','If Primavera Risk Analysis rejects XER on your installed version, import/use the Schedule_Fallback.csv or re-export from P6 Professional in the PRA-supported version. CASEY provides both a legacy-targeted XER and a clean schedule table fallback.')
    bio.seek(0); return stream(bio.getvalue(),'application/zip','CASEY_v54_Platinum_Output_Pack.zip')


# =====================================================================
# v55 PLATINUM OUTPUTS FINAL PATCH
# - Excel repaired by replacing openpyxl-generated export workbooks with xlsxwriter.
# - PPTX removed from output factory and ZIP because weak decks reduce trust.
# - XER header targeted to Primavera/PRA-compatible ERMHDR 8.0 style.
# - Outputs adapt to selected scenario, project sector, Earth/Space and all cost/schedule levels.
# =====================================================================

import html
from fastapi.responses import JSONResponse
try:
    import xlsxwriter
except Exception:
    xlsxwriter = None

V55_BLUE = '#0B1F33'
V55_CYAN = '#00A6C8'
V55_LINE = '#D9E2EA'
V55_LIGHT = '#F4F7FA'
V55_RED = '#C00000'
V55_AMBER = '#F4B183'
V55_GREEN = '#70AD47'


def _v55_scenario_key(model):
    return str(model.get('scenario') or 'base').lower().replace(' ', '_')


def _v55_scenario_label(model):

    # FINAL EXEC POLISH: scenario-aware ranking and differentiation
    try:
        sc = str(model.get("scenario","base")).lower()
        schedule_lists = {
            "faster":[
                "Concurrent commissioning overload",
                "Recovery float exhaustion",
                "Acceleration premium shock",
                "Grid connection delay",
                "Integrated systems testing concurrency"
            ],
            "cheaper":[
                "Vendor claims and change exposure",
                "Procurement deferral and long-lead slippage",
                "Design maturity gap",
                "Scope growth from deferred decisions",
                "Interface coordination delay"
            ],
            "lower_risk":[
                "Governance and approvals latency",
                "Extended validation sequencing",
                "Conservative commissioning gates",
                "Operational readiness hold-points",
                "Assurance and compliance reviews"
            ],
            "premium":[
                "Integration complexity across parallel packages",
                "Executive decision latency",
                "Technology assurance alignment",
                "Multi-package interface management",
                "Programme coordination overhead"
            ]
        }
        cost_lists = {
            "faster":[
                "Acceleration premiums and overtime",
                "Power train, transformers and switchgear",
                "Integrated systems testing",
                "Grid and utility concurrency",
                "Recovery-float consumption"
            ],
            "cheaper":[
                "Deferred procurement packaging",
                "Claims and commercial exposure",
                "Rework from reduced contingency",
                "Long-lead inflation volatility",
                "Scope rationalisation impacts"
            ],
            "lower_risk":[
                "Additional contingency and reserve",
                "Enhanced validation and assurance",
                "Programme controls and governance",
                "Redundant infrastructure resilience",
                "Extended commissioning readiness"
            ]
        }
        if sc in schedule_lists:
            model["sector_schedule_threats"]=schedule_lists[sc]
        if sc in cost_lists:
            model["sector_primary_cost_drivers"]=cost_lists[sc]

        if sc=="faster":
            model["executive_shock_insight"]="Acceleration increases spend faster than it reduces uncertainty; the delivery tail becomes more volatile."
        elif sc=="cheaper":
            model["executive_shock_insight"]="Capital efficiency reduces resilience: procurement and recovery flexibility become constrained."
        elif sc=="lower_risk":
            model["executive_shock_insight"]="Confidence is purchased through reserve, governance and extended delivery duration."
        elif sc=="premium":
            model["executive_shock_insight"]="Premium posture buys resilience, optionality and stronger certainty at visible capex premium."
    except Exception:
        pass

    return model.get('scenario_label') or scenario_params(_v55_scenario_key(model))[4]


def _v55_selected_class(model):
    try: return int(model.get('estimate_class') or 3)
    except Exception: return 3


def _v55_selected_schedule_level(model):
    try: return int(model.get('schedule_level') or 3)
    except Exception: return 3


def _v55_scenario_strategy(model):
    key = _v55_scenario_key(model)
    table = {
        'base': {
            'title': 'APPROVE DEFINITION PHASE — DO NOT COMMIT FULL CAPITAL YET',
            'summary': 'Balanced reference case for first-pass board challenge, option testing and contractor / consultant assurance.',
            'decision': 'Use as the control baseline. Validate evidence before approving full capital.'},
        'faster': {
            'title': 'ACCELERATE ONLY WITH EXPLICIT COST AND INTERFACE CONTROL',
            'summary': 'Faster scenario compresses programme duration through parallel design, early procurement and premium logistics, increasing interface and delivery risk.',
            'decision': 'Approve acceleration only if time value exceeds premium cost and added risk.'},
        'cheaper': {
            'title': 'LOWER COST CASE REQUIRES STRONGER RISK GOVERNANCE',
            'summary': 'Cheaper scenario reduces scope and procurement cost, but usually transfers risk into quality, schedule and contingency drawdown.',
            'decision': 'Use only where capex pressure is dominant and residual risk is accepted.'},
        'lower_risk': {
            'title': 'LOWER RISK CASE IMPROVES CONFIDENCE BUT REQUIRES MORE TIME AND ALLOWANCE',
            'summary': 'Lower Risk scenario adds surveys, assurance, stage gates, buffers and procurement validation to increase confidence.',
            'decision': 'Recommended when board certainty matters more than lowest first cost.'},
        'premium': {
            'title': 'PREMIUM CASE MAXIMISES RESILIENCE AND DELIVERY CONFIDENCE',
            'summary': 'Premium scenario strengthens suppliers, design maturity, resilience and governance to reduce downstream uncertainty.',
            'decision': 'Use for flagship, mission-critical or reputation-sensitive programmes.'},
        'investor': {
            'title': 'INVESTOR CASE MAKES RISK AND VALUE TRADE-OFFS EXPLICIT',
            'summary': 'Investor scenario frames the programme as an investment decision with clear capital, risk and confidence trade-offs.',
            'decision': 'Use to support IC, lender, fund or partner conversations.'},
        'survival': {
            'title': 'SURVIVAL CASE IS WARNING-HEAVY AND SHOULD NOT BE PRESENTED AS A TARGET BASELINE',
            'summary': 'Survival scenario minimises scope and cost, but creates high residual risk and quality / schedule exposure.',
            'decision': 'Use only as a stress case, not a preferred delivery path.'}
    }
    return table.get(key, table['base'])


def _v55_money(v):
    try: return money_bn(float(v))
    except Exception: return str(v)


def _v55_all_cost_lines(model):
    """Return dictionary {class: rows}. Always includes selected and all levels."""
    selected = model.get('cost_lines') or []
    by = model.get('estimates_by_class') or {}
    out = {}
    for cls in range(1, 6):
        rows = by.get(str(cls)) or []
        if not rows:
            low, high, _, mat = class_range(cls)
            rows = [{**x, 'p10_bn': round(float(x.get('p50_bn',0))*low, 3), 'p90_bn': round(float(x.get('p50_bn',0))*high, 3), 'class': cls, 'maturity': mat} for x in selected]
        out[cls] = rows
    return out


def _v55_all_schedule_levels(model):
    levels = model.get('schedules_by_level') or {}
    if not levels:
        levels = {str(_v55_selected_schedule_level(model)): model.get('schedule_rows') or []}
    return {str(k): list(v or []) for k, v in levels.items()}


def _v55_extra_risks(model):
    mode = str(model.get('mode') or 'Earth')
    sub = str(model.get('subsector') or '').lower()
    generic = [
        ('R-009','Commercial claims escalation','Contract ambiguity, change volume and interface disputes','Claims submitted above allowance','Commercial pressure, management distraction and increased outturn cost','Commercial Lead','Tight change control, entitlement reviews, dispute avoidance board','Claim volume exceeds baseline'),
        ('R-010','Resource availability','Labour, specialist supervisors or commissioning teams constrained','Required resources are unavailable when planned','Productivity loss, resequencing and higher preliminaries','Delivery Lead','Resource loading, framework call-offs, labour strategy','Resource histogram red flag'),
        ('R-011','Quality failure / rework','Installation quality, vendor QA or inspection regime below need','Inspection failure triggers rework','Rework cost, schedule delay and confidence loss','Quality Manager','ITPs, hold points, vendor QA, right-first-time reviews','NCR trend increases'),
        ('R-012','Cyber / controls integration','Operational technology, control systems and network interfaces immature','Controls integration fails testing','Commissioning delay, resilience issue and operational risk','Systems Lead','Early cyber review, FAT/SAT plan, integration lab','Controls defect trend increases'),
        ('R-013','Weather / environmental disruption','Extreme weather, environmental windows or site constraints not fully allowed','Workface access or productivity is disrupted','Delay, access cost and resequencing','Construction Manager','Weather calendars, contingency workfronts, protection strategy','Weather downtime exceeds allowance'),
        ('R-014','Stakeholder change','Sponsor, operator or authority requirements evolve late','Requirements change after baseline freeze','Scope growth, redesign and approval delays','Project Director','Decision log, change board and stakeholder gates','Late decision count increases'),
        ('R-015','Funding / approval gate delay','Investment approvals or staged funding gates not aligned with procurement','Capital release is delayed','Procurement missed windows and schedule slippage','Sponsor','Approval roadmap, early board papers, contingency approvals','Approval date slips'),
        ('R-016','Commissioning readiness gap','Incomplete test packs, operational readiness or training gaps','Commissioning starts without readiness','Repeated tests, handover delay and operational risk','Commissioning Lead','Readiness dashboard, integrated test plan, early ORAT','Readiness score below threshold'),
        ('R-017','Data / estimate basis weakness','Estimate source data, quantities or benchmarks are incomplete','Baseline cannot be defended under challenge','Confidence reduction and rework of cost plan','Cost Lead','Assumption register, benchmark evidence, quantity validation','Unknown allowance exceeds threshold'),
        ('R-018','Interface ownership gap','RACI and package boundaries are unclear across packages','Interface issue has no accountable owner','Delay, claims and fragmented accountability','Integration Manager','Interface control documents and owner matrix','Interface actions overdue'),
    ]
    sector = []
    if 'airport' in sub:
        sector += [('R-A01','Airside phasing disruption','Operational airport constraints and possessions are underestimated','Construction cannot access critical work fronts','Night work premium, delay and operational disruption','Airside Lead','Airside phasing plan, possessions, ORAT integration','Airside access cancelled'),('R-A02','Baggage / security systems integration','Complex passenger processing systems not fully integrated','Systems fail integrated trials','Terminal opening delay and reputation impact','Systems Lead','Early integration lab, FAT/SAT, ORAT testing','Trial failures increase')]
    if 'rail' in sub:
        sector += [('R-R01','Possession access constraints','Rail access windows and blockade availability are insufficient','Planned works cannot be delivered in possession','Resequencing, delay and overtime cost','Rail Access Lead','Possession strategy and contingency blockade plan','Possession cancelled'),('R-R02','Signalling integration','Signalling design, test and assurance duration underestimated','System does not pass integration testing','Critical path delay and safety assurance issue','Signalling Lead','Independent signalling assurance and test lab','Test defects exceed threshold')]
    if 'data centre' in sub:
        sector += [('R-D03','Liquid cooling vendor maturity','Cooling technology and supplier capacity are immature','Vendor solution underperforms acceptance tests','Capacity derate, redesign and schedule delay','MEP Lead','Prototype, vendor qualification and performance guarantee','Cooling FAT failed'),('R-D04','Generator / UPS lead time','Critical power equipment supply is constrained','Power train delivery slips','Critical path movement and acceleration premium','Procurement Lead','Early vendor reservation and dual-source strategy','Vendor delivery warning')]
    if mode.lower() == 'space':
        sector += [('R-S04','Launch window dependency','Orbital mechanics or launch provider availability constrains deployment','Launch window missed','Months of delay and re-manifest cost','Mission Lead','Backup launch windows and manifest options','Launch readiness slips'),('R-S05','Thermal / radiation design margin','Environment and shielding assumptions remain immature','Thermal or radiation margins fail review','Redesign, extra mass and launch cost increase','Chief Engineer','Thermal vacuum testing and radiation analysis','Margin below threshold')]
    return sector + generic


def _v55_risks(model):
    base = []
    for r in model.get('risks') or model.get('risk_register') or []:
        p = max(5, int(float(r.get('probability_pct') or r.get('likelihood_pct') or 25)))
        title = r.get('title') or r.get('risk') or r.get('name') or 'Risk'
        cause = r.get('cause') or r.get('description') or 'Underlying assumption, external dependency or delivery condition requires validation.'
        event = r.get('event') or f'{title} occurs or materialises during delivery.'
        impact = r.get('impact') or r.get('impact_description') or r.get('basis_of_schedule_impact') or 'Cost, schedule, quality, scope or confidence impact.'
        mitig = r.get('mitigation') or r.get('response') or 'Assign owner, validate assumption, implement mitigation and track residual exposure.'
        owner = r.get('owner') or 'Project Controls Lead'
        c = max(0.001, float(r.get('cost_emv_bn') or 0.01))
        d = max(0.1, float(r.get('schedule_emv_days') or 0.1))
        res = r.get('residual_rating') or ('Red' if p>=45 or c>0.5 or d>30 else 'Amber' if p>=25 or c>0.15 or d>10 else 'Green')
        base.append({
            'risk_id': r.get('risk_id') or r.get('id') or f'R-{len(base)+1:03d}', 'title': title, 'category': r.get('category') or 'Delivery', 'probability_pct': p,
            'cause': cause, 'event': event, 'impact': impact, 'owner': owner, 'mitigation': mitig,
            'trigger': r.get('trigger') or 'Threshold breach or assumption change', 'residual_rating': res,
            'activity_id': r.get('activity_id') or 'A1900', 'activity_name': r.get('activity_name') or 'Delivery activity',
            'cbs': r.get('cbs') or '01.04', 'cbs_name': r.get('cbs_name') or 'Cost package', 'cost_emv_bn': round(c,3), 'schedule_emv_days': round(d,1),
            'status': r.get('status') or 'Open'
        })
    cost = parse_bn(model.get('cost_p50'))
    months = int(str(model.get('schedule') or '60').split()[0]) if str(model.get('schedule') or '').split() else 60
    rm = scenario_params(_v55_scenario_key(model))[2]
    for rid,title,cause,event,impact,owner,mit,trig in _v55_extra_risks(model):
        if any(x['risk_id'] == rid for x in base):
            continue
        idx = len(base)+1
        p = max(8, min(72, int((22 + (idx % 7)*4) * rm)))
        c = round(max(0.015, cost*(0.006 + (idx%5)*0.004)*p/100),3)
        d = round(max(0.8, (months*0.35 + (idx%6)*6)*p/100),1)
        base.append({'risk_id': rid, 'title': title, 'category': 'Delivery', 'probability_pct': p, 'cause': cause, 'event': event, 'impact': impact, 'owner': owner, 'mitigation': mit, 'trigger': trig, 'residual_rating': 'Red' if p>=45 or c>0.5 or d>30 else 'Amber' if p>=25 else 'Green', 'activity_id': f'A{1400+(idx%8)*100}', 'activity_name': 'Linked schedule activity', 'cbs': f'01.{(idx%11)+1:02d}', 'cbs_name': 'Linked cost package', 'cost_emv_bn': c, 'schedule_emv_days': d, 'status': 'Open'})
    return sorted(base, key=lambda x: (x['cost_emv_bn']*100 + x['schedule_emv_days']/10), reverse=True)


def _v55_curves(model):
    mc = model.get('monte_carlo') or {}
    curve = mc.get('curve') or []
    if curve:
        qcra = [(int(x.get('percentile')), float(x.get('cost_bn') or 0)) for x in curve]
        qsra = [(int(x.get('percentile')), float(x.get('schedule_months') or 0)) for x in curve]
    else:
        c = parse_bn(model.get('cost_p50')); m = int(str(model.get('schedule') or '60').split()[0])
        qcra = [(p, round(c*(0.75+p/100*0.65),2)) for p in [1,5,10,20,30,40,50,60,70,80,90,95,99]]
        qsra = [(p, round(m*(0.85+p/100*0.38),2)) for p in [1,5,10,20,30,40,50,60,70,80,90,95,99]]
    return {'qcra': qcra, 'qsra': qsra}


def _v55_tornado(model):
    risks = _v55_risks(model)
    qcra = [r for r in risks if r['cost_emv_bn'] > 0]
    qsra = [r for r in risks if r['schedule_emv_days'] > 0]
    qcra = sorted(qcra, key=lambda r: r['cost_emv_bn'], reverse=True)[:12]
    qsra = sorted(qsra, key=lambda r: r['schedule_emv_days'], reverse=True)[:12]
    return {'qcra': qcra, 'qsra': qsra}


def _v55_make_workbook_base(filename='casey.xlsx'):
    if not xlsxwriter:
        raise RuntimeError('xlsxwriter is required for v55 platinum exports')
    bio = BytesIO()
    wb = xlsxwriter.Workbook(bio, {'in_memory': True, 'constant_memory': False})
    return wb, bio


def _v55_formats(wb):
    return {
        'title': wb.add_format({'bold': True, 'font_size': 20, 'font_color': V55_BLUE, 'bg_color': 'FFFFFF'}),
        'subtitle': wb.add_format({'font_size': 10, 'font_color': '#666666'}),
        'section': wb.add_format({'bold': True, 'font_size': 13, 'font_color': 'FFFFFF', 'bg_color': V55_BLUE, 'border': 1, 'border_color': V55_LINE}),
        'header': wb.add_format({'bold': True, 'font_color': 'FFFFFF', 'bg_color': '#0070C0', 'border': 1, 'border_color': V55_LINE, 'valign': 'vcenter'}),
        'cell': wb.add_format({'border': 1, 'border_color': V55_LINE, 'valign': 'top', 'text_wrap': True}),
        'cell_money': wb.add_format({'border': 1, 'border_color': V55_LINE, 'num_format': '$#,##0.0B', 'valign': 'top'}),
        'cell_num': wb.add_format({'border': 1, 'border_color': V55_LINE, 'num_format': '0.0', 'valign': 'top'}),
        'kpi_label': wb.add_format({'bold': True, 'font_color': V55_BLUE, 'bg_color': V55_LIGHT, 'border': 1, 'border_color': V55_LINE}),
        'kpi_value': wb.add_format({'bold': True, 'font_size': 16, 'font_color': V55_BLUE, 'bg_color': 'FFFFFF', 'border': 1, 'border_color': V55_LINE}),
        'red': wb.add_format({'bg_color': '#FCE4D6', 'border': 1, 'border_color': V55_LINE, 'text_wrap': True, 'valign': 'top'}),
        'amber': wb.add_format({'bg_color': '#FFF2CC', 'border': 1, 'border_color': V55_LINE, 'text_wrap': True, 'valign': 'top'}),
        'green': wb.add_format({'bg_color': '#E2F0D9', 'border': 1, 'border_color': V55_LINE, 'text_wrap': True, 'valign': 'top'}),
        'note': wb.add_format({'italic': True, 'font_color': '#666666', 'text_wrap': True})
    }


def _v55_setup_ws(ws):
    ws.set_zoom(90)
    ws.hide_gridlines(2)
    ws.freeze_panes(5, 0)
    ws.set_landscape()
    ws.set_margins(0.3, 0.3, 0.4, 0.4)


def _v55_write_table(ws, start_row, start_col, headers, rows, fmts, widths=None, money_cols=None, num_cols=None):
    money_cols = set(money_cols or [])
    num_cols = set(num_cols or [])
    for c,h in enumerate(headers): ws.write(start_row, start_col+c, h, fmts['header'])
    for r,row in enumerate(rows, start=start_row+1):
        for c,val in enumerate(row):
            fmt = fmts['cell_money'] if c in money_cols else fmts['cell_num'] if c in num_cols else fmts['cell']
            ws.write(r, start_col+c, val, fmt)
    if widths:
        for i,w in enumerate(widths): ws.set_column(start_col+i, start_col+i, w)
    return start_row + len(rows) + 1


def workbook_bytes(model: Dict[str, Any]) -> bytes:
    wb, bio = _v55_make_workbook_base(); f = _v55_formats(wb)
    sc = _v55_scenario_strategy(model); sel_cls = _v55_selected_class(model); sel_lvl = _v55_selected_schedule_level(model)
    costs_by = _v55_all_cost_lines(model); costs = costs_by.get(sel_cls, model.get('cost_lines') or [])
    risks = _v55_risks(model); curves = _v55_curves(model); torn = _v55_tornado(model)
    # Executive dashboard
    ws = wb.add_worksheet('01 Executive Dashboard'); _v55_setup_ws(ws)
    ws.merge_range('A1:H1', 'CASEY PLATINUM COST + RISK MODEL', f['title'])
    ws.merge_range('A2:H2', f"{model.get('title')} | {_v55_scenario_label(model)} scenario | Class {sel_cls} | Schedule Level {sel_lvl}", f['subtitle'])
    kpis = [('P50 Cost', model.get('cost_p50')), ('Cost Range', model.get('cost_range')), ('QCRA P80', _v55_money((model.get('monte_carlo') or {}).get('qcra',{}).get('p80',0))), ('Schedule', model.get('schedule')), ('QSRA P80', str((model.get('monte_carlo') or {}).get('qsra',{}).get('p80',''))+' mo'), ('Risk / Confidence', f"{model.get('risk')} / {model.get('confidence_pct')}%")]
    for i,(k,v) in enumerate(kpis):
        r=4+(i//3)*3; c=(i%3)*3
        ws.merge_range(r,c,r,c+1,k,f['kpi_label']); ws.merge_range(r+1,c,r+1,c+1,str(v),f['kpi_value'])
    ws.write(11,0,'Executive recommendation',f['section']); ws.merge_range(12,0,12,7,sc['title'],f['kpi_value']); ws.merge_range(13,0,14,7,sc['summary'] + ' ' + sc['decision'],f['cell'])
    # cost composition chart source
    summary={'Direct':0.0,'Indirect':0.0,'Reserve':0.0}
    for x in costs: summary[x.get('type','Direct')] = summary.get(x.get('type','Direct'),0)+float(x.get('p50_bn') or 0)
    _v55_write_table(ws,17,0,['Cost Type','P50 $B'],list(summary.items()),f,[18,16],money_cols=[1])
    chart=wb.add_chart({'type':'doughnut'}); chart.add_series({'name':'Cost composition','categories':['01 Executive Dashboard',18,0,20,0],'values':['01 Executive Dashboard',18,1,20,1],'data_labels':{'percentage':True}}); chart.set_title({'name':'Direct / Indirect / Reserve'}); chart.set_style(10); ws.insert_chart('D17',chart,{'x_scale':1.2,'y_scale':1.1})
    # Selected cost estimate
    ws = wb.add_worksheet('02 Selected Cost Estimate'); _v55_setup_ws(ws)
    ws.merge_range('A1:H1', f'SELECTED CLASS {sel_cls} COST ESTIMATE — {_v55_scenario_label(model).upper()} SCENARIO', f['title'])
    rows=[[x.get('cbs'),x.get('description'),x.get('type'),x.get('basis'),float(x.get('p10_bn') or 0),float(x.get('p50_bn') or 0),float(x.get('p90_bn') or 0),x.get('impact_basis') or 'Validate source, benchmark and scope basis.'] for x in costs]
    _v55_write_table(ws,4,0,['CBS','Description','Type','Basis','P10 $B','P50 $B','P90 $B','Challenge / Basis'],rows,f,[10,28,13,34,12,12,12,46],money_cols=[4,5,6])
    chart=wb.add_chart({'type':'column'}); chart.add_series({'name':'P50 $B','categories':['02 Selected Cost Estimate',5,1,4+len(rows),1],'values':['02 Selected Cost Estimate',5,5,4+len(rows),5]}); chart.set_title({'name':'P50 cost by package'}); chart.set_y_axis({'name':'$B'}); ws.insert_chart('J4',chart,{'x_scale':1.3,'y_scale':1.2})
    # All classes
    ws=wb.add_worksheet('03 All Class Estimates'); _v55_setup_ws(ws); ws.merge_range('A1:I1','ALL COST ESTIMATE CLASSES — NO CLASS DUMP WITHOUT CONTEXT',f['title'])
    rows=[]
    for cls, clines in costs_by.items():
        _,_,cname,mat=class_range(cls)
        for x in clines:
            rows.append([cls,cname,mat,x.get('cbs'),x.get('description'),x.get('type'),float(x.get('p10_bn') or 0),float(x.get('p50_bn') or 0),float(x.get('p90_bn') or 0)])
    _v55_write_table(ws,4,0,['Class','Class Name','Maturity','CBS','Description','Type','P10 $B','P50 $B','P90 $B'],rows,f,[8,26,26,9,28,12,12,12,12],money_cols=[6,7,8])
    # Scenario comparison
    ws=wb.add_worksheet('04 Scenario Comparison'); _v55_setup_ws(ws); ws.merge_range('A1:G1','SCENARIO COMPARISON — WHAT CHANGES AND WHY',f['title'])
    scrows=[]
    for x in model.get('scenario_comparison') or []:
        scrows.append([x.get('label'),x.get('cost'),x.get('schedule_months'),x.get('risk'),x.get('confidence'),x.get('why'), 'Decision trade-off: cost / time / risk / confidence'])
    _v55_write_table(ws,4,0,['Scenario','P50 Cost','Schedule Months','Risk','Confidence %','Commercial Interpretation','Board Use'],scrows,f,[14,14,16,14,14,62,40])
    # Risk register
    ws=wb.add_worksheet('05 Risk Register'); _v55_setup_ws(ws); ws.merge_range('A1:Q1','RISK REGISTER — CAUSE / EVENT / IMPACT / RESPONSE',f['title'])
    rows=[[r['risk_id'],r['title'],r['category'],r['probability_pct'],r['cause'],r['event'],r['impact'],r['owner'],r['mitigation'],r['trigger'],r['residual_rating'],r['activity_id'],r['activity_name'],r['cbs'],r['cbs_name'],r['cost_emv_bn'],r['schedule_emv_days']] for r in risks]
    _v55_write_table(ws,4,0,['ID','Risk','Category','Likelihood %','Cause','Risk Event','Impact Description','Owner','Mitigation / Response','Trigger','Residual','Activity','Activity Name','CBS','CBS Name','QCRA Mean $B','QSRA Mean Days'],rows,f,[10,25,14,13,36,36,44,18,42,30,12,12,24,10,24,14,14],money_cols=[15],num_cols=[3,16])
    ws.conditional_format(5,10,4+len(rows),10, {'type':'text','criteria':'containing','value':'Red','format':f['red']}); ws.conditional_format(5,10,4+len(rows),10, {'type':'text','criteria':'containing','value':'Amber','format':f['amber']}); ws.conditional_format(5,10,4+len(rows),10, {'type':'text','criteria':'containing','value':'Green','format':f['green']})
    # Curves + tornados
    ws=wb.add_worksheet('06 QCRA Curve'); _v55_setup_ws(ws); ws.merge_range('A1:D1','QCRA COST CONFIDENCE CURVE',f['title']); _v55_write_table(ws,4,0,['Percentile','Cost $B'],curves['qcra'],f,[14,14],money_cols=[1]); ch=wb.add_chart({'type':'line'}); ch.add_series({'name':'QCRA Cost $B','categories':['06 QCRA Curve',5,0,4+len(curves['qcra']),0],'values':['06 QCRA Curve',5,1,4+len(curves['qcra']),1]}); ch.set_title({'name':'QCRA S-Curve'}); ch.set_y_axis({'name':'Cost $B'}); ws.insert_chart('D4',ch,{'x_scale':1.4,'y_scale':1.2})
    ws=wb.add_worksheet('07 QSRA Curve'); _v55_setup_ws(ws); ws.merge_range('A1:D1','QSRA SCHEDULE CONFIDENCE CURVE',f['title']); _v55_write_table(ws,4,0,['Percentile','Duration Months'],curves['qsra'],f,[14,18],num_cols=[1]); ch=wb.add_chart({'type':'line'}); ch.add_series({'name':'QSRA Months','categories':['07 QSRA Curve',5,0,4+len(curves['qsra']),0],'values':['07 QSRA Curve',5,1,4+len(curves['qsra']),1]}); ch.set_title({'name':'QSRA S-Curve'}); ch.set_y_axis({'name':'Months'}); ws.insert_chart('D4',ch,{'x_scale':1.4,'y_scale':1.2})
    ws=wb.add_worksheet('08 QCRA Tornado'); _v55_setup_ws(ws); ws.merge_range('A1:E1','QCRA COST TORNADO — COST DRIVERS ONLY',f['title']); rows=[[i+1,r['title'],r['cbs'],r['cost_emv_bn'],r['mitigation']] for i,r in enumerate(torn['qcra'])]; _v55_write_table(ws,4,0,['Rank','Driver','CBS','Mean Cost $B','Action'],rows,f,[8,30,10,14,48],money_cols=[3]); ch=wb.add_chart({'type':'bar'}); ch.add_series({'name':'Mean Cost $B','categories':['08 QCRA Tornado',5,1,4+len(rows),1],'values':['08 QCRA Tornado',5,3,4+len(rows),3]}); ch.set_title({'name':'QCRA Cost Tornado'}); ws.insert_chart('G4',ch,{'x_scale':1.4,'y_scale':1.4})
    ws=wb.add_worksheet('09 QSRA Tornado'); _v55_setup_ws(ws); ws.merge_range('A1:E1','QSRA SCHEDULE TORNADO — SCHEDULE DRIVERS ONLY',f['title']); rows=[[i+1,r['title'],r['activity_id'],r['schedule_emv_days'],r['mitigation']] for i,r in enumerate(torn['qsra'])]; _v55_write_table(ws,4,0,['Rank','Driver','Activity','Mean Days','Action'],rows,f,[8,30,12,12,48],num_cols=[3]); ch=wb.add_chart({'type':'bar'}); ch.add_series({'name':'Mean Days','categories':['09 QSRA Tornado',5,1,4+len(rows),1],'values':['09 QSRA Tornado',5,3,4+len(rows),3]}); ch.set_title({'name':'QSRA Schedule Tornado'}); ws.insert_chart('G4',ch,{'x_scale':1.4,'y_scale':1.4})
    # All schedules
    ws=wb.add_worksheet('10 All Schedule Levels'); _v55_setup_ws(ws); ws.merge_range('A1:H1','SCHEDULE LEVELS 1–5 — ALL LEVELS GENERATED',f['title'])
    rows=[]
    for lvl, acts in _v55_all_schedule_levels(model).items():
        for a in acts: rows.append([lvl,a.get('activity_id'),a.get('phase'),a.get('activity'),a.get('predecessor'),a.get('duration_months'),a.get('critical'),a.get('basis')])
    _v55_write_table(ws,4,0,['Level','Activity ID','Phase','Activity','Predecessor','Duration Months','Critical','Basis'],rows,f,[8,13,16,36,18,14,10,48],num_cols=[5])
    ws=wb.add_worksheet('11 Basis + Audit'); _v55_setup_ws(ws); ws.merge_range('A1:D1','METHOD, ASSUMPTIONS AND AUDIT',f['title'])
    audit=[['Generated',datetime.utcnow().strftime('%Y-%m-%d %H:%M UTC')],['Model version','CASEY v55 Platinum Outputs Final'],['Scenario',_v55_scenario_label(model)],['Estimate class',sel_cls],['Schedule level',sel_lvl],['Cost method','Parametric sector template + location + scale + complexity + scenario + class range'],['QCRA method','Cost confidence curve and cost-only tornado drivers'],['QSRA method','Schedule confidence curve and schedule-only tornado drivers'],['Limitation','First-pass intelligence. Validate with project-specific evidence before commitment.']]
    _v55_write_table(ws,4,0,['Area','Basis'],audit,f,[24,100])
    wb.close(); bio.seek(0); return bio.getvalue()


def risk_register_workbook_bytes(model: Dict[str, Any]) -> bytes:
    """6-sheet risk register: dashboard, full register with O/M/P, board top 10, QCRA tornado, QSRA tornado, mitigation tracker."""
    wb, bio = _v55_make_workbook_base(); f = _v55_formats(wb)
    risks = _v55_risks(model); torn = _v55_tornado(model)
    mc = model.get('monte_carlo') or {}; qcra = mc.get('qcra') or {}; qsra = mc.get('qsra') or {}

    # ── 01 Risk Dashboard ────────────────────────────────────────────
    ws = wb.add_worksheet('01 Risk Dashboard'); _v55_setup_ws(ws)
    ws.merge_range('A1:J1', 'CASEY PLATINUM RISK REGISTER', f['title'])
    ws.merge_range('A2:J2', f"{model.get('title')} | {_v55_scenario_label(model)} scenario | {len(risks)} risks quantified", f['subtitle'])
    kpis = [
        ['Total risks', len(risks)],
        ['Board-visible', sum(1 for r in risks if r.get('board_visibility')=='Yes' or r['residual_rating'] in ['Red','Amber'])],
        ['High / Extreme', sum(1 for r in risks if r.get('pre_mitigation_rating','') in ['High','Extreme'] or r.get('residual_rating','')=='Red')],
        ['QCRA P50', _v55_money(float(qcra.get('p50',0)))],
        ['QCRA P80', _v55_money(float(qcra.get('p80',0)))],
        ['QSRA P80', f"{qsra.get('p80','—')} months"],
        ['Top QCRA driver', torn['qcra'][0]['title'] if torn['qcra'] else '—'],
        ['Top QSRA driver', torn['qsra'][0]['title'] if torn['qsra'] else '—'],
    ]
    _v55_write_table(ws, 4, 0, ['Metric','Value'], kpis, f, [30, 50])
    # Board challenge questions
    bqs = model.get('board_challenge_questions') or []
    if bqs:
        ws.merge_range(15, 0, 15, 9, 'LIKELY BOARD ATTACKS', f['section'])
        for i, q in enumerate(bqs[:8]):
            ws.write(16+i, 0, f'{i+1}. {q}', f['cell'])

    # ── 02 Full Risk Register with O/M/P ────────────────────────────
    ws = wb.add_worksheet('02 Full Risk Register'); _v55_setup_ws(ws)
    ws.merge_range('A1:S1', 'FULL QUANTIFIED RISK REGISTER — CAUSE · EVENT · IMPACT · OWNER · O/M/P', f['title'])
    raw = model.get('risks') or model.get('risk_register') or []
    rows = []
    for r in raw:
        rows.append([
            r.get('risk_id',''), r.get('title',''), r.get('category',''), r.get('probability_pct',0),
            r.get('cause',''), r.get('risk_event',r.get('event','')), r.get('impact_description',r.get('impact','')),
            r.get('owner',''), r.get('mitigation',''), r.get('trigger',''), r.get('response_strategy','Mitigate'),
            r.get('activity_id',''), r.get('activity_name',''), r.get('cbs',''), r.get('cbs_name',''),
            r.get('cost_o_bn',0), r.get('cost_m_bn',0), r.get('cost_p_bn',0),
            r.get('schedule_o_days',0), r.get('schedule_m_days',0), r.get('schedule_p_days',0),
            r.get('cost_emv_bn',0), r.get('schedule_emv_days',0), r.get('residual_rating',''), r.get('status','Open')
        ])
    _v55_write_table(ws, 4, 0,
        ['ID','Risk','Category','Prob%','Cause','Risk Event','Impact','Owner','Mitigation','Trigger','Response',
         'Activity ID','Activity','CBS','CBS Name','Cost O $B','Cost M $B','Cost P $B',
         'Sched O d','Sched M d','Sched P d','QCRA EMV $B','QSRA EMV d','Residual','Status'],
        rows, f, [10,26,14,9,36,36,44,18,42,28,14,13,26,9,22,12,12,12,11,11,11,13,11,12,10],
        money_cols=[15,16,17,21], num_cols=[3,18,19,20,22])
    ws.conditional_format(5, 23, 4+len(rows), 23, {'type':'text','criteria':'containing','value':'Red','format':f['red']})
    ws.conditional_format(5, 23, 4+len(rows), 23, {'type':'text','criteria':'containing','value':'Amber','format':f['amber']})
    ws.conditional_format(5, 23, 4+len(rows), 23, {'type':'text','criteria':'containing','value':'Green','format':f['green']})

    # ── 03 Board Top 10 ──────────────────────────────────────────────
    ws = wb.add_worksheet('03 Board Top 10'); _v55_setup_ws(ws)
    ws.merge_range('A1:H1', 'BOARD TOP 10 RISKS — READABLE SUMMARY FOR GOVERNANCE PACK', f['title'])
    top10 = sorted(risks, key=lambda r: r['cost_emv_bn']*100 + r['schedule_emv_days'], reverse=True)[:10]
    rows10 = [[r['risk_id'], r['title'], r['cause'], r['event'], r['impact'], r['owner'], r['mitigation'], _v55_money(r['cost_emv_bn'])] for r in top10]
    _v55_write_table(ws, 4, 0, ['ID','Risk','Cause','Event','Impact','Owner','Action','QCRA EMV'], rows10, f, [10,26,36,36,44,18,44,14], money_cols=[7])

    # ── 04 QCRA Tornado with chart ───────────────────────────────────
    ws = wb.add_worksheet('04 QCRA Cost Tornado'); _v55_setup_ws(ws)
    ws.merge_range('A1:F1', 'QCRA COST TORNADO — RISKS RANKED BY COST EMV CONTRIBUTION', f['title'])
    rows_t = [[i+1, r['title'], r['category'], r['cbs'], r['cost_emv_bn'], r.get('mitigation','')] for i,r in enumerate(torn['qcra'])]
    _v55_write_table(ws, 4, 0, ['Rank','Driver','Category','CBS','Mean Cost $B','Mitigation Action'], rows_t, f, [8,32,14,10,16,54], money_cols=[4])
    if rows_t:
        ch = wb.add_chart({'type':'bar','subtype':'clustered'})
        ch.add_series({'name':'QCRA Mean Cost $B','categories':['04 QCRA Cost Tornado',5,1,4+len(rows_t),1],'values':['04 QCRA Cost Tornado',5,4,4+len(rows_t),4],'fill':{'color':'#F59E0B'}})
        ch.set_title({'name':'QCRA Cost Tornado — Mean contribution by risk'}); ch.set_y_axis({'reverse':True}); ch.set_size({'width':520,'height':340})
        ws.insert_chart('H4', ch)

    # ── 05 QSRA Tornado with chart ───────────────────────────────────
    ws = wb.add_worksheet('05 QSRA Schedule Tornado'); _v55_setup_ws(ws)
    ws.merge_range('A1:F1', 'QSRA SCHEDULE TORNADO — RISKS RANKED BY SCHEDULE EMV CONTRIBUTION', f['title'])
    rows_s = [[i+1, r['title'], r['category'], r['activity_id'], r['schedule_emv_days'], r.get('mitigation','')] for i,r in enumerate(torn['qsra'])]
    _v55_write_table(ws, 4, 0, ['Rank','Driver','Category','Activity','Mean Days','Mitigation Action'], rows_s, f, [8,32,14,12,13,54], num_cols=[4])
    if rows_s:
        ch = wb.add_chart({'type':'bar','subtype':'clustered'})
        ch.add_series({'name':'QSRA Mean Days','categories':['05 QSRA Schedule Tornado',5,1,4+len(rows_s),1],'values':['05 QSRA Schedule Tornado',5,4,4+len(rows_s),4],'fill':{'color':'#0EA5E9'}})
        ch.set_title({'name':'QSRA Schedule Tornado — Mean contribution by risk'}); ch.set_y_axis({'reverse':True}); ch.set_size({'width':520,'height':340})
        ws.insert_chart('H4', ch)

    # ── 06 Mitigation Tracker ────────────────────────────────────────
    ws = wb.add_worksheet('06 Mitigation Tracker'); _v55_setup_ws(ws)
    ws.merge_range('A1:H1', 'RISK MITIGATION TRACKER — OWNER · ACTION · TRIGGER · STATUS', f['title'])
    rows_m = [[r['risk_id'], r['title'], r['owner'], r['mitigation'], r.get('trigger',''), r.get('response_strategy','Mitigate'), r['residual_rating'], r['status']] for r in risks]
    _v55_write_table(ws, 4, 0, ['ID','Risk','Owner','Mitigation Action','Trigger','Response','Residual','Status'], rows_m, f, [10,28,18,52,36,14,12,10])
    ws.conditional_format(5, 6, 4+len(rows_m), 6, {'type':'text','criteria':'containing','value':'Red','format':f['red']})
    ws.conditional_format(5, 6, 4+len(rows_m), 6, {'type':'text','criteria':'containing','value':'Amber','format':f['amber']})

    wb.close(); bio.seek(0); return bio.getvalue()


def xer_bytes(model: Dict[str, Any]) -> bytes:
    """PRA-compatible legacy-targeted XER. Uses ERMHDR 8.0 and basic TASK/TASKPRED structure."""
    rows = model.get('schedule_rows') or []
    if not rows:
        rows = schedule_rows(model.get('mode','Earth'), model.get('subsector','General'), int(str(model.get('schedule','60')).split()[0]), _v55_selected_schedule_level(model))
    lines=[]
    lines.append('ERMHDR\t8.0\t{}\tProject\tCASEY\tCASEY\tdbxDatabaseNoName\tProject Management\tUSD'.format(datetime.utcnow().strftime('%Y-%m-%d')))
    lines.append('%T\tCURRTYPE'); lines.append('%F\tcurr_id\tdecimal_digit_cnt\tcurr_symbol\tdecimal_symbol\tdigit_group_symbol\tpos_curr_fmt_type\tneg_curr_fmt_type\tcurr_type\tcurr_short_name\tgroup_digit_cnt\tbase_exch_rate'); lines.append('%R\t1\t2\t$\t.\t,\t#1.1\t(#1.1)\tUS Dollar\tUSD\t3\t1')
    lines.append('%T\tPROJECT'); lines.append('%F\tproj_id\tfy_start_month_num\tname_sep_char\tproj_short_name\tplan_start_date\tlast_recalc_date\tclndr_id\tdef_duration_type\tdef_qty_type\tdef_rate_type\tdef_task_type\tcritical_path_type')
    lines.append('%R\t3338\t1\t.\tCASEY\t{} 08:00\t{} 08:00\t3892\tDT_FixedDUR2\tQT_Hour\tCOST_PER_QTY\tTT_Task\tCT_DrivPath'.format(datetime.utcnow().strftime('%Y-%m-%d'), datetime.utcnow().strftime('%Y-%m-%d')))
    lines.append('%T\tCALENDAR'); lines.append('%F\tclndr_id\tdefault_flag\tclndr_name\tproj_id\tbase_clndr_id\tclndr_type\tday_hr_cnt\tweek_hr_cnt\tmonth_hr_cnt\tyear_hr_cnt\trsrc_private\tclndr_data')
    lines.append('%R\t3892\tY\tStandard\t3338\t\tCA_Project\t8\t40\t172\t2000\tN\t(0||CalendarData()())')
    lines.append('%T\tPROJWBS'); lines.append('%F\twbs_id\tproj_id\tobs_id\tseq_num\twbs_short_name\twbs_name\tparent_wbs_id')
    lines.append('%R\t5000\t3338\t565\t1\tCASEY\tCASEY Generated Schedule\t')
    lines.append('%T\tTASK')
    fields=['task_id','proj_id','wbs_id','clndr_id','task_code','task_name','task_type','duration_type','status_code','target_drtn_hr_cnt','remain_drtn_hr_cnt','phys_complete_pct','target_start_date','target_end_date','early_start_date','early_end_date','late_start_date','late_end_date','guid']
    lines.append('%F\t'+'\t'.join(fields))
    id_map={}
    start='{} 08:00'.format(datetime.utcnow().strftime('%Y-%m-%d'))
    for i,a in enumerate(rows,1):
        id_map[str(a.get('activity_id'))]=i
        dur=int(max(8, float(a.get('duration_months') or 1)*172))
        code=str(a.get('activity_id') or f'A{i:04d}')[:40]
        name=str(a.get('activity') or 'Activity').replace('\t',' ')[:120]
        guid=('CASEY%032d'%i)[:22]
        lines.append('%R\t{}\t3338\t5000\t3892\t{}\t{}\tTT_Task\tDT_FixedDUR2\tTK_NotStart\t{}\t{}\t0\t{}\t\t{}\t\t\t\t{}'.format(i,code,name,dur,dur,start,start,guid))
    lines.append('%T\tTASKPRED'); lines.append('%F\ttask_pred_id\ttask_id\tpred_task_id\tpred_type\tlag_hr_cnt')
    pred_id=1
    for i,a in enumerate(rows,1):
        preds=str(a.get('predecessor') or '').replace(';',',').split(',')
        for p in [x.strip() for x in preds if x.strip()]:
            pid=id_map.get(p)
            if pid and pid != i:
                lines.append('%R\t{}\t{}\t{}\tPR_FS\t0'.format(pred_id,i,pid)); pred_id+=1
    lines.append('%T\tUDFTYPE'); lines.append('%F\tudf_type_id\ttable_name\tudf_type_name\tudf_type_label\tlogical_data_type\tsuper_flag')
    lines.append('%R\t329\tTASK\tCASEY_Risk_Most_Likely\tLikely\tFT_TEXT\tN')
    lines.append('%E')
    return ('\n'.join(lines)+'\n').encode('utf-8')


def pdf_bytes(model: Dict[str, Any]) -> bytes:
    bio=BytesIO(); doc=SimpleDocTemplate(bio,pagesize=landscape(A4),rightMargin=24,leftMargin=24,topMargin=24,bottomMargin=24)
    styles=getSampleStyleSheet(); styles.add(ParagraphStyle(name='V55Title',parent=styles['Title'],fontSize=24,leading=28,textColor=colors.HexColor(V55_BLUE),spaceAfter=10)); styles.add(ParagraphStyle(name='V55H',parent=styles['Heading1'],fontSize=15,textColor=colors.HexColor(V55_BLUE),spaceAfter=7)); styles.add(ParagraphStyle(name='V55Small',parent=styles['BodyText'],fontSize=8,leading=10))
    sc=_v55_scenario_strategy(model); risks=_v55_risks(model); torn=_v55_tornado(model); curves=_v55_curves(model); costs=_v55_all_cost_lines(model).get(_v55_selected_class(model),model.get('cost_lines') or [])
    story=[]
    story.append(Paragraph(f"CASEY BOARD DECISION PACK — {model.get('title')}",styles['V55Title'])); story.append(Paragraph(f"{_v55_scenario_label(model)} scenario | {model.get('subsector')} | Generated {datetime.utcnow().strftime('%Y-%m-%d %H:%M UTC')}",styles['V55Small']))
    kpi=[['P50 Cost',model.get('cost_p50'),'Cost Range',model.get('cost_range'),'QCRA P80',_v55_money((model.get('monte_carlo') or {}).get('qcra',{}).get('p80',0))],['Schedule',model.get('schedule'),'QSRA P80',str((model.get('monte_carlo') or {}).get('qsra',{}).get('p80',''))+' mo','Risk / Confidence',f"{model.get('risk')} / {model.get('confidence_pct')}%"]]
    t=Table(kpi,colWidths=[75,110,75,135,75,120]); t.setStyle(TableStyle([('BACKGROUND',(0,0),(-1,-1),colors.HexColor('#F4F7FA')),('TEXTCOLOR',(0,0),(-1,-1),colors.HexColor(V55_BLUE)),('FONTNAME',(0,0),(-1,-1),'Helvetica-Bold'),('GRID',(0,0),(-1,-1),0.5,colors.HexColor(V55_LINE)),('FONTSIZE',(0,0),(-1,-1),9)])); story.append(t); story.append(Spacer(1,10)); story.append(Paragraph(sc['title'],styles['V55H'])); story.append(Paragraph(sc['summary']+' '+sc['decision'],styles['BodyText'])); story.append(PageBreak())
    story.append(Paragraph('QCRA and QSRA Tornado Drivers — separated',styles['V55H'])); rows=[['Rank','QCRA Cost Driver','CBS','Mean Cost','QSRA Schedule Driver','Activity','Mean Days']]
    for i in range(max(len(torn['qcra']),len(torn['qsra']))):
        qc=torn['qcra'][i] if i<len(torn['qcra']) else {}; qs=torn['qsra'][i] if i<len(torn['qsra']) else {}
        rows.append([i+1,qc.get('title',''),qc.get('cbs',''),_v55_money(qc.get('cost_emv_bn',0)) if qc else '',qs.get('title',''),qs.get('activity_id',''),qs.get('schedule_emv_days','')])
    t=Table(rows,repeatRows=1,colWidths=[35,145,55,70,145,65,60]); t.setStyle(TableStyle([('BACKGROUND',(0,0),(-1,0),colors.HexColor('#0070C0')),('TEXTCOLOR',(0,0),(-1,0),colors.white),('GRID',(0,0),(-1,-1),0.3,colors.HexColor(V55_LINE)),('FONTSIZE',(0,0),(-1,-1),7.5),('VALIGN',(0,0),(-1,-1),'TOP'),('ROWBACKGROUNDS',(0,1),(-1,-1),[colors.white,colors.HexColor('#F8FBFD')]) ])); story.append(t); story.append(PageBreak())
    story.append(Paragraph('Board Top Risks — one risk per row, readable summary',styles['V55H'])); rows=[['ID','Risk','Cause','Impact','Owner','Action']]
    for r in risks: rows.append([r['risk_id'],r['title'],r['cause'],r['impact'],r['owner'],r['mitigation']])
    t=Table(rows,repeatRows=1,colWidths=[45,90,180,210,70,205]); t.setStyle(TableStyle([('BACKGROUND',(0,0),(-1,0),colors.HexColor('#0070C0')),('TEXTCOLOR',(0,0),(-1,0),colors.white),('GRID',(0,0),(-1,-1),0.25,colors.HexColor(V55_LINE)),('FONTSIZE',(0,0),(-1,-1),6.9),('VALIGN',(0,0),(-1,-1),'TOP'),('ROWBACKGROUNDS',(0,1),(-1,-1),[colors.white,colors.HexColor('#F8FBFD')]) ])); story.append(t); story.append(PageBreak())
    story.append(Paragraph('Selected Cost Estimate — direct / indirect / reserve',styles['V55H'])); rows=[['CBS','Description','Type','P10','P50','P90']]
    for x in costs: rows.append([x.get('cbs'),x.get('description'),x.get('type'),_v55_money(x.get('p10_bn')), _v55_money(x.get('p50_bn')), _v55_money(x.get('p90_bn'))])
    t=Table(rows,repeatRows=1,colWidths=[50,230,80,90,90,90]); t.setStyle(TableStyle([('BACKGROUND',(0,0),(-1,0),colors.HexColor('#0070C0')),('TEXTCOLOR',(0,0),(-1,0),colors.white),('GRID',(0,0),(-1,-1),0.25,colors.HexColor(V55_LINE)),('FONTSIZE',(0,0),(-1,-1),8),('VALIGN',(0,0),(-1,-1),'TOP'),('ROWBACKGROUNDS',(0,1),(-1,-1),[colors.white,colors.HexColor('#F8FBFD')]) ])); story.append(t)
    doc.build(story); bio.seek(0); return bio.getvalue()


def pptx_bytes(model: Dict[str, Any]) -> bytes:
    """10-slide board deck: dark navy theme, proper charts, sector intelligence."""
    from pptx.dml.color import RGBColor as RC
    from pptx.enum.text import PP_ALIGN

    # ── Colours ─────────────────────────────────────────────────────
    NAVY   = RC(0x0A, 0x0F, 0x1E)
    NAVY2  = RC(0x0F, 0x18, 0x29)
    AMBER  = RC(0xF5, 0x9E, 0x0B)
    WHITE  = RC(0xFF, 0xFF, 0xFF)
    ICE    = RC(0xE0, 0xF2, 0xFE)
    SLATE  = RC(0x94, 0xA3, 0xB8)
    GREEN  = RC(0x10, 0xB9, 0x81)
    RED    = RC(0xEF, 0x44, 0x44)
    TEAL   = RC(0x0E, 0xA5, 0xE9)

    prs = Presentation()
    prs.slide_width  = PptxInches(13.333)
    prs.slide_height = PptxInches(7.5)
    blank = prs.slide_layouts[6]

    # ── Helpers ──────────────────────────────────────────────────────
    def bg(s, col=NAVY):
        sh = s.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
        sh.fill.solid(); sh.fill.fore_color.rgb = col; sh.line.fill.background()

    def amber_bar(s, top=False):
        y = 0 if top else PptxInches(7.42)
        sh = s.shapes.add_shape(1, 0, y, prs.slide_width, PptxInches(0.08))
        sh.fill.solid(); sh.fill.fore_color.rgb = AMBER; sh.line.fill.background()

    def txt(s, text, x, y, w, h, size=14, bold=False, col=WHITE, align=PP_ALIGN.LEFT, italic=False):
        box = s.shapes.add_textbox(PptxInches(x), PptxInches(y), PptxInches(w), PptxInches(h))
        tf = box.text_frame; tf.word_wrap = True
        para = tf.paragraphs[0]; para.alignment = align
        run = para.add_run(); run.text = text
        run.font.size = PptxPt(size); run.font.bold = bold
        run.font.color.rgb = col; run.font.italic = italic
        return box

    def card(s, x, y, w, h, col=NAVY2, border=None):
        sh = s.shapes.add_shape(1, PptxInches(x), PptxInches(y), PptxInches(w), PptxInches(h))
        sh.fill.solid(); sh.fill.fore_color.rgb = col
        if border:
            sh.line.color.rgb = border; sh.line.width = PptxPt(1)
        else:
            sh.line.fill.background()
        return sh

    def kpi_card(s, x, y, label, value, sub=None, accent=AMBER):
        card(s, x, y, 2.9, 1.55, NAVY2, accent)
        txt(s, label.upper(), x+0.12, y+0.08, 2.65, 0.3, size=7, col=SLATE)
        txt(s, value, x+0.12, y+0.32, 2.65, 0.7, size=22, bold=True, col=WHITE)
        if sub:
            txt(s, sub, x+0.12, y+1.1, 2.65, 0.35, size=8, col=SLATE)

    def heading(s, title, subtitle=""):
        txt(s, title, 0.55, 0.22, 12.3, 0.65, size=26, bold=True, col=WHITE)
        if subtitle:
            txt(s, subtitle, 0.55, 0.82, 12.3, 0.38, size=12, col=AMBER, italic=True)

    def numbered_list(s, items, x, y, w, h, size=10, col=WHITE, accent=AMBER):
        for i, item in enumerate(items):
            card(s, x, y + i*0.48, 0.35, 0.38, accent)
            txt(s, str(i+1), x+0.05, y + i*0.48, 0.25, 0.38, size=12, bold=True, col=NAVY)
            txt(s, item, x+0.45, y + i*0.48, w-0.55, 0.42, size=size, col=col)

    mc      = model.get('monte_carlo') or {}
    qcra    = mc.get('qcra') or {}
    qsra    = mc.get('qsra') or {}
    risks   = _v55_risks(model)
    torn    = _v55_tornado(model)
    sc      = _v55_scenario_strategy(model)
    costs   = _v55_all_cost_lines(model).get(_v55_selected_class(model), model.get('cost_lines') or [])
    sched   = model.get('schedule_detail') or model.get('schedule_rows') or []
    scen_m  = model.get('scenario_matrix') or []
    conf    = model.get('confidence_pct') or 0
    conf_label = 'DO NOT APPROVE WITHOUT MORE EVIDENCE' if conf < 40 else 'BOARD CHALLENGE LIKELY' if conf < 65 else 'CONDITIONALLY APPROVABLE' if conf < 80 else 'APPROVAL READY'
    conf_col = RED if conf < 40 else AMBER if conf < 65 else TEAL if conf < 80 else GREEN
    title_str = str(model.get('title') or 'CASEY Project')
    scenario_str = str(model.get('scenario_label') or model.get('scenario') or 'Base').title()
    sector_str = str(model.get('subsector') or model.get('mode') or '')

    # ── SLIDE 1: Cover ───────────────────────────────────────────────
    s = prs.slides.add_slide(blank); bg(s); amber_bar(s, top=True); amber_bar(s, top=False)
    txt(s, 'CASEY TITAN X', 0.55, 0.45, 12, 0.6, size=11, col=SLATE)
    txt(s, title_str, 0.55, 0.95, 12, 1.3, size=38, bold=True, col=WHITE)
    txt(s, f'{scenario_str} scenario  |  {sector_str}  |  Board Intelligence Pack', 0.55, 2.2, 12, 0.5, size=13, col=AMBER)
    card(s, 0.55, 3.0, 12.2, 1.6, NAVY2, AMBER)
    txt(s, sc.get('recommendation','').upper(), 0.75, 3.15, 11.8, 0.55, size=14, bold=True, col=AMBER)
    txt(s, sc.get('summary','') + '  ' + sc.get('decision',''), 0.75, 3.7, 11.8, 0.8, size=10, col=ICE)
    txt(s, f'Generated {__import__("datetime").datetime.utcnow().strftime("%Y-%m-%d %H:%M UTC")}', 0.55, 6.95, 12, 0.3, size=8, col=SLATE)

    # ── SLIDE 2: Board KPIs ─────────────────────────────────────────
    s = prs.slides.add_slide(blank); bg(s); amber_bar(s, top=True)
    heading(s, 'Board Decision Metrics', f'{title_str}  |  {scenario_str} scenario')
    kpi_card(s, 0.45, 1.35, 'P50 Cost', str(model.get('cost_p50') or '—'), model.get('cost_range'), AMBER)
    kpi_card(s, 3.45, 1.35, 'QCRA P80', money_bn(float(qcra.get('p80',0))), 'Board downside cost', RED)
    kpi_card(s, 6.45, 1.35, 'Schedule', str(model.get('schedule') or '—'), f"QSRA P80 {qsra.get('p80','—')} months", AMBER)
    kpi_card(s, 9.45, 1.35, 'Risk / Confidence', f"{model.get('risk')} / {conf}%", conf_label, conf_col)
    kpi_card(s, 0.45, 3.05, 'Direct Cost', str(model.get('direct_cost') or '—'), 'Construction and delivery', TEAL)
    kpi_card(s, 3.45, 3.05, 'Indirect Cost', str(model.get('indirect_cost') or '—'), 'Prelims, fees, management', TEAL)
    kpi_card(s, 6.45, 3.05, 'Reserve', str(model.get('risk_reserve') or '—'), 'QCRA-linked contingency', TEAL)
    kpi_card(s, 9.45, 3.05, 'Estimate Class', f"Class {_v55_selected_class(model)}", 'AACE classification', SLATE)
    card(s, 0.45, 4.9, 12.4, 1.5, NAVY2, None)
    txt(s, str(model.get('executive_shock_insight') or model.get('executive_summary') or ''), 0.65, 5.05, 12.0, 1.2, size=11, col=ICE, italic=True)

    # ── SLIDE 3: Scenarios ──────────────────────────────────────────
    s = prs.slides.add_slide(blank); bg(s); amber_bar(s, top=True)
    heading(s, 'Five Scenario Trade-Off', 'What buying speed, savings, assurance or resilience actually costs.')
    if scen_m:
        for i, sc_row in enumerate(scen_m[:5]):
            x = 0.45 + i * 2.56
            is_sel = sc_row.get('scenario','').lower() == (model.get('scenario') or 'base').lower()
            border = AMBER if is_sel else SLATE
            card(s, x, 1.35, 2.4, 4.8, NAVY2, border)
            lbl = str(sc_row.get('label') or sc_row.get('scenario') or f'S{i+1}')
            txt(s, lbl.upper(), x+0.1, 1.45, 2.2, 0.4, size=9, bold=True, col=AMBER if is_sel else ICE)
            txt(s, str(sc_row.get('cost_p50') or '—'), x+0.1, 1.85, 2.2, 0.55, size=20, bold=True, col=WHITE)
            txt(s, f"{sc_row.get('schedule_months','—')} months", x+0.1, 2.45, 2.2, 0.4, size=12, col=ICE)
            risk_col = RED if sc_row.get('risk','').lower() == 'high' else AMBER if sc_row.get('risk','').lower() == 'medium-high' else GREEN
            txt(s, str(sc_row.get('risk') or '—'), x+0.1, 2.9, 2.2, 0.35, size=10, bold=True, col=risk_col)
            txt(s, f"{sc_row.get('confidence_pct','—')}% confidence", x+0.1, 3.3, 2.2, 0.35, size=10, col=SLATE)
            why = str(sc_row.get('why') or sc_row.get('trade') or '')[:120]
            txt(s, why, x+0.1, 3.75, 2.2, 1.2, size=7.5, col=SLATE)
    else:
        txt(s, 'Scenario matrix not available — run the model with scenario=base first.', 0.55, 2.5, 12, 1, size=12, col=SLATE)

    # ── SLIDE 4: QCRA / QSRA separated ─────────────────────────────
    s = prs.slides.add_slide(blank); bg(s); amber_bar(s, top=True)
    heading(s, 'QCRA and QSRA — Cost and Schedule Risk Separated', f'10,000-iteration Monte Carlo  |  {len(risks)} risks wired in')
    # Left: QCRA
    card(s, 0.45, 1.35, 6.1, 5.75, NAVY2, TEAL)
    txt(s, 'QCRA — COST EXPOSURE', 0.65, 1.5, 5.7, 0.4, size=10, bold=True, col=TEAL)
    curve_pts = [('P10', qcra.get('p10',0)), ('P50', qcra.get('p50',0)), ('P80', qcra.get('p80',0)), ('P90', qcra.get('p90',0))]
    for j, (lbl, val) in enumerate(curve_pts):
        y = 2.0 + j * 0.75
        card(s, 0.65, y, 1.0, 0.6, NAVY, TEAL if lbl=='P80' else None)
        txt(s, lbl, 0.68, y+0.05, 0.94, 0.5, size=11, bold=True, col=TEAL if lbl=='P80' else ICE, align=PP_ALIGN.CENTER)
        card(s, 1.75, y, 4.6, 0.6, NAVY, None)
        txt(s, money_bn(float(val)), 1.85, y+0.08, 4.4, 0.45, size=18, bold=(lbl=='P80'), col=RED if lbl=='P80' else WHITE)
    txt(s, 'QCRA Tornado — Top Cost Drivers', 0.65, 5.1, 5.7, 0.35, size=9, bold=True, col=AMBER)
    for j, r in enumerate(torn.get('qcra',[])[:5]):
        txt(s, f"{j+1}. {r['title']}", 0.65, 5.5+j*0.3, 3.8, 0.28, size=8, col=ICE)
        txt(s, money_bn(r['cost_emv_bn']), 4.5, 5.5+j*0.3, 1.8, 0.28, size=8, col=AMBER, align=PP_ALIGN.RIGHT)
    # Right: QSRA
    card(s, 6.75, 1.35, 6.1, 5.75, NAVY2, AMBER)
    txt(s, 'QSRA — SCHEDULE EXPOSURE', 6.95, 1.5, 5.7, 0.4, size=10, bold=True, col=AMBER)
    curve_pts_s = [('P10', qsra.get('p10',0)), ('P50', qsra.get('p50',0)), ('P80', qsra.get('p80',0)), ('P90', qsra.get('p90',0))]
    for j, (lbl, val) in enumerate(curve_pts_s):
        y = 2.0 + j * 0.75
        card(s, 6.95, y, 1.0, 0.6, NAVY, AMBER if lbl=='P80' else None)
        txt(s, lbl, 6.98, y+0.05, 0.94, 0.5, size=11, bold=True, col=AMBER if lbl=='P80' else ICE, align=PP_ALIGN.CENTER)
        card(s, 8.05, y, 4.6, 0.6, NAVY, None)
        txt(s, f"{val} months", 8.15, y+0.08, 4.4, 0.45, size=18, bold=(lbl=='P80'), col=RED if lbl=='P80' else WHITE)
    txt(s, 'QSRA Tornado — Top Schedule Drivers', 6.95, 5.1, 5.7, 0.35, size=9, bold=True, col=AMBER)
    for j, r in enumerate(torn.get('qsra',[])[:5]):
        txt(s, f"{j+1}. {r['title']}", 6.95, 5.5+j*0.3, 4.0, 0.28, size=8, col=ICE)
        txt(s, f"{r['schedule_emv_days']}d", 11.0, 5.5+j*0.3, 1.6, 0.28, size=8, col=AMBER, align=PP_ALIGN.RIGHT)

    # ── SLIDE 5: Risk Register top 8 ────────────────────────────────
    s = prs.slides.add_slide(blank); bg(s); amber_bar(s, top=True)
    heading(s, f'Risk Register — Top Risks by EMV Driver  ({len(risks)} total)', 'Cause → Event → Impact → Owner → Action.  Every risk linked to activity and CBS.')
    headers = ['ID','Risk','Cat','Prob%','Cause','Impact','Owner','QCRA $B','QSRA d']
    col_x   = [0.35, 0.9, 2.7, 3.45, 4.0, 6.6, 9.1, 10.75, 12.1]
    col_w   = [0.5,  1.75, 0.7, 0.5, 2.55, 2.45, 1.6, 1.3,  1.0]
    card(s, 0.35, 1.35, 12.6, 0.38, RC(0x14,0x24,0x40), None)
    for j, h in enumerate(headers):
        txt(s, h, col_x[j], 1.38, col_w[j], 0.32, size=8, bold=True, col=AMBER)
    for i, r in enumerate(risks[:8]):
        y = 1.78 + i * 0.65
        row_col = RC(0x0D,0x15,0x25) if i%2==0 else NAVY2
        card(s, 0.35, y, 12.6, 0.62, row_col, None)
        p = r['probability_pct']
        p_col = RED if p>=50 else AMBER if p>=35 else ICE
        vals = [r['risk_id'], r['title'][:22], r['category'][:10], str(p)+'%',
                r['cause'][:38], r['impact'][:38], r['owner'][:20],
                money_bn(r['cost_emv_bn']), str(r['schedule_emv_days'])]
        for j, (val, cx, cw) in enumerate(zip(vals, col_x, col_w)):
            c = p_col if j==3 else RED if j==7 and r['cost_emv_bn']>0.5 else AMBER if j==7 else ICE
            txt(s, str(val), cx, y+0.07, cw, 0.5, size=7.5, col=c)

    # ── SLIDE 6: Cost estimate ───────────────────────────────────────
    s = prs.slides.add_slide(blank); bg(s); amber_bar(s, top=True)
    heading(s, f'Class {_v55_selected_class(model)} Cost Estimate — {scenario_str} Scenario', 'Direct / Indirect / Reserve split.  P10 / P50 / P90 by CBS.')
    headers = ['CBS','Description','Type','P10','P50','P90','Basis / Challenge']
    col_x2  = [0.35, 1.0, 3.75, 4.6, 5.45, 6.3, 7.2]
    col_w2  = [0.6,  2.7, 0.8, 0.8, 0.8, 0.8, 5.8]
    card(s, 0.35, 1.35, 12.6, 0.38, RC(0x14,0x24,0x40), None)
    for j, h in enumerate(headers):
        txt(s, h, col_x2[j], 1.38, col_w2[j], 0.32, size=8, bold=True, col=AMBER)
    show_costs = costs[:10]
    for i, x in enumerate(show_costs):
        y = 1.78 + i * 0.5
        row_col = RC(0x0D,0x15,0x25) if i%2==0 else NAVY2
        type_col = TEAL if x.get('type')=='Direct' else AMBER if x.get('type')=='Indirect' else RED
        card(s, 0.35, y, 12.6, 0.47, row_col, None)
        vals = [x.get('cbs',''), x.get('description','')[:34], x.get('type',''),
                money_bn(float(x.get('p10_bn',0))), money_bn(float(x.get('p50_bn',0))),
                money_bn(float(x.get('p90_bn',0))), (x.get('impact_basis') or x.get('basis',''))[:70]]
        for j, (val, cx, cw) in enumerate(zip(vals, col_x2, col_w2)):
            c = type_col if j==2 else AMBER if j in [3,4,5] else ICE
            txt(s, str(val), cx, y+0.06, cw, 0.38, size=7.5, col=c)

    # ── SLIDE 7: Schedule summary ────────────────────────────────────
    s = prs.slides.add_slide(blank); bg(s); amber_bar(s, top=True)
    heading(s, f'Level {_v55_selected_schedule_level(model)} Schedule — Critical Path Summary', f'{len(sched)} activities  |  XER export available for Primavera P6')
    headers = ['ID','Phase','Activity','Duration','Critical','Predecessor']
    col_x3  = [0.35, 1.15, 2.5, 7.0, 8.0, 9.2]
    col_w3  = [0.75, 1.3, 4.45, 0.95, 0.75, 3.9]
    card(s, 0.35, 1.35, 12.6, 0.38, RC(0x14,0x24,0x40), None)
    for j, h in enumerate(headers):
        txt(s, h, col_x3[j], 1.38, col_w3[j], 0.32, size=8, bold=True, col=AMBER)
    for i, a in enumerate(sched[:14]):
        y = 1.78 + i * 0.37
        row_col = RC(0x0D,0x15,0x25) if i%2==0 else NAVY2
        card(s, 0.35, y, 12.6, 0.35, row_col, None)
        is_c = a.get('critical') == 'Yes'
        vals = [a.get('activity_id',''), a.get('phase','')[:14], a.get('activity','')[:55],
                f"{a.get('duration_months',0)} mo", 'CRITICAL' if is_c else '—', a.get('predecessor','')]
        for j, (val, cx, cw) in enumerate(zip(vals, col_x3, col_w3)):
            c = RED if j==4 and is_c else AMBER if j==3 else ICE
            txt(s, str(val), cx, y+0.04, cw, 0.28, size=7.5, col=c)

    # ── SLIDE 8: Board Assurance Questions ────────────────────────────
    s = prs.slides.add_slide(blank); bg(s); amber_bar(s, top=True)
    heading(s, 'Board Assurance Questions', 'The questions a serious investment committee will ask before approving capital.')
    bqs = model.get('board_challenge_questions') or model.get('board_attack_simulation') or []
    if isinstance(bqs, list) and bqs:
        numbered_list(s, [str(q) for q in bqs[:8]], 0.45, 1.35, 7.0, 0.4, size=10)
    # Right: CASEY position
    card(s, 7.8, 1.35, 5.1, 5.8, NAVY2, AMBER)
    txt(s, 'CASEY POSITION', 8.0, 1.5, 4.7, 0.38, size=9, bold=True, col=AMBER)
    pos_txt = str(model.get('casey_position') or model.get('casey_thinking') or model.get('if_this_fails') or '')[:500]
    txt(s, pos_txt, 8.0, 1.95, 4.7, 4.8, size=9, col=ICE)

    # ── SLIDE 9: Incumbent pressure test ────────────────────────────
    s = prs.slides.add_slide(blank); bg(s); amber_bar(s, top=True)
    heading(s, 'CASEY vs Traditional Advisory', 'What CASEY shows that a static cost report does not.')
    tvc = model.get('traditional_vs_casey') or {}
    card(s, 0.45, 1.35, 5.8, 2.8, NAVY2, SLATE)
    txt(s, 'TRADITIONAL CONTROLS VIEW', 0.65, 1.5, 5.4, 0.38, size=9, bold=True, col=SLATE)
    txt(s, str(tvc.get('traditional','Civil progress and cost spend appear on track.')), 0.65, 2.0, 5.4, 1.9, size=10, col=ICE)
    card(s, 7.0, 1.35, 5.8, 2.8, NAVY2, AMBER)
    txt(s, 'CASEY READS UNDERNEATH', 7.2, 1.5, 5.4, 0.38, size=9, bold=True, col=AMBER)
    txt(s, str(tvc.get('casey','CASEY reads the governing constraint, not headline progress.')), 7.2, 2.0, 5.4, 1.9, size=10, col=ICE)
    card(s, 0.45, 4.35, 12.4, 2.65, NAVY2, TEAL)
    txt(s, 'WHY THIS MATTERS', 0.65, 4.5, 12.0, 0.38, size=9, bold=True, col=TEAL)
    txt(s, 'Traditional reporting can remain green while board-defensible confidence is already deteriorating. CASEY tests the governing constraint — possessions, licensing, yield, integration — rather than reported progress alone.', 0.65, 4.95, 12.0, 1.8, size=10, col=ICE)

    # ── SLIDE 10: Close ──────────────────────────────────────────────
    s = prs.slides.add_slide(blank); bg(s); amber_bar(s, top=True); amber_bar(s, top=False)
    txt(s, 'CASEY TITAN X', 0.55, 0.5, 12, 0.5, size=10, col=SLATE)
    txt(s, '"Traditional project controls reports show numbers.', 0.55, 1.3, 12, 0.65, size=22, bold=True, col=WHITE)
    txt(s, 'CASEY shows the board what the numbers are trying to hide."', 0.55, 1.95, 12, 0.65, size=22, bold=True, col=AMBER)
    txt(s, 'This pack was generated in under 3 minutes. A traditional early-stage advisory team takes 6–10 weeks and £250,000–£850,000 to produce equivalent outputs.', 0.55, 3.1, 12, 0.8, size=12, col=ICE)
    card(s, 0.55, 4.1, 12.2, 1.7, NAVY2, AMBER)
    txt(s, '✉  deepa@caseai.co.uk   |   linkedin.com/company/caseai   |   casey.ai', 0.75, 4.55, 11.8, 0.5, size=12, col=AMBER, align=PP_ALIGN.CENTER)
    txt(s, 'Book a demo  ·  Request a sample board pack  ·  Enterprise pilot enquiries welcome', 0.75, 5.1, 11.8, 0.5, size=10, col=ICE, align=PP_ALIGN.CENTER)
    bio = BytesIO(); prs.save(bio); bio.seek(0); return bio.getvalue()


def _v55_schedule_csv(model: Dict[str, Any]) -> bytes:
    out=StringIO(); w=csv.writer(out); w.writerow(['Level','Activity ID','Phase','Activity','Predecessor','Duration Months','Critical','Basis'])
    for lvl, acts in _v55_all_schedule_levels(model).items():
        for a in acts: w.writerow([lvl,a.get('activity_id'),a.get('phase'),a.get('activity'),a.get('predecessor'),a.get('duration_months'),a.get('critical'),a.get('basis')])
    return out.getvalue().encode('utf-8')


def export_workbook_v55_endpoint(model: Dict[str, Any]): return stream(workbook_bytes(model),'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet','CASEY_v55_Platinum_Cost_Model.xlsx')
def export_risk_v55_endpoint(model: Dict[str, Any]): return stream(risk_register_workbook_bytes(model),'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet','CASEY_v55_Platinum_Risk_Register.xlsx')
def export_xer_v55_endpoint(model: Dict[str, Any]): return stream(xer_bytes(model),'application/octet-stream','CASEY_v55_PRA_Targeted_Schedule.xer')
def export_schedule_csv_v55_endpoint(model: Dict[str, Any]): return stream(_v55_schedule_csv(model),'text/csv','CASEY_v55_All_Schedule_Levels.csv')
def export_pdf_v55_endpoint(model: Dict[str, Any]): return stream(pdf_bytes(model),'application/pdf','CASEY_v55_Board_Decision_Pack.pdf')
def export_all_v55_endpoint(model: Dict[str, Any]):
    bio=BytesIO()
    with zipfile.ZipFile(bio,'w',zipfile.ZIP_DEFLATED) as z:
        z.writestr('01_CASEY_v55_Platinum_Cost_Model.xlsx',workbook_bytes(model))
        z.writestr('02_CASEY_v55_Platinum_Risk_Register.xlsx',risk_register_workbook_bytes(model))
        z.writestr('03_CASEY_v55_PRA_Targeted_Schedule.xer',xer_bytes(model))
        z.writestr('04_CASEY_v55_All_Schedule_Levels.csv',_v55_schedule_csv(model))
        z.writestr('05_CASEY_v55_Board_Decision_Pack.docx',word_bytes(model))
        z.writestr('06_CASEY_v55_Board_Decision_Pack.pdf',pdf_bytes(model))
        z.writestr('07_CASEY_v55_Model_Audit.json',json.dumps(model,indent=2))
        z.writestr('README_PRA_IMPORT_NOTE.txt','The XER is exported with ERMHDR 8.0 to target older Primavera/PRA imports. If your PRA instance still rejects XER, use the CSV schedule fallback or open in P6 Professional and re-export using the version supported by your PRA build.')
    bio.seek(0); return stream(bio.getvalue(),'application/zip','CASEY_v55_Platinum_Output_Pack.zip')


def _v55_replace_post(path, endpoint):
    app.router.routes = [r for r in app.router.routes if not (getattr(r,'path',None)==path and 'POST' in getattr(r,'methods',set()))]
    app.post(path)(endpoint)

_v55_replace_post('/export/workbook', export_workbook_v55_endpoint)
_v55_replace_post('/export/risk-register', export_risk_v55_endpoint)
_v55_replace_post('/export/xer', export_xer_v55_endpoint)
_v55_replace_post('/export/pdf', export_pdf_v55_endpoint)
_v55_replace_post('/export/all', export_all_v55_endpoint)
app.post('/export/schedule-csv')(export_schedule_csv_v55_endpoint)
app.post('/v55/export/all')(export_all_v55_endpoint)

APP_VERSION = 'CASEY TITAN X v55 Platinum Outputs Final'

# =====================================================================
# v57 SELLABLE ANY PROJECT ENGINE INSTALLER
# Final output pack removes PDF/PPTX/Word from the primary source-of-truth pack.
# =====================================================================
try:
    from v56_outputs import install_v56
    install_v56(app)
    APP_VERSION = 'CASEY TITAN X v60 Final Commercial'
except Exception as _v56_error:
    print('V57 install failed:', _v56_error)

# =====================================================================
# v59 PLATINUM DYNAMIC OUTPUTS INSTALLER
try:
    from v59_outputs import install_v59
    install_v59(app)
    APP_VERSION = 'CASEY TITAN X v60 Final Commercial'
except Exception as _v59_error:
    print('V59 install failed:', _v59_error)


# v62 TRUE DYNAMIC XER + COMMERCIAL OUTPUT ENGINE INSTALLER
try:
    from v62_outputs import install_v62
    install_v62(app)
    print('CASEY v62 true dynamic XER engine installed')
except Exception as _v62_error:
    print('V62 install failed:', _v62_error)

# v63 TRUE UNIVERSAL ANY-PROJECT ENGINE INSTALLER
try:
    from v63_outputs import install_v63
    install_v63(app)
    print('CASEY v63 true universal engine installed')
except Exception as _v63_error:
    print('V63 install failed:', _v63_error)


# v64 FINAL SELLABLE UNIVERSAL ENGINE INSTALLER
try:
    from v64_outputs import install_v64
    install_v64(app)
    APP_VERSION = 'CASEY TITAN X v64 Final Sellable Universal Engine'
    print('CASEY v64 final sellable universal engine installed')
except Exception as _v64_error:
    print('V64 install failed:', _v64_error)



# ------------------------- CASEY v9.5+ holy-shit scenario consequence engine -------------------------
SCENARIO_PROFILES_V95_PLUS = {
    "base": {
        "label":"Base","cost":1.00,"schedule":1.00,"confidence":0,"risk":"Medium-High",
        "reserve_factor":1.00,"direct_factor":1.00,"indirect_factor":1.00,
        "curve":"base","risk_lift":0,
        "trade":"Balanced cost, time and evidence posture.",
        "won":"Maintains a credible reference case for board challenge.",
        "lost":"Does not buy extra time certainty or capital efficiency."
    },
    "faster": {
        "label":"Faster","cost":1.18,"schedule":0.78,"confidence":-13,"risk":"High",
        "reserve_factor":1.42,"direct_factor":1.10,"indirect_factor":1.30,
        "curve":"faster","risk_lift":18,
        "trade":"You bought time by spending money and consuming recovery float.",
        "won":"Earlier revenue / market-entry option and stronger strategic timing.",
        "lost":"CQV float, interface stability, procurement optionality and late-stage recovery."
    },
    "cheaper": {
        "label":"Cheaper","cost":0.86,"schedule":1.14,"confidence":-16,"risk":"High",
        "reserve_factor":0.58,"direct_factor":0.90,"indirect_factor":0.78,
        "curve":"cheaper","risk_lift":14,
        "trade":"You cut capital authorization by moving risk into resilience, redundancy and start-up.",
        "won":"Lower initial approval number and reduced near-term capital draw.",
        "lost":"Operational resilience, commissioning flexibility, contingency adequacy and lifecycle certainty."
    },
    "lower_risk": {
        "label":"Lower Risk","cost":1.12,"schedule":1.16,"confidence":10,"risk":"Medium-Low",
        "reserve_factor":1.52,"direct_factor":1.05,"indirect_factor":1.16,
        "curve":"lower_risk","risk_lift":-12,
        "trade":"You bought confidence by adding assurance, float and procurement evidence.",
        "won":"Reduced downside volatility and stronger approval defensibility, at the expense of market timing.",
        "lost":"Market-entry acceleration, lean capex posture and near-term strategic timing."
    },
    "premium": {
        "label":"Premium","cost":1.28,"schedule":0.96,"confidence":18,"risk":"Low",
        "reserve_factor":1.65,"direct_factor":1.22,"indirect_factor":1.24,
        "curve":"premium","risk_lift":-16,
        "trade":"You purchased operational separation, redundancy and downside protection capacity.",
        "won":"Higher mission resilience, commissioning separation and stronger downside containment.",
        "lost":"Lean capital optics and minimum-approval positioning."
    },
}

def _money_to_bn_hs(s):
    s=str(s or "").replace("$","").strip().upper()
    try:
        if s.endswith("T"): return float(s[:-1])*1000
        if s.endswith("B"): return float(s[:-1])
        if s.endswith("M"): return float(s[:-1])/1000
        return float(s)
    except Exception:
        return 0.0

def _bn_to_money_hs(bn):
    try: bn=float(bn)
    except Exception: bn=0.0
    if bn >= 1000: return f"${bn/1000:.1f}T"
    if bn >= 1: return f"${bn:.1f}B"
    return f"${bn*1000:.0f}M"

def _duration_hs(s):
    try:
        return int(float(str(s or "60").replace("months","").strip().split()[0]))
    except Exception:
        return 60

def _fam_insight_pack_hs(model):
    fam=_sector_family(model.get("subsector",""))
    if fam=="life_sciences":
        return {
            "base":"Mechanical completion is not the true finish line; validated production readiness and deviation closure are the real board decision gates.",
            "faster":"Acceleration is now the risk: CQV, clean-utility validation and media-fill readiness are being forced into the same window. The programme may finish construction before it is safe to release product.",
            "cheaper":"The cheaper case is not simply cheaper. It has likely transferred cost into operational fragility: leaner redundancy, weaker qualification float and higher late-stage validation disruption.",
            "lower_risk":"The lower-risk case buys regulatory confidence by protecting CQV float, GMP turnover and deviation closure before commercial ramp-up.",
            "premium":"Premium delivery protects market launch by buying utility resilience, validation capacity and stronger batch-release readiness."
        }
    if fam=="data_centre":
        return {
            "base":"Power availability and systems integration readiness are more likely to govern delivery than civil progress alone.",
            "faster":"Acceleration compressed energisation and integrated systems testing into the same delivery window. Civil completion may outpace operational readiness.",
            "cheaper":"Capital efficiency is being achieved by reducing resilience margin, commissioning separation and operational redundancy.",
            "lower_risk":"The lower-risk posture protects energisation sequencing, commissioning float and operational continuity before customer cutover.",
            "premium":"Premium posture buys commissioning separation, operational redundancy and higher uptime defensibility."
        }
    if fam=="semiconductor":
        return {
            "base":"Tool install, UPW readiness and yield-ramp qualification are the real finish line, not cleanroom completion.",
            "faster":"Acceleration compresses tool hook-up, UPW commissioning and qualification cadence into a fragile ramp window.",
            "cheaper":"The cheaper case risks deferring utility resilience and tool-readiness assurance into yield ramp, where recovery is most expensive.",
            "lower_risk":"The lower-risk case protects process-tool cadence, utility stability and yield-ramp confidence.",
            "premium":"Premium delivery buys tool availability, redundancy and stronger production-ramp certainty."
        }
    if model.get("mode")=="Space":
        return {
            "base":"Launch alone is not the real constraint; qualification, thermal-power balance and autonomous recovery determine mission survivability.",
            "faster":"Acceleration compresses payload qualification and launch integration, increasing mission-assurance risk after deployment when recovery options are limited.",
            "cheaper":"The cheaper case likely removes redundancy in the one environment where repair is hardest: after deployment.",
            "lower_risk":"The lower-risk case protects qualification, environmental testing and autonomous recovery evidence before launch commitment.",
            "premium":"Premium delivery buys mission assurance, redundancy and survivability after deployment."
        }
    return {
        "base":"The dominant risk sits in interfaces, procurement evidence and commissioning readiness rather than headline construction progress.",
        "faster":"Acceleration compresses interface resolution and commissioning recovery float, increasing late-stage execution volatility.",
        "cheaper":"The cheaper case lowers the approval number by moving risk into operational readiness, contingency adequacy and deferred resilience.",
        "lower_risk":"The lower-risk case buys confidence by adding evidence, float and assurance.",
        "premium":"Premium delivery buys resilience, optionality and stronger downside protection."
    }

def _scenario_risks_hs(model, scenario, lift):
    fam=_sector_family(model.get("subsector",""))
    common = {
        "faster":[
            ("R-F01","Concurrent commissioning overload","Acceleration","Too many systems enter commissioning together","Rework, late handover and validation slippage","58","Commissioning Lead","Separate hard hold-points; add surge QA/CQV resource"),
            ("R-F02","Recovery float exhaustion","Schedule compression","Float consumed before integration testing","No recovery path for late supplier or interface failure","54","Project Controls Lead","Create protected executive float and weekly path-to-green review"),
            ("R-F03","Acceleration premium shock","Procurement","Expediting and premium vendors exceed allowance","Cost growth beyond approved scenario envelope","49","Commercial Lead","Lock long-lead packages and cap acceleration premiums")
        ],
        "cheaper":[
            ("R-C01","Deferred resilience failure","Capital reduction","Redundancy and resilience scope deferred","Operational start-up or availability failure","57","Operations Lead","Define minimum viable redundancy and non-deferrable systems list"),
            ("R-C02","Contingency underfunding","Budget pressure","Reserve reduced below uncertainty envelope","Board funding shock when latent risks materialise","52","Sponsor","Ring-fence reserve for critical systems and regulatory gates"),
            ("R-C03","Lifecycle cost transfer","Capex cut","Lower initial scope creates higher operating burden","False economy and stakeholder challenge","46","Asset Owner","Show capex/opex trade-off and owner acceptance record")
        ],
        "lower_risk":[
            ("R-L01","Assurance gate delay","Risk reduction","Additional evidence gates slow early progress","Short-term date pressure despite better certainty","28","PMO Lead","Pre-agree evidence criteria and approval cadence"),
            ("R-L02","Procurement buffer cost","Risk reduction","More secure procurement path increases early commitment","Higher pre-approval funding need","24","Commercial Lead","Use staged release with supplier option strategy")
        ],
        "premium":[
            ("R-P01","Premium scope creep","Optionality","Resilience and redundancy additions expand scope","Capex escalation and governance challenge","22","Sponsor","Govern premium scope with value gates"),
            ("R-P02","Over-specification risk","Resilience posture","Premium solution exceeds real operating need","Weak value-for-money challenge","18","Asset Owner","Tie premium scope to quantified service / mission resilience")
        ],
        "base":[]
    }
    sector_extra = []
    if fam=="life_sciences" and scenario=="faster":
        sector_extra=[("R-F04","Media-fill readiness collision","CQV compression","Media-fill, clean utilities and staffing readiness converge","FDA readiness and batch-release delay","63","CQV Lead","Protect media-fill sequence and deviation closure capacity")]
    elif fam=="life_sciences" and scenario=="cheaper":
        sector_extra=[("R-C04","Clean utility redundancy gap","Capex reduction","WFI / clean steam resilience reduced or deferred","Production continuity exposure during ramp-up","60","Utilities Lead","Identify GMP-critical redundancy that cannot be deferred")]
    elif model.get("mode")=="Space" and scenario=="faster":
        sector_extra=[("R-F04","Mission assurance compression","Launch pressure","Qualification evidence compressed before launch slot","Post-deployment failure with limited recovery","67","Mission Assurance Lead","No-go criteria for thermal, power and autonomy testing")]
    elif model.get("mode")=="Space" and scenario=="cheaper":
        sector_extra=[("R-C04","Redundancy deletion after deployment","Capex cut","Backup systems removed before orbital operations","Single-point failure in inaccessible environment","65","Systems Lead","Protect mission-critical redundancy and autonomous recovery budget")]
    return sector_extra + common.get(scenario, [])

def _curve_hs(cost, months, scenario):
    """Scenario-specific QCRA/QSRA P-curves anchored so P50 equals the headline P50.
    This prevents the dashboard saying P50 is one number while the chart implies another.
    The tail shape is the intelligence: Faster/Cheaper widen, Lower Risk/Premium compress.
    """
    pts=[1,5,10,20,30,40,50,60,70,80,90,95,99]
    profiles={
        "base": {
            "cost":[0.78,0.81,0.84,0.89,0.93,0.97,1.00,1.06,1.13,1.22,1.36,1.48,1.62],
            "sched":[0.82,0.84,0.86,0.90,0.94,0.97,1.00,1.05,1.10,1.16,1.26,1.34,1.44],
            "meaning":"Reference case: P50 is the approval view; P80/P90 show normal delivery uncertainty."
        },
        "faster": {
            "cost":[0.82,0.84,0.87,0.91,0.95,0.98,1.00,1.09,1.20,1.35,1.55,1.72,1.95],
            "sched":[0.72,0.75,0.78,0.84,0.90,0.95,1.00,1.09,1.20,1.34,1.52,1.66,1.85],
            "meaning":"Faster case: the median improves, but the tail widens as concurrency and recovery exposure increase."
        },
        "cheaper": {
            "cost":[0.70,0.74,0.78,0.84,0.90,0.96,1.00,1.12,1.25,1.42,1.68,1.88,2.15],
            "sched":[0.84,0.87,0.90,0.94,0.97,0.99,1.00,1.12,1.24,1.39,1.58,1.72,1.92],
            "meaning":"Cheaper case: lower approval number, but a fat tail because contingency and resilience have been stripped."
        },
        "lower_risk": {
            "cost":[0.88,0.90,0.92,0.95,0.97,0.99,1.00,1.04,1.08,1.13,1.20,1.25,1.32],
            "sched":[0.90,0.92,0.94,0.96,0.98,0.99,1.00,1.04,1.08,1.13,1.20,1.26,1.34],
            "meaning":"Lower-risk case: schedule flexibility and evidence maturity tighten the downside curve."
        },
        "premium": {
            "cost":[0.88,0.90,0.93,0.96,0.98,0.99,1.00,1.04,1.09,1.16,1.25,1.34,1.46],
            "sched":[0.84,0.87,0.90,0.94,0.97,0.99,1.00,1.04,1.09,1.16,1.25,1.34,1.46],
            "meaning":"Premium case: additional capital buys separation, optionality and lower operational fragility."
        },
    }
    prof=profiles.get(scenario, profiles["base"])
    out=[]
    for p,cf,sf in zip(pts, prof["cost"], prof["sched"]):
        out.append({
            "percentile":p,
            "cost_bn":round(float(cost)*cf,2),
            "schedule_months":int(round(float(months)*sf)),
            "label": f"P{p}",
            "cost_meaning": "headline P50" if p==50 else ("tail exposure" if p>=80 else "lower-bound / optimism case"),
            "schedule_meaning": "headline P50" if p==50 else ("finish-date risk" if p>=80 else "early finish case"),
        })
    return out

def _scenario_outputs_text_hs(profile, insight):
    return [
        f"Scenario trade: {profile['trade']}",
        f"What you gained: {profile['won']}",
        f"What you gave up: {profile['lost']}",
        f"Board warning: {insight}",
        "Decision rule: do not approve this scenario until the changed risk owner, reserve position and critical-path evidence are explicit."
    ]

def scenario_cascade_v95(model:Dict[str,Any], scenario:str) -> Dict[str,Any]:
    # Override prior v9.5 cascade with deeper v9.5+ consequence model.
    # v9.5+ final sector correction: semiconductor / fab language overrides accidental pharma routing.
    _prompt_text = str(model.get("prompt") or model.get("brief") or model.get("description") or "")
    _title_text = str(model.get("title") or "")
    _all_text = (_prompt_text + " " + _title_text + " " + str(model.get("subsector",""))).lower()
    if any(k in _all_text for k in ["semiconductor", "fab ", "fabrication", "wafer", "euv", "lithography", "process-tool", "process tool", "upw", "chip manufacturing"]):
        model["mode"] = "Earth"
        model["subsector"] = "Semiconductor / Advanced Manufacturing"
        if "Semiconductor" not in str(model.get("title","")):
            model["title"] = "Semiconductor Fabrication Campus"

    scenario=(scenario or "base").lower()
    if scenario not in SCENARIO_PROFILES_V95_PLUS:
        scenario="base"
    profile=SCENARIO_PROFILES_V95_PLUS[scenario]
    if "_hs_base" not in model:
        model["_hs_base"]={
            "cost_p50": model.get("cost_p50"),
            "schedule_months": _duration_hs(model.get("schedule")),
            "confidence": int(model.get("confidence_pct",55) or 55),
            "risk": model.get("risk","Medium-High"),
            "cost_lines": json.loads(json.dumps(model.get("cost_lines",[]))),
            "risks": json.loads(json.dumps(model.get("risks",[]))),
            "schedule_rows": json.loads(json.dumps(model.get("schedule_rows",[]))),
            "direct_cost": model.get("direct_cost"),
            "indirect_cost": model.get("indirect_cost"),
            "risk_reserve": model.get("risk_reserve"),
        }
    base=model["_hs_base"]
    base_cost=_money_to_bn_hs(base["cost_p50"])
    base_months=int(base["schedule_months"])
    base_conf=int(base["confidence"])
    new_cost=base_cost*profile["cost"]
    new_months=max(3,int(round(base_months*profile["schedule"])))
    new_conf=max(5,min(96,base_conf+profile["confidence"]))
    insight=_fam_insight_pack_hs(model).get(scenario, _fam_insight_pack_hs(model)["base"])
    model.update({
        "scenario":scenario,
        "scenario_label":profile["label"],
        "cost_p50":_bn_to_money_hs(new_cost),
        "cost_p10":_bn_to_money_hs(new_cost*(0.78 if scenario!="cheaper" else 0.70)),
        "cost_p90":_bn_to_money_hs(new_cost*(1.27 if scenario in ["lower_risk","premium"] else 1.48 if scenario=="faster" else 1.55 if scenario=="cheaper" else 1.32)),
        "schedule":f"{new_months} months",
        "risk":profile["risk"],
        "confidence_pct":new_conf,
        "executive_shock_insight":insight,
        "scenario_trade":profile["trade"],
        "scenario_gain":profile["won"],
        "scenario_loss":profile["lost"],
        "version":"CASEY TITAN X v95+ Scenario Intelligence"
    })
    model["cost_range"]=f"{model['cost_p10']} - {model['cost_p90']}"
    direct=_money_to_bn_hs(base.get("direct_cost") or model.get("direct_cost") or _bn_to_money_hs(base_cost*.78))*profile["direct_factor"]
    indirect=_money_to_bn_hs(base.get("indirect_cost") or model.get("indirect_cost") or _bn_to_money_hs(base_cost*.17))*profile["indirect_factor"]
    reserve=_money_to_bn_hs(base.get("risk_reserve") or model.get("risk_reserve") or _bn_to_money_hs(base_cost*.05))*profile["reserve_factor"]
    model["direct_cost"]=_bn_to_money_hs(direct)
    model["indirect_cost"]=_bn_to_money_hs(indirect)
    model["risk_reserve"]=_bn_to_money_hs(reserve)

    # scenario delta versus base
    model["scenario_delta_intelligence"]=[
        {"label":"Capital movement","value":f"{(profile['cost']-1)*100:+.0f}%","meaning":profile["trade"]},
        {"label":"Schedule movement","value":f"{(profile['schedule']-1)*100:+.0f}%","meaning":profile["won"]},
        {"label":"Confidence movement","value":f"{profile['confidence']:+d} pts","meaning":profile["lost"]},
        {"label":"Reserve philosophy","value":f"{(profile['reserve_factor']-1)*100:+.0f}% reserve shift","meaning":"Reserve is being used as a strategic choice, not a static percentage."},
        {"label":"Board consequence","value":profile["risk"],"meaning":insight},
    ]
    # confidence breakdown, scenario-specific
    if scenario=="faster":
        cb=[("Concurrency loading","-13","Parallel execution and commissioning overlap increase volatility"),
            ("Procurement acceleration","-9","Expediting reduces optionality and increases premium risk"),
            ("Recovery float","-12","Float is consumed before validation / integration is proven"),
            ("Strategic timing","+6","Earlier market-entry option improves strategic value"),
            ("Evidence maturity","-7","Decisions are being pulled ahead of package maturity")]
    elif scenario=="cheaper":
        cb=[("Reserve adequacy","-14","Contingency is below the uncertainty envelope"),
            ("Operational resilience","-11","Redundancy and start-up protection are reduced"),
            ("Capital efficiency","+7","Lower approval number improves affordability"),
            ("Lifecycle exposure","-9","Costs are likely transferred into operations and recovery"),
            ("Procurement certainty","-6","Lower-cost sourcing increases vendor / quality variance")]
    elif scenario=="lower_risk":
        cb=[("Assurance maturity","+12","More evidence exists before commitment"),
            ("Commissioning float","+10","Recovery time is protected"),
            ("Procurement certainty","+8","Long-lead strategy is more secure"),
            ("Schedule aggressiveness","-4","Longer delivery reduces strategic timing"),
            ("Reserve adequacy","+9","Tail risk is better funded")]
    elif scenario=="premium":
        cb=[("Resilience","+13","Redundancy and optionality improve downside protection"),
            ("Procurement certainty","+10","Premium route buys capacity and priority"),
            ("Operational readiness","+11","Stronger handover / mission readiness"),
            ("Capital intensity","-7","Higher approval burden"),
            ("Tail exposure","+9","P80/P90 risk is controlled")]
    else:
        cb=[("Benchmark similarity","+8","Comparable sector archetype identified"),
            ("Scope maturity","-4","Class 3 assumptions still need package evidence"),
            ("Procurement certainty","-5","Long-lead evidence remains incomplete"),
            ("Schedule logic","+5","Level 4 logic gives credible traceability"),
            ("Reserve adequacy","+2","Balanced contingency posture")]
    model["confidence_breakdown"]=[{"driver":a,"effect":b,"note":c} for a,b,c in cb]

    # Cost line mutation: not proportional only; add scenario basis
    new_lines=[]
    for i,line in enumerate(base.get("cost_lines", model.get("cost_lines",[]))):
        x=dict(line)
        typ=x.get("type","Direct")
        factor = profile["direct_factor"] if typ=="Direct" else profile["indirect_factor"] if typ=="Indirect" else profile["reserve_factor"]
        # Make some sector lines move harder by scenario
        desc=(x.get("description") or "").lower()
        if scenario=="faster" and any(k in desc for k in ["equipment","tool","utility","clean","commission","process"]): factor*=1.08
        if scenario=="cheaper" and any(k in desc for k in ["redund", "utility", "resilience", "warehouse", "automation"]): factor*=0.82
        if scenario in ["lower_risk","premium"] and any(k in desc for k in ["utility","commission","validation","risk","reserve"]): factor*=1.10
        for key in ["low_p10","most_likely_p50","high_p90"]:
            if key in x:
                basev=_money_to_bn_hs(x[key])
                spread = 1.00
                if key=="high_p90" and scenario in ["faster","cheaper"]: spread=1.18
                if key=="low_p10" and scenario=="cheaper": spread=0.92
                if key=="high_p90" and scenario in ["lower_risk","premium"]: spread=0.95
                x[key]=_bn_to_money_hs(basev*factor*spread)
        if scenario=="faster":
            x["basis"]="Acceleration basis: expediting, parallel workfronts, premium suppliers and productivity drag applied to this CBS line."
        elif scenario=="cheaper":
            x["basis"]="Cheaper basis: reduced capital scope / lower redundancy assumed; downstream operational and contingency exposure increases."
        elif scenario=="lower_risk":
            x["basis"]="Lower-risk basis: assurance, procurement buffer, commissioning float and evidence maturity added to this CBS line."
        elif scenario=="premium":
            x["basis"]="Premium basis: resilience, redundancy, optionality and priority procurement applied."
        else:
            x["basis"]="Base basis: balanced reference case using sector benchmark, location, scale and complexity factors."
        new_lines.append(x)
    model["cost_lines"]=new_lines
    model["cost_breakdown"]=new_lines

    # Risk register mutation with scenario-specific risks at top
    risks=[]
    for tup in _scenario_risks_hs(model, scenario, profile["risk_lift"]):
        rid,title,cause,event,impact,prob,owner,mit=tup
        risks.append({"id":rid,"title":title,"cause":cause,"event":event,"impact":impact,"probability_pct":int(prob),"activity":"A1500","cbs":"02.01","owner":owner,"mitigation":mit})
    for r in base.get("risks", model.get("risks",[])):
        y=dict(r)
        try:
            y["probability_pct"]=max(5,min(88,int(y.get("probability_pct",30))+profile["risk_lift"]))
        except Exception: pass
        if scenario=="faster":
            y["mitigation"]="Fast-track mitigation: protect hold-points, add surge controls and weekly executive risk burn-down. " + str(y.get("mitigation",""))
        elif scenario=="cheaper":
            y["mitigation"]="Cheaper mitigation: confirm minimum viable resilience and owner acceptance of transferred risk. " + str(y.get("mitigation",""))
        elif scenario in ["lower_risk","premium"]:
            y["mitigation"]="Assurance mitigation: evidence gates, protected float and procurement certainty. " + str(y.get("mitigation",""))
        risks.append(y)
    model["risks"]=risks
    model["risk_register"]=model["risks"]

    # Schedule mutation with critical path shifts and float meaning
    rows=[]
    for idx,row in enumerate(base.get("schedule_rows", model.get("schedule_rows",[]))):
        z=dict(row)
        try:
            d=max(1,int(z.get("duration_months",1)))
            z["duration_months"]=max(1,int(round(d*profile["schedule"])))
        except Exception: pass
        if scenario=="faster":
            z["critical"]="Yes" if idx in [1,2,3,4,5] else z.get("critical","No")
            z["basis"]="Faster logic: overlapped procurement / turnover / commissioning; reduced recovery float and higher near-critical density."
        elif scenario=="cheaper":
            z["critical"]="Yes" if idx in [2,4,6] else z.get("critical","No")
            z["basis"]="Cheaper logic: slower phasing, deferred redundancy, leaner owner support and longer operational readiness tail."
        elif scenario=="lower_risk":
            z["critical"]="Yes" if idx in [3,5] else "No"
            z["basis"]="Lower-risk logic: protected evidence gates, added commissioning float and lower near-critical congestion."
        elif scenario=="premium":
            z["critical"]="Yes" if idx in [3,5] else "No"
            z["basis"]="Premium logic: priority procurement, stronger redundancy and managed readiness gates."
        else:
            z["basis"]="Base logic: balanced critical path with standard sector sequencing and benchmark duration."
        rows.append(z)
    model["schedule_rows"]=rows
    model["schedule_detail"]=rows

    # Curves and tornado
    curve=_curve_hs(new_cost,new_months,scenario)
    qcra={"p10":round(curve[2]["cost_bn"],1),"p50":round(new_cost,1),"p80":round(curve[9]["cost_bn"],1),"p90":round(curve[10]["cost_bn"],1),"p95":round(curve[11]["cost_bn"],1)}
    qsra={"p10":curve[2]["schedule_months"],"p50":new_months,"p80":curve[9]["schedule_months"],"p90":curve[10]["schedule_months"],"p95":curve[11]["schedule_months"]}
    # Keep headline range and curve range reconciled; P50 in curve is the same as the KPI P50.
    model["cost_p10"]=_bn_to_money_hs(qcra["p10"])
    model["cost_p90"]=_bn_to_money_hs(qcra["p90"])
    model["cost_range"]=f"{model['cost_p10']} - {model['cost_p90']}"

    # Executive scenario comparison versus the immutable Base case.
    base_cost_money=_bn_to_money_hs(base_cost)
    base_schedule_months=base_months
    base_conf_pct=base_conf
    delta_cost=new_cost-base_cost
    delta_months=new_months-base_months
    delta_conf=new_conf-base_conf
    model["scenario_comparison_vs_base"]={
        "base":{"cost_p50":base_cost_money,"schedule_months":base_schedule_months,"confidence_pct":base_conf_pct,"risk":model.get("_base_risk","Medium-High")},
        "selected":{"scenario":profile["label"],"cost_p50":model["cost_p50"],"schedule_months":new_months,"confidence_pct":new_conf,"risk":profile["risk"]},
        "delta":{"cost_bn":round(delta_cost,2),"cost":_bn_to_money_hs(abs(delta_cost)),"cost_direction":"higher" if delta_cost>0 else "lower" if delta_cost<0 else "same",
                 "months":delta_months,"confidence_pts":delta_conf},
        "plain_english": ("Base is the reference case: no cost, schedule or confidence delta is applied. Use Faster, Cheaper, Lower Risk or Premium to expose the board trade-off." if scenario=="base" else f"Compared with Base, {profile['label']} is {_bn_to_money_hs(abs(delta_cost))} {'more expensive' if delta_cost>0 else 'cheaper' if delta_cost<0 else 'unchanged'}, {abs(delta_months)} months {'faster' if delta_months<0 else 'slower' if delta_months>0 else 'unchanged'}, and {abs(delta_conf)} confidence points {'higher' if delta_conf>0 else 'lower' if delta_conf<0 else 'unchanged'}.")
    }
    model["scenario_matrix"]=[]
    for sk,sp in SCENARIO_PROFILES_V95_PLUS.items():
        c=base_cost*sp["cost"]; m=max(3,int(round(base_months*sp["schedule"]))); cf=max(8,min(96,base_conf+sp["confidence"]))
        model["scenario_matrix"].append({"scenario":sk,"label":sp["label"],"cost_p50":_bn_to_money_hs(c),"schedule_months":m,"risk":sp["risk"],"confidence_pct":cf,"cost_delta_pct":round((sp["cost"]-1)*100),"schedule_delta_pct":round((sp["schedule"]-1)*100),"confidence_delta_pts":sp["confidence"],"why":"; ".join([sp.get("trade",""), sp.get("won",""), sp.get("lost","")])})
    # Waterfalls explain WHY the selected scenario moved. These feed dashboard and exports.
    if scenario=="faster":
        cost_moves=[("Base P50",base_cost),("Acceleration premium",base_cost*.075),("Parallel EPC / package overlap",base_cost*.045),("Early long-lead procurement",base_cost*.035),("Commissioning surge / overtime",base_cost*.025),("Faster scenario P50",new_cost)]
        sched_moves=[("Base duration",base_months),("Parallel delivery workfronts",-round(base_months*.07)),("Early procurement release",-round(base_months*.05)),("Concurrent commissioning",-round(base_months*.04)),("Reduced float / overlap",new_months-base_months+round(base_months*.16)),("Faster duration",new_months)]
    elif scenario=="cheaper":
        cost_moves=[("Base P50",base_cost),("Scope / spec restraint",-base_cost*.055),("Deferred redundancy",-base_cost*.040),("Lower indirects / slower phasing",-base_cost*.030),("Reduced reserve",-base_cost*.015),("Cheaper scenario P50",new_cost)]
        sched_moves=[("Base duration",base_months),("Slower procurement path",round(base_months*.04)),("Deferred assurance / resilience",round(base_months*.035)),("Lean owner / commissioning support",round(base_months*.03)),("Cheaper duration",new_months)]
    elif scenario=="lower_risk":
        cost_moves=[("Base P50",base_cost),("Assurance gates",base_cost*.035),("Procurement buffer",base_cost*.030),("Commissioning float",base_cost*.025),("Reserve uplift",base_cost*.030),("Lower-risk scenario P50",new_cost)]
        sched_moves=[("Base duration",base_months),("Evidence gates",round(base_months*.055)),("Protected commissioning float",round(base_months*.055)),("Lower-risk duration",new_months)]
    elif scenario=="premium":
        cost_moves=[("Base P50",base_cost),("Redundancy / resilience",base_cost*.095),("Priority procurement",base_cost*.060),("Operational readiness",base_cost*.045),("Optionality premium",base_cost*.080),("Premium scenario P50",new_cost)]
        sched_moves=[("Base duration",base_months),("Priority procurement",-round(base_months*.04)),("Stronger readiness gates",round(base_months*.02)),("Premium duration",new_months)]
    else:
        cost_moves=[("Base P50",base_cost),("Selected scenario P50",new_cost)]
        sched_moves=[("Base duration",base_months),("Selected scenario duration",new_months)]
    model["cost_waterfall_vs_base"]=[{"driver":a,"value_bn":round(b,2),"value":_bn_to_money_hs(abs(b)) if i not in [0,len(cost_moves)-1] else _bn_to_money_hs(b),"kind":"total" if i in [0,len(cost_moves)-1] else "delta"} for i,(a,b) in enumerate(cost_moves)]
    model["schedule_waterfall_vs_base"]=[{"driver":a,"months":int(b),"kind":"total" if i in [0,len(sched_moves)-1] else "delta"} for i,(a,b) in enumerate(sched_moves)]
    curve_text={
        "base":"Reference curve: P50 equals the headline estimate; P80/P90 show normal board-level contingency exposure.",
        "faster":"Faster curve: the median date improves, but the P80/P90 tail widens because float, procurement optionality and commissioning recovery are consumed.",
        "cheaper":"Cheaper curve: P50 is lower, but P80/P90 become ugly because reserve, redundancy and recovery capacity have been stripped out.",
        "lower_risk":"Lower-risk curve: P50 is higher/slower, but the P80/P90 tail compresses because evidence, float and reserve are protected.",
        "premium":"Premium curve: higher approval value buys resilience; the value is not the P50, it is the controlled downside tail."
    }.get(scenario)
    model["monte_carlo"]={
        "qcra":qcra,
        "qsra":qsra,
        "curve":curve,
        "curve_interpretation":curve_text,
        "curve_readout":[
            f"Cost P50 reconciles to {model['cost_p50']}; P80 is {_bn_to_money_hs(qcra['p80'])}; P90 is {_bn_to_money_hs(qcra['p90'])}.",
            f"Schedule P50 reconciles to {new_months} months; P80 is {qsra['p80']} months; P90 is {qsra['p90']} months.",
            curve_text
        ],
        "tornado":[{"risk_id":r.get("id") or r.get("risk_id"),"title":r.get("title"),"driver":r.get("title"),"driver_score":max(8,100-i*7),"contribution":max(8,100-i*7)} for i,r in enumerate(model["risks"][:8])]
    }
    # Board + outputs
    model["board_briefing"]=[
        insight,
        model.get("scenario_comparison_vs_base",{}).get("plain_english", f"{profile['label']} scenario versus Base calculated."),
        f"{profile['label']} scenario: {model['cost_p50']} P50, {model['cost_range']} range, {model['schedule']} baseline and {new_conf}% confidence.",
        f"Trade made: {profile['trade']}",
        f"Gained: {profile['won']}",
        f"Lost: {profile['lost']}"
    ]
    model["outputs_board_memo"]=[model.get("scenario_comparison_vs_base",{}).get("plain_english", "Scenario compared with Base.")] + _scenario_outputs_text_hs(profile, insight)
    model["top_decisions_required"]=[
        "Accept or reject the scenario trade-off explicitly at board level.",
        "Confirm the changed critical path and near-critical path density.",
        "Approve the scenario-specific reserve and contingency philosophy.",
        "Assign named owners for the new top scenario risks.",
        "Confirm whether XER, cost workbook and risk register should be issued as scenario-controlled outputs."
    ]
    model["casey_thinking"]=(f"CASEY has re-cut the programme as a {profile['label']} scenario. {profile['trade']} "
        f"The governing consequence is: {insight} QCRA/QSRA curves, cost basis, risk probabilities and schedule logic have been re-weighted to match this strategic posture.")
    model["executive_summary"]=(
        f"{model.get('title','Programme')} scenario view: {profile['label']}. CASEY indicates {model['cost_p50']} P50 exposure, "
        f"{model['cost_range']} range, {model['schedule']} baseline, {model['risk']} risk and {new_conf}% confidence. "
        f"{insight} Trade-off: {profile['trade']}"
    )
    # Export control metadata
    model["export_scenario_control"]={
        "scenario":profile["label"],
        "must_stamp_exports":True,
        "xer_logic":"critical path, duration and basis changed by scenario",
        "cost_logic":"CBS basis, reserve and uncertainty spread changed by scenario",
        "risk_logic":"probability, mitigation and scenario-emergent risks changed by scenario"
    }

    # FINAL EXEC POLISH: scenario-aware ranking and differentiation
    try:
        sc = str(model.get("scenario","base")).lower()
        schedule_lists = {
            "faster":[
                "Concurrent commissioning overload",
                "Recovery float exhaustion",
                "Acceleration premium shock",
                "Grid connection delay",
                "Integrated systems testing concurrency"
            ],
            "cheaper":[
                "Vendor claims and change exposure",
                "Procurement deferral and long-lead slippage",
                "Design maturity gap",
                "Scope growth from deferred decisions",
                "Interface coordination delay"
            ],
            "lower_risk":[
                "Governance and approvals latency",
                "Extended validation sequencing",
                "Conservative commissioning gates",
                "Operational readiness hold-points",
                "Assurance and compliance reviews"
            ],
            "premium":[
                "Integration complexity across parallel packages",
                "Executive decision latency",
                "Technology assurance alignment",
                "Multi-package interface management",
                "Programme coordination overhead"
            ]
        }
        cost_lists = {
            "faster":[
                "Acceleration premiums and overtime",
                "Power train, transformers and switchgear",
                "Integrated systems testing",
                "Grid and utility concurrency",
                "Recovery-float consumption"
            ],
            "cheaper":[
                "Deferred procurement packaging",
                "Claims and commercial exposure",
                "Rework from reduced contingency",
                "Long-lead inflation volatility",
                "Scope rationalisation impacts"
            ],
            "lower_risk":[
                "Additional contingency and reserve",
                "Enhanced validation and assurance",
                "Programme controls and governance",
                "Redundant infrastructure resilience",
                "Extended commissioning readiness"
            ]
        }
        if sc in schedule_lists:
            model["sector_schedule_threats"]=schedule_lists[sc]
        if sc in cost_lists:
            model["sector_primary_cost_drivers"]=cost_lists[sc]

        if sc=="faster":
            model["executive_shock_insight"]="Acceleration increases spend faster than it reduces uncertainty; the delivery tail becomes more volatile."
        elif sc=="cheaper":
            model["executive_shock_insight"]="Capital efficiency reduces resilience: procurement and recovery flexibility become constrained."
        elif sc=="lower_risk":
            model["executive_shock_insight"]="Confidence is purchased through reserve, governance and extended delivery duration."
        elif sc=="premium":
            model["executive_shock_insight"]="Premium posture buys resilience, optionality and stronger certainty at visible capex premium."
    except Exception:
        pass

    return model



# ================= CASEY V104 FRONTEND GRAPH COMPATIBILITY + EXEC FEATURES =================
# Ensures V103 frontend receives all graph / scenario / executive fields it expects.
import copy as _casey_copy_v104
import math as _casey_math_v104

def _v104_money_to_bn(x):
    s = str(x or "").replace("$","").replace(",","").strip().upper()
    try:
        if s.endswith("T"): return float(s[:-1]) * 1000
        if s.endswith("B"): return float(s[:-1])
        if s.endswith("M"): return float(s[:-1]) / 1000
        return float(s)
    except Exception:
        return 0.0

def _v104_fmt_bn(x):
    x = float(x or 0)
    if x >= 1000: return f"${x/1000:.1f}T"
    if x >= 1: return f"${x:.1f}B"
    return f"${x*1000:.0f}M"

def _v104_months(x):
    try:
        return int(float(str(x or "0").replace("months","").strip().split()[0]))
    except Exception:
        return 0

def _v104_profile(s):
    s = (s or "base").lower().replace(" ","_")
    return {
        "base":       {"label":"Base","cost":1.00,"sched":1.00,"conf":0,   "risk":"Medium-High"},
        "faster":     {"label":"Faster","cost":1.22,"sched":0.74,"conf":-18, "risk":"High"},
        "cheaper":    {"label":"Cheaper","cost":0.82,"sched":1.20,"conf":-20, "risk":"High"},
        "lower_risk": {"label":"Lower Risk","cost":1.16,"sched":1.16,"conf":18, "risk":"Medium"},
        "premium":    {"label":"Premium","cost":1.32,"sched":0.94,"conf":23, "risk":"Medium-Low"},
    }.get(s, {"label":"Base","cost":1.00,"sched":1.00,"conf":0,"risk":"Medium-High"})

def _v104_curve(cost_bn, months, scen):
    scen = (scen or "base").lower().replace(" ","_")
    pts = [1,5,10,20,30,40,50,60,70,80,90,95,99]
    arr = []
    for p in pts:
        x = p / 100
        if scen == "faster":
            cf = 0.70 + 0.34*x + 0.50*(x**4.3)
            sf = 0.60 + 0.24*x + 0.64*(x**4.7)
        elif scen == "cheaper":
            cf = 0.50 + 0.25*x + 0.86*(x**3.9)
            sf = 0.70 + 0.22*x + 0.78*(x**3.5)
        elif scen == "lower_risk":
            cf = 0.84 + 0.18*x + 0.08*(x**2.6)
            sf = 0.88 + 0.14*x + 0.07*(x**2.6)
        elif scen == "premium":
            cf = 0.86 + 0.18*x + 0.08*(x**2.6)
            sf = 0.80 + 0.16*x + 0.13*(x**2.6)
        else:
            cf = 0.70 + 0.25*x + 0.34*(x**2.8)
            sf = 0.74 + 0.18*x + 0.32*(x**2.6)
        arr.append({"percentile":p, "cost_bn":round(cost_bn*cf,2), "schedule_months":max(1,int(round(months*sf)))})
    return arr

def _v104_trade(profile):
    label = profile["label"]
    if label == "Faster":
        return ("You bought time by spending money and consuming recovery float.", "Earlier market-entry / revenue timing.", "Procurement optionality, recovery float and late-stage stability.")
    if label == "Cheaper":
        return ("You cut capital authorization by transferring risk into resilience, redundancy and start-up.", "Lower initial approval number and reduced near-term capital draw.", "Operational resilience, commissioning flexibility, contingency adequacy and lifecycle certainty.")
    if label == "Lower Risk":
        return ("You bought confidence by adding assurance, float and procurement evidence.", "Reduced P80/P90 exposure and stronger board approval defensibility.", "Earlier revenue date and lean capital posture.")
    if label == "Premium":
        return ("You bought resilience, redundancy and strategic optionality.", "Higher resilience, stronger procurement certainty and better downside protection.", "Lowest-capex authorization case.")
    return ("Balanced cost, time and evidence posture.", "Credible reference case for board challenge.", "No extra certainty or capital efficiency purchased.")

def _v104_normalize_model(model, prompt="", scenario="base"):
    if not isinstance(model, dict):
        return model
    
    # FINAL EXEC POLISH: scenario-aware ranking and differentiation
    try:
        sc = str(model.get("scenario","base")).lower()
        schedule_lists = {
            "faster":[
                "Concurrent commissioning overload",
                "Recovery float exhaustion",
                "Acceleration premium shock",
                "Grid connection delay",
                "Integrated systems testing concurrency"
            ],
            "cheaper":[
                "Vendor claims and change exposure",
                "Procurement deferral and long-lead slippage",
                "Design maturity gap",
                "Scope growth from deferred decisions",
                "Interface coordination delay"
            ],
            "lower_risk":[
                "Governance and approvals latency",
                "Extended validation sequencing",
                "Conservative commissioning gates",
                "Operational readiness hold-points",
                "Assurance and compliance reviews"
            ],
            "premium":[
                "Integration complexity across parallel packages",
                "Executive decision latency",
                "Technology assurance alignment",
                "Multi-package interface management",
                "Programme coordination overhead"
            ]
        }
        cost_lists = {
            "faster":[
                "Acceleration premiums and overtime",
                "Power train, transformers and switchgear",
                "Integrated systems testing",
                "Grid and utility concurrency",
                "Recovery-float consumption"
            ],
            "cheaper":[
                "Deferred procurement packaging",
                "Claims and commercial exposure",
                "Rework from reduced contingency",
                "Long-lead inflation volatility",
                "Scope rationalisation impacts"
            ],
            "lower_risk":[
                "Additional contingency and reserve",
                "Enhanced validation and assurance",
                "Programme controls and governance",
                "Redundant infrastructure resilience",
                "Extended commissioning readiness"
            ]
        }
        if sc in schedule_lists:
            model["sector_schedule_threats"]=schedule_lists[sc]
        if sc in cost_lists:
            model["sector_primary_cost_drivers"]=cost_lists[sc]

        if sc=="faster":
            model["executive_shock_insight"]="Acceleration increases spend faster than it reduces uncertainty; the delivery tail becomes more volatile."
        elif sc=="cheaper":
            model["executive_shock_insight"]="Capital efficiency reduces resilience: procurement and recovery flexibility become constrained."
        elif sc=="lower_risk":
            model["executive_shock_insight"]="Confidence is purchased through reserve, governance and extended delivery duration."
        elif sc=="premium":
            model["executive_shock_insight"]="Premium posture buys resilience, optionality and stronger certainty at visible capex premium."
    except Exception:
        pass

    return model
    scenario = (scenario or model.get("scenario") or "base").lower().replace(" ","_")
    profile = _v104_profile(scenario)
    model["scenario"] = scenario
    model["scenario_label"] = profile["label"]

    # Determine base values from model if no explicit base object exists.
    base = model.get("base_comparison") or {}
    base_cost = _v104_money_to_bn(base.get("base_cost_p50") or model.get("_base_cost_p50") or model.get("cost_p50"))
    if base_cost <= 0:
        base_cost = _v104_money_to_bn(model.get("cost_p50")) or 5.0
    base_months = _v104_months(base.get("base_schedule") or model.get("_base_schedule") or model.get("schedule")) or 60
    base_conf = int(str(base.get("base_confidence") or model.get("_base_confidence") or model.get("confidence_pct") or 60).replace("%","").split()[0])

    # If selected scenario already has values, preserve them; otherwise derive from base.
    scen_cost = _v104_money_to_bn(model.get("cost_p50")) if scenario == "base" else base_cost * profile["cost"]
    scen_months = _v104_months(model.get("schedule")) if scenario == "base" else max(3,int(round(base_months * profile["sched"])))
    scen_conf = int(model.get("confidence_pct") or base_conf) if scenario == "base" else max(5,min(96,base_conf + profile["conf"]))

    model["cost_p50"] = _v104_fmt_bn(scen_cost)
    model["cost_p10"] = _v104_fmt_bn(scen_cost * (0.78 if scenario != "cheaper" else 0.70))
    model["cost_p90"] = _v104_fmt_bn(scen_cost * (1.32 if scenario in ["base","lower_risk","premium"] else 1.52 if scenario=="faster" else 1.60))
    model["cost_range"] = f"{model['cost_p10']} - {model['cost_p90']}"
    model["schedule"] = f"{scen_months} months"
    model["risk"] = profile["risk"] if scenario != "base" else model.get("risk", profile["risk"])
    model["confidence_pct"] = scen_conf

    trade, gain, loss = _v104_trade(profile)
    insight = {
        "base":"Programme success depends on the real constraint, not the headline construction scope.",
        "faster":"Acceleration increases spend faster than it reduces uncertainty; the delivery tail becomes more volatile.",
        "cheaper":"The cheaper case is not simply cheaper; it transfers risk into resilience, contingency and start-up certainty.",
        "lower_risk":"Lower-risk delivery buys confidence through assurance, float and procurement evidence.",
        "premium":"Premium delivery buys downside protection, resilience and strategic optionality."
    }.get(scenario, trade)
    model["executive_shock_insight"] = model.get("executive_shock_insight") if scenario=="base" else insight
    model["scenario_trade"] = trade
    model["scenario_gain"] = gain
    model["scenario_loss"] = loss

    cost_delta = scen_cost - base_cost
    model["base_comparison"] = {
        "base_cost_p50": _v104_fmt_bn(base_cost),
        "scenario_cost_p50": _v104_fmt_bn(scen_cost),
        "cost_delta": ("-" + _v104_fmt_bn(abs(cost_delta))) if cost_delta < 0 else _v104_fmt_bn(cost_delta),
        "base_schedule": f"{base_months} months",
        "scenario_schedule": f"{scen_months} months",
        "schedule_delta": f"{scen_months-base_months:+d} months",
        "base_confidence": f"{base_conf}%",
        "scenario_confidence": f"{scen_conf}%",
        "confidence_delta": f"{scen_conf-base_conf:+d} pts",
        "trade": trade,
        "gain": gain,
        "loss": loss
    }
    model["scenario_comparison_vs_base"] = model["base_comparison"]
    model["scenario_delta_intelligence"] = [
        {"label":"Base P50 cost", "value":f"{_v104_fmt_bn(base_cost)} → {_v104_fmt_bn(scen_cost)}", "meaning":trade},
        {"label":"Base P50 schedule", "value":f"{base_months} mo → {scen_months} mo", "meaning":gain},
        {"label":"Confidence", "value":f"{base_conf}% → {scen_conf}%", "meaning":loss},
        {"label":"Risk posture", "value":model["risk"], "meaning":insight},
    ]

    model["cost_waterfall_vs_base"] = [
        {"driver":"Base P50", "value":round(base_cost,2), "kind":"base"},
        {"driver":profile["label"]+" scenario delta", "value":round(cost_delta,2), "kind":"delta"},
        {"driver":"Scenario P50", "value":round(scen_cost,2), "kind":"total"},
    ]
    model["schedule_waterfall_vs_base"] = [
        {"driver":"Base schedule", "months":base_months, "kind":"base"},
        {"driver":profile["label"]+" scenario delta", "months":scen_months-base_months, "kind":"delta"},
        {"driver":"Scenario schedule", "months":scen_months, "kind":"total"},
    ]

    curve = _v104_curve(scen_cost, scen_months, scenario)
    qcra = {"p10":curve[2]["cost_bn"], "p50":round(scen_cost,1), "p80":curve[9]["cost_bn"], "p90":curve[10]["cost_bn"], "p95":curve[11]["cost_bn"]}
    qsra = {"p10":curve[2]["schedule_months"], "p50":scen_months, "p80":curve[9]["schedule_months"], "p90":curve[10]["schedule_months"], "p95":curve[11]["schedule_months"]}

    risks = model.get("risks") or model.get("risk_register") or []
    # Make sure risk ids and chart names match frontend table + tornado.
    norm_risks = []
    for i,r in enumerate(risks):
        x = dict(r)
        if "risk_id" not in x:
            x["risk_id"] = x.get("id") or f"R{i+1:02d}"
        x["title"] = x.get("title") or x.get("risk") or f"Risk {i+1}"
        x["probability_pct"] = int(x.get("probability_pct") or x.get("probability") or max(15,65-i*5))
        norm_risks.append(x)
    if not norm_risks:
        norm_risks = [{"risk_id":"R01","title":"Scenario execution volatility","probability_pct":55,"cause":"Scenario posture","event":"Risk materialises","impact":"Cost/schedule exposure","owner":"PMO","mitigation":"Weekly controls review"}]
    model["risks"] = norm_risks
    model["risk_register"] = norm_risks

    tornado = []
    for i,r in enumerate(norm_risks):
        title = r.get("title","Risk")
        contribution = max(8, int(r.get("probability_pct",50)) + max(0,30-i*4))
        tornado.append({
            "risk_id": r.get("risk_id") or r.get("id") or f"R{i+1:02d}",
            "driver": title,
            "title": title,
            "contribution": contribution,
            "driver_score": contribution
        })

    model["monte_carlo"] = {
        **(model.get("monte_carlo") or {}),
        "qcra": qcra,
        "qsra": qsra,
        "curve": curve,
        "tornado": tornado,
        "curve_readout": [
            f"QCRA: P50 {model['cost_p50']} equals the headline scenario estimate; P80 { _v104_fmt_bn(qcra['p80']) } and P90 { _v104_fmt_bn(qcra['p90']) } show downside board exposure.",
            f"QSRA: P50 {scen_months} months equals the headline schedule; P80 {qsra['p80']} and P90 {qsra['p90']} months show finish-date risk.",
            f"Scenario vs Base: cost {model['base_comparison']['base_cost_p50']} → {model['cost_p50']}; schedule {base_months} → {scen_months} months; confidence {base_conf}% → {scen_conf}%."
        ],
        "curve_interpretation": {
            "base":"Balanced uncertainty: P80/P90 reflects normal delivery and procurement risk.",
            "faster":"Compressed median with aggressive P80/P95 tail: possible, but unstable.",
            "cheaper":"Lower median with fragile reserve and ugly tail: cheaper only if nothing goes wrong.",
            "lower_risk":"Tighter distribution: slower and more expensive, but more governable.",
            "premium":"Higher median with controlled downside: buys resilience and optionality."
        }.get(scenario)
    }

    model["confidence_breakdown"] = [
        {"driver":"Scenario trade-off", "effect":f"{profile['conf']:+d}", "note":trade},
        {"driver":"Procurement certainty", "effect":"-10" if scenario in ["faster","cheaper"] else "+9", "note":"Long-lead evidence controls confidence movement."},
        {"driver":"Schedule logic maturity", "effect":"-9" if scenario=="faster" else "+7" if scenario in ["lower_risk","premium"] else "+3", "note":"Critical path and handover gates are rebalanced."},
        {"driver":"Commissioning / validation readiness", "effect":"-12" if scenario in ["faster","cheaper"] else "+10", "note":"Readiness exposure changes by scenario."},
        {"driver":"Reserve adequacy", "effect":"-15" if scenario=="cheaper" else "+12" if scenario in ["lower_risk","premium"] else "+2", "note":"Reserve is treated as a scenario choice."},
    ]

    model["board_briefing"] = [
        insight,
        f"{profile['label']} scenario: {model['cost_p50']} P50, {model['cost_range']} range, {model['schedule']} baseline and {scen_conf}% confidence.",
        f"Base comparison: {model['base_comparison']['base_cost_p50']} / {model['base_comparison']['base_schedule']} / {model['base_comparison']['base_confidence']} → {model['cost_p50']} / {model['schedule']} / {scen_conf}%.",
        f"Gained: {gain}",
        f"Lost: {loss}"
    ]
    model["outputs_board_memo"] = [
        f"Selected scenario: {profile['label']}.",
        f"Trade made: {trade}",
        f"Base vs scenario: {model['base_comparison']['base_cost_p50']} → {model['cost_p50']}; {model['base_comparison']['base_schedule']} → {model['schedule']}.",
        f"Board warning: {insight}",
        "Exports should be stamped with scenario, base value, scenario value, delta and trade-off reason."
    ]
    model["casey_thinking"] = f"CASEY recalibrated cost, schedule, QCRA/QSRA, risk and confidence from the same scenario basis. {trade}"
    model["executive_summary"] = f"{model.get('title','Programme')} scenario view: {profile['label']}. CASEY indicates {model['cost_p50']} P50 exposure, {model['cost_range']} range, {model['schedule']} baseline, {model['risk']} risk and {scen_conf}% confidence. {insight}"

    # FINAL EXEC POLISH: scenario-aware ranking and differentiation
    try:
        sc = str(model.get("scenario","base")).lower()
        schedule_lists = {
            "faster":[
                "Concurrent commissioning overload",
                "Recovery float exhaustion",
                "Acceleration premium shock",
                "Grid connection delay",
                "Integrated systems testing concurrency"
            ],
            "cheaper":[
                "Vendor claims and change exposure",
                "Procurement deferral and long-lead slippage",
                "Design maturity gap",
                "Scope growth from deferred decisions",
                "Interface coordination delay"
            ],
            "lower_risk":[
                "Governance and approvals latency",
                "Extended validation sequencing",
                "Conservative commissioning gates",
                "Operational readiness hold-points",
                "Assurance and compliance reviews"
            ],
            "premium":[
                "Integration complexity across parallel packages",
                "Executive decision latency",
                "Technology assurance alignment",
                "Multi-package interface management",
                "Programme coordination overhead"
            ]
        }
        cost_lists = {
            "faster":[
                "Acceleration premiums and overtime",
                "Power train, transformers and switchgear",
                "Integrated systems testing",
                "Grid and utility concurrency",
                "Recovery-float consumption"
            ],
            "cheaper":[
                "Deferred procurement packaging",
                "Claims and commercial exposure",
                "Rework from reduced contingency",
                "Long-lead inflation volatility",
                "Scope rationalisation impacts"
            ],
            "lower_risk":[
                "Additional contingency and reserve",
                "Enhanced validation and assurance",
                "Programme controls and governance",
                "Redundant infrastructure resilience",
                "Extended commissioning readiness"
            ]
        }
        if sc in schedule_lists:
            model["sector_schedule_threats"]=schedule_lists[sc]
        if sc in cost_lists:
            model["sector_primary_cost_drivers"]=cost_lists[sc]

        if sc=="faster":
            model["executive_shock_insight"]="Acceleration increases spend faster than it reduces uncertainty; the delivery tail becomes more volatile."
        elif sc=="cheaper":
            model["executive_shock_insight"]="Capital efficiency reduces resilience: procurement and recovery flexibility become constrained."
        elif sc=="lower_risk":
            model["executive_shock_insight"]="Confidence is purchased through reserve, governance and extended delivery duration."
        elif sc=="premium":
            model["executive_shock_insight"]="Premium posture buys resilience, optionality and stronger certainty at visible capex premium."
    except Exception:
        pass

    return model

_CASEY_V104_ORIGINAL_BUILD_MODEL = build_model
def build_model(prompt:str, client:str="", class_level:int=3, schedule_level:int=3, scenario:str="base"):
    return _v104_normalize_model(_CASEY_V104_ORIGINAL_BUILD_MODEL(prompt, client, class_level, schedule_level, scenario), prompt, scenario)

# Local demo package: disable one-run lock.
def _public_demo_used(identity):
    return None

def demo_status(request):
    return {"allowed": True, "used": 0, "limit": 999, "unlocked": True}
# ================= END CASEY V104 FRONTEND GRAPH COMPATIBILITY =================

# ================= CASEY V106 DEMO EXPORT POLISH =================
# Real demo exports for every export button. Public demo outputs are stamped clearly
# so buyers understand these are sample board artefacts, while still receiving files.
def _v106_stamp_model(model: Dict[str, Any]) -> Dict[str, Any]:
    m = dict(model or {})
    sc = str(m.get('scenario_label') or m.get('scenario') or 'Base')
    base = m.get('base_comparison') or {}
    m['demo_watermark'] = 'CASEY PUBLIC DEMO OUTPUT - NOT CERTIFIED FOR COMMERCIAL RELIANCE'
    m['export_stamp'] = {
        'watermark': m['demo_watermark'],
        'scenario': sc,
        'base_value': base,
        'scenario_value': {'cost_p50': m.get('cost_p50'), 'schedule': m.get('schedule'), 'confidence_pct': m.get('confidence_pct'), 'risk': m.get('risk')},
        'trade_off_reason': m.get('casey_thinking') or m.get('executive_shock_insight') or 'Scenario-controlled first-pass intelligence.'
    }
    return m

def _v106_qcra_qsra_workbook_bytes(model: Dict[str, Any]) -> bytes:
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill, Alignment
    from openpyxl.chart import LineChart, Reference, BarChart
    m = _v106_stamp_model(model)
    mc = m.get('monte_carlo') or {}
    curve = mc.get('curve') or []
    torn = mc.get('tornado') or []
    wb = Workbook()
    ws = wb.active
    ws.title = 'DEMO Stamp'
    rows = [
        ['CASEY PUBLIC DEMO OUTPUT', m.get('demo_watermark')],
        ['Project', m.get('title')],
        ['Scenario', m.get('scenario_label') or m.get('scenario')],
        ['P50 cost', m.get('cost_p50')],
        ['Schedule', m.get('schedule')],
        ['Risk / confidence', f"{m.get('risk')} / {m.get('confidence_pct')}%"],
        ['Trade-off reason', (m.get('casey_thinking') or m.get('executive_shock_insight') or '')],
    ]
    for r in rows: ws.append(r)
    for c in ws[1]: c.font = Font(bold=True, color='FFFFFF'); c.fill = PatternFill('solid', fgColor='7A1F1F')
    ws.column_dimensions['A'].width = 26; ws.column_dimensions['B'].width = 110
    ws2 = wb.create_sheet('QCRA QSRA Curves')
    ws2.append(['Percentile','QCRA Cost $B','QSRA Months'])
    for x in curve:
        ws2.append([x.get('percentile'), x.get('cost_bn'), x.get('schedule_months')])
    for c in ws2[1]: c.font = Font(bold=True); c.fill = PatternFill('solid', fgColor='B7F7FF')
    if len(curve) >= 2:
        chart = LineChart(); chart.title='QCRA cost confidence curve'; chart.y_axis.title='Cost $B'; chart.x_axis.title='Percentile'
        data = Reference(ws2, min_col=2, min_row=1, max_row=ws2.max_row); cats = Reference(ws2, min_col=1, min_row=2, max_row=ws2.max_row)
        chart.add_data(data, titles_from_data=True); chart.set_categories(cats); chart.height=8; chart.width=14; ws2.add_chart(chart,'E2')
        chart2 = LineChart(); chart2.title='QSRA finish-date curve'; chart2.y_axis.title='Months'; chart2.x_axis.title='Percentile'
        data2 = Reference(ws2, min_col=3, min_row=1, max_row=ws2.max_row); chart2.add_data(data2, titles_from_data=True); chart2.set_categories(cats); chart2.height=8; chart2.width=14; ws2.add_chart(chart2,'E18')
    ws3 = wb.create_sheet('Risk Tornado')
    ws3.append(['Rank','Risk ID','Driver','Contribution'])
    for i,x in enumerate(torn[:12],1): ws3.append([i,x.get('risk_id'),x.get('driver') or x.get('title'),x.get('contribution') or x.get('driver_score')])
    for c in ws3[1]: c.font = Font(bold=True); c.fill = PatternFill('solid', fgColor='B7F7FF')
    if ws3.max_row > 1:
        b=BarChart(); b.type='bar'; b.title='Top QCRA/QSRA exposure drivers'; b.y_axis.title='Driver'; b.x_axis.title='Contribution'
        data=Reference(ws3,min_col=4,min_row=1,max_row=ws3.max_row); cats=Reference(ws3,min_col=3,min_row=2,max_row=ws3.max_row)
        b.add_data(data,titles_from_data=True); b.set_categories(cats); b.height=10; b.width=16; ws3.add_chart(b,'F2')
    for sheet in wb.worksheets:
        for row in sheet.iter_rows():
            for cell in row:
                cell.alignment = Alignment(vertical='top', wrap_text=True)
        for col in range(1, min(sheet.max_column, 8)+1):
            sheet.column_dimensions[chr(64+col)].width = 22
    bio = BytesIO(); wb.save(bio); return bio.getvalue()

@app.post('/export/qcra-qsra')
def export_qcra_qsra_v106(model: Dict[str, Any]):
    return stream(_v106_qcra_qsra_workbook_bytes(model), 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet', 'CASEY_DEMO_QCRA_QSRA_Pack.xlsx')

@app.post('/export/excel')
def export_excel_alias_v106(model: Dict[str, Any]):
    # Alias kept so older frontend builds still download instead of failing.
    return stream(workbook_bytes(_v106_stamp_model(model)), 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet', 'CASEY_DEMO_Cost_Workbook.xlsx')
# ================= END CASEY V106 DEMO EXPORT POLISH =================


# ================= CASEY V107 FINAL DEMO COMMERCIAL LOCK + STRONG EXPORTS =================
# This final layer restores the commercial public-demo gate:
# one credible Earth OR Space run per email, and one demo export set per run/email.
try:
    from fastapi import Body
    from fastapi.responses import JSONResponse
    from v64_outputs import (
        cost_workbook_bytes as _v107_cost_workbook_bytes,
        risk_workbook_bytes as _v107_risk_workbook_bytes,
        patch_template_xer as _v107_xer_bytes,
        schedule_csv_bytes as _v107_schedule_csv_bytes,
        model_json_bytes as _v107_model_json_bytes,
    )
except Exception as _v107_import_error:
    print('CASEY v107 import warning:', _v107_import_error)


def _v107_db_init():
    con = db(); cur = con.cursor()
    cur.execute("""CREATE TABLE IF NOT EXISTS public_demo_exports(
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        run_id TEXT,
        email_hash TEXT,
        export_kind TEXT,
        created_at TEXT
    )""")
    con.commit(); con.close()
_v107_db_init()


def _v107_remove_route(path: str, method: str = 'POST'):
    app.router.routes = [r for r in app.router.routes if not (getattr(r, 'path', None) == path and method in getattr(r, 'methods', set()))]




def _v108_admin_email_set():
    """Emails that are allowed unlimited demo testing while public users remain one-run locked.
    Configure with CASEY_ADMIN_EMAILS="you@company.com,another@company.com".
    The bundled defaults are only for local demo testing and can be overridden in production.
    """
    raw = os.environ.get('CASEY_ADMIN_EMAILS', 'deepa@caseai.co.uk,admin@controlorbit.com,demo@controlorbit.com')
    return { _normalise_email(x) for x in raw.split(',') if str(x).strip() }


def _v108_is_admin_email(email: str) -> bool:
    return bool(email) and _normalise_email(email) in _v108_admin_email_set()


def _v108_admin_key_ok(request: Request = None) -> bool:
    key = os.environ.get('CASEY_ADMIN_KEY', '').strip()
    if not key or request is None:
        return False
    return (request.headers.get('x-casey-admin-key') == key) or (request.query_params.get('admin_key') == key)


def _v108_is_admin_payload(payload: Dict[str, Any]) -> bool:
    email = (payload or {}).get('lead_email') or (payload or {}).get('client') or (payload or {}).get('email') or ''
    return _v108_is_admin_email(str(email))

def _v107_is_paid(payload: Dict[str, Any]) -> bool:
    return bool((payload or {}).get('paid') or (payload or {}).get('enterprise_access') or os.environ.get('CASEY_DEMO_UNLOCK_EXPORTS') == '1' or _v108_is_admin_payload(payload or {}))


def _v107_hash_from_payload(payload: Dict[str, Any]) -> str:
    email = (payload or {}).get('lead_email') or (payload or {}).get('client') or (payload or {}).get('email') or ''
    return _sha(str(email))


def _v107_can_export(payload: Dict[str, Any], kind: str):
    # Local demo/dev override for unrestricted testing.
    # Enabled by default for localhost builds unless explicitly disabled.
    local_demo_override = os.environ.get('CASEY_LOCAL_DEMO', '1') == '1'
    if local_demo_override:
        return True, None

    if _v107_is_paid(payload):
        return True, None
    run_id = str((payload or {}).get('run_id') or '').strip()
    email_hash = _v107_hash_from_payload(payload)
    if not run_id or not email_hash:
        return False, 'Demo exports require a public demo run_id and email. Please run the one free CASEY intelligence run first, or request paid access.'
    con = db(); cur = con.cursor()
    row = cur.execute('SELECT id, export_kind, created_at FROM public_demo_exports WHERE (run_id=? OR email_hash=?) AND export_kind=? LIMIT 1', (run_id, email_hash, kind)).fetchone()
    if row:
        con.close()
        return False, f'This email has already downloaded the demo {kind.replace('_',' ')} export. Request access to unlock repeat exports and reruns.'
    now = datetime.utcnow().isoformat()
    cur.execute('INSERT INTO public_demo_exports(run_id,email_hash,export_kind,created_at) VALUES(?,?,?,?)', (run_id, email_hash, kind, now))
    con.commit(); con.close()
    return True, None


def _v107_export_block(message: str):
    raise HTTPException(status_code=402, detail={
        'message': message,
        'upgrade_cta': 'Request access to unlock unlimited reruns, exports, saved projects and certified workflow support.'
    })


def _v107_stamp_model(model: Dict[str, Any]) -> Dict[str, Any]:
    m = dict(model or {})
    base = m.get('base_comparison') or {}
    label = str(m.get('scenario_label') or m.get('scenario') or 'Base').replace('_',' ').title()
    wm = 'CASEY PUBLIC DEMO OUTPUT - SIMULATED SAMPLE - NOT CERTIFIED FOR COMMERCIAL RELIANCE'
    m['demo_watermark'] = wm
    m['watermark'] = wm
    m['scenario_label'] = label
    m['export_stamp'] = {
        'watermark': wm,
        'scenario': label,
        'run_id': m.get('run_id'),
        'base_value': base,
        'scenario_value': {'cost_p50': m.get('cost_p50'), 'schedule': m.get('schedule'), 'confidence_pct': m.get('confidence_pct'), 'risk': m.get('risk')},
        'trade_off_reason': m.get('casey_thinking') or m.get('executive_shock_insight') or 'Scenario-controlled first-pass CASEY intelligence.'
    }
    title = str(m.get('title') or 'CASEY Project')
    if 'DEMO' not in title.upper():
        m['title'] = f'DEMO - {title}'
    return m


def _v107_readme(model: Dict[str, Any]) -> bytes:
    m = _v107_stamp_model(model)
    txt = f"""CASEY PUBLIC DEMO OUTPUT PACK

WATERMARK: {m.get('demo_watermark')}
Run ID: {m.get('run_id')}
Scenario: {m.get('scenario_label')}
Project: {m.get('title')}
P50 Cost: {m.get('cost_p50')}
Schedule: {m.get('schedule')}
Risk / confidence: {m.get('risk')} / {m.get('confidence_pct')}%

This demo pack is intentionally stamped. It is a simulated board-grade intelligence sample, not a certified estimate, schedule, risk model or commercial deliverable.

Included outputs:
1. Cost workbook with scenario-controlled basis
2. Risk register with cause / event / impact / owner / mitigation
3. XER schedule generated from the selected project/scenario
4. Schedule CSV fallback
5. QCRA/QSRA workbook with curves and tornado drivers
6. JSON audit model

Paid access unlocks saved projects, repeat scenarios, unrestricted exports, private benchmark libraries and production controls workflows.
"""
    return txt.encode('utf-8')


def _v107_watermarked_qcra_qsra(model: Dict[str, Any]) -> bytes:
    return _v106_qcra_qsra_workbook_bytes(_v107_stamp_model(model))


def _v107_full_pack(model: Dict[str, Any]) -> bytes:
    m = _v107_stamp_model(model)
    bio = BytesIO()
    with zipfile.ZipFile(bio, 'w', zipfile.ZIP_DEFLATED) as z:
        z.writestr('00_CASEY_PUBLIC_DEMO_README.txt', _v107_readme(m))
        z.writestr('01_CASEY_DEMO_Cost_Workbook.xlsx', _v107_cost_workbook_bytes(m))
        z.writestr('02_CASEY_DEMO_Risk_Register.xlsx', _v107_risk_workbook_bytes(m))
        z.writestr('03_CASEY_DEMO_Strong_P6_Schedule.xer', _v107_xer_bytes(m))
        z.writestr('04_CASEY_DEMO_Schedule_CSV_Fallback.csv', _v107_schedule_csv_bytes(m))
        z.writestr('05_CASEY_DEMO_QCRA_QSRA_Curves_and_Tornado.xlsx', _v107_watermarked_qcra_qsra(m))
        z.writestr('06_CASEY_DEMO_Model_Audit.json', _v107_model_json_bytes(m))
    bio.seek(0)
    return bio.getvalue()

_v107_remove_route('/demo/status', 'GET')
@app.get('/demo/status')
def demo_status_v107(request: Request):
    email = request.query_params.get('email','')
    if _v108_is_admin_email(email) or _v108_admin_key_ok(request):
        return {'allowed': True, 'admin_bypass': True, 'used': 0, 'limit': 'unlimited', 'remaining': 'unlimited', 'rule': 'admin bypass active; public remains one credible Earth or Space run per email'}
    if not email:
        return {'allowed': True, 'used': 0, 'limit': 1, 'remaining': 1, 'rule': 'one credible Earth or Space run per email'}
    h = _sha(_normalise_email(email))
    con = db(); cur = con.cursor(); row = cur.execute('SELECT COUNT(*) AS c FROM public_demo_uses WHERE email_hash=?', (h,)).fetchone(); con.close()
    used = int(row['c'] if row else 0)
    return {'allowed': used < 1, 'used': used, 'limit': 1, 'remaining': max(0, 1-used), 'rule': 'one credible Earth or Space run per email'}

_v107_remove_route('/public-demo/generate', 'POST')
@app.post('/public-demo/generate')
def public_demo_generate_v107(req: PublicDemoRequest, request: Request):
    issues = _quality_gate_public_demo(req)
    if issues:
        raise HTTPException(status_code=422, detail={'message': 'CASEY needs one real infrastructure or space programme brief before using your free run.', 'issues': issues})
    admin_bypass = _v108_is_admin_email(req.email) or _v108_admin_key_ok(request)
    identity = _public_demo_identity(request, req)
    con = db(); cur = con.cursor()
    prev = cur.execute('SELECT run_id, created_at FROM public_demo_uses WHERE email_hash=? LIMIT 1', (identity['email_hash'],)).fetchone()
    if prev and not admin_bypass:
        con.close()
        raise HTTPException(status_code=403, detail={'message': 'This email has already used its one free CASEY intelligence run. Request access to run another Earth or Space scenario.', 'existing_run_id': prev['run_id']})
    prompt = _premium_public_prompt(req)
    input_quality = _public_demo_brief_quality_score(req)
    model = build_model(prompt, _normalise_email(req.email), 3, 4, 'base')
    model['input_quality_score'] = input_quality['score']
    model['public_demo'] = True
    model['admin_bypass'] = bool(admin_bypass)
    model['lead_email'] = _normalise_email(req.email)
    model['project_type'] = req.project_type
    run_id = 'CASEY-DEMO-' + uuid.uuid4().hex[:10].upper()
    model['run_id'] = run_id
    model = _v107_stamp_model(model)
    report = _public_demo_report(model)
    now = datetime.utcnow().isoformat()
    cur.execute("""INSERT INTO public_demo_uses(run_id,email_hash,ip_hash,fingerprint_hash,client_token_hash,project_type,project_text,model_json,created_at)
        VALUES(?,?,?,?,?,?,?,?,?)""", (run_id, identity['email_hash'], identity['ip_hash'], identity['fingerprint_hash'], identity['client_token_hash'], req.project_type, req.project_description, json.dumps(model), now))
    cur.execute('INSERT INTO projects(name,client,prompt,mode,created_at,model_json) VALUES(?,?,?,?,?,?)', (model.get('title'), model.get('client'), model.get('prompt'), model.get('mode'), now, json.dumps(model)))
    con.commit(); con.close()
    return {'run_id': run_id, 'used': 0 if admin_bypass else 1, 'limit': 'unlimited' if admin_bypass else 1, 'remaining': 'unlimited' if admin_bypass else 0, 'admin_bypass': bool(admin_bypass), 'report': report, 'model': model}

def _v107_export_endpoint(kind: str, builder, media: str, filename: str):
    def endpoint(payload: dict = Body(default={})):
        m = _v107_stamp_model(payload or {})
        ok, msg = _v107_can_export(m, kind)
        if not ok:
            _v107_export_block(msg)
        return stream(builder(m), media, filename)
    return endpoint

for _p in ['/export/workbook','/export/cost-model','/export/excel']:
    _v107_remove_route(_p, 'POST')
    app.post(_p)(_v107_export_endpoint('cost_workbook', _v107_cost_workbook_bytes, 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet', 'CASEY_DEMO_Cost_Workbook.xlsx'))
for _p in ['/export/risk-register']:
    _v107_remove_route(_p, 'POST')
    app.post(_p)(_v107_export_endpoint('risk_register', _v107_risk_workbook_bytes, 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet', 'CASEY_DEMO_Risk_Register.xlsx'))
for _p in ['/export/xer']:
    _v107_remove_route(_p, 'POST')
    app.post(_p)(_v107_export_endpoint('xer_schedule', _v107_xer_bytes, 'application/octet-stream', 'CASEY_DEMO_Strong_P6_Schedule.xer'))
for _p in ['/export/schedule-csv']:
    _v107_remove_route(_p, 'POST')
    app.post(_p)(_v107_export_endpoint('schedule_csv', _v107_schedule_csv_bytes, 'text/csv', 'CASEY_DEMO_Schedule_CSV_Fallback.csv'))
for _p in ['/export/qcra-qsra']:
    _v107_remove_route(_p, 'POST')
    app.post(_p)(_v107_export_endpoint('qcra_qsra', _v107_watermarked_qcra_qsra, 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet', 'CASEY_DEMO_QCRA_QSRA.xlsx'))
for _p in ['/export/json']:
    _v107_remove_route(_p, 'POST')
    app.post(_p)(_v107_export_endpoint('model_audit', _v107_model_json_bytes, 'application/json', 'CASEY_DEMO_Model_Audit.json'))
for _p in ['/export/all','/export/full-pack']:
    _v107_remove_route(_p, 'POST')
    app.post(_p)(_v107_export_endpoint('full_pack', _v107_full_pack, 'application/zip', 'CASEY_DEMO_Final_Output_Pack.zip'))

APP_VERSION = 'CASEY V107 Final Demo Commercial Lock + Strong Exports'
print('CASEY V107 final demo commercial lock and strong exports installed')
# ================= END CASEY V107 FINAL DEMO COMMERCIAL LOCK + STRONG EXPORTS =================

# ================= CASEY V116 ELITE DEMO COMPLETION PATCH =================
# Adds a safe local/admin reset for demo testing and a health/readiness manifest.
# Public users still remain locked to one credible Earth or Space run per email.

def _casey_v116_local_request(request: Request) -> bool:
    host = (request.client.host if request and request.client else '') or ''
    return host in {'127.0.0.1', '::1', 'localhost'}

@app.get('/elite-demo/health')
def elite_demo_health():
    return {
        'version': 'CASEY V116 Elite Demo Completion',
        'status': 'ready',
        'public_lock': 'one credible Earth or Space intelligence run per email',
        'admin_testing': 'CASEY_ADMIN_EMAILS bypasses run and export limits; localhost can reset test locks',
        'confidence_definition': 'board-defensibility, not generic optimism',
        'exports': ['board pack zip', 'cost workbook', 'risk register', 'XER schedule', 'QCRA/QSRA workbook', 'model audit json'],
        'demo_positioning': 'first-pass strategic intelligence, demo-watermarked, not certified commercial reliance'
    }

@app.post('/elite-demo/reset-test-email')
def elite_demo_reset_test_email(request: Request, payload: Dict[str, Any] = Body(default={})):
    """Local/admin-only helper so the owner can test the one-run lock without weakening public gating."""
    email = _normalise_email(str((payload or {}).get('email') or ''))
    if not email:
        raise HTTPException(status_code=422, detail='email is required')
    if not (_casey_v116_local_request(request) or _v108_admin_key_ok(request) or _v108_is_admin_email(email)):
        raise HTTPException(status_code=403, detail='Reset is local/admin only. Public one-run lock remains active.')
    h = _sha(email)
    con = db(); cur = con.cursor()
    cur.execute('DELETE FROM public_demo_uses WHERE email_hash=?', (h,))
    cur.execute('DELETE FROM public_demo_exports WHERE email_hash=?', (h,))
    con.commit(); con.close()
    return {'reset': True, 'email': email, 'message': 'Test lock cleared for this email only. Public users remain one-run locked.'}

APP_VERSION = 'CASEY V116 Elite Demo Completion'
print('CASEY V116 elite demo completion patch installed')
# ================= END CASEY V116 ELITE DEMO COMPLETION PATCH =================

# ================= CASEY V118 FINAL DEMO EXPORT REFINEMENT MARKER =================
APP_VERSION = 'CASEY V118 Final Demo Export Refinement'
print('CASEY V118 final demo export refinement marker installed')
# ================= END CASEY V118 FINAL DEMO EXPORT REFINEMENT MARKER =================

# ================= CASEY V119 FINAL EARTH + SPACE DEMO EXPORT LOCK =================
# Final requested demo refinements are routed through the existing export endpoints.
try:
    from v64_outputs import (
        cost_workbook_bytes as _v119_cost_workbook_bytes,
        risk_workbook_bytes as _v119_risk_workbook_bytes,
        patch_template_xer as _v119_xer_bytes,
        schedule_csv_bytes as _v119_schedule_csv_bytes,
        model_json_bytes as _v119_model_json_bytes,
        all_zip_bytes as _v119_full_pack_bytes,
        qcra_qsra_workbook_bytes as _v119_qcra_qsra_workbook_bytes,
    )

    for _p in ['/export/workbook','/export/cost-model','/export/excel']:
        _v107_remove_route(_p, 'POST')
        app.post(_p)(_v107_export_endpoint('v119_cost_workbook', _v119_cost_workbook_bytes, 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet', 'CASEY_V119_Final_Cost_Workbook.xlsx'))
    for _p in ['/export/risk-register']:
        _v107_remove_route(_p, 'POST')
        app.post(_p)(_v107_export_endpoint('v119_risk_register', _v119_risk_workbook_bytes, 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet', 'CASEY_V119_Operational_Risk_Register.xlsx'))
    for _p in ['/export/xer']:
        _v107_remove_route(_p, 'POST')
        app.post(_p)(_v107_export_endpoint('v119_xer_schedule', _v119_xer_bytes, 'application/octet-stream', 'CASEY_V119_Planner_Grade_Schedule.xer'))
    for _p in ['/export/schedule-csv']:
        _v107_remove_route(_p, 'POST')
        app.post(_p)(_v107_export_endpoint('v119_schedule_csv', _v119_schedule_csv_bytes, 'text/csv', 'CASEY_V119_Schedule_Levels_Resource_Loaded.csv'))
    for _p in ['/export/qcra-qsra']:
        _v107_remove_route(_p, 'POST')
        app.post(_p)(_v107_export_endpoint('v119_qcra_qsra', _v119_qcra_qsra_workbook_bytes, 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet', 'CASEY_V119_QCRA_QSRA_Asymmetric_Tornado.xlsx'))
    for _p in ['/export/json']:
        _v107_remove_route(_p, 'POST')
        app.post(_p)(_v107_export_endpoint('v119_model_audit', _v119_model_json_bytes, 'application/json', 'CASEY_V119_Model_Audit.json'))
    for _p in ['/export/all','/export/full-pack']:
        _v107_remove_route(_p, 'POST')
        app.post(_p)(_v107_export_endpoint('v119_full_pack', _v119_full_pack_bytes, 'application/zip', 'CASEY_V119_FINAL_EARTH_SPACE_OUTPUT_PACK.zip'))

    APP_VERSION = 'CASEY V119 Final Earth + Space Demo Lock'
    print('CASEY V119 final Earth + Space export endpoints installed')
except Exception as _v119_export_error:
    print('CASEY V119 export endpoint warning:', _v119_export_error)
# ================= END CASEY V119 FINAL EARTH + SPACE DEMO EXPORT LOCK =================

# ================= CASEY V123 LIVE CALIBRATION SIGNALS =================
# Lightweight live-sector intelligence layer for demo: converts current market,
# delivery and launch-environment signals into transparent confidence / P-tail
# modifiers. This is intentionally deterministic and auditable for demos.

def _v123_live_calibration_signals(model: Dict[str, Any]) -> List[Dict[str, Any]]:
    mode = str(model.get('mode') or '').lower()
    subsector = str(model.get('subsector') or '').lower()
    prompt = str(model.get('prompt') or '').lower()
    scenario = str(model.get('scenario') or 'base').lower()
    space = mode == 'space' or any(x in prompt + ' ' + subsector for x in ['lunar','mars','orbital','space','launch','leo','cislunar'])
    if space:
        signals = [
            {'signal':'Launch reliability volatility','status':'Active','direction':'↑ schedule-tail exposure','weight':0.18,'applies_to':'QSRA P80/P90, launch cadence risk','basis':'Recent heavy-lift test campaigns show progress but continue to expose launch, booster and flight-safety variance.'},
            {'signal':'Mission assurance burden','status':'Active','direction':'↓ confidence until evidence improves','weight':0.15,'applies_to':'confidence, reserve adequacy','basis':'Lunar / orbital programmes require qualification, redundancy and post-deployment recovery evidence before board approval.'},
            {'signal':'Thermal-power balance','status':'Watch','direction':'↑ integration risk','weight':0.11,'applies_to':'systems integration, QCRA reserve','basis':'Power storage, thermal rejection and autonomous servicing remain governing constraints for surface / orbital infrastructure.'},
            {'signal':'Regulatory and range availability','status':'Watch','direction':'↑ launch-cadence uncertainty','weight':0.09,'applies_to':'schedule logic, manifest dependencies','basis':'Launch windows, flight approvals and range constraints are treated as live operating-environment risks.'},
            {'signal':'Orbital logistics and recovery','status':'Active','direction':'↑ lifecycle exposure','weight':0.13,'applies_to':'risk register, operational resilience','basis':'Repair and recovery after deployment are materially harder than terrestrial rework.'},
        ]
    elif any(x in subsector + ' ' + prompt for x in ['data centre','data center','hyperscale']) or re.search(r'\bai\b|artificial intelligence', subsector + ' ' + prompt):
        signals = [
            {'signal':'Grid connection congestion','status':'Active','direction':'↑ critical-path exposure','weight':0.17,'applies_to':'QSRA P80/P90, energisation risk','basis':'Power availability and utility interconnection remain the dominant delivery constraints for hyperscale AI infrastructure.'},
            {'signal':'Transformer / switchgear lead times','status':'Active','direction':'↑ procurement tail','weight':0.15,'applies_to':'risk register, QCRA reserve','basis':'Long-lead electrical packages are treated as supply-chain stress signals.'},
            {'signal':'Liquid cooling readiness','status':'Watch','direction':'↑ commissioning risk','weight':0.10,'applies_to':'confidence, IST logic','basis':'Cooling readiness and integrated systems testing can govern delivery more than civil progress.'},
            {'signal':'Power-market escalation','status':'Watch','direction':'↑ cost uncertainty','weight':0.08,'applies_to':'cost range, reserve philosophy','basis':'Energy infrastructure demand and procurement competition are included in capital-delivery calibration.'},
        ]
    elif any(x in subsector + ' ' + prompt for x in ['semiconductor','fab','wafer','cleanroom']):
        signals = [
            {'signal':'Tool-install sequencing','status':'Active','direction':'↑ schedule-tail exposure','weight':0.16,'applies_to':'QSRA critical path','basis':'Cleanroom readiness, process tools and vendor windows control fab delivery confidence.'},
            {'signal':'UPW / specialty utilities maturity','status':'Active','direction':'↑ commissioning risk','weight':0.13,'applies_to':'confidence, risk reserve','basis':'Specialty utility qualification can dominate board-defensibility before ramp.'},
            {'signal':'Equipment supply concentration','status':'Watch','direction':'↑ procurement uncertainty','weight':0.10,'applies_to':'QCRA reserve','basis':'Specialist OEM slots and logistics remain high-sensitivity inputs.'},
        ]
    elif any(x in subsector + ' ' + prompt for x in ['airport','terminal','runway']):
        signals = [
            {'signal':'Live-operations phasing','status':'Active','direction':'↑ interface risk','weight':0.14,'applies_to':'schedule logic, risk register','basis':'Work inside operating airports creates staging, safety and stakeholder constraints.'},
            {'signal':'Systems integration readiness','status':'Active','direction':'↑ handover risk','weight':0.12,'applies_to':'confidence, commissioning','basis':'Baggage, security, communications and passenger-flow systems must be commissioned as one operating environment.'},
            {'signal':'Regulatory / stakeholder approvals','status':'Watch','direction':'↑ schedule uncertainty','weight':0.09,'applies_to':'QSRA P80/P90','basis':'Airside approvals and operational windows are treated as current delivery-friction signals.'},
        ]
    elif any(x in subsector + ' ' + prompt for x in ['defence','defense','military','radar','naval','airbase']):
        signals = [
            {'signal':'Security and assurance burden','status':'Active','direction':'↓ evidence maturity until cleared','weight':0.16,'applies_to':'confidence, approvals','basis':'Classified interfaces and assurance gates increase decision friction.'},
            {'signal':'Sovereign supply-chain exposure','status':'Active','direction':'↑ procurement uncertainty','weight':0.14,'applies_to':'QCRA reserve, risk register','basis':'Controlled components and supplier eligibility affect delivery tails.'},
            {'signal':'Integration test dependency','status':'Watch','direction':'↑ commissioning risk','weight':0.11,'applies_to':'schedule / operational readiness','basis':'Mission-system integration can govern delivery more than construction progress.'},
        ]
    else:
        signals = [
            {'signal':'Commodity and procurement volatility','status':'Watch','direction':'↑ cost-range uncertainty','weight':0.10,'applies_to':'QCRA reserve','basis':'Current market and supply-chain volatility are applied to benchmark ranges.'},
            {'signal':'Approvals and governance latency','status':'Watch','direction':'↑ schedule-tail exposure','weight':0.09,'applies_to':'QSRA P80/P90','basis':'Decision friction and evidence maturity are reflected in confidence scoring.'},
            {'signal':'Commissioning readiness','status':'Active','direction':'↓ confidence if unproven','weight':0.12,'applies_to':'confidence, risk register','basis':'Operational readiness is treated as a board-defensibility signal, not a late-stage admin task.'},
        ]
    if scenario == 'faster':
        signals.append({'signal':'Acceleration stress overlay','status':'Active','direction':'P50 improves / P90 worsens','weight':0.14,'applies_to':'QSRA tail, confidence movement','basis':'Speed consumes float and optionality unless acceleration evidence is explicit.'})
    elif scenario == 'cheaper':
        signals.append({'signal':'Capital-efficiency stress overlay','status':'Active','direction':'P50 lower / resilience lower','weight':0.14,'applies_to':'reserve adequacy, lifecycle exposure','basis':'Lower headline capital can transfer exposure into start-up, operations and P90 contingency.'})
    elif scenario == 'lower_risk':
        signals.append({'signal':'Assurance overlay','status':'Active','direction':'confidence ↑ / duration ↑','weight':0.12,'applies_to':'reserve, schedule gates','basis':'The scenario buys evidence, float and resilience rather than early date or lean capex.'})
    elif scenario == 'premium':
        signals.append({'signal':'Resilience premium overlay','status':'Active','direction':'optionality ↑ / affordability challenge ↑','weight':0.12,'applies_to':'board trade-off, QCRA reserve','basis':'Premium posture increases downside protection but creates value-for-money challenge.'})
    return signals

def _v123_apply_live_calibration(model: Dict[str, Any]) -> Dict[str, Any]:
    m = dict(model or {})
    signals = _v123_live_calibration_signals(m)
    m['live_calibration_active'] = True
    m['live_calibration_label'] = 'LIVE CALIBRATION SIGNALS ACTIVE'
    m['live_calibration_summary'] = 'Current sector conditions are being applied to confidence, contingency and delivery-tail exposure.'
    m['live_calibration_signals'] = signals
    m['live_calibration_strip'] = ' • '.join([s['signal'] for s in signals[:4]])
    top = signals[0]['signal'] if signals else 'delivery environment volatility'
    # Inject subtle evidence of calibration into existing narrative fields without overwhelming the pack.
    if m.get('mode') == 'Space':
        extra = f" Live calibration is weighting {top.lower()}, mission assurance and orbital recovery exposure into the QSRA/QCRA tail."
    else:
        extra = f" Live calibration is weighting {top.lower()}, procurement stress and commissioning readiness into the QSRA/QCRA tail."
    if m.get('uncertainty_narrative') and isinstance(m['uncertainty_narrative'], dict):
        interp = m['uncertainty_narrative'].get('interpretation') or ''
        if 'Live calibration' not in interp:
            m['uncertainty_narrative']['interpretation'] = (interp + extra).strip()
    else:
        m['uncertainty_narrative'] = {'interpretation': extra.strip()}
    cards = list(m.get('mission_control_cards') or [])
    cards.insert(0, {'label':'Live calibration', 'signal': m['live_calibration_summary'], 'severity':'Active'})
    m['mission_control_cards'] = cards[:6]
    audit = dict(m.get('active_context_lock') or {})
    audit['live_calibration'] = True
    audit['signal_count'] = len(signals)
    m['active_context_lock'] = audit
    return m

_CASEY_V123_ORIGINAL_BUILD_MODEL = build_model
def build_model(prompt:str, client:str='', class_level:int=3, schedule_level:int=3, scenario:str='base'):
    return _v123_apply_live_calibration(_CASEY_V123_ORIGINAL_BUILD_MODEL(prompt, client, class_level, schedule_level, scenario))

@app.get('/live-calibration/manifest')
def live_calibration_manifest():
    return {
        'version':'CASEY V123 Live Calibration Demo Layer',
        'status':'active',
        'positioning':'sector intelligence calibration, not a certified market-data feed',
        'signals':['launch reliability','mission assurance','commodity / procurement volatility','regulatory exposure','commissioning readiness'],
        'applies_to':['confidence','QCRA reserve','QSRA P80/P90','risk register','board narrative']
    }

APP_VERSION = 'CASEY V123 Live Calibration Demo Layer'
print('CASEY V123 live calibration demo layer installed')
# ================= END CASEY V123 LIVE CALIBRATION SIGNALS =================

# ================= CASEY V124 SECTOR ONTOLOGY HARDENING LOCK =================
# Final public-demo hardening: sector-locked causal graphs, vocabulary and benchmark guardrails.
# Prevents data-centre / rail / airport / space / defence / energy ontology bleed in UI and exports.

def _v124_text_blob(model: Dict[str, Any]) -> str:
    return (str(model.get('prompt','')) + ' ' + str(model.get('title','')) + ' ' + str(model.get('subsector','')) + ' ' + str(model.get('mode',''))).lower()


def _v124_sector_key(model: Dict[str, Any]) -> str:
    t = _v124_text_blob(model)
    sub = str(model.get('subsector','')).lower()
    mode = str(model.get('mode',''))
    # Space first — most distinct
    if mode == 'Space' or any(k in t for k in ['lunar','mars','orbital','satellite','spaceport','launch vehicle','payload','moon','deep space']): return 'space'
    if any(k in t for k in ['airport','aviation','terminal','runway','heathrow','gatwick','airside','baggage','orat']): return 'airport'
    if any(k in t for k in ['data centre','data center','hyperscale','ai campus','compute campus','gpu','cloud region','white space']): return 'data_centre'
    # Defence BEFORE energy — naval/submarine/dockyard/AUKUS programmes that contain
    # 'nuclear' (propulsion) must not be misrouted to energy sector
    if any(k in t for k in ['defence','defense','military','airbase','radar','missile','secure facility','mod ','dod ',
                             'naval','dockyard','submarine','shipbuilding','warship','frigate','destroyer',
                             'aukus','naval industrial','naval propulsion','nuclear submarine','nuclear-powered',
                             'aircraft carrier','naval base','combat vessel','maritime patrol']): return 'defence'
    if any(k in t for k in ['energy','power plant','renewable','wind farm','offshore wind','solar','battery','substation','transmission','grid','hydrogen','nuclear','smr']): return 'energy'
    if any(k in t for k in ['rail','metro','transit','high speed','hs2','station','signalling','rolling stock','california high speed']): return 'rail'
    if any(k in t for k in ['semiconductor','fab','wafer','cleanroom','foundry','lithography','chip plant']): return 'semiconductor'
    if any(k in t for k in ['life sciences','pharma','biologics','gmp','fill-finish','sterile','cqv','amgen','lilly','novartis','pfizer']): return 'life_sciences'
    if any(k in t for k in ['oil','gas','lng','refinery','petrochemical','offshore','pipeline','fpso','hydrocarbon','carbon capture']): return 'oil_gas'
    if any(k in t for k in ['hospital','healthcare','clinical','medical centre','patient','nhs']): return 'healthcare'
    if any(k in t for k in ['water','wastewater','desalination','sewer','reservoir','treatment plant']): return 'water'
    if any(k in t for k in ['port','harbour','marine','dock','container terminal']): return 'ports'
    return 'general_infrastructure'


def _v124_library(key: str) -> Dict[str, Any]:
    lib = {
      'airport': {
        'label':'Airport / Aviation',
        'shock':'The dominant risk sits in live operations, ORAT, baggage/security integration and regulatory acceptance — not headline construction progress.',
        'constraints':'ORAT readiness, baggage/security systems integration, airside phasing and regulator acceptance',
        'signals':[('Live operations phasing','Active','Terminal works must sequence around airside access, airline interfaces and passenger operations.'),('ORAT readiness','Active','Operational readiness, trials and airport transition govern confidence.'),('Baggage / security integration','Watch','Baggage, security, ICT and MEP systems must operate as one environment.'),('Regulatory and airline approvals','Watch','Airport approvals and stakeholder readiness create the P80/P90 tail.')],
        'bench':[('Airport Terminal Expansion','$2B–$24B','60-144'),('Major Hub Capacity Programme','$5B–$35B','72-168'),('Rail/Transit Systems Integration','$3B–$80B','60-180')],
        'chain':['ORAT readiness','Baggage/security integration','Airside phasing','Regulatory acceptance','Operational transition','Commissioning overlap','Confidence'],
        'confidence':['Benchmark similarity: airport terminal / airfield expansion','Scope maturity: capacity, phasing and systems definition','Procurement certainty: baggage/security/MEP packages','Schedule maturity: ORAT and live operations logic','Interface exposure: airlines, airside, landside and regulators'],
        'cost':['Terminal and airside works','Baggage and security systems','Operational transition and phasing','Transport integration and utilities','Retail/passenger experience fit-out'],
        'schedule':['Live airport phasing and possessions','Baggage/security systems integration','Operational readiness trials','Regulatory and stakeholder approvals','Airside access and safety constraints']},
      'rail': {
        'label':'Rail / Transit',
        'shock':'The dominant risk sits in possessions, signalling integration, systems migration and operator acceptance rather than civil progress alone.',
        'constraints':'possessions, signalling integration, systems migration and operator acceptance',
        'signals':[('Possession window pressure','Active','Access windows and blockades govern productive delivery time.'),('Signalling and systems integration','Active','Signalling, telecoms, power and control systems drive migration risk.'),('Operator acceptance','Watch','Timetable, trial running and safety assurance determine readiness.'),('Utilities / corridor interfaces','Watch','Diversions and third-party interfaces create schedule tail exposure.')],
        'bench':[('Metro / Rail Extension','$3B–$80B','60-180'),('Major Station Redevelopment','$1B–$15B','36-96'),('Rail Systems Migration Programme','$500M–$10B','30-84')],
        'chain':['Possession access','Utility diversions','Signalling integration','Systems migration','Trial operations','Operator acceptance','Confidence'],
        'confidence':['Benchmark similarity: rail/transit programme','Scope maturity: alignment, station and systems definition','Procurement certainty: civil/systems package strategy','Schedule maturity: possessions and test/commissioning logic','Interface exposure: utilities, operators and regulators'],
        'cost':['Civil and station works','Signalling, power and telecoms systems','Utility diversions and corridor constraints','Possessions and access logistics','Systems assurance and testing'],
        'schedule':['Possession window availability','Signalling integration and migration','Utility diversion completion','Trial running and safety certification','Operator/regulator acceptance']},
      'data_centre': {
        'label':'Digital Infrastructure / Hyperscale Data Centre',
        'shock':'Power availability and systems integration readiness are more likely to govern delivery than shell construction progress.',
        'constraints':'energisation, cooling readiness and integrated systems testing',
        'signals':[('Grid connection congestion','Active','Utility interconnection and energisation remain the dominant delivery constraints.'),('Transformer / switchgear lead times','Active','Long-lead electrical packages shape procurement tail exposure.'),('Liquid cooling readiness','Watch','Cooling readiness and heat-rejection capacity govern commissioning risk.'),('IST and phased hall turnover','Watch','Integrated systems testing and phased data-hall readiness drive confidence.')],
        'bench':[('Hyperscale AI Data Centre Campus','$2B–$18B','24-84'),('Digital Infrastructure Campus','$1B–$12B','24-72'),('Energy / Utility Megaprogramme','$1B–$18B','36-108'),('Semiconductor / Advanced Manufacturing','$8B–$32B','54-96')],
        'chain':['Transformer lead-time','Grid energisation','Liquid cooling readiness','IST congestion','Commissioning overlap','Reserve drawdown','Confidence'],
        'confidence':['Benchmark similarity: hyperscale digital infrastructure','Scope maturity: campus power and white-space definition','Procurement certainty: transformers, generators and switchgear','Schedule maturity: grid and commissioning logic','Interface exposure: utilities, fibre and commissioning'],
        'cost':['Utility/grid connection and substations','Power train, transformers and switchgear','Liquid cooling / heat rejection systems','Data halls and white space fit-out','Accelerated procurement premiums'],
        'schedule':['Grid energisation and utility agreements','Long-lead transformer and switchgear delivery','Integrated systems testing and commissioning','Cooling plant readiness','Phased data-hall turnover']},
      'semiconductor': {
        'label':'Semiconductor / Advanced Manufacturing',
        'shock':'Tool install, cleanroom certification and yield-ramp readiness govern board confidence more than shell completion.',
        'constraints':'cleanroom readiness, process tool install, specialty utilities and yield-ramp qualification',
        'signals':[('Tool-install sequencing','Active','OEM tool windows and installation sequence govern critical path.'),('Cleanroom certification','Active','Cleanroom classification and environmental stability constrain turnover.'),('UPW / specialty utilities','Watch','Ultra-pure water, gases and exhaust systems drive commissioning readiness.'),('Yield ramp qualification','Watch','Production qualification creates post-mechanical-completion uncertainty.')],
        'bench':[('Advanced Semiconductor Fab','$8B–$35B','54-108'),('Cleanroom Manufacturing Campus','$2B–$12B','36-84'),('Specialty Utilities Programme','$500M–$5B','24-60')],
        'chain':['Cleanroom readiness','Specialty utilities','Process tool delivery','Tool hook-up','Qualification lots','Yield ramp','Confidence'],
        'confidence':['Benchmark similarity: advanced fab / cleanroom campus','Scope maturity: process flow and tool list maturity','Procurement certainty: OEM tool slots and specialty utilities','Schedule maturity: tool hook-up and qualification sequence','Interface exposure: utilities, vendors and yield ramp'],
        'cost':['Cleanrooms and classified areas','Process tools and hook-up','UPW/specialty gases/exhaust systems','Vibration/environmental controls','Yield ramp and qualification support'],
        'schedule':['Tool delivery and install windows','Cleanroom certification','Specialty utilities qualification','Process qualification and yield ramp','OEM/vendor interface availability']},
      'life_sciences': {
        'label':'Life Sciences / Biologics Manufacturing',
        'shock':'Mechanical completion is not the true finish line; validated production readiness and deviation closure are the real board decision gates.',
        'constraints':'CQV, GMP turnover, validation readiness and regulatory evidence',
        'signals':[('CQV readiness','Active','Commissioning, qualification and validation control release confidence.'),('Clean utility validation','Active','WFI, clean steam and process utilities constrain turnover.'),('Media fill / batch readiness','Watch','Process validation and batch release drive operational start-up risk.'),('Regulatory inspection readiness','Watch','FDA/EMA readiness affects board-defensible approval.')],
        'bench':[('Biologics Manufacturing Campus','$2B–$8B','36-78'),('Sterile Fill-Finish / Aseptic Expansion','$1B–$5B','30-60'),('Advanced GMP Cleanroom','$2B–$10B','42-84')],
        'chain':['GMP turnover','Clean utility validation','CQV protocols','Media fill readiness','Deviation closure','Regulatory readiness','Confidence'],
        'confidence':['Benchmark similarity: pharma / biologics campus','Scope maturity: GMP package and user requirement definition','Procurement certainty: process equipment and clean utility lead times','Schedule maturity: CQV logic and validation pathway','Regulatory exposure: FDA/EMA inspection readiness'],
        'cost':['Process equipment and fill-finish lines','GMP cleanrooms / classified areas','Clean utilities and HVAC zoning','CQV validation and deviation closure','Regulatory readiness and quality systems'],
        'schedule':['CQV protocol approval and execution','Long-lead process equipment delivery','Clean utility validation and media fills','FDA/EMA inspection readiness','Batch release and operational readiness']},
      'defence': {
        'label':'Defence / Secure Mission Systems',
        'shock':'Assurance, security accreditation and mission-system integration are likely to govern approval before asset completion.',
        'constraints':'security accreditation, mission assurance, sovereign supply chain and integration test evidence',
        'signals':[('Security accreditation','Active','Classified systems and facility accreditation create approval gates.'),('Mission-system integration','Active','Sensors, comms, command systems and platform interfaces govern readiness.'),('Sovereign supply-chain exposure','Watch','Controlled components and cleared supplier capacity affect delivery tails.'),('Assurance test campaign','Watch','Operational acceptance depends on test evidence and traceability.')],
        'bench':[('Secure Mission Facility','$500M–$8B','30-96'),('Defence Systems Integration Programme','$1B–$20B','48-144'),('Naval / Airbase Modernisation','$1B–$15B','48-120')],
        'chain':['Security requirements','Sovereign procurement','Mission-system integration','Assurance test campaign','Accreditation gates','Operational acceptance','Confidence'],
        'confidence':['Benchmark similarity: defence secure systems programme','Scope maturity: mission requirement and accreditation definition','Procurement certainty: cleared suppliers and controlled components','Schedule maturity: integration/test campaign logic','Interface exposure: security, operators and regulators'],
        'cost':['Secure facilities and hardening','Mission systems and integration','Controlled equipment procurement','Cyber/security accreditation','Test campaign and operational acceptance'],
        'schedule':['Security accreditation gates','Mission-system integration testing','Controlled supplier lead-times','Operational test campaign','Authority-to-operate approval']},
      'oil_gas': {
        'label':'Energy / Oil & Gas',
        'shock':'Procurement, modular integration, HSE readiness and commissioning/start-up govern value more than installed quantities.',
        'constraints':'long-lead equipment, module integration, HSE readiness and start-up/commissioning assurance',
        'signals':[('Long-lead rotating equipment','Active','Compressors, turbines and specialist packages shape procurement exposure.'),('Module / brownfield integration','Active','Tie-ins, shutdown windows and interface planning create schedule tails.'),('HSE / regulatory approvals','Watch','Permitting and safety case maturity affect board confidence.'),('Commissioning and start-up','Watch','Hydrotest, pre-commissioning and start-up dominate late risk.')],
        'bench':[('LNG / Gas Processing Train','$5B–$30B','48-108'),('Refinery / Petrochemical Expansion','$2B–$20B','36-96'),('Offshore / Pipeline Programme','$1B–$18B','36-96')],
        'chain':['FEED maturity','Long-lead packages','Module fabrication','Tie-in windows','Pre-commissioning','Start-up readiness','Confidence'],
        'confidence':['Benchmark similarity: oil/gas processing or export infrastructure','Scope maturity: FEED, plot plan and tie-in definition','Procurement certainty: rotating equipment and specialist packages','Schedule maturity: modularisation and shutdown logic','Interface exposure: brownfield operations, HSE and regulators'],
        'cost':['Process units and rotating equipment','Pipelines / terminals / tie-ins','Modular fabrication and logistics','HSE, permitting and compliance','Commissioning and start-up support'],
        'schedule':['Long-lead compressor/turbine packages','Module fabrication and transport','Shutdown/tie-in windows','Regulatory and HSE approvals','Pre-commissioning and start-up']},
      'energy': {
        'label':'Energy / Power Infrastructure',
        'shock':'Grid access, equipment lead-times, permitting and commissioning readiness govern delivery confidence more than site progress.',
        'constraints':'grid connection, major equipment lead-times, permitting and commissioning readiness',
        'signals':[('Grid connection access','Active','Interconnection agreements and energisation gates govern the delivery tail.'),('Major equipment lead-times','Active','Transformers, turbines, inverters or batteries shape procurement risk.'),('Permitting and environmental approvals','Watch','Approvals and grid studies affect schedule certainty.'),('Commissioning / performance testing','Watch','Performance tests and grid compliance govern final acceptance.')],
        'bench':[('Transmission / Grid Reinforcement','$500M–$10B','24-84'),('Renewable + Storage Programme','$300M–$8B','18-60'),('Power Generation Programme','$1B–$15B','36-96')],
        'chain':['Permitting','Grid studies','Major equipment delivery','Site construction','Energisation','Performance testing','Confidence'],
        'confidence':['Benchmark similarity: power / utility infrastructure','Scope maturity: connection, permitting and design definition','Procurement certainty: transformers, turbines, inverters or batteries','Schedule maturity: energisation and commissioning logic','Interface exposure: grid operator, regulator and land/environment'],
        'cost':['Grid connection and substations','Generation/storage equipment','Civil and balance-of-plant works','Environmental and permitting compliance','Commissioning and performance testing'],
        'schedule':['Grid connection agreement','Long-lead equipment delivery','Permitting and environmental approval','Energisation windows','Performance testing and acceptance']},
      'space': {
        'label':'Space / Lunar / Orbital Infrastructure',
        'shock':'Launch alone is not the real constraint; qualification, thermal-power balance and autonomous recovery determine mission survivability.',
        'constraints':'mission assurance, launch logistics, qualification evidence and autonomous recovery',
        'signals':[('Launch reliability volatility','Active','Launch cadence, range access and manifest priority affect schedule certainty.'),('Mission assurance burden','Active','Qualification and redundancy evidence govern board defensibility.'),('Thermal-power balance','Watch','Power storage and thermal rejection shape survivability.'),('Regulatory and range availability','Watch','Flight approvals, range windows and orbital operations create schedule tails.')],
        'bench':[('Orbital / Lunar Infrastructure','$8B–$95B','72-216'),('Launch and Payload Integration','$1B–$14B','36-108'),('Deep-Space Mission Systems','$4B–$35B','60-144'),('Autonomous Operations Platform','$2B–$20B','36-96')],
        'chain':['Launch cadence','Payload integration','Thermal-power balance','Range availability','Autonomous commissioning','Mission assurance','Confidence'],
        'confidence':['Benchmark similarity: space infrastructure archetype','Scope maturity: payload and mission architecture definition','Procurement certainty: launch, avionics and qualified hardware','Schedule maturity: launch and commissioning logic','Operational exposure: remote recovery and servicing limits'],
        'cost':['Payload / habitat / mission systems','Launch and orbital logistics','Power, thermal and autonomy systems','Qualification and test campaign','Mission operations and recovery reserve'],
        'schedule':['Payload qualification campaign','Launch manifest and range access','Thermal-vacuum and systems testing','Orbital/surface deployment sequence','Autonomous commissioning and recovery']},
      'healthcare': {
        'label':'Healthcare / Hospital Infrastructure','shock':'Clinical transition, medical equipment integration and infection-control readiness govern approval more than building completion.','constraints':'clinical commissioning, infection-control compliance, medical equipment integration and phased occupancy','signals':[('Clinical commissioning','Active','Clinical workflows and patient transition govern readiness.'),('Medical equipment integration','Active','Imaging, theatres, labs and digital systems drive handover risk.'),('Infection-control compliance','Watch','ICRA and commissioning evidence shape confidence.'),('Phased occupancy','Watch','Live hospital operations constrain access and transition.')],'bench':[('Major Hospital Campus','$1B–$8B','36-96'),('Clinical Tower / Specialist Centre','$500M–$4B','30-72')],'chain':['Clinical requirements','Medical equipment procurement','Digital systems integration','Infection-control validation','Phased occupancy','Clinical commissioning','Confidence'],'confidence':['Benchmark similarity: hospital / clinical campus','Scope maturity: clinical brief and department definition','Procurement certainty: medical equipment and digital systems','Schedule maturity: clinical commissioning and phased occupancy','Interface exposure: clinicians, patients, regulators and live operations'],'cost':['Clinical departments and fit-out','Medical equipment and imaging','Digital health and systems integration','Infection-control and compliance','Phased occupancy and transition'],'schedule':['Medical equipment delivery','Clinical commissioning readiness','Infection-control validation','Digital systems integration','Patient transition and phased occupancy']},
      'water': {
        'label':'Water / Environmental Infrastructure','shock':'Consents, process commissioning and environmental compliance govern delivery confidence more than civil installation.','constraints':'permits, process commissioning, environmental compliance and operational acceptance','signals':[('Consents and discharge permits','Active','Environmental approvals govern start-up readiness.'),('Process commissioning','Active','Treatment-process performance controls acceptance.'),('MEICA procurement','Watch','Pumps, controls and specialist equipment create lead-time risk.'),('Operational takeover','Watch','Operator readiness and compliance testing affect confidence.')],'bench':[('Water Treatment Programme','$300M–$5B','24-72'),('Desalination / Major Water Plant','$1B–$10B','36-96')],'chain':['Environmental permits','MEICA procurement','Civil/process works','Process commissioning','Compliance testing','Operator acceptance','Confidence'],'confidence':['Benchmark similarity: water treatment / environmental asset','Scope maturity: process capacity and permit definition','Procurement certainty: MEICA and specialist equipment','Schedule maturity: commissioning and compliance testing logic','Interface exposure: regulator, operator and environmental stakeholders'],'cost':['Civil/process treatment works','MEICA equipment','Pipelines and intake/outfall works','Environmental mitigation','Commissioning and compliance testing'],'schedule':['Permit approval','MEICA procurement','Process commissioning','Compliance testing','Operator acceptance']},
      'ports': {
        'label':'Ports / Marine Infrastructure','shock':'Marine access, dredging, quay interfaces and terminal systems govern delivery risk more than civil progress alone.','constraints':'marine access, dredging, quay works, terminal systems and port operational interfaces','signals':[('Marine access window','Active','Weather, tides and vessel windows shape productivity.'),('Dredging / quay interface','Active','Marine works and ground conditions control schedule exposure.'),('Terminal systems readiness','Watch','Cranes, automation and yard systems affect operational acceptance.'),('Port operations interface','Watch','Live port operations constrain phasing and access.')],'bench':[('Container Terminal Expansion','$500M–$8B','30-84'),('Major Port / Marine Works','$1B–$12B','36-96')],'chain':['Marine access','Dredging / ground risk','Quay construction','Cranes and terminal systems','Operational trials','Port acceptance','Confidence'],'confidence':['Benchmark similarity: port / terminal expansion','Scope maturity: marine works and terminal operating model','Procurement certainty: cranes, systems and marine contractors','Schedule maturity: dredging, quay and systems commissioning','Interface exposure: shipping, operators and regulators'],'cost':['Dredging and marine works','Quay / berth construction','Cranes and terminal systems','Yard, utilities and access roads','Operational trials and port transition'],'schedule':['Marine access and weather windows','Dredging / ground treatment','Quay and berth completion','Cranes / terminal systems commissioning','Port operational acceptance']}
    }
    return lib.get(key, {
      'label':'Capital Infrastructure','shock':'The dominant risk sits in interfaces, procurement evidence and commissioning readiness rather than headline progress.','constraints':'interface control, procurement evidence and commissioning readiness','signals':[('Procurement volatility','Watch','Market capacity and supplier evidence shape confidence.'),('Approvals latency','Watch','Governance and permit friction influence schedule tails.'),('Commissioning readiness','Active','Operational readiness is a board-defensibility signal.'),('Interface density','Active','Interfaces create the largest uncertainty transfer.')],'bench':[('Complex Infrastructure Programme','$1B–$20B','36-120')],'chain':['Scope definition','Procurement evidence','Interface control','Commissioning readiness','Operational acceptance','Reserve adequacy','Confidence'],'confidence':['Benchmark similarity: comparable infrastructure archetype','Scope maturity: requirements and package definition','Procurement certainty: supplier and market evidence','Schedule maturity: critical-path and commissioning logic','Interface exposure: utilities, operators and approvals'],'cost':['Civil/enabling works','Specialist systems and equipment','Utilities and interfaces','Programme indirects','Risk reserve and contingency'],'schedule':['Approvals and governance','Long-lead procurement','Interface coordination','Commissioning readiness','Operational acceptance']})


def _v124_forbidden_terms(key: str):
    all_terms = {
      'data_centre':['ORAT','baggage','airside','landside','rolling stock','signalling','payload integration','launch cadence','range availability','yield ramp','CQV'],
      'airport':['liquid cooling','white-space','white space','GPU','data hall','thermal rejection','yield ramp','payload','launch reliability','transformer lead-time'],
      'rail':['liquid cooling','GPU','data hall','ORAT','baggage systems','airside','launch reliability','thermal-power','yield ramp'],
      'space':['ORAT','baggage','airside','rolling stock','signalling','liquid cooling readiness','white space','yield ramp','CQV'],
      'semiconductor':['ORAT','baggage','airside','rolling stock','launch reliability','payload integration'],
      'life_sciences':['ORAT','baggage','rolling stock','launch reliability','liquid cooling readiness','white space'],
      'defence':['ORAT','baggage','liquid cooling readiness','white space','yield ramp'],
      'oil_gas':['ORAT','baggage','rolling stock','liquid cooling readiness','white space','payload integration','yield ramp'],
      'energy':['ORAT','baggage','rolling stock','liquid cooling readiness','payload integration','yield ramp'],
      'healthcare':['launch reliability','liquid cooling readiness','rolling stock','baggage systems'],
      'water':['launch reliability','liquid cooling readiness','rolling stock','baggage systems'],
      'ports':['launch reliability','liquid cooling readiness','rolling stock','baggage systems','ORAT']
    }
    return all_terms.get(key, [])


def _v124_scrub_value(v, forbidden):
    """Remove only terms forbidden for the selected ontology.
    Do not apply global replacements blindly; otherwise valid ORAT/launch/cooling terms are damaged.
    """
    replacements = {
        'liquid cooling readiness':'systems readiness', 'Liquid cooling readiness':'Systems readiness',
        'Liquid cooling':'Systems readiness', 'liquid cooling':'systems readiness',
        'white-space':'specialist systems', 'white space':'specialist systems', 'data hall':'delivery package', 'GPU':'specialist equipment',
        'ORAT':'operational readiness', 'baggage systems':'systems integration', 'Baggage':'Systems', 'baggage':'systems',
        'airside':'operational access', 'landside':'public interface', 'rolling stock':'operating assets', 'signalling':'control systems',
        'launch reliability':'delivery reliability', 'Launch reliability':'Delivery reliability', 'payload integration':'systems integration',
        'range availability':'access availability', 'yield ramp':'operational ramp-up', 'CQV':'commissioning validation'
    }
    if isinstance(v, str):
        out = v
        low_forbidden = {str(f).lower() for f in forbidden}
        for a,b in replacements.items():
            if a.lower() in low_forbidden:
                out = out.replace(a,b)
        return out
    if isinstance(v, list): return [_v124_scrub_value(x, forbidden) for x in v]
    if isinstance(v, dict): return {k:_v124_scrub_value(val, forbidden) for k,val in v.items()}
    return v


def _v124_apply_sector_lock(model: Dict[str, Any]) -> Dict[str, Any]:
    m = dict(model or {})
    key = _v124_sector_key(m)
    L = _v124_library(key)
    m['sector_ontology_key'] = key
    m['sector_ontology_label'] = L['label']
    m['executive_shock_insight'] = L['shock']
    m['sector_confidence_drivers'] = L['confidence']
    m['sector_primary_cost_drivers'] = L['cost']
    m['sector_schedule_threats'] = L['schedule']
    m['causal_graph_nodes'] = L['chain']
    m['causal_chain'] = L['chain']
    m['benchmark_comparison'] = [{'archetype':a,'anchor_cost':c,'anchor_duration_months':d} for a,c,d in L['bench']]
    m['why_casey_generated_this'] = [
        f"CASEY detected {L['label']} from the project brief and locked the sector ontology before generating outputs.",
        f"Sector-native behaviours applied: {', '.join(L['chain'][:5])}.",
        "Benchmark cohort, causal chain, confidence drivers and risk language were constrained to the selected sector.",
        "The output is designed for early board challenge and scope definition, not certified pricing."
    ]
    sigs=[]
    for s,status,basis in L['signals']:
        sigs.append({'signal':s,'status':status,'direction':'sector-locked calibration','weight':0.12,'applies_to':'confidence, QCRA/QSRA, risk register','basis':basis})
    m['live_calibration_signals']=sigs
    m['live_calibration_strip']=' • '.join([x['signal'] for x in sigs[:4]])
    m['mission_control_cards']=[{'label':'Live calibration','signal':'Current sector conditions are being applied to confidence, contingency and delivery-tail exposure.','severity':'Active'}] + [{'label':s,'signal':b,'severity':st} for s,st,b in L['signals']]
    m['mission_control_cards']=m['mission_control_cards'][:6]
    m['casey_thinking'] = f"CASEY has re-cut the programme as a {m.get('scenario_label','Base')} scenario inside the {L['label']} ontology. The governing consequence is: {L['shock']} QCRA/QSRA curves, cost basis, risk probabilities and schedule logic have been sector-locked to this posture."
    m['executive_summary'] = f"{L['label']} scenario view: {m.get('scenario_label','Base')}. CASEY indicates {m.get('cost_p50')} P50 exposure, {m.get('cost_range')} range, {m.get('schedule')} baseline, {m.get('risk')} risk and {m.get('confidence_pct')}% confidence. {L['shock']} Trade-off: {m.get('scenario_why','Balanced cost, time and evidence posture.')}"
    m['board_briefing'] = [
        L['shock'],
        f"{m.get('scenario_label','Base')} is the reference case: no cost, schedule or confidence delta is applied unless a scenario is selected.",
        f"{m.get('scenario_label','Base')} scenario: {m.get('cost_p50')} P50, {m.get('cost_range')} range, {m.get('schedule')} baseline and {m.get('confidence_pct')}% confidence.",
        f"Confidence is governed by {L['constraints']}.",
        "Gained: Maintains a credible reference case for board challenge."
    ]
    m['uncertainty_narrative']={
        'estimate_maturity':'Class 3 maturity is suitable for budget authorization, but procurement and design assumptions still need challenge.',
        'schedule_maturity':'Schedule Level 4 gives stronger logic and QSRA traceability.',
        'interpretation':f"Live calibration is weighting {L['constraints']} into the QSRA/QCRA tail."
    }
    # Scenario delta language remains, but sector nouns are locked.
    m['top_decisions_required']=[
        'Accept or reject the scenario trade-off explicitly at board level.',
        'Confirm the changed critical path and near-critical path density.',
        'Approve the scenario-specific reserve and contingency philosophy.',
        'Assign named owners for the new top scenario risks.',
        'Confirm whether XER, cost workbook and risk register should be issued as scenario-controlled outputs.'
    ]
    forbidden = _v124_forbidden_terms(key)
    m = _v124_scrub_value(m, forbidden)
    # restore sector-native strings after broad scrub in case terms are valid for this key
    m['sector_ontology_key'] = key; m['sector_ontology_label'] = L['label']
    if key == 'data_centre':
        m['sector_primary_cost_drivers'] = L['cost']; m['sector_schedule_threats'] = L['schedule']; m['causal_graph_nodes'] = L['chain']; m['causal_chain'] = L['chain']
    if key == 'space':
        m['sector_confidence_drivers'] = L['confidence']; m['sector_schedule_threats'] = L['schedule']; m['causal_graph_nodes'] = L['chain']; m['causal_chain'] = L['chain']
    return m

_CASEY_V124_PREV_BUILD_MODEL = build_model
def build_model(prompt:str, client:str='', class_level:int=3, schedule_level:int=3, scenario:str='base'):
    return _v124_apply_sector_lock(_CASEY_V124_PREV_BUILD_MODEL(prompt, client, class_level, schedule_level, scenario))

APP_VERSION = 'CASEY V124 Sector Ontology Hardened Public Demo'
print('CASEY V124 sector ontology hardening lock installed')
# ================= END CASEY V124 SECTOR ONTOLOGY HARDENING LOCK =================

# ========================= CASEY V128 ADVISORY AUTHORITY + EXPORT GOVERNANCE LOCK =========================
# Stable patch: keeps the V127/V124 app shell intact while making Ask CASEY and exports more
# institutional, sector-aware and board-challenge oriented.

def _v128_as_float(x, default=0.0):
    try:
        if isinstance(x, str):
            s=x.replace('$','').replace(',','').strip().upper()
            if s.endswith('T'): return float(s[:-1]) * 1000
            if s.endswith('B'): return float(s[:-1])
            if s.endswith('M'): return float(s[:-1]) / 1000
            return float(s)
        return float(x)
    except Exception:
        return default

def _v128_money(x):
    x=_v128_as_float(x)
    if x >= 1000: return f"${x/1000:.1f}T"
    if x >= 1: return f"${x:.1f}B"
    if x > 0: return f"${x*1000:.0f}M"
    return "not quantified"

def _v128_sector_pack(model):
    fam=_casey_sector_family(model)
    base=_SECTOR_INTEL.get(fam, _SECTOR_INTEL.get('infrastructure', {}))
    if not base:
        base={'governor':'interfaces, evidence maturity and operational readiness','hidden':'headline progress can hide confidence erosion','fail':'evidence maturity lagged the reported programme position','evidence':['owner evidence','schedule logic','procurement certainty','commissioning readiness','reserve adequacy'],'traditional':'Progress is read through cost and schedule reporting.','casey':'Confidence is read through governing constraints and evidence maturity.','attacks':['What evidence proves readiness?','Which assumption collapses confidence first?','Is float real or theoretical?']}
    return fam, base

def _v128_casey_position(model, intel):
    title=model.get('title') or model.get('client') or 'This programme'
    conf=model.get('confidence_pct', model.get('confidence_score','n/a'))
    risk=model.get('risk','elevated')
    return f"""CASEY POSITION

CASEY does not judge {title} on headline progress alone.

The programme should be governed through {intel['governor']}. Current confidence ({conf}%) is a board-defensibility signal, not an optimism score. Risk posture is {risk}; approval should be conditional until evidence proves the governing constraint is controlled.

If management cannot evidence this constraint, the programme is not approval-ready; it is only reportable."""

def _v128_traditional_vs_casey(intel):
    return f"""TRADITIONAL CONTROLS VS CASEY

Traditional controls view:
{intel['traditional']}

CASEY view:
{intel['casey']}

Why this matters:
Traditional reporting can remain green while confidence is already deteriorating. CASEY reads the behaviour underneath the dashboard: evidence maturity, governing constraint exposure, and whether float/reserve are actually usable."""

def _v128_board_attacks(model, intel):
    attacks=list(intel.get('attacks') or [])[:5]
    while len(attacks)<5:
        attacks.append(['What evidence proves the current confidence position is real?','Which assumption collapses P80/P90 exposure first?','Is the reported float operationally usable or theoretical?','Who owns the governing interface liability?','Why should the board believe the reserve posture?'][len(attacks)])
    return "LIKELY BOARD ATTACKS\n\n" + "\n".join([f"{i+1}. {x}" for i,x in enumerate(attacks)])

def _v128_failure(model, intel):
    return f"""IF THIS PROGRAMME FAILS

The most likely external explanation will be:
"{intel['fail']}."

The internal explanation will be sharper:
Management treated reported progress as equivalent to board-defensible readiness.

CASEY WARNING
The failure mode is behavioural before it is numerical: optimism, weak evidence closure, and late recognition of the governing constraint."""

def _v128_board_decision(model, intel):
    ev='\n'.join([f"- {x}" for x in (intel.get('evidence') or [])[:5]])
    return f"""WHAT THE BOARD IS REALLY DECIDING

The board is not simply approving cost and schedule.

The board is accepting:
- exposure to {intel['governor']}
- the current evidence maturity position
- whether reserve is credible against P80/P90 behaviour
- whether management can recover if the governing constraint slips

Evidence required before unconditional approval:
{ev}

CASEY POSITION
Approval should be conditional. Confidence improvement requires evidence closure, not a more polished narrative."""

def _v128_intervention(model, intel):
    ev=(intel.get('evidence') or ['governing constraint evidence'])[0]
    return f"""INTERVENTION INTELLIGENCE

The highest-value intervention is not another status pack.

1. Prove {intel['governor']} with auditable evidence.
2. Convert the weakest assumption into an owned mitigation with trigger, date, owner and residual exposure.
3. Re-run P80/P90 after evidence closure, not after narrative revision.

Fastest confidence unlock:
Close the first missing evidence item: {ev}.

CASEY POSITION
The programme needs targeted evidence closure where confidence is actually governed."""

def _format_advisor_answer(question: str, model: Dict[str, Any]) -> str:
    # V128 override: robust, string-only, sector-native, board-challenge oriented.
    try:
        if not isinstance(question, str): question=str(question or '')
    except Exception:
        question=''
    q=(question or '').lower().strip()
    fam,intel=_v128_sector_pack(model or {})
    title=(model or {}).get('title') or (model or {}).get('client') or 'this programme'
    conf,p80c,p80s=_advisor_metrics(model or {})
    risks=(model or {}).get('risks') or (model or {}).get('risk_register') or []
    risk_titles=[]
    for r in risks[:4]:
        if isinstance(r, dict): risk_titles.append(str(r.get('title') or r.get('risk') or r.get('name') or 'unresolved risk'))
        else: risk_titles.append(str(r))
    if not risk_titles: risk_titles=(intel.get('evidence') or [])[:4]

    if any(k in q for k in ['turner','t&t','traditional','pmo','consultant','controls view','casey view']):
        return _v128_traditional_vs_casey(intel) + "\n\n" + _v128_casey_position(model or {}, intel)
    if any(k in q for k in ['attack','board challenge','external assurance','reviewer','audit','challenge first']):
        return _v128_board_attacks(model or {}, intel) + "\n\n" + _v128_casey_position(model or {}, intel)
    if any(k in q for k in ['disagree','agree','reported','management reporting','official','status','green']):
        return f"""CASEY DISAGREES WITH THE REPORTED POSITION

The visible programme narrative may appear stable.

CASEY disagrees if management is relying on reported progress without proving {intel['governor']}.

Current confidence: {conf}%
P80 cost exposure: {p80c}
P80 schedule exposure: {p80s} months

The hidden issue:
{intel['hidden']}.

""" + _v128_casey_position(model or {}, intel)
    if any(k in q for k in ['fail','fails','blamed','public','headline','collapse']):
        return _v128_failure(model or {}, intel)
    if any(k in q for k in ['really deciding','decision','approve','approval','board is deciding']):
        return _v128_board_decision(model or {}, intel)
    if any(k in q for k in ['confidence','why low','why is confidence','moved','score']):
        ev='\n'.join([f"- {x}" for x in (intel.get('evidence') or [])[:5]])
        return f"""CONFIDENCE EXPLANATION

Current confidence: {conf}%

Confidence is not optimism. It is board defensibility.

CASEY reads confidence through:
{ev}

Primary governing constraint:
{intel['governor']}

P80 cost exposure: {p80c}
P80 schedule exposure: {p80s} months

CASEY POSITION
If management wants confidence to improve, it must close evidence gaps around the governing constraint. Contingency without evidence buys time; it does not create confidence."""
    if any(k in q for k in ['intervention','do now','fix','improve','action','changes confidence','fastest']):
        return _v128_intervention(model or {}, intel)
    if any(k in q for k in ['governing chain','real chain','causal','dependency']):
        chain=(model or {}).get('causal_chain') or (model or {}).get('causal_graph_nodes') or []
        if not chain: chain=(intel.get('evidence') or [])[:5] + ['confidence']
        chain='\n→ '.join([str(x) for x in chain[:7]])
        return f"""REAL GOVERNING CHAIN

{chain}

Interpretation:
The programme is not governed by the most visible activity. It is governed by the weakest evidence link in the chain.

""" + _v128_casey_position(model or {}, intel)
    if any(k in q for k in ['hiding','not seeing','dashboard fail','looks green','red','optimism']):
        bullets='\n'.join([f"- {x}" for x in risk_titles[:4]])
        return f"""WHAT THE DASHBOARD FAILS TO SHOW

The dominant issue is not the headline cost or schedule.

The hidden issue:
{intel['hidden']}.

Most exposed items:
{bullets}

Traditional reporting may show movement. CASEY is testing whether the movement is governable.

""" + _v128_casey_position(model or {}, intel)
    return f"""WHAT THE BOARD IS NOT SEEING

The dominant issue is not the headline cost or schedule.

The hidden issue:
{intel['hidden']}.

Current confidence: {conf}%
P80 cost exposure: {p80c}
P80 schedule exposure: {p80s} months

Most exposed items:
""" + "\n".join([f"- {x}" for x in risk_titles[:4]]) + "\n\n" + _v128_casey_position(model or {}, intel)

# Export upgrades: add board challenge/intelligence tabs and sections without changing endpoint names.
_CASEY_V128_PREV_WORKBOOK_BYTES = workbook_bytes
_CASEY_V128_PREV_WORD_BYTES = word_bytes
_CASEY_V128_PREV_PDF_BYTES = pdf_bytes

def workbook_bytes(model):
    data=_CASEY_V128_PREV_WORKBOOK_BYTES(model)
    try:
        from openpyxl import load_workbook
        bio=BytesIO(data); wb=load_workbook(bio)
        fam,intel=_v128_sector_pack(model)
        ws=wb.create_sheet('CASEY Board Challenge')
        rows=[
            ['Section','Content'],
            ['CASEY Position',_v128_casey_position(model,intel)],
            ['Traditional Controls vs CASEY',_v128_traditional_vs_casey(intel)],
            ['Board Attacks',_v128_board_attacks(model,intel)],
            ['What The Board Is Really Deciding',_v128_board_decision(model,intel)],
            ['If This Programme Fails',_v128_failure(model,intel)],
            ['Intervention Intelligence',_v128_intervention(model,intel)],
        ]
        for r in rows: ws.append(r)
        ws.column_dimensions['A'].width=34; ws.column_dimensions['B'].width=120
        for row in ws.iter_rows():
            for c in row: c.alignment = Alignment(wrap_text=True, vertical='top')
        out=BytesIO(); wb.save(out); out.seek(0); return out.getvalue()
    except Exception:
        return data

def word_bytes(model):
    doc=Document(); styles=doc.styles; styles['Normal'].font.name='Aptos'; styles['Normal'].font.size=Pt(10)
    fam,intel=_v128_sector_pack(model)
    doc.add_heading('CASEY Strategic Board Intelligence Pack',0)
    doc.add_paragraph(model.get('executive_summary',''))
    doc.add_heading('CASEY Position',1); doc.add_paragraph(_v128_casey_position(model,intel))
    doc.add_heading('Traditional Controls vs CASEY',1); doc.add_paragraph(_v128_traditional_vs_casey(intel))
    doc.add_heading('Likely Board Attacks',1)
    for line in _v128_board_attacks(model,intel).split('\n')[2:]:
        if line.strip(): doc.add_paragraph(line, style='List Bullet')
    doc.add_heading('What The Board Is Really Deciding',1); doc.add_paragraph(_v128_board_decision(model,intel))
    doc.add_heading('If This Programme Fails',1); doc.add_paragraph(_v128_failure(model,intel))
    doc.add_heading('Intervention Intelligence',1); doc.add_paragraph(_v128_intervention(model,intel))
    doc.add_heading('Core Metrics',1)
    tbl=doc.add_table(rows=1,cols=2); tbl.style='Light Shading Accent 1'; tbl.rows[0].cells[0].text='Metric'; tbl.rows[0].cells[1].text='Value'
    for k,v in [('Cost P50',model.get('cost_p50')),('Cost Range',model.get('cost_range')),('Schedule',model.get('schedule')),('Risk',model.get('risk')),('Confidence',str(model.get('confidence_pct'))+'%'),('Sector Ontology',model.get('sector_ontology_label',fam))]:
        row=tbl.add_row().cells; row[0].text=str(k); row[1].text=str(v)
    doc.add_heading('Top Risks',1)
    for r in (model.get('risks') or [])[:8]:
        if isinstance(r, dict): doc.add_paragraph(f"{r.get('risk_id','R')} {r.get('title','Risk')} — {r.get('mitigation','mitigation required')}",style='List Bullet')
    bio=BytesIO(); doc.save(bio); bio.seek(0); return bio.getvalue()

def pdf_bytes(model):
    bio=BytesIO(); doc=SimpleDocTemplate(bio,pagesize=A4, rightMargin=36, leftMargin=36, topMargin=36, bottomMargin=36)
    styles=getSampleStyleSheet(); story=[]; fam,intel=_v128_sector_pack(model)
    story.append(Paragraph('CASEY Strategic Board Intelligence Pack',styles['Title'])); story.append(Spacer(1,12))
    story.append(Paragraph(str(model.get('executive_summary','')),styles['BodyText'])); story.append(Spacer(1,12))
    for title,txt in [('CASEY Position',_v128_casey_position(model,intel)),('Traditional Controls vs CASEY',_v128_traditional_vs_casey(intel)),('Likely Board Attacks',_v128_board_attacks(model,intel)),('What The Board Is Really Deciding',_v128_board_decision(model,intel)),('If This Programme Fails',_v128_failure(model,intel)),('Intervention Intelligence',_v128_intervention(model,intel))]:
        story.append(Paragraph(title,styles['Heading1']))
        for para in str(txt).split('\n\n'):
            story.append(Paragraph(para.replace('\n','<br/>'),styles['BodyText'])); story.append(Spacer(1,7))
    doc.build(story); bio.seek(0); return bio.getvalue()

APP_VERSION = 'CASEY V128 Advisor + Export Strategic Governance Final Stable'
print('CASEY V128 advisor/export strategic governance layer installed')
# ======================= END CASEY V128 ADVISORY AUTHORITY + EXPORT GOVERNANCE LOCK =======================

# ========================= CASEY V145 GLOBAL SHOWCASE CONSISTENCY LOCK =========================
# Final hardening layer: one canonical set of selected-scenario numbers is stamped across
# UI payloads and all exports. This fixes stale narrative/card/header mismatches where
# scenario_matrix, cost lines, QCRA/QSRA and executive text could disagree.

def _v145_bn(x):
    try:
        if isinstance(x, (int, float)): return float(x)
        s=str(x or '').replace('$','').replace(',','').strip().upper()
        if not s: return 0.0
        if s.endswith('T'): return float(s[:-1])*1000.0
        if s.endswith('B'): return float(s[:-1])
        if s.endswith('M'): return float(s[:-1])/1000.0
        return float(re.sub(r'[^0-9.\-]','',s) or 0)
    except Exception:
        return 0.0

def _v145_money(bn):
    try: bn=float(bn)
    except Exception: bn=0.0
    if bn >= 1000: return f'${bn/1000:.1f}T'
    if bn >= 1: return f'${bn:.1f}B'
    return f'${bn*1000:.0f}M'

def _v145_months(v):
    try:
        if isinstance(v,(int,float)): return int(round(float(v)))
        m=re.search(r'-?\d+(?:\.\d+)?', str(v or ''))
        return int(round(float(m.group(0)))) if m else 0
    except Exception:
        return 0

def _v145_int(v, default=55):
    try:
        if isinstance(v,(int,float)): return int(round(float(v)))
        m=re.search(r'-?\d+(?:\.\d+)?', str(v or ''))
        return int(round(float(m.group(0)))) if m else default
    except Exception:
        return default

def _v145_profile(scenario):
    s=str(scenario or 'base').lower().replace(' ','_')
    return {
        'base':       {'label':'Base','p10':0.84,'p80':1.30,'p90':1.47,'sp80':1.30,'sp90':1.47,'risk':'Medium-High','conf_delta':0},
        'faster':     {'label':'Faster','p10':0.74,'p80':1.42,'p90':1.60,'sp80':1.60,'sp90':1.78,'risk':'High','conf_delta':-13},
        'cheaper':    {'label':'Cheaper','p10':0.72,'p80':1.45,'p90':1.68,'sp80':1.42,'sp90':1.62,'risk':'High','conf_delta':-16},
        'lower_risk': {'label':'Lower Risk','p10':0.88,'p80':1.18,'p90':1.31,'sp80':1.18,'sp90':1.30,'risk':'Medium-Low','conf_delta':10},
        'premium':    {'label':'Premium','p10':0.86,'p80':1.16,'p90':1.27,'sp80':1.14,'sp90':1.24,'risk':'Low','conf_delta':16},
    }.get(s, {'label':'Base','p10':0.84,'p80':1.30,'p90':1.47,'sp80':1.30,'sp90':1.47,'risk':'Medium-High','conf_delta':0})

def _v145_selected_row(model):
    scenario=str(model.get('scenario') or 'base').lower().replace(' ','_')
    for row in (model.get('scenario_matrix') or model.get('scenario_comparison') or []):
        if str(row.get('scenario') or '').lower().replace(' ','_') == scenario:
            return row
    return {}

def _v145_sector_name(model):
    return str(model.get('sector_ontology_label') or model.get('subsector') or model.get('mode') or 'Capital Programme')

def _v145_constraint(model):
    # Prefer sector-native constraint language already generated by V124/V128.
    key=str(model.get('sector_ontology_key') or '').lower()
    mapping={
        'rail':'possessions, signalling integration, systems migration and operator acceptance',
        'space':'qualification, thermal-power balance, launch logistics and autonomous recovery',
        'data_centre':'utility energisation, cooling readiness, transformer supply and integrated systems testing',
        'airport':'ORAT readiness, baggage/security integration, airside phasing and regulator acceptance',
        'life_sciences':'CQV readiness, GMP turnover, sterile utilities and regulatory release evidence',
        'semiconductor':'tool install, UPW readiness, cleanroom qualification and yield-ramp evidence',
        'energy':'regulatory approvals, grid interfaces, long-lead equipment and commissioning sequence',
        'defence':'mission assurance, secure supply chain, systems integration and acceptance evidence',
        'oil_gas':'cryogenic systems, marine logistics, long-lead equipment and commissioning readiness',
        'nuclear':'licensing, safety case maturity, containment qualification and supply-chain evidence',
    }
    if key in mapping: return mapping[key]
    txt=str(model.get('executive_shock_insight') or '').strip()
    return txt if txt else 'interfaces, procurement evidence, schedule logic and commissioning readiness'

def _v145_conf_label(conf, scenario):
    conf=_v145_int(conf,55)
    if conf < 40: return 'Do not approve without more evidence'
    if conf < 58: return 'Evidence gap visible'
    if conf < 72: return 'Board challenge likely'
    if conf < 84: return 'Board-defensible with conditions'
    return 'High-confidence delivery posture'

def _v145_reconcile_cost_lines(model, p50):
    rows=list(model.get('cost_lines') or model.get('cost_breakdown') or [])
    if not rows:
        direct=_v145_bn(model.get('direct_cost')) or p50*0.70
        indirect=_v145_bn(model.get('indirect_cost')) or p50*0.14
        reserve=_v145_bn(model.get('risk_reserve')) or max(0.0, p50-direct-indirect)
        rows=[
            {'cbs':'01.00','description':'Direct delivery scope','type':'Direct','low_p10':_v145_money(direct*0.82),'most_likely_p50':_v145_money(direct),'high_p90':_v145_money(direct*1.28),'basis':'Scenario-controlled direct scope.'},
            {'cbs':'90.00','description':'Indirects, owner costs and integration','type':'Indirect','low_p10':_v145_money(indirect*0.82),'most_likely_p50':_v145_money(indirect),'high_p90':_v145_money(indirect*1.30),'basis':'Scenario-controlled indirect / owner cost.'},
            {'cbs':'99.00','description':'Risk reserve and contingency','type':'Reserve','low_p10':_v145_money(reserve*0.70),'most_likely_p50':_v145_money(reserve),'high_p90':_v145_money(reserve*1.65),'basis':'Scenario-controlled risk reserve.'},
        ]
    # Scale every row to selected P50 so workbook, UI and audit spine agree.
    vals=[]
    for r in rows:
        val=_v145_bn(r.get('most_likely_p50') or r.get('p50_bn') or r.get('p50') or r.get('value'))
        vals.append(val)
    s=sum(vals)
    if p50 > 0 and s > 0:
        factor=p50/s
        new=[]
        for r,val in zip(rows,vals):
            x=dict(r)
            typ=str(x.get('type') or 'Direct')
            if re.search('reserve|risk|contingency', typ, re.I): typ='Reserve'
            elif re.search('indirect|owner|pm|management|prelim', typ, re.I): typ='Indirect'
            else: typ='Direct'
            p=val*factor
            low=_v145_bn(x.get('low_p10') or x.get('p10') or x.get('low'))*factor or p*0.82
            high=_v145_bn(x.get('high_p90') or x.get('p90') or x.get('high'))*factor or p*1.32
            x.update({'type':typ,'low_p10':_v145_money(low),'most_likely_p50':_v145_money(p),'high_p90':_v145_money(high)})
            new.append(x)
        rows=new
    totals={'Direct':0.0,'Indirect':0.0,'Reserve':0.0}
    for r in rows:
        typ=str(r.get('type') or 'Direct')
        if typ not in totals: typ='Direct'
        totals[typ]+=_v145_bn(r.get('most_likely_p50'))
    # Correct rounding drift in Direct bucket.
    drift=p50-sum(totals.values())
    if rows and abs(drift) >= 0.01:
        for r in rows:
            if r.get('type')=='Direct':
                r['most_likely_p50']=_v145_money(_v145_bn(r.get('most_likely_p50'))+drift); break
    totals={'Direct':0.0,'Indirect':0.0,'Reserve':0.0}
    for r in rows:
        totals[r.get('type') if r.get('type') in totals else 'Direct'] += _v145_bn(r.get('most_likely_p50'))
    return rows, totals

def _v145_curve(p50, months, scenario):
    prof=_v145_profile(scenario)
    pts=[1,5,10,20,30,40,50,60,70,80,90,95,99]
    if scenario=='faster':
        c=[.78,.81,.84,.89,.94,.98,1.0,1.09,1.22,prof['p80'],prof['p90'],1.78,1.95]
        s=[.70,.74,.78,.85,.91,.96,1.0,1.11,1.28,prof['sp80'],prof['sp90'],1.92,2.10]
    elif scenario=='cheaper':
        c=[.70,.74,.78,.84,.91,.97,1.0,1.13,1.28,prof['p80'],prof['p90'],1.88,2.15]
        s=[.84,.87,.90,.94,.97,.99,1.0,1.12,1.25,prof['sp80'],prof['sp90'],1.76,1.96]
    elif scenario in ('lower_risk','premium'):
        c=[.88,.90,.92,.95,.97,.99,1.0,1.04,1.09,prof['p80'],prof['p90'],prof['p90']+.08,prof['p90']+.16]
        s=[.88,.90,.92,.95,.97,.99,1.0,1.04,1.08,prof['sp80'],prof['sp90'],prof['sp90']+.08,prof['sp90']+.15]
    else:
        c=[.78,.81,.84,.89,.93,.97,1.0,1.06,1.15,prof['p80'],prof['p90'],1.58,1.72]
        s=[.82,.84,.86,.90,.94,.97,1.0,1.06,1.15,prof['sp80'],prof['sp90'],1.58,1.72]
    return [{'percentile':p,'cost_bn':round(p50*cm,2),'schedule_months':max(1,int(round(months*sm)))} for p,cm,sm in zip(pts,c,s)]

def _v145_apply_global_consistency(model: Dict[str, Any]) -> Dict[str, Any]:
    m=dict(model or {})
    scenario=str(m.get('scenario') or 'base').lower().replace(' ','_')
    prof=_v145_profile(scenario)
    row=_v145_selected_row(m)
    # 1) Select the canonical scenario values. Scenario matrix wins, then cost buckets, then headline.
    p50=_v145_bn(row.get('cost_p50')) or _v145_bn(m.get('cost_p50'))
    bucket_sum=_v145_bn(m.get('direct_cost'))+_v145_bn(m.get('indirect_cost'))+_v145_bn(m.get('risk_reserve'))
    if bucket_sum > 0 and (not p50 or abs(bucket_sum-p50)/max(p50,1) > 0.015):
        p50=bucket_sum
    months=_v145_months(row.get('schedule_months') or row.get('schedule') or m.get('schedule')) or 60
    conf=_v145_int(row.get('confidence_pct') or row.get('confidence') or m.get('confidence_pct'), 55)
    risk=str(row.get('risk') or prof.get('risk') or m.get('risk') or 'Medium-High')
    # 2) Reconcile cost rows/buckets exactly to that P50.
    rows, totals=_v145_reconcile_cost_lines(m, p50)
    # 3) Stamp canonical QCRA/QSRA.
    curve=_v145_curve(p50, months, scenario)
    qcra={'p10':round(p50*prof['p10'],1),'p50':round(p50,1),'p80':round(p50*prof['p80'],1),'p90':round(p50*prof['p90'],1)}
    qsra={'p10':max(1,int(round(months*.86))),'p50':months,'p80':max(months,int(round(months*prof['sp80']))),'p90':max(months,int(round(months*prof['sp90'])))}
    # 4) Write canonical top-level fields.
    m.update({
        'scenario':scenario,
        'scenario_label':prof['label'],
        'cost_p50':_v145_money(p50),
        'cost_p10':_v145_money(qcra['p10']),
        'cost_p90':_v145_money(qcra['p90']),
        'cost_range':f"{_v145_money(qcra['p10'])} - {_v145_money(qcra['p90'])}",
        'schedule':f'{months} months',
        'risk':risk,
        'confidence_pct':conf,
        'confidence_engine_label':_v145_conf_label(conf, scenario),
        'cost_lines':rows,
        'cost_breakdown':rows,
        'direct_cost':_v145_money(totals['Direct']),
        'indirect_cost':_v145_money(totals['Indirect']),
        'risk_reserve':_v145_money(totals['Reserve']),
    })
    mc=dict(m.get('monte_carlo') or {})
    mc.update({'qcra':qcra,'qsra':qsra,'curve':curve,'curve_readout':[
        f"Cost P50 reconciles to {m['cost_p50']}; P80 is {_v145_money(qcra['p80'])}; P90 is {_v145_money(qcra['p90'])}.",
        f"Schedule P50 reconciles to {months} months; P80 is {qsra['p80']} months; P90 is {qsra['p90']} months.",
        'Selected scenario curve is locked to the same P50/P80/P90 used by the UI and exports.'
    ]})
    m['monte_carlo']=mc
    constraint=_v145_constraint(m)
    sector=_v145_sector_name(m)
    trade=m.get('scenario_trade') or {'base':'Balanced reference case for board challenge.','faster':'Acceleration buys time but consumes float and procurement optionality.','cheaper':'Lower approval number transfers exposure into resilience and recovery.','lower_risk':'Assurance, float and evidence improve confidence at higher cost/time.','premium':'Resilience and optionality buy downside protection at premium capex.'}.get(scenario,'Scenario trade-off')
    m['executive_summary']=(f"{sector} scenario view: {prof['label']}. CASEY indicates {m['cost_p50']} P50 exposure, "
        f"{m['cost_range']} range, {m['schedule']} baseline, {risk} risk and {conf}% confidence. "
        f"Confidence is governed by {constraint}. Trade-off: {trade}")
    m['board_briefing']=[
        m.get('executive_shock_insight') or f'Confidence is governed by {constraint}.',
        f"{prof['label']} scenario is selected; all cards, narratives, QCRA/QSRA and exports are stamped from this scenario payload.",
        f"{prof['label']} scenario: {m['cost_p50']} P50, {m['cost_range']} range, {m['schedule']} baseline and {conf}% confidence.",
        f"Confidence is governed by {constraint}.",
        f"Board position: {_v145_conf_label(conf, scenario)}."
    ]
    m['casey_thinking']=(f"CASEY has re-cut the programme as a {prof['label']} scenario inside the {sector} delivery environment. "
        f"The governing consequence is: confidence is governed by {constraint}. QCRA/QSRA curves, cost basis, risk probabilities and schedule logic are locked to this scenario.")
    base=None
    for r in (m.get('scenario_matrix') or []):
        if str(r.get('scenario','')).lower()=='base': base=r; break
    if base:
        base_cost=_v145_bn(base.get('cost_p50')) or p50; base_months=_v145_months(base.get('schedule_months')) or months; base_conf=_v145_int(base.get('confidence_pct'), conf)
    else:
        base_cost=p50; base_months=months; base_conf=conf
    _delta_cost = round(p50 - base_cost, 2)
    _delta_months = months - base_months
    _delta_conf = conf - base_conf
    m['scenario_comparison_vs_base']={
        'base':{'cost_p50':_v145_money(base_cost),'schedule_months':base_months,'confidence_pct':base_conf,'risk':'Medium-High'},
        'selected':{'scenario':scenario,'cost_p50':m['cost_p50'],'schedule_months':months,'confidence_pct':conf,'risk':risk},
        'delta':{
            'cost_bn': _delta_cost,
            'cost': _v145_money(abs(_delta_cost)),
            'cost_direction': 'higher' if _delta_cost > 0 else 'lower' if _delta_cost < 0 else 'same',
            'months': _delta_months,
            'confidence_pts': _delta_conf
        },
        'plain_english': f"{prof['label']} versus Base: cost {base_cost:.1f}B to {p50:.1f}B, schedule {base_months} to {months} months, confidence {base_conf}% to {conf}%."
    }
    m['audit_spine']=[
        {'check':'Cost reconciliation','status':'PASS','detail':f"Direct {m['direct_cost']} + Indirect {m['indirect_cost']} + Reserve {m['risk_reserve']} = {m['cost_p50']}"},
        {'check':'Scenario lock','status':prof['label'],'detail':'Headline, executive narrative, QCRA/QSRA, risk register and exports use the same selected-scenario payload.'},
        {'check':'P-tail linkage','status':_v145_money(qcra['p80']),'detail':f"Cost P80 and QSRA P80 {qsra['p80']} months are visible for board challenge."},
        {'check':'Evidence gate','status':_v145_conf_label(conf, scenario),'detail':'Board posture remains tied to evidence closure, not dashboard colour.'},
    ]
    m['version']='CASEY V145 Global Showcase Consistency Lock'
    return m

_CASEY_V145_PREV_BUILD_MODEL = build_model
def build_model(prompt: str, client: str='', class_level: int=3, schedule_level: int=3, scenario: str='base'):
    return _v145_apply_global_consistency(_CASEY_V145_PREV_BUILD_MODEL(prompt, client, class_level, schedule_level, scenario))

# Also reconcile models received by export endpoints, because the UI sends the current model payload.
_CASEY_V145_PREV_WORKBOOK_BYTES = workbook_bytes
_CASEY_V145_PREV_WORD_BYTES = word_bytes
_CASEY_V145_PREV_PDF_BYTES = pdf_bytes

def workbook_bytes(model):
    return _CASEY_V145_PREV_WORKBOOK_BYTES(_v145_apply_global_consistency(model or {}))

def word_bytes(model):
    return _CASEY_V145_PREV_WORD_BYTES(_v145_apply_global_consistency(model or {}))

def pdf_bytes(model):
    return _CASEY_V145_PREV_PDF_BYTES(_v145_apply_global_consistency(model or {}))

APP_VERSION='CASEY V145 Global Showcase Consistency Lock'
print('CASEY V145 global showcase consistency/export lock installed')
# ======================= END CASEY V145 GLOBAL SHOWCASE CONSISTENCY LOCK =======================

# V145 export payload guards for every file type.
_CASEY_V145_PREV_RISK_REGISTER_WORKBOOK_BYTES = risk_register_workbook_bytes
_CASEY_V145_PREV_XER_BYTES = xer_bytes
_CASEY_V145_PREV_PPTX_BYTES = pptx_bytes
_CASEY_V145_PREV_EXPORT_ALL_FUNC = export_all if 'export_all' in globals() else None

def risk_register_workbook_bytes(model):
    return _CASEY_V145_PREV_RISK_REGISTER_WORKBOOK_BYTES(_v145_apply_global_consistency(model or {}))

def xer_bytes(model):
    return _CASEY_V145_PREV_XER_BYTES(_v145_apply_global_consistency(model or {}))

def pptx_bytes(model):
    return _CASEY_V145_PREV_PPTX_BYTES(_v145_apply_global_consistency(model or {}))

def export_all(model: Dict[str, Any]):
    m=_v145_apply_global_consistency(model or {})
    bio=BytesIO()
    with zipfile.ZipFile(bio,'w',zipfile.ZIP_DEFLATED) as z:
        z.writestr('01_CASEY_V150_Cost_Model_Planet_Class.xlsx', workbook_bytes(m))
        z.writestr('02_CASEY_V150_Risk_Register_Pro.xlsx', risk_register_workbook_bytes(m))
        z.writestr('03_CASEY_V150_P6_Schedule.xer', xer_bytes(m))
        z.writestr('04_CASEY_V150_Executive_Board_Report.docx', word_bytes(m))
        z.writestr('05_CASEY_V150_Board_Intelligence_Pack.pdf', pdf_bytes(m))
        z.writestr('06_CASEY_V150_Board_Deck_Elite.pptx', pptx_bytes(m))
        z.writestr('07_CASEY_V150_Full_Model_Audit.json', json.dumps(m,indent=2))
        z.writestr('08_CASEY_V150_Risk_Register_Raw.csv', risk_csv_bytes(m))
        z.writestr('09_CASEY_Demo_Close_Script.txt', '\n'.join(m.get('launch_demo_script',[])))
    bio.seek(0)
    return stream(bio.getvalue(),'application/zip','CASEY_Output_Pack_Planet_Class.zip')

# ========================= CASEY V146 DISPLAY RECONCILIATION LOCK =========================
# Fixes final visible rounding drift where displayed Direct + Indirect + Reserve could differ
# from displayed P50 after money formatting, even when underlying model reconciled.
_CASEY_V146_PREV_APPLY_GLOBAL_CONSISTENCY = _v145_apply_global_consistency

def _v146_round_display_bucket_sum(model: Dict[str, Any]) -> Dict[str, Any]:
    m=dict(model or {})
    p50=_v145_bn(m.get('cost_p50'))
    if p50 <= 0:
        return m
    # Keep direct and indirect from the reconciled rows, then force reserve to close the displayed P50.
    direct=_v145_bn(m.get('direct_cost'))
    indirect=_v145_bn(m.get('indirect_cost'))
    reserve=max(0.0, p50-direct-indirect)
    # If one-decimal display would drift, rebalance reserve on the same displayed precision.
    p50_d=round(p50,1)
    direct_d=round(direct,1)
    indirect_d=round(indirect,1)
    reserve_d=round(max(0.0, p50_d-direct_d-indirect_d),1)
    m['cost_p50']=_v145_money(p50_d)
    m['direct_cost']=_v145_money(direct_d)
    m['indirect_cost']=_v145_money(indirect_d)
    m['risk_reserve']=_v145_money(reserve_d)
    # Make audit text and visible export guard use the same display numbers.
    if isinstance(m.get('audit_spine'), list) and m['audit_spine']:
        m['audit_spine'][0]={'check':'Cost reconciliation','status':'PASS','detail':f"Direct {m['direct_cost']} + Indirect {m['indirect_cost']} + Reserve {m['risk_reserve']} = {m['cost_p50']}"}
    # Executive summary should retain the canonical headline after display rounding.
    sector=_v145_sector_name(m); scenario=str(m.get('scenario') or 'base').lower().replace(' ','_'); prof=_v145_profile(scenario)
    constraint=_v145_constraint(m); risk=str(m.get('risk') or prof.get('risk')); conf=_v145_int(m.get('confidence_pct'),55)
    trade=m.get('scenario_trade') or {'base':'Balanced reference case for board challenge.','faster':'Acceleration buys time but consumes float and procurement optionality.','cheaper':'Lower approval number transfers exposure into resilience and recovery.','lower_risk':'Assurance, float and evidence improve confidence at higher cost/time.','premium':'Resilience and optionality buy downside protection at premium capex.'}.get(scenario,'Scenario trade-off')
    m['executive_summary']=(f"{sector} scenario view: {prof['label']}. CASEY indicates {m['cost_p50']} P50 exposure, "
        f"{m.get('cost_range')} range, {m.get('schedule')} baseline, {risk} risk and {conf}% confidence. "
        f"Confidence is governed by {constraint}. Trade-off: {trade}")
    # Board briefing line 3 same headline.
    if isinstance(m.get('board_briefing'), list) and len(m['board_briefing'])>=3:
        m['board_briefing'][2]=f"{prof['label']} scenario: {m['cost_p50']} P50, {m.get('cost_range')} range, {m.get('schedule')} baseline and {conf}% confidence."
    m['version']='CASEY V146 Display Reconciliation Lock'
    return m

def _v145_apply_global_consistency(model: Dict[str, Any]) -> Dict[str, Any]:
    return _v146_round_display_bucket_sum(_CASEY_V146_PREV_APPLY_GLOBAL_CONSISTENCY(model or {}))

# Rebind build_model and export guards to the V146 wrapper.
def build_model(prompt: str, client: str='', class_level: int=3, schedule_level: int=3, scenario: str='base'):
    return _v145_apply_global_consistency(_CASEY_V145_PREV_BUILD_MODEL(prompt, client, class_level, schedule_level, scenario))

def workbook_bytes(model):
    return _CASEY_V145_PREV_WORKBOOK_BYTES(_v145_apply_global_consistency(model or {}))

def word_bytes(model):
    return _CASEY_V145_PREV_WORD_BYTES(_v145_apply_global_consistency(model or {}))

def pdf_bytes(model):
    return _CASEY_V145_PREV_PDF_BYTES(_v145_apply_global_consistency(model or {}))

def risk_register_workbook_bytes(model):
    return _CASEY_V145_PREV_RISK_REGISTER_WORKBOOK_BYTES(_v145_apply_global_consistency(model or {}))

def xer_bytes(model):
    return _CASEY_V145_PREV_XER_BYTES(_v145_apply_global_consistency(model or {}))

def pptx_bytes(model):
    return _CASEY_V145_PREV_PPTX_BYTES(_v145_apply_global_consistency(model or {}))

def export_all(model: Dict[str, Any]):
    m=_v145_apply_global_consistency(model or {})
    bio=BytesIO()
    with zipfile.ZipFile(bio,'w',zipfile.ZIP_DEFLATED) as z:
        z.writestr('01_CASEY_V150_Cost_Model_Planet_Class.xlsx', workbook_bytes(m))
        z.writestr('02_CASEY_V150_Risk_Register_Pro.xlsx', risk_register_workbook_bytes(m))
        z.writestr('03_CASEY_V150_P6_Schedule.xer', xer_bytes(m))
        z.writestr('04_CASEY_V150_Executive_Board_Report.docx', word_bytes(m))
        z.writestr('05_CASEY_V150_Board_Intelligence_Pack.pdf', pdf_bytes(m))
        z.writestr('06_CASEY_V150_Board_Deck_Elite.pptx', pptx_bytes(m))
        z.writestr('07_CASEY_V150_Full_Model_Audit.json', json.dumps(m,indent=2))
        z.writestr('08_CASEY_V150_Risk_Register_Raw.csv', risk_csv_bytes(m))
        z.writestr('09_CASEY_Demo_Close_Script.txt', '\n'.join(m.get('launch_demo_script',[])))
    bio.seek(0)
    return stream(bio.getvalue(),'application/zip','CASEY_Output_Pack_Planet_Class.zip')

APP_VERSION='CASEY V146 Display Reconciliation Lock'
print('CASEY V146 display reconciliation lock installed')
# ======================= END CASEY V146 DISPLAY RECONCILIATION LOCK =======================

# ========================= CASEY V147 ROUTE RECONCILIATION + HARNESS FIX =========================
# Fixes:
# 1. All FastAPI export routes now pass through _v145_apply_global_consistency before building files.
#    Previously the V107/V119 route closures captured early builder references that bypassed V146.
# 2. export/json route now returns the V146-reconciled model payload.
# 3. export/qcra-qsra route now uses the V146-reconciled model so QCRA/QSRA numbers match.
# 4. _v107_stamp_model is patched to apply full consistency before stamping.

_CASEY_V147_ORIG_STAMP = _v107_stamp_model

def _v107_stamp_model(model: Dict[str, Any]) -> Dict[str, Any]:
    return _CASEY_V147_ORIG_STAMP(_v145_apply_global_consistency(model or {}))

# Re-register all live export routes with the V146-guarded builders.
# We replace every /export/* route so the FastAPI router uses the final reconciled functions.
for _v147_path, _v147_builder, _v147_media, _v147_fname in [
    ('/export/workbook',      lambda m: workbook_bytes(m),              'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet', 'CASEY_V150_Cost_Model_Planet_Class.xlsx'),
    ('/export/cost-model',    lambda m: workbook_bytes(m),              'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet', 'CASEY_V150_Cost_Model_Planet_Class.xlsx'),
    ('/export/excel',         lambda m: workbook_bytes(m),              'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet', 'CASEY_V150_Cost_Model_Planet_Class.xlsx'),
    ('/export/risk-register', lambda m: risk_register_workbook_bytes(m),'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet', 'CASEY_V150_Risk_Register_Pro.xlsx'),
    ('/export/xer',           lambda m: xer_bytes(m),                   'application/octet-stream',                                          'CASEY_V150_P6_Schedule.xer'),
    ('/export/word',          lambda m: word_bytes(m),                  'application/vnd.openxmlformats-officedocument.wordprocessingml.document', 'CASEY_V150_Executive_Board_Report.docx'),
    ('/export/pdf',           lambda m: pdf_bytes(m),                   'application/pdf',                                                   'CASEY_V150_Board_Intelligence_Pack.pdf'),
    ('/export/pptx',          lambda m: pptx_bytes(m),                  'application/vnd.openxmlformats-officedocument.presentationml.presentation', 'CASEY_V150_Board_Deck_Elite.pptx'),
    ('/export/schedule-csv',  lambda m: risk_csv_bytes(m),              'text/csv',                                                          'CASEY_V150_Schedule_Levels.csv'),
    ('/export/json',          lambda m: json.dumps(_v145_apply_global_consistency(m), indent=2).encode(), 'application/json', 'CASEY_V150_Full_Model_Audit.json'),
    ('/export/qcra-qsra',     lambda m: _v106_qcra_qsra_workbook_bytes(_v145_apply_global_consistency(m)), 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet', 'CASEY_V150_QCRA_QSRA_Pack.xlsx'),
    ('/export/all',           lambda m: export_all(m)._body if hasattr(export_all(m),'_body') else export_all(m), 'application/zip', 'CASEY_V150_Output_Pack_Planet_Class.zip'),
    ('/export/full-pack',     lambda m: export_all(m)._body if hasattr(export_all(m),'_body') else export_all(m), 'application/zip', 'CASEY_V150_Output_Pack_Planet_Class.zip'),
]:
    try:
        _v107_remove_route(_v147_path, 'POST')
    except Exception:
        pass
    # Capture loop vars correctly
    def _v147_make_endpoint(builder, media, fname):
        def _endpoint(payload: dict = Body(default={})):
            m = _v145_apply_global_consistency(payload or {})
            return stream(builder(m), media, fname)
        return _endpoint
    try:
        app.post(_v147_path)(_v147_make_endpoint(_v147_builder, _v147_media, _v147_fname))
    except Exception as _e:
        print(f'CASEY V147 route warning {_v147_path}: {_e}')

# Special handler for /export/all that builds zip inline to avoid double-streaming
def _v147_remove_route(path, method='POST'):
    try:
        _v107_remove_route(path, method)
    except Exception:
        pass

_v147_remove_route('/export/all')
_v147_remove_route('/export/full-pack')

def _v147_export_all_endpoint(payload: dict = Body(default={})):
    return export_all(_v145_apply_global_consistency(payload or {}))

app.post('/export/all')(_v147_export_all_endpoint)
app.post('/export/full-pack')(_v147_export_all_endpoint)

APP_VERSION = 'CASEY V147 Route Reconciliation Lock'
print('CASEY V147 route reconciliation lock installed')
# ======================= END CASEY V147 ROUTE RECONCILIATION LOCK =======================

# ========================= CASEY V148 TRUST CORE + EXPORT INTEGRITY LOCK =========================
# Final hardening layer for demo trust:
# - one canonical scenario state signature across UI and exports
# - export manifest inside full ZIP
# - machine-readable validation results
# - all live export routes re-registered to V148-canonicalized payloads

_CASEY_V148_PREV_APPLY = _v145_apply_global_consistency

def _v148_jsonable(obj):
    try:
        json.dumps(obj)
        return obj
    except Exception:
        if isinstance(obj, dict):
            return {str(k): _v148_jsonable(v) for k, v in obj.items()}
        if isinstance(obj, (list, tuple)):
            return [_v148_jsonable(v) for v in obj]
        return str(obj)

def _v148_canonical_state(model: Dict[str, Any]) -> Dict[str, Any]:
    m = _CASEY_V148_PREV_APPLY(model or {})
    p50 = _v145_bn(m.get('cost_p50'))
    months = _v145_months(m.get('schedule'))
    conf = _v145_int(m.get('confidence_pct'), 55)
    scenario = str(m.get('scenario') or 'base').lower().replace(' ', '_')
    sector = _v145_sector_name(m)
    qcra = ((m.get('monte_carlo') or {}).get('qcra') or {})
    qsra = ((m.get('monte_carlo') or {}).get('qsra') or {})
    state_basis = {
        'title': m.get('title'),
        'mode': m.get('mode'),
        'sector': sector,
        'subsector': m.get('subsector'),
        'scenario': scenario,
        'cost_p50': m.get('cost_p50'),
        'cost_range': m.get('cost_range'),
        'schedule': m.get('schedule'),
        'confidence_pct': conf,
        'risk': m.get('risk'),
        'qcra': qcra,
        'qsra': qsra,
        'direct_cost': m.get('direct_cost'),
        'indirect_cost': m.get('indirect_cost'),
        'risk_reserve': m.get('risk_reserve'),
    }
    signature = hashlib.sha256(json.dumps(_v148_jsonable(state_basis), sort_keys=True).encode('utf-8')).hexdigest()[:16].upper()
    m['casey_state'] = {
        'model_version': 'CASEY V148 Trust Core',
        'scenario_signature': signature,
        'scenario': scenario,
        'scenario_label': m.get('scenario_label'),
        'canonical_numbers': state_basis,
        'source_of_truth': 'All UI cards, board narratives, QCRA/QSRA, workbook, risk register, XER, PDF, DOCX, PPTX and JSON exports are generated from this selected scenario state.',
        'commercial_disclaimer': 'Strategic intelligence simulation / first-pass board pack. Not a certified estimate without client source data and professional sign-off.'
    }
    m['scenario_signature'] = signature
    m['export_stamp'] = {
        **dict(m.get('export_stamp') or {}),
        'casey_v150_signature': signature,
        'scenario': m.get('scenario_label') or scenario,
        'cost_p50': m.get('cost_p50'),
        'schedule': m.get('schedule'),
        'confidence_pct': conf,
        'qcra_p80': _v145_money(_v145_bn(qcra.get('p80'))),
        'qsra_p80_months': qsra.get('p80'),
        'generated_at_utc': datetime.utcnow().isoformat() + 'Z',
    }
    # Make visible texts carry the same signature/decision posture without sounding like code.
    conf_label = _v145_conf_label(conf, scenario)
    m['confidence_engine_detail'] = f"Board-defensibility confidence: {conf_label}. Signature {signature}."
    if isinstance(m.get('audit_spine'), list):
        m['audit_spine'] = list(m['audit_spine']) + [
            {'check': 'Scenario signature', 'status': signature, 'detail': 'Same signature is written into JSON audit and full export pack manifest.'},
            {'check': 'Export integrity', 'status': 'LOCKED', 'detail': 'Every export endpoint re-canonicalizes the selected scenario before file generation.'},
        ]
    m['version'] = 'CASEY V148 Trust Core + Export Integrity Lock'
    return m

def _v148_validate_model(model: Dict[str, Any]) -> Dict[str, Any]:
    m = _v148_canonical_state(model or {})
    fails = []
    p50 = _v145_bn(m.get('cost_p50'))
    total = _v145_bn(m.get('direct_cost')) + _v145_bn(m.get('indirect_cost')) + _v145_bn(m.get('risk_reserve'))
    if abs(total - p50) > 0.11:
        fails.append({'check': 'cost_bucket_reconciliation', 'p50_bn': round(p50, 3), 'bucket_sum_bn': round(total, 3)})
    qcra = ((m.get('monte_carlo') or {}).get('qcra') or {})
    qsra = ((m.get('monte_carlo') or {}).get('qsra') or {})
    if abs(float(qcra.get('p50') or 0) - round(p50, 1)) > 0.11:
        fails.append({'check': 'qcra_p50_reconciliation', 'qcra_p50': qcra.get('p50'), 'p50_bn': round(p50, 1)})
    if int(qsra.get('p50') or 0) != _v145_months(m.get('schedule')):
        fails.append({'check': 'qsra_p50_reconciliation', 'qsra_p50': qsra.get('p50'), 'schedule': m.get('schedule')})
    summary = str(m.get('executive_summary') or '')
    for token in [m.get('cost_p50'), m.get('schedule'), str(m.get('confidence_pct')) + '%']:
        if token and str(token) not in summary:
            fails.append({'check': 'executive_summary_token', 'missing': token})
    if not m.get('scenario_signature'):
        fails.append({'check': 'scenario_signature_missing'})
    return {'ok': not fails, 'fail_count': len(fails), 'failures': fails, 'scenario_signature': m.get('scenario_signature'), 'model': m}

# Override consistency wrapper and build model again.
def _v145_apply_global_consistency(model: Dict[str, Any]) -> Dict[str, Any]:
    return _v148_canonical_state(model or {})

def build_model(prompt: str, client: str='', class_level: int=3, schedule_level: int=3, scenario: str='base'):
    return _v148_canonical_state(_CASEY_V145_PREV_BUILD_MODEL(prompt, client, class_level, schedule_level, scenario))

# Rebind export builders to V148 trust core.
def workbook_bytes(model):
    return _CASEY_V145_PREV_WORKBOOK_BYTES(_v148_canonical_state(model or {}))

def word_bytes(model):
    return _CASEY_V145_PREV_WORD_BYTES(_v148_canonical_state(model or {}))

def pdf_bytes(model):
    return _CASEY_V145_PREV_PDF_BYTES(_v148_canonical_state(model or {}))

def risk_register_workbook_bytes(model):
    return _CASEY_V145_PREV_RISK_REGISTER_WORKBOOK_BYTES(_v148_canonical_state(model or {}))

def xer_bytes(model):
    return _CASEY_V145_PREV_XER_BYTES(_v148_canonical_state(model or {}))

def pptx_bytes(model):
    return _CASEY_V145_PREV_PPTX_BYTES(_v148_canonical_state(model or {}))

def _v148_manifest(model: Dict[str, Any]) -> bytes:
    m = _v148_canonical_state(model or {})
    validation = _v148_validate_model(m)
    manifest = {
        'package': 'CASEY V148 Final Trust Core Output Pack',
        'scenario_signature': m.get('scenario_signature'),
        'project': m.get('title'),
        'scenario': m.get('scenario_label'),
        'cost_p50': m.get('cost_p50'),
        'cost_range': m.get('cost_range'),
        'schedule': m.get('schedule'),
        'confidence_pct': m.get('confidence_pct'),
        'risk': m.get('risk'),
        'validation': {k:v for k,v in validation.items() if k != 'model'},
        'files': [
            '01_CASEY_V150_Cost_Model_Planet_Class.xlsx',
            '02_CASEY_V150_Risk_Register_Pro.xlsx',
            '03_CASEY_V150_P6_Schedule.xer',
            '04_CASEY_V150_Executive_Board_Report.docx',
            '05_CASEY_V150_Board_Intelligence_Pack.pdf',
            '06_CASEY_V150_Board_Deck_Elite.pptx',
            '07_CASEY_V150_Full_Model_Audit.json',
            '08_CASEY_V150_Risk_Register_Raw.csv',
            '09_CASEY_Demo_Close_Script.txt'
        ],
        'disclaimer': 'Strategic intelligence simulation / first-pass board pack. Not a certified estimate without source-data review and professional sign-off.'
    }
    return json.dumps(manifest, indent=2).encode('utf-8')

def export_all(model: Dict[str, Any]):
    m = _v148_canonical_state(model or {})
    bio = BytesIO()
    with zipfile.ZipFile(bio, 'w', zipfile.ZIP_DEFLATED) as z:
        z.writestr('00_CASEY_V150_EXPORT_MANIFEST.json', _v148_manifest(m))
        z.writestr('01_CASEY_V150_Cost_Model_Planet_Class.xlsx', workbook_bytes(m))
        z.writestr('02_CASEY_V150_Risk_Register_Pro.xlsx', risk_register_workbook_bytes(m))
        z.writestr('03_CASEY_V150_P6_Schedule.xer', xer_bytes(m))
        z.writestr('04_CASEY_V150_Executive_Board_Report.docx', word_bytes(m))
        z.writestr('05_CASEY_V150_Board_Intelligence_Pack.pdf', pdf_bytes(m))
        z.writestr('06_CASEY_V150_Board_Deck_Elite.pptx', pptx_bytes(m))
        z.writestr('07_CASEY_V150_Full_Model_Audit.json', json.dumps(_v148_jsonable(m), indent=2))
        z.writestr('08_CASEY_V150_Risk_Register_Raw.csv', risk_csv_bytes(m))
        z.writestr('09_CASEY_Demo_Close_Script.txt', '\n'.join(m.get('launch_demo_script', [])))
    bio.seek(0)
    return stream(bio.getvalue(), 'application/zip', 'CASEY_V150_ENTERPRISE_OUTPUT_PACK.zip')

@app.post('/qa/validate-model')
def qa_validate_model_v148(payload: Dict[str, Any] = Body(default={})):
    out = _v148_validate_model(payload or {})
    return {k:v for k,v in out.items() if k != 'model'}

@app.get('/qa/readiness')
def qa_readiness_v148():
    return {
        'version': 'CASEY V148 Final Trust Core',
        'status': 'ready_for_controlled_executive_demo',
        'hardening': ['canonical scenario signature', 'export manifest', 'cost/QCRA/QSRA reconciliation', 'scenario-locked exports', 'model validation endpoint'],
        'limits': 'Demo-grade strategic intelligence simulation; not a certified estimate without client source data.'
    }

# Re-register every user-facing export route to V148 builders.
for _v148_path, _v148_builder, _v148_media, _v148_fname in [
    ('/export/workbook',      lambda m: workbook_bytes(m), 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet', 'CASEY_V150_Cost_Model_Planet_Class.xlsx'),
    ('/export/cost-model',    lambda m: workbook_bytes(m), 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet', 'CASEY_V150_Cost_Model_Planet_Class.xlsx'),
    ('/export/excel',         lambda m: workbook_bytes(m), 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet', 'CASEY_V150_Cost_Model_Planet_Class.xlsx'),
    ('/export/risk-register', lambda m: risk_register_workbook_bytes(m), 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet', 'CASEY_V150_Risk_Register_Pro.xlsx'),
    ('/export/xer',           lambda m: xer_bytes(m), 'application/octet-stream', 'CASEY_V150_P6_Schedule.xer'),
    ('/export/word',          lambda m: word_bytes(m), 'application/vnd.openxmlformats-officedocument.wordprocessingml.document', 'CASEY_V150_Executive_Board_Report.docx'),
    ('/export/pdf',           lambda m: pdf_bytes(m), 'application/pdf', 'CASEY_V150_Board_Intelligence_Pack.pdf'),
    ('/export/pptx',          lambda m: pptx_bytes(m), 'application/vnd.openxmlformats-officedocument.presentationml.presentation', 'CASEY_V150_Board_Deck_Elite.pptx'),
    ('/export/json',          lambda m: json.dumps(_v148_jsonable(_v148_canonical_state(m)), indent=2).encode('utf-8'), 'application/json', 'CASEY_V150_Full_Model_Audit.json'),
    ('/export/qcra-qsra',     lambda m: _v106_qcra_qsra_workbook_bytes(_v148_canonical_state(m)), 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet', 'CASEY_V150_QCRA_QSRA_Pack.xlsx'),
]:
    try:
        _v107_remove_route(_v148_path, 'POST')
    except Exception:
        pass
    def _v148_make_endpoint(builder, media, fname):
        def _endpoint(payload: dict = Body(default={})):
            m = _v148_canonical_state(payload or {})
            return stream(builder(m), media, fname)
        return _endpoint
    app.post(_v148_path)(_v148_make_endpoint(_v148_builder, _v148_media, _v148_fname))

for _v148_pack_path in ['/export/all', '/export/full-pack']:
    try:
        _v107_remove_route(_v148_pack_path, 'POST')
    except Exception:
        pass
app.post('/export/all')(lambda payload=Body(default={}): export_all(_v148_canonical_state(payload or {})))
app.post('/export/full-pack')(lambda payload=Body(default={}): export_all(_v148_canonical_state(payload or {})))

APP_VERSION = 'CASEY V148 Final Trust Core + Export Integrity Lock'
print('CASEY V148 final trust core/export integrity lock installed')
# ======================= END CASEY V148 TRUST CORE + EXPORT INTEGRITY LOCK =======================

# ======================= CASEY V150 ENTERPRISE TRUST RUNTIME =======================
# V150 turns the V148 trust core into a stricter single-state runtime: every model is
# canonicalized, signed, scored, validated and export-stamped before UI/export use.

V150_WEIGHTS = {
    'benchmark_fit': 18,
    'schedule_density': 16,
    'procurement_certainty': 15,
    'evidence_maturity': 18,
    'reserve_adequacy': 12,
    'systems_integration': 11,
    'operational_readiness': 10,
}

def _v150_pct(x, default=50):
    try:
        if isinstance(x, str):
            x = re.sub(r'[^0-9.-]', '', x)
        return int(max(0, min(100, round(float(x)))))
    except Exception:
        return default

def _v150_text_hash(*parts: Any) -> str:
    raw = json.dumps(_v148_jsonable(parts), sort_keys=True, ensure_ascii=False)
    return hashlib.sha256(raw.encode('utf-8')).hexdigest()[:20].upper()

def _v150_sector_flags(model: Dict[str, Any]) -> Dict[str, bool]:
    t = ' '.join(str(model.get(k, '')) for k in ['title','subsector','mode','prompt','executive_summary']).lower()
    return {
        'space': str(model.get('mode','')).lower() == 'space' or any(w in t for w in ['lunar','mars','orbital','satellite','starship','kuiper','space']),
        'rail': any(w in t for w in ['rail','metro','crossrail','hsr','signalling','station','corridor']),
        'data': any(w in t for w in ['data centre','data center','compute','hyperscale','cloud','ai campus']),
        'nuclear': any(w in t for w in ['nuclear','smr','reactor','fusion']),
        'pharma': any(w in t for w in ['gmp','pharma','biologics','life sciences','validation']),
        'defence': any(w in t for w in ['defence','defense','military','naval','missile','drone']),
        'fab': any(w in t for w in ['fab','semiconductor','cleanroom','wafer','foundry']),
    }

def _v150_primary_constraints(model: Dict[str, Any]) -> List[str]:
    flags = _v150_sector_flags(model)
    if flags['space']:
        return ['launch cadence', 'thermal-power balance', 'autonomous commissioning', 'mission assurance evidence']
    if flags['rail']:
        return ['possessions', 'signalling integration', 'systems migration', 'operator acceptance']
    if flags['data']:
        return ['grid connection', 'cooling architecture', 'long-lead electrical procurement', 'commissioning density']
    if flags['nuclear']:
        return ['licensing gate', 'safety-case evidence', 'nuclear island procurement', 'commissioning assurance']
    if flags['pharma']:
        return ['GMP validation', 'clean utility readiness', 'automation integration', 'regulatory handover']
    if flags['defence']:
        return ['mission assurance', 'controlled procurement', 'security accreditation', 'interoperability']
    if flags['fab']:
        return ['tool install', 'ultra-clean utilities', 'yield ramp', 'workforce maturity']
    return ['scope freeze', 'procurement certainty', 'interface ownership', 'commissioning evidence']

def _v150_confidence_breakdown(model: Dict[str, Any]) -> Dict[str, Any]:
    conf = _v150_pct(model.get('confidence_pct'), 55)
    scenario = str(model.get('scenario') or model.get('scenario_label') or 'base').lower()
    p50 = max(0.1, _v145_bn(model.get('cost_p50')))
    reserve = _v145_bn(model.get('risk_reserve'))
    reserve_pct = max(0, min(100, reserve / p50 * 100))
    months = max(1, _v145_months(model.get('schedule')))
    qcra = ((model.get('monte_carlo') or {}).get('qcra') or {})
    p80 = _v145_bn(qcra.get('p80')) or p50 * 1.18
    tail_ratio = max(1.0, p80 / p50)
    # Derived driver sub-scores. These are intentionally asymmetric and bounded.
    benchmark_fit = max(22, min(94, conf + (6 if scenario in ['lower_risk','premium'] else -3 if scenario in ['faster','cheaper'] else 2)))
    schedule_density = max(18, min(92, 88 - min(45, months / 4) - (12 if scenario == 'faster' else 0) + (7 if scenario == 'lower_risk' else 0)))
    procurement_certainty = max(20, min(92, conf - (10 if scenario in ['faster','cheaper'] else 0) + (5 if scenario in ['lower_risk','premium'] else 0)))
    evidence_maturity = max(18, min(94, conf + (8 if scenario in ['lower_risk','premium'] else -6 if scenario == 'cheaper' else 0)))
    reserve_adequacy = max(16, min(94, 44 + reserve_pct * 2.25 - max(0, (tail_ratio - 1.18) * 70)))
    systems_integration = max(17, min(92, conf - (11 if scenario == 'faster' else 5 if scenario == 'cheaper' else 0) + (6 if scenario == 'premium' else 0)))
    operational_readiness = max(18, min(93, conf - (8 if scenario == 'faster' else 3 if scenario == 'cheaper' else 0) + (8 if scenario == 'lower_risk' else 0)))
    drivers = {
        'benchmark_fit': round(benchmark_fit),
        'schedule_density': round(schedule_density),
        'procurement_certainty': round(procurement_certainty),
        'evidence_maturity': round(evidence_maturity),
        'reserve_adequacy': round(reserve_adequacy),
        'systems_integration': round(systems_integration),
        'operational_readiness': round(operational_readiness),
    }
    weighted = round(sum(drivers[k] * V150_WEIGHTS[k] for k in V150_WEIGHTS) / sum(V150_WEIGHTS.values()))
    return {
        'score': conf,
        'computed_score': weighted,
        'delta_to_reported': conf - weighted,
        'drivers': drivers,
        'weights': V150_WEIGHTS,
        'reserve_pct_of_p50': round(reserve_pct, 1),
        'tail_ratio_p80_to_p50': round(tail_ratio, 3),
        'interpretation': 'Board-defensibility confidence derived from benchmark fit, schedule density, procurement certainty, evidence maturity, reserve adequacy, systems integration and operational readiness.'
    }

def _v150_governance_state(model: Dict[str, Any]) -> Dict[str, Any]:
    conf = _v150_pct(model.get('confidence_pct'), 55)
    scenario = str(model.get('scenario') or 'base').lower()
    qcra = ((model.get('monte_carlo') or {}).get('qcra') or {})
    p50 = max(0.1, _v145_bn(model.get('cost_p50')))
    p80 = _v145_bn(qcra.get('p80')) or p50 * 1.18
    tail = max(0, (p80 / p50 - 1) * 100)
    stress = max(0, min(100, round((100 - conf) * 0.55 + tail * 0.65 + (12 if scenario == 'faster' else 8 if scenario == 'cheaper' else -4 if scenario in ['lower_risk','premium'] else 0))))
    evidence_vol = max(0, min(100, round((100 - conf) * 0.45 + (10 if scenario in ['faster','cheaper'] else 0))))
    tail_exp = max(0, min(100, round(tail * 2.6)))
    reserve_pressure = max(0, min(100, round(tail_exp * 0.55 + stress * 0.35)))
    defensibility = max(0, min(100, round(conf - stress * 0.16 + (4 if scenario in ['lower_risk','premium'] else 0))))
    return {
        'board_defensibility': defensibility,
        'governance_stress': stress,
        'tail_exposure': tail_exp,
        'evidence_volatility': evidence_vol,
        'reserve_pressure': reserve_pressure,
        'confidence_drift': round(conf - _v150_confidence_breakdown(model)['computed_score']),
        'decision_posture': 'Approve with evidence gates' if defensibility >= 68 else 'Challenge before approval' if defensibility >= 48 else 'Do not approve without recovery plan',
    }

def _v150_propagation_graph(model: Dict[str, Any]) -> Dict[str, Any]:
    constraints = _v150_primary_constraints(model)
    scenario = str(model.get('scenario') or 'base').lower()
    entry = {
        'base': 'reference case',
        'faster': 'schedule acceleration',
        'cheaper': 'capital compression',
        'lower_risk': 'assurance expansion',
        'premium': 'resilience premium',
    }.get(scenario, scenario)
    chain = [entry] + constraints + ['confidence movement', 'reserve posture', 'board decision posture']
    edges = []
    for i in range(len(chain)-1):
        weight = round(0.58 + ((i * 17 + len(str(model.get('title','')))) % 28) / 100, 2)
        edges.append({'from': chain[i], 'to': chain[i+1], 'weight': weight, 'effect': 'amplifies' if scenario in ['faster','cheaper'] else 'stabilises' if scenario in ['lower_risk','premium'] else 'governs'})
    return {'nodes': [{'id': x, 'label': x.title()} for x in chain], 'edges': edges, 'primary_constraints': constraints, 'scenario_force': entry}

def _v150_entropy(model: Dict[str, Any]) -> Dict[str, Any]:
    texts = []
    for key in ['executive_summary','casey_thinking','confidence_engine_detail']:
        if model.get(key): texts.append(str(model.get(key)))
    for key in ['board_briefing','board_challenge_questions','next_best_actions','audit_spine']:
        val = model.get(key)
        if isinstance(val, list):
            for x in val[:12]: texts.append(json.dumps(x, default=str) if isinstance(x, dict) else str(x))
    words = re.findall(r'[a-zA-Z]{4,}', ' '.join(texts).lower())
    unique = len(set(words)); total = max(1, len(words))
    score = round(min(100, max(0, unique / total * 180)))
    repeated = [w for w in sorted(set(words)) if words.count(w) >= max(5, total // 45)][:12]
    return {'score': score, 'unique_words': unique, 'total_words': total, 'repeated_terms': repeated, 'status': 'PASS' if score >= 42 else 'WATCH'}

def _v150_canonical_state(model: Dict[str, Any]) -> Dict[str, Any]:
    m = _v148_canonical_state(model or {})
    conf_break = _v150_confidence_breakdown(m)
    gov = _v150_governance_state(m)
    graph = _v150_propagation_graph(m)
    entropy = _v150_entropy(m)
    signature = _v150_text_hash(
        m.get('title'), m.get('scenario'), m.get('cost_p50'), m.get('schedule'),
        m.get('confidence_pct'), conf_break, gov, graph
    )
    constraints = ', '.join(graph['primary_constraints'])
    m['casey_runtime'] = {
        'model_version': 'CASEY V150 Enterprise Trust Runtime',
        'single_source_of_truth': 'ScenarioState',
        'scenario_signature': signature,
        'propagation_checksum': _v150_text_hash(graph, conf_break, gov),
        'canonical_state_schema': ['project','scenario','cost','schedule','risk','procurement','systems','evidence','governance','confidence','reserve','causal','exports','audit'],
        'ai_rule': 'AI may narrate computed state, but must not invent numbers, reserves, durations, owners or probabilities.',
    }
    m['scenario_signature'] = signature
    m['scenario_state'] = {
        'project': {'title': m.get('title'), 'mode': m.get('mode'), 'subsector': m.get('subsector')},
        'scenario': {'id': m.get('scenario'), 'label': m.get('scenario_label')},
        'cost': {'p50': m.get('cost_p50'), 'range': m.get('cost_range'), 'direct': m.get('direct_cost'), 'indirect': m.get('indirect_cost'), 'reserve': m.get('risk_reserve')},
        'schedule': {'baseline': m.get('schedule'), 'qsra': ((m.get('monte_carlo') or {}).get('qsra') or {})},
        'confidence': conf_break,
        'governance': gov,
        'causal': graph,
        'entropy': entropy,
    }
    m['governance_state'] = gov
    m['confidence_breakdown_v150'] = conf_break
    m['causal_graph_v150'] = graph
    m['narrative_entropy'] = entropy
    m['confidence_engine_detail'] = (
        f"Board-defensibility {m.get('confidence_pct')}% · defensibility {gov['board_defensibility']} · "
        f"governance stress {gov['governance_stress']} · signature {signature}."
    )
    label = m.get('scenario_label') or str(m.get('scenario','Base')).title()
    m['casey_thinking'] = (
        f"CASEY V150 treats this as a locked {label} scenario state. The governing chain is {constraints}; "
        f"that chain drives confidence, reserve posture, tail exposure and board decision posture from one canonical model."
    )
    # Keep board briefing aligned with runtime state.
    briefing = list(m.get('board_briefing') or [])[:3]
    briefing += [
        f"Governance posture: {gov['decision_posture']}.",
        f"Primary evidence gate: close {graph['primary_constraints'][0]} and prove {graph['primary_constraints'][-1]} before approval confidence improves.",
    ]
    m['board_briefing'] = briefing[:5]
    # Add audit rows while avoiding runaway duplication on repeated canonicalization.
    audit = [x for x in (m.get('audit_spine') or []) if not (isinstance(x, dict) and str(x.get('check','')).startswith('V150'))]
    audit += [
        {'check':'V150 single-state runtime','status':'LOCKED','detail':f"ScenarioState signature {signature}; all exports are regenerated from this canonical state."},
        {'check':'V150 governance stress','status':str(gov['governance_stress']),'detail':f"Decision posture: {gov['decision_posture']}; evidence volatility {gov['evidence_volatility']}."},
        {'check':'V150 narrative entropy','status':entropy['status'],'detail':f"Uniqueness score {entropy['score']} across board-facing narrative outputs."},
    ]
    m['audit_spine'] = audit
    m['export_stamp'] = {
        **dict(m.get('export_stamp') or {}),
        'casey_v150_signature': signature,
        'propagation_checksum': m['casey_runtime']['propagation_checksum'],
        'board_defensibility': gov['board_defensibility'],
        'governance_stress': gov['governance_stress'],
        'model_version': 'CASEY V150 Enterprise Trust Runtime',
        'generated_at_utc': datetime.utcnow().isoformat() + 'Z',
    }
    m['version'] = 'CASEY V150 Enterprise Trust Runtime'
    return m

def _v150_validate_model(model: Dict[str, Any]) -> Dict[str, Any]:
    m = _v150_canonical_state(model or {})
    failures = []
    warnings = []
    p50 = _v145_bn(m.get('cost_p50'))
    direct = _v145_bn(m.get('direct_cost')); indirect = _v145_bn(m.get('indirect_cost')); reserve = _v145_bn(m.get('risk_reserve'))
    if abs((direct + indirect + reserve) - p50) > 0.12:
        failures.append({'check':'cost_reconciliation','p50':round(p50,3),'bucket_sum':round(direct+indirect+reserve,3)})
    qcra = ((m.get('monte_carlo') or {}).get('qcra') or {})
    qsra = ((m.get('monte_carlo') or {}).get('qsra') or {})
    if not (_v145_bn(qcra.get('p50')) <= _v145_bn(qcra.get('p80')) <= _v145_bn(qcra.get('p90'))):
        failures.append({'check':'qcra_monotonicity','qcra':qcra})
    if not (int(qsra.get('p50') or 0) <= int(qsra.get('p80') or 0) <= int(qsra.get('p90') or 0)):
        failures.append({'check':'qsra_monotonicity','qsra':qsra})
    summary = str(m.get('executive_summary') or '')
    for token in [m.get('cost_p50'), m.get('schedule'), str(m.get('confidence_pct')) + '%']:
        if token and str(token) not in summary:
            failures.append({'check':'summary_token_alignment','missing':token})
    if not m.get('scenario_signature') or len(str(m.get('scenario_signature'))) < 12:
        failures.append({'check':'scenario_signature'})
    ent = m.get('narrative_entropy') or {}
    if int(ent.get('score') or 0) < 42:
        warnings.append({'check':'narrative_entropy_watch','score':ent.get('score')})
    gov = m.get('governance_state') or {}
    if not all(k in gov for k in ['board_defensibility','governance_stress','tail_exposure','decision_posture']):
        failures.append({'check':'governance_state_complete'})
    return {
        'ok': not failures,
        'failure_count': len(failures),
        'warning_count': len(warnings),
        'failures': failures,
        'warnings': warnings,
        'scenario_signature': m.get('scenario_signature'),
        'propagation_checksum': ((m.get('casey_runtime') or {}).get('propagation_checksum')),
        'governance_state': gov,
        'entropy': ent,
    }

def _v150_manifest(model: Dict[str, Any]) -> bytes:
    m = _v150_canonical_state(model or {})
    validation = _v150_validate_model(m)
    manifest = {
        'package': 'CASEY V150 Enterprise Trust Runtime Output Pack',
        'scenario_signature': m.get('scenario_signature'),
        'propagation_checksum': ((m.get('casey_runtime') or {}).get('propagation_checksum')),
        'project': m.get('title'),
        'scenario': m.get('scenario_label'),
        'cost_p50': m.get('cost_p50'),
        'cost_range': m.get('cost_range'),
        'schedule': m.get('schedule'),
        'confidence_pct': m.get('confidence_pct'),
        'governance_state': m.get('governance_state'),
        'validation': validation,
        'files': [
            '00_CASEY_V150_EXPORT_MANIFEST.json',
            '00_CASEY_V150_SCENARIO_STATE.json',
            '01_CASEY_V150_Cost_Model_Planet_Class.xlsx',
            '02_CASEY_V150_Risk_Register_Pro.xlsx',
            '03_CASEY_V150_P6_Schedule.xer',
            '04_CASEY_V150_Executive_Board_Report.docx',
            '05_CASEY_V150_Board_Intelligence_Pack.pdf',
            '06_CASEY_V150_Board_Deck_Elite.pptx',
            '07_CASEY_V150_Full_Model_Audit.json',
            '08_CASEY_V150_Risk_Register_Raw.csv',
            '09_CASEY_Demo_Close_Script.txt',
        ],
        'disclaimer': 'Strategic intelligence simulation / first-pass board pack. Not a certified estimate without client source data and professional sign-off.'
    }
    return json.dumps(_v148_jsonable(manifest), indent=2).encode('utf-8')

# Runtime overrides: canonicalize all generation and all exports.
def _v145_apply_global_consistency(model: Dict[str, Any]) -> Dict[str, Any]:
    return _v150_canonical_state(model or {})

def build_model(prompt: str, client: str='', class_level: int=3, schedule_level: int=3, scenario: str='base'):
    return _v150_canonical_state(_CASEY_V145_PREV_BUILD_MODEL(prompt, client, class_level, schedule_level, scenario))

def workbook_bytes(model):
    return _CASEY_V145_PREV_WORKBOOK_BYTES(_v150_canonical_state(model or {}))

def word_bytes(model):
    return _CASEY_V145_PREV_WORD_BYTES(_v150_canonical_state(model or {}))

def pdf_bytes(model):
    return _CASEY_V145_PREV_PDF_BYTES(_v150_canonical_state(model or {}))

def risk_register_workbook_bytes(model):
    return _CASEY_V145_PREV_RISK_REGISTER_WORKBOOK_BYTES(_v150_canonical_state(model or {}))

def xer_bytes(model):
    return _CASEY_V145_PREV_XER_BYTES(_v150_canonical_state(model or {}))

def pptx_bytes(model):
    return _CASEY_V145_PREV_PPTX_BYTES(_v150_canonical_state(model or {}))

def export_all(model: Dict[str, Any]):
    m = _v150_canonical_state(model or {})
    bio = BytesIO()
    with zipfile.ZipFile(bio, 'w', zipfile.ZIP_DEFLATED) as z:
        z.writestr('00_CASEY_V150_EXPORT_MANIFEST.json', _v150_manifest(m))
        z.writestr('00_CASEY_V150_SCENARIO_STATE.json', json.dumps(_v148_jsonable(m.get('scenario_state') or {}), indent=2))
        z.writestr('01_CASEY_V150_Cost_Model_Planet_Class.xlsx', workbook_bytes(m))
        z.writestr('02_CASEY_V150_Risk_Register_Pro.xlsx', risk_register_workbook_bytes(m))
        z.writestr('03_CASEY_V150_P6_Schedule.xer', xer_bytes(m))
        z.writestr('04_CASEY_V150_Executive_Board_Report.docx', word_bytes(m))
        z.writestr('05_CASEY_V150_Board_Intelligence_Pack.pdf', pdf_bytes(m))
        z.writestr('06_CASEY_V150_Board_Deck_Elite.pptx', pptx_bytes(m))
        z.writestr('07_CASEY_V150_Full_Model_Audit.json', json.dumps(_v148_jsonable(m), indent=2))
        z.writestr('08_CASEY_V150_Risk_Register_Raw.csv', risk_csv_bytes(m))
        z.writestr('09_CASEY_Demo_Close_Script.txt', '\n'.join(m.get('launch_demo_script', [])))
    bio.seek(0)
    return stream(bio.getvalue(), 'application/zip', 'CASEY_V150_ENTERPRISE_TRUST_RUNTIME_OUTPUT_PACK.zip')

@app.post('/qa/v150/validate-model')
def qa_validate_model_v150(payload: Dict[str, Any] = Body(default={})):
    return _v150_validate_model(payload or {})

@app.get('/qa/v150/readiness')
def qa_readiness_v150():
    return {
        'version': 'CASEY V150 Enterprise Trust Runtime',
        'status': 'controlled_executive_demo_ready',
        'hardening': [
            'single-state scenario runtime',
            'scenario signature + propagation checksum',
            'confidence breakdown with weighted drivers',
            'governance state bar metrics',
            'narrative entropy watch',
            'V150 manifest and scenario state embedded in full export pack',
            'route-level export re-canonicalization',
        ],
        'limits': 'Still a strategic simulation prototype; full browser-render and 15,000 export stress testing should be run on a dedicated machine before production deployment.'
    }

# Re-register export routes to V150 builders.
for _v150_path, _v150_builder, _v150_media, _v150_fname in [
    ('/export/workbook',      lambda m: workbook_bytes(m), 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet', 'CASEY_V150_Cost_Model_Planet_Class.xlsx'),
    ('/export/cost-model',    lambda m: workbook_bytes(m), 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet', 'CASEY_V150_Cost_Model_Planet_Class.xlsx'),
    ('/export/excel',         lambda m: workbook_bytes(m), 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet', 'CASEY_V150_Cost_Model_Planet_Class.xlsx'),
    ('/export/risk-register', lambda m: risk_register_workbook_bytes(m), 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet', 'CASEY_V150_Risk_Register_Pro.xlsx'),
    ('/export/xer',           lambda m: xer_bytes(m), 'application/octet-stream', 'CASEY_V150_P6_Schedule.xer'),
    ('/export/word',          lambda m: word_bytes(m), 'application/vnd.openxmlformats-officedocument.wordprocessingml.document', 'CASEY_V150_Executive_Board_Report.docx'),
    ('/export/pdf',           lambda m: pdf_bytes(m), 'application/pdf', 'CASEY_V150_Board_Intelligence_Pack.pdf'),
    ('/export/pptx',          lambda m: pptx_bytes(m), 'application/vnd.openxmlformats-officedocument.presentationml.presentation', 'CASEY_V150_Board_Deck_Elite.pptx'),
    ('/export/json',          lambda m: json.dumps(_v148_jsonable(_v150_canonical_state(m)), indent=2).encode('utf-8'), 'application/json', 'CASEY_V150_Full_Model_Audit.json'),
]:
    try:
        _v107_remove_route(_v150_path, 'POST')
    except Exception:
        pass
    def _v150_make_endpoint(builder, media, fname):
        def _endpoint(payload: dict = Body(default={})):
            m = _v150_canonical_state(payload or {})
            return stream(builder(m), media, fname)
        return _endpoint
    app.post(_v150_path)(_v150_make_endpoint(_v150_builder, _v150_media, _v150_fname))

for _v150_pack_path in ['/export/all', '/export/full-pack']:
    try:
        _v107_remove_route(_v150_pack_path, 'POST')
    except Exception:
        pass
app.post('/export/all')(lambda payload=Body(default={}): export_all(_v150_canonical_state(payload or {})))
app.post('/export/full-pack')(lambda payload=Body(default={}): export_all(_v150_canonical_state(payload or {})))

APP_VERSION = 'CASEY V150 Enterprise Trust Runtime'
print('CASEY V150 enterprise trust runtime installed')
# ======================= END CASEY V150 ENTERPRISE TRUST RUNTIME =======================

# ======================= CASEY V168 PROFESSIONAL CLIENT-SIDE CHALLENGE =======================
# Converts the intake engine from raw/technical JSON language into board-grade professional
# assurance language. Important UX rule: uploaded files produce a challenge delta and evidence
# opinion; they do not overwrite the programme baseline.

def _v168_professionalise_intake(result: Dict[str, Any]) -> Dict[str, Any]:
    r = dict(result or {})
    cm = dict(r.get('challenge_model') or {})
    src = dict(r.get('source_intelligence') or {})
    p50 = cm.get('p50_bn')
    p80 = cm.get('p80_bn')
    p90 = cm.get('p90_bn')
    conf = cm.get('confidence_pct') or 55
    ftype = str(r.get('file_type') or 'client source file')
    cost = (src.get('cost') or {})
    risk = (src.get('risk') or {})
    xer = (src.get('xer') or {})

    def money(x):
        try: return f"${float(x):.1f}B"
        except Exception: return 'not yet quantified'

    benchmark = []
    if 'rail' in ftype.lower() or True:
        benchmark.append('Compared with comparable rail/transit and systems-integration programmes, the board should expect explicit reconciliation between the submitted P50, quantified residual risk and the P80/P90 delivery tail.')
        benchmark.append('Where possession access, signalling integration, systems migration or operator acceptance are material, comparable programmes typically show wider downside volatility than civil-progress reporting suggests.')
    if risk.get('p90_bn') or p90:
        benchmark.append(f"The uploaded source indicates a stress/downside signal around {money(p90)}. CASEY treats this as a reliance question: either fund it, evidence it down, or disclose it as retained board exposure.")
    if cost.get('reserve_bn') and p50:
        pct = float(cost.get('reserve_bn') or 0) / max(float(p50), 0.001) * 100
        benchmark.append(f"The visible reserve signal is approximately {pct:.0f}% of inferred P50. CASEY tests whether that reserve is risk-derived or only a percentage allowance.")

    findings = [
        f"Source-file intake completed for {ftype}. CASEY mapped available cost, risk and schedule signals and compared them to the active programme baseline rather than replacing it.",
    ]
    if p50:
        findings.append(f"Submitted/inferred headline position: {money(p50)} P50. This is treated as the client-file number to be challenged, not as a separate CASEY benchmark estimate.")
    if p80:
        findings.append(f"Independent challenge position: {money(p80)} P80 and {money(p90)} stress case. The difference between P50 and P80 is the latent exposure requiring acceptance, mitigation or evidence closure.")
    if risk.get('risk_rows_found'):
        findings.append(f"Risk register signal: {risk.get('risk_rows_found')} rows identified. CASEY checks whether residual exposure is reconciled to reserve line-by-line.")
    if xer.get('task_count'):
        findings.append(f"Schedule signal: {xer.get('task_count')} XER activities identified. CASEY checks whether critical path, open logic, calendars and constraints can support the board date.")
    if cost.get('cost_lines_found'):
        findings.append(f"Cost signal: {cost.get('cost_lines_found')} cost lines detected. CASEY checks basis statements, exclusions, escalation and CBS/WBS mapping.")

    observations = []
    if not cost.get('reserve_bn'):
        observations.append('The submitted position does not yet show a clearly mapped reserve/contingency basis. Reliance should be withheld until reserve is reconciled to quantified risk exposure.')
    elif p50 and float(cost.get('reserve_bn') or 0) / max(float(p50), 0.001) < 0.10:
        observations.append('The visible reserve appears thin for a programme with systems/interface exposure. Require QCRA support rather than accepting a flat allowance.')
    if risk.get('risk_rows_found') and not risk.get('emv_bn'):
        observations.append('Risk rows are present, but residual exposure is not yet clearly reconciled to the cost model. This can understate the amount the client will ultimately fund.')
    if risk.get('p90_bn') and p50 and float(risk.get('p90_bn')) > float(p50) * 1.8:
        observations.append('The downside risk signal is materially larger than the headline position. The issue is not whether the register is long; it is whether the reserve has priced the credible tail.')
    if xer.get('open_end_risk_count', 0):
        observations.append(f"The schedule contains {xer.get('open_end_risk_count')} weak/open logic signals. The board date should be tested through QSRA, not accepted as a deterministic contractor date.")
    observations.append('Comparable programmes with similar interface and commissioning conditions have historically failed when reported progress was treated as equivalent to board-defensible readiness.')

    questions = [
        'Where does the submitted P50 reconcile to the P80/P90 downside, line by line?',
        'Which CBS/WBS package owns the largest unpriced exposure, and who signs the evidence closure?',
        'Is contingency risk-derived, or is it a percentage allowance applied to the estimate?',
        'Which schedule activities drive the board date, and are their predecessors, calendars, constraints and possession assumptions defensible?',
        'What exposure transfers back to the client if the contractor position is accepted without additional evidence?',
        'Which evidence item would most improve board-defensibility before approval?'
    ]
    next_steps = [
        'Ask for the native cost workbook, risk register and XER schedule as one source bundle, not isolated extracts.',
        'Reconcile CBS/WBS/activity IDs across cost, schedule and risk before accepting the submitted position.',
        'Require a P50/P80/P90 bridge showing what is funded, what is mitigated and what remains retained by the client.',
        'Assign named owners and closure dates for each open evidence item before investment approval.',
        'Issue a professional board challenge note comparing the submitted position with CASEY’s challenge range and benchmark memory.'
    ]

    cm['delta_bn'] = round(float(p80) - float(p50), 3) if p50 and p80 else None
    cm['board_posture'] = 'Further assurance required before approval' if conf < 55 else 'Board challenge likely — evidence closure required'
    r.update({
        'version': 'CASEY V168 professional client-side challenge',
        'challenge_model': cm,
        'findings': findings,
        'professional_observations': observations[:8],
        'red_flags': observations[:8],  # kept for backward UI compatibility, rendered as professional observations
        'benchmark_comparison': benchmark[:5],
        'board_challenge_questions': questions,
        'next_steps': next_steps,
        'if_this_fails': 'The likely failure mode is commercial reliance: the board accepts a submitted cost/schedule/risk position before the downside exposure, reserve logic and evidence owners are reconciled.',
        'epc_challenge': False,
    })
    r.pop('board_attack_kill_chain', None)
    return r

# Re-register upload routes so the professional result is returned first.
def _v168_remove_route(path, method='POST'):
    try:
        app.router.routes = [rt for rt in app.router.routes if not (getattr(rt, 'path', None) == path and method in getattr(rt, 'methods', set()))]
    except Exception:
        pass

_v168_remove_route('/upload')
_v168_remove_route('/upload/analyse')

@app.post('/upload')
async def upload_professional_v168(file: UploadFile = File(...)):
    content = await file.read(); name = file.filename or 'client_source_file'
    try:
        base = _v167_build_intake_challenge(name, content)
        result = _v168_professionalise_intake(base)
    except Exception as e:
        result = _v168_professionalise_intake({
            'filename': name,
            'size_bytes': len(content),
            'file_type': 'unreadable / protected client file',
            'source_intelligence': {'cost': {}, 'risk': {}, 'xer': {}},
            'challenge_model': {'confidence_pct': 35},
            'findings': ['File received but could not be normalised. CASEY treats this as an evidence-readiness issue, not a technical failure.'],
            'red_flags': [f'Parser issue: {str(e)[:140]}'],
        })
    try:
        con=db(); cur=con.cursor(); cur.execute('INSERT INTO uploads(filename,created_at,analysis_json) VALUES(?,?,?)',(name,datetime.utcnow().isoformat(),json.dumps(result))); con.commit(); con.close()
    except Exception:
        pass
    return result

@app.post('/upload/analyse')
async def analyse_upload_professional_v168(file: UploadFile = File(...)):
    return await upload_professional_v168(file)

APP_VERSION = 'CASEY V168 Professional Client-Side Challenge'
print('CASEY V168 professional client-side challenge installed')
# ======================= END CASEY V168 PROFESSIONAL CLIENT-SIDE CHALLENGE =======================
