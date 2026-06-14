"""
CASEY main.py — COMPLETE ADDITIONS
====================================
Add all of these to your existing main.py.
Place AFTER your existing /export/pdf route.

STEP 1: Add these imports near the top of main.py:

    import subprocess, tempfile, json, os, base64
    from fastapi import Request, UploadFile, File, Form
    from fastapi.responses import Response as FastAPIResponse

STEP 2: Install Python dependencies:
    pip install anthropic pdfplumber python-pptx python-docx openpyxl --break-system-packages

STEP 3: npm install in backend folder (for pptxgenjs):
    npm install pptxgenjs

STEP 4: Set Render environment variable:
    ANTHROPIC_API_KEY = sk-ant-api03-...
"""

import subprocess, tempfile, json, os, base64
from fastapi import Request, UploadFile, File, Form
from fastapi.responses import Response as FastAPIResponse


# ════════════════════════════════════════════════════════════════
# ROUTE 1: BOARD PACK PPTX
# 13-slide world-class PowerPoint — any sector, country, currency
# ════════════════════════════════════════════════════════════════

@app.post("/export/board-pack-pptx")
async def export_board_pack_pptx(request: Request):
    try:
        payload = await request.json()
        with tempfile.NamedTemporaryFile(mode='w', suffix='.json', delete=False, encoding='utf-8') as f:
            json.dump(payload, f, ensure_ascii=False, default=str)
            model_path = f.name
        output_path = model_path.replace('.json', '.pptx')
        script_path = os.path.join(os.path.dirname(__file__), 'generate_board_pack.js')
        if not os.path.exists(script_path):
            return {"error": "generate_board_pack.js not found in backend folder"}
        result = subprocess.run(
            ['node', script_path, model_path, output_path],
            capture_output=True, text=True, timeout=60,
            cwd=os.path.dirname(script_path)
        )
        if result.returncode != 0:
            return {"error": f"PPTX generation failed: {result.stderr or result.stdout}"}
        if not os.path.exists(output_path):
            return {"error": "PPTX file was not created"}
        with open(output_path, 'rb') as f:
            content = f.read()
        try:
            os.unlink(model_path)
            os.unlink(output_path)
        except Exception:
            pass
        title = str(payload.get('title') or payload.get('subsector') or 'CASEY')
        safe = ''.join(c if c.isalnum() or c in '-_ ' else '_' for c in title)[:50]
        return FastAPIResponse(
            content=content,
            media_type='application/vnd.openxmlformats-officedocument.presentationml.presentation',
            headers={
                'Content-Disposition': f'attachment; filename="CASEY_{safe}_Board_Pack.pptx"',
                'Content-Length': str(len(content)),
            }
        )
    except subprocess.TimeoutExpired:
        return {"error": "Board pack timed out — check Node.js is available on Render"}
    except Exception as e:
        return {"error": f"Board pack failed: {str(e)}"}


# ════════════════════════════════════════════════════════════════
# ROUTE 2: DOCUMENT CHALLENGE
# Upload any board pack PDF/PPTX/DOCX — CASEY reads it and
# returns what's missing, what the board will challenge,
# and a rebuilt confidence score
# ════════════════════════════════════════════════════════════════

@app.post("/advisor/challenge-document")
async def challenge_document(request: Request):
    """
    Accepts JSON with:
      - file_b64: base64-encoded file content
      - file_name: original filename
      - file_type: MIME type
      - model: optional CASEY model context (title, cost, confidence etc)
    Returns:
      - summary: one paragraph of findings
      - critical_gaps: list of {category, finding, recommendation}
      - board_questions: list of strings
      - missing_elements: list of strings
    """
    try:
        import anthropic
        payload = await request.json()
        file_b64 = payload.get('file_b64', '')
        file_name = payload.get('file_name', 'document')
        file_type = payload.get('file_type', 'application/pdf')
        model_ctx = payload.get('model', {})

        # Extract text from document
        doc_text = ""
        if file_b64:
            file_bytes = base64.b64decode(file_b64)
            fname_lower = file_name.lower()

            if fname_lower.endswith('.pdf'):
                try:
                    import pdfplumber, io
                    with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
                        pages = [p.extract_text() or '' for p in pdf.pages[:20]]
                        doc_text = '\n'.join(pages)[:12000]
                except ImportError:
                    doc_text = f"[PDF: {file_name} — pdfplumber not installed]"

            elif fname_lower.endswith(('.pptx', '.ppt')):
                try:
                    from pptx import Presentation
                    import io
                    prs = Presentation(io.BytesIO(file_bytes))
                    texts = []
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, 'text'):
                                texts.append(shape.text)
                    doc_text = '\n'.join(texts)[:12000]
                except ImportError:
                    doc_text = f"[PPTX: {file_name} — python-pptx not installed]"

            elif fname_lower.endswith(('.docx', '.doc')):
                try:
                    from docx import Document
                    import io
                    doc = Document(io.BytesIO(file_bytes))
                    doc_text = '\n'.join(p.text for p in doc.paragraphs)[:12000]
                except ImportError:
                    doc_text = f"[DOCX: {file_name} — python-docx not installed]"

        # Build context string from model
        model_context = ""
        if model_ctx:
            model_context = f"""
The programme being reviewed:
- Title: {model_ctx.get('title') or model_ctx.get('subsector', 'Unknown')}
- CASEY P50 estimate: {model_ctx.get('cost_p50', 'Unknown')}
- CASEY confidence score: {model_ctx.get('confidence_pct', 'Unknown')}%
- Schedule: {model_ctx.get('schedule', 'Unknown')}
"""

        # Challenge prompt — CASEY acting as hostile board examiner
        prompt = f"""You are CASEY, the world's most rigorous programme intelligence system. 
A user has uploaded their existing board pack for you to challenge.
{model_context}

DOCUMENT CONTENT:
{doc_text if doc_text else '[No text extracted — responding based on document metadata only]'}

Your job: act as a hostile, expert investment committee examiner. 
Review this document and identify exactly what will fail in a board room.

Return ONLY valid JSON in this exact structure (no markdown, no preamble):
{{
  "summary": "Two sentences summarising the document and the overall board-readiness verdict.",
  "critical_gaps": [
    {{
      "category": "Category name (e.g. QCRA/QSRA, Governing Constraint, Benchmark, Confidence Score, Evidence Gate)",
      "finding": "What is missing or wrong in one sentence.",
      "recommendation": "What CASEY would require to fix it."
    }}
  ],
  "board_questions": [
    "Question the board will ask that this pack cannot answer"
  ],
  "missing_elements": [
    "Element missing from this pack"
  ],
  "confidence_gap": "One sentence on how far this pack is from 75% board confidence threshold"
}}

Generate 3-5 critical gaps, 5 board questions, 3-4 missing elements.
Be specific, technical and unsparing. Do not be polite."""

        client = anthropic.Anthropic()
        msg = client.messages.create(
            model="claude-sonnet-4-6",
            max_tokens=2000,
            messages=[{"role": "user", "content": prompt}]
        )
        raw = msg.content[0].text.strip()
        # Strip markdown fences if present
        if raw.startswith('```'):
            raw = raw.split('```')[1]
            if raw.startswith('json'):
                raw = raw[4:]
        result = json.loads(raw)
        return result

    except ImportError:
        return {"error": "anthropic not installed — pip install anthropic"}
    except json.JSONDecodeError as e:
        return {"error": f"Failed to parse Claude response: {str(e)}", "raw": raw[:500] if 'raw' in dir() else ""}
    except Exception as e:
        return {"error": f"Challenge failed: {str(e)}"}


# ════════════════════════════════════════════════════════════════
# ROUTE 3: MONTHLY ACTUALS INGESTION
# Accepts monthly progress data → returns updated model with
# actual_progress_t (0-1) so the timeline head moves to real position
# ════════════════════════════════════════════════════════════════

@app.post("/actuals/ingest")
async def ingest_actuals(request: Request):
    """
    Accepts JSON with:
      - model: the CASEY model
      - actuals: list of {date, spend_bn, milestone, confidence, note}
    Returns:
      - updated_model: model with actual_progress_t, actual_spend_bn, actual_confidence
      - timeline_head: float 0-1 for where to place the timeline head
      - narrative: one sentence describing the actual vs planned position
    """
    try:
        payload = await request.json()
        model = payload.get('model', {})
        actuals = payload.get('actuals', [])

        if not actuals:
            return {"error": "No actuals provided", "timeline_head": 0}

        # Get the latest actual
        latest = actuals[0]  # Already sorted newest-first from frontend

        # Parse the total programme
        total_months_str = str(model.get('schedule', '') or model.get('schedule_months', '') or '24')
        total_months_match = __import__('re').search(r'\d+', total_months_str)
        total_months = int(total_months_match.group()) if total_months_match else 24

        # Get start date
        start_str = model.get('start_date', '')
        from datetime import datetime
        try:
            if start_str:
                start = datetime.fromisoformat(start_str.replace('Z', ''))
            else:
                start = datetime.now().replace(day=1)
        except Exception:
            start = datetime.now().replace(day=1)

        # Calculate months elapsed to the actual report date
        try:
            report_date_str = latest.get('date', '')
            if report_date_str:
                # Handle YYYY-MM format from month input
                if len(report_date_str) == 7:
                    report_date_str += '-01'
                report_date = datetime.fromisoformat(report_date_str)
                months_elapsed = (report_date.year - start.year) * 12 + (report_date.month - start.month)
                months_elapsed = max(0, min(months_elapsed, total_months))
            else:
                months_elapsed = 0
        except Exception:
            months_elapsed = 0

        # Calculate progress as fraction of total programme
        timeline_head = months_elapsed / max(total_months, 1)

        # Build updated model
        updated = dict(model)
        updated['actual_progress_t'] = round(timeline_head, 4)
        updated['actual_spend_bn'] = float(latest.get('spend_bn', 0))
        updated['actual_confidence'] = int(latest.get('confidence', model.get('confidence_pct', 60)))
        updated['actual_milestone'] = latest.get('milestone', '')
        updated['actual_report_date'] = latest.get('date', '')
        updated['_actualsLoaded'] = True

        # Simple narrative
        planned_pct = round(timeline_head * 100)
        spend_bn = float(latest.get('spend_bn', 0))
        p50_bn = float(model.get('cost_p50_bn') or 1)
        spend_pct = round(spend_bn / p50_bn * 100) if p50_bn > 0 else 0

        if spend_pct > planned_pct + 10:
            narrative = f"Programme is {months_elapsed} months in ({planned_pct}% of baseline) but has consumed {spend_pct}% of P50 budget — cost overrun signal."
        elif spend_pct < planned_pct - 10:
            narrative = f"Programme is {months_elapsed} months in ({planned_pct}% of baseline) with {spend_pct}% budget consumed — under-spend may indicate programme delay, not efficiency."
        else:
            narrative = f"Programme is {months_elapsed} months in ({planned_pct}% of baseline). Spend and schedule appear broadly aligned at {spend_pct}% of P50 budget consumed."

        return {
            "updated_model": updated,
            "timeline_head": timeline_head,
            "months_elapsed": months_elapsed,
            "narrative": narrative,
            "latest_actual": latest,
        }

    except Exception as e:
        return {"error": f"Actuals ingestion failed: {str(e)}", "timeline_head": 0}


# ════════════════════════════════════════════════════════════════
# ROUTE 4: WORKBOOK WITH COVER TAB + RISK HEATMAP
# Returns XLSX with a proper cover tab and risk heatmap tab
# ════════════════════════════════════════════════════════════════

@app.post("/export/workbook-with-cover")
async def export_workbook_with_cover(request: Request):
    """
    Generates an XLSX workbook with:
      Tab 1 — Cover (programme name, P50/P80/P90, confidence, RAG, date, CASEY branding)
      Tab 2 — Risk Heatmap (5×5 grid with risks plotted by probability × impact)
      Tab 3+ — Existing cost model data
    """
    try:
        import openpyxl
        from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
        from openpyxl.utils import get_column_letter

        payload = await request.json()
        model = payload
        curr = model.get('currency_symbol', '£')

        wb = openpyxl.Workbook()

        # ── COVER TAB ──────────────────────────────────────────────────
        ws_cover = wb.active
        ws_cover.title = "Cover"

        # Colours
        BG_DARK  = "03060C"  # CASEY dark background
        CYAN     = "8DF7FF"
        AMBER    = "F59E0B"
        GREEN    = "10B981"
        RED      = "EF4444"
        WHITE    = "FFFFFF"
        GREY     = "94A3B8"
        MIDGREY  = "1E293B"

        conf = int(model.get('confidence_pct', 60))
        rag = "GREEN" if conf >= 75 else "AMBER" if conf >= 55 else "RED"
        rag_color = GREEN if conf >= 75 else AMBER if conf >= 55 else RED
        rag_verdict = "Ready for capital commitment" if conf >= 75 else f"Conditional — {75-conf}pts below 75% board threshold" if conf >= 55 else "Do not approve without recovery plan"

        def cell(ws, row, col, value, bold=False, size=11, color=WHITE, bg=None, align='left', wrap=False):
            c = ws.cell(row=row, column=col, value=value)
            c.font = Font(name='Calibri', bold=bold, size=size, color=color)
            c.alignment = Alignment(horizontal=align, vertical='center', wrap_text=wrap)
            if bg:
                c.fill = PatternFill("solid", fgColor=bg)
            return c

        # Set column widths
        ws_cover.column_dimensions['A'].width = 28
        ws_cover.column_dimensions['B'].width = 36
        ws_cover.column_dimensions['C'].width = 22
        ws_cover.column_dimensions['D'].width = 18
        for row in range(1, 40):
            ws_cover.row_dimensions[row].height = 16

        # Header
        ws_cover.row_dimensions[1].height = 48
        cell(ws_cover, 1, 1, "CASEY PROGRAMME INTELLIGENCE", bold=True, size=16, color=CYAN, bg=BG_DARK, align='center')
        ws_cover.merge_cells('A1:D1')

        ws_cover.row_dimensions[2].height = 28
        title = str(model.get('title') or model.get('subsector') or 'Programme')
        cell(ws_cover, 2, 1, title, bold=True, size=13, color=WHITE, bg=MIDGREY, align='center')
        ws_cover.merge_cells('A2:D2')

        # RAG verdict
        ws_cover.row_dimensions[3].height = 32
        cell(ws_cover, 3, 1, rag + " — " + rag_verdict, bold=True, size=12, color=BG_DARK, bg=rag_color, align='center')
        ws_cover.merge_cells('A3:D3')

        # Key metrics
        row = 5
        metrics = [
            ("P50 Cost (Baseline)", str(model.get('cost_p50', '—'))),
            ("P80 Cost (Board approve at)", str(model.get('cost_p80', '—'))),
            ("P90 Cost (Stress case)", str(model.get('cost_p90', '—'))),
            ("Schedule", str(model.get('schedule', '—'))),
            ("Confidence Score", f"{conf}%"),
            ("Board Threshold", "75%"),
            ("Gap to Threshold", f"{max(0, 75-conf)} pts" if conf < 75 else "None — board ready"),
            ("Estimate Class", str(model.get('estimate_class_name', f"Class {model.get('estimate_class', 3)}"))),
            ("Schedule Level", str(model.get('schedule_level_name', f"Level {model.get('schedule_level', 4)}"))),
            ("Scenario Active", str(model.get('scenario_label', 'Base'))),
            ("Location", str(model.get('location', '—'))),
            ("Sector", str(model.get('subsector', model.get('mode', '—')))),
            ("Report Date", __import__('datetime').datetime.now().strftime('%d %b %Y')),
            ("Prepared by", "CASEY Programme Intelligence"),
            ("Governing Constraint", str(model.get('governing_constraint_prominent', '—'))[:80]),
        ]
        for label, value in metrics:
            cell(ws_cover, row, 1, label, bold=True, size=10, color=GREY, bg=BG_DARK)
            cell(ws_cover, row, 2, value, size=11, color=WHITE, bg=MIDGREY, wrap=True)
            row += 1

        # ── RISK HEATMAP TAB ──────────────────────────────────────────
        ws_heat = wb.create_sheet("Risk Heatmap")
        ws_heat.column_dimensions['A'].width = 6
        for col in ['B','C','D','E','F']:
            ws_heat.column_dimensions[col].width = 20
        for r in range(1, 30):
            ws_heat.row_dimensions[r].height = 18

        PROB_LABELS  = ['Almost Certain', 'Likely', 'Possible', 'Unlikely', 'Rare']
        IMPACT_LABELS= ['Negligible', 'Minor', 'Moderate', 'Major', 'Critical']
        CELL_COLORS  = [  # row=prob (high to low), col=impact (low to high)
            ["C6EFCE","C6EFCE","FFEB9C","FFC7CE","FFC7CE"],  # Almost certain
            ["C6EFCE","C6EFCE","FFEB9C","FFC7CE","FFC7CE"],  # Likely
            ["C6EFCE","FFEB9C","FFEB9C","FFC7CE","FFC7CE"],  # Possible
            ["C6EFCE","C6EFCE","FFEB9C","FFEB9C","FFC7CE"],  # Unlikely
            ["C6EFCE","C6EFCE","C6EFCE","FFEB9C","FFEB9C"],  # Rare
        ]

        # Header row
        cell(ws_heat, 1, 1, "Risk Heatmap", bold=True, size=12, color=WHITE, bg=BG_DARK)
        ws_heat.merge_cells('A1:F1')
        for ci, il in enumerate(IMPACT_LABELS):
            cell(ws_heat, 2, ci+2, il, bold=True, size=9, color=WHITE, bg=MIDGREY, align='center')

        # Grid
        for ri, pl in enumerate(PROB_LABELS):
            cell(ws_heat, ri+3, 1, pl, bold=True, size=8, color=WHITE, bg=MIDGREY, align='right')
            for ci in range(5):
                bg_hex = CELL_COLORS[ri][ci]
                ws_heat.cell(row=ri+3, column=ci+2).fill = PatternFill("solid", fgColor=bg_hex)

        # Plot risks on heatmap
        def get_prob_idx(prob_str):
            p = str(prob_str).lower()
            if any(x in p for x in ['almost','very high','certain']): return 0
            if any(x in p for x in ['high','likely']): return 1
            if any(x in p for x in ['medium','moderate','possible']): return 2
            if any(x in p for x in ['low','unlikely']): return 3
            return 4

        def get_impact_idx(imp_str):
            im = str(imp_str).lower()
            if any(x in im for x in ['critical','catastrophic']): return 4
            if any(x in im for x in ['major','significant','high']): return 3
            if any(x in im for x in ['moderate','medium']): return 2
            if any(x in im for x in ['minor','low']): return 1
            return 0

        risks = model.get('risks', model.get('risk_register', []))
        cell_texts = [[[] for _ in range(5)] for _ in range(5)]
        for risk in risks:
            pi = get_prob_idx(risk.get('probability', ''))
            ii = get_impact_idx(risk.get('impact', risk.get('consequence', '')))
            title_r = str(risk.get('title', risk.get('risk', '?')))[:20]
            cell_texts[pi][ii].append(title_r)

        for ri in range(5):
            for ci in range(5):
                if cell_texts[ri][ci]:
                    existing = ws_heat.cell(row=ri+3, column=ci+2).value or ''
                    ws_heat.cell(row=ri+3, column=ci+2).value = '\n'.join(cell_texts[ri][ci])
                    ws_heat.cell(row=ri+3, column=ci+2).font = Font(size=7, bold=True)
                    ws_heat.cell(row=ri+3, column=ci+2).alignment = Alignment(horizontal='center', vertical='center', wrap_text=True)
                    ws_heat.row_dimensions[ri+3].height = max(18, 14*len(cell_texts[ri][ci]))

        # ── RISK REGISTER TAB ──────────────────────────────────────────
        ws_risk = wb.create_sheet("Risk Register")
        headers = ['#','Risk Title','Probability','Impact','Schedule (wk)','Cost EMV (Bn)','Owner','Mitigation','Trigger','RAG']
        for ci, h in enumerate(headers, 1):
            c = ws_risk.cell(row=1, column=ci, value=h)
            c.font = Font(bold=True, size=9, color=WHITE)
            c.fill = PatternFill("solid", fgColor=BG_DARK)
            c.alignment = Alignment(horizontal='center')

        for ri, risk in enumerate(risks, 2):
            conf_r = str(risk.get('probability','')).lower()
            rag_r = "RED" if any(x in conf_r for x in ['high','likely','almost','critical']) else "AMBER" if any(x in conf_r for x in ['medium','moderate','possible']) else "GREEN"
            rag_c_r = FFC7CE if rag_r=="RED" else FFEB9C if rag_r=="AMBER" else C6EFCE
            row_vals = [
                ri-1,
                str(risk.get('title',risk.get('risk','?')))[:50],
                str(risk.get('probability','?')),
                str(risk.get('impact','?')),
                str(risk.get('schedule_impact_weeks','?')),
                str(risk.get('cost_emv_bn',risk.get('cost_impact_bn','?'))),
                str(risk.get('owner','?')),
                str(risk.get('mitigation','?'))[:80],
                str(risk.get('trigger','?'))[:60],
                rag_r,
            ]
            for ci, val in enumerate(row_vals, 1):
                c = ws_risk.cell(row=ri, column=ci, value=val)
                c.font = Font(size=9)
                c.alignment = Alignment(wrap_text=True, vertical='top')
                if ci == len(row_vals):
                    c.fill = PatternFill("solid", fgColor=rag_c_r if rag_r else "FFFFFF")

        for col in range(1, len(headers)+1):
            ws_risk.column_dimensions[get_column_letter(col)].width = 16

        # Save
        with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as f:
            xlsx_path = f.name
        wb.save(xlsx_path)
        with open(xlsx_path, 'rb') as f:
            content = f.read()
        try:
            os.unlink(xlsx_path)
        except Exception:
            pass

        title_s = str(model.get('title') or model.get('subsector') or 'CASEY')
        safe = ''.join(c if c.isalnum() or c in '-_ ' else '_' for c in title_s)[:50]

        return FastAPIResponse(
            content=content,
            media_type='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
            headers={
                'Content-Disposition': f'attachment; filename="CASEY_{safe}_Workbook.xlsx"',
                'Content-Length': str(len(content)),
            }
        )

    except ImportError:
        return {"error": "openpyxl not installed — pip install openpyxl"}
    except Exception as e:
        import traceback
        return {"error": f"Workbook failed: {str(e)}", "trace": traceback.format_exc()[-500:]}


# ════════════════════════════════════════════════════════════════
# ROUTE 5: AI TIMELINE NARRATIVE (optional)
# Called by the advisor what-if bar on the timeline
# ════════════════════════════════════════════════════════════════

@app.post("/ai/timeline-narrative")
async def timeline_narrative(request: Request):
    try:
        import anthropic
        payload = await request.json()
        prompt_text = payload.get('prompt', '')
        if not prompt_text:
            return {"text": ""}
        client = anthropic.Anthropic()
        msg = client.messages.create(
            model="claude-sonnet-4-6",
            max_tokens=150,
            messages=[{"role": "user", "content": prompt_text}]
        )
        return {"text": msg.content[0].text if msg.content else ""}
    except Exception as e:
        return {"text": "", "error": str(e)}


# ════════════════════════════════════════════════════════════════
# ROUTE 6: TEST NODE (verify pptxgenjs available — run once then remove)
# ════════════════════════════════════════════════════════════════

# @app.get("/test-node")
# async def test_node():
#     r = subprocess.run(['node', '--version'], capture_output=True, text=True)
#     r2 = subprocess.run(['node', '-e', 'require("pptxgenjs"); console.log("pptxgenjs OK")'], capture_output=True, text=True, cwd=os.path.dirname(__file__))
#     return {"node": r.stdout.strip(), "pptxgenjs": r2.stdout.strip() or r2.stderr.strip(), "ok": r.returncode == 0}


# ════════════════════════════════════════════════════════════════
# ROUTE 7: ADVISOR MEMORY — server-side storage
# Stores advisor conversations server-side, keyed by email + programme
# Called by the CASEY_Features useAdvisorMemory hook
# Requires: pip install anthropic (already done)
# ════════════════════════════════════════════════════════════════

import hashlib
from pathlib import Path

MEMORY_DIR = Path("/tmp/casey_memories")
MEMORY_DIR.mkdir(exist_ok=True)

def memory_path(email: str, programme_id: str) -> Path:
    key = hashlib.sha256(f"{email}:{programme_id}".encode()).hexdigest()[:16]
    return MEMORY_DIR / f"mem_{key}.json"

@app.post("/advisor/memory/save")
async def save_advisor_memory(request: Request):
    try:
        payload = await request.json()
        email = payload.get('email', 'anon')
        programme_id = payload.get('programme_id', 'global')
        messages = payload.get('messages', [])
        path = memory_path(email, programme_id)
        data = {"email": email, "programme_id": programme_id, "messages": messages[:40], "updated": __import__('datetime').datetime.now().isoformat()}
        path.write_text(json.dumps(data))
        return {"saved": True, "count": len(messages)}
    except Exception as e:
        return {"saved": False, "error": str(e)}

@app.post("/advisor/memory/load")
async def load_advisor_memory(request: Request):
    try:
        payload = await request.json()
        email = payload.get('email', 'anon')
        programme_id = payload.get('programme_id', 'global')
        path = memory_path(email, programme_id)
        if not path.exists():
            return {"messages": [], "found": False}
        data = json.loads(path.read_text())
        return {"messages": data.get("messages", []), "found": True, "updated": data.get("updated")}
    except Exception as e:
        return {"messages": [], "found": False, "error": str(e)}


# ════════════════════════════════════════════════════════════════
# ROUTE 8: VERSION HISTORY — server-side audit trail
# Stores every model version server-side, keyed by programme id
# This closes T&T's counter-argument about audit trails
# ════════════════════════════════════════════════════════════════

VERSIONS_DIR = Path("/tmp/casey_versions")
VERSIONS_DIR.mkdir(exist_ok=True)

@app.post("/versions/save")
async def save_version(request: Request):
    try:
        payload = await request.json()
        programme_id = str(payload.get('programme_id', 'unknown'))
        event = payload.get('event', 'Model updated')
        model_snapshot = payload.get('model', {})
        safe_id = ''.join(c if c.isalnum() or c in '-_' else '_' for c in programme_id)[:40]
        path = VERSIONS_DIR / f"v_{safe_id}.json"
        existing = []
        if path.exists():
            try: existing = json.loads(path.read_text())
            except: existing = []
        entry = {
            "id": __import__('time').time_ns(),
            "timestamp": __import__('datetime').datetime.now().isoformat(),
            "event": event,
            "confidence_pct": model_snapshot.get('confidence_pct'),
            "cost_p50": model_snapshot.get('cost_p50'),
            "cost_p80": model_snapshot.get('cost_p80'),
            "schedule": model_snapshot.get('schedule'),
            "scenario": model_snapshot.get('scenario_label') or model_snapshot.get('scenario', 'Base'),
            "estimate_class": model_snapshot.get('estimate_class_name') or f"Class {model_snapshot.get('estimate_class', 3)}",
            "programme_id": programme_id,
        }
        versions = [entry, *existing][:40]
        path.write_text(json.dumps(versions))
        return {"saved": True, "total": len(versions)}
    except Exception as e:
        return {"saved": False, "error": str(e)}

@app.get("/versions/{programme_id}")
async def get_versions(programme_id: str):
    try:
        safe_id = ''.join(c if c.isalnum() or c in '-_' else '_' for c in programme_id)[:40]
        path = VERSIONS_DIR / f"v_{safe_id}.json"
        if not path.exists():
            return {"versions": [], "found": False}
        versions = json.loads(path.read_text())
        return {"versions": versions, "found": True, "count": len(versions)}
    except Exception as e:
        return {"versions": [], "found": False, "error": str(e)}


# ════════════════════════════════════════════════════════════════
# ROUTE 9: PORTFOLIO — server-side saved programmes
# Backs up saved projects to server so they survive browser clears
# ════════════════════════════════════════════════════════════════

PORTFOLIO_DIR = Path("/tmp/casey_portfolios")
PORTFOLIO_DIR.mkdir(exist_ok=True)

@app.post("/portfolio/save")
async def save_portfolio(request: Request):
    try:
        payload = await request.json()
        email = str(payload.get('email', 'anon'))
        projects = payload.get('projects', [])
        safe = hashlib.sha256(email.encode()).hexdigest()[:16]
        path = PORTFOLIO_DIR / f"p_{safe}.json"
        data = {"email": email, "projects": projects[:30], "updated": __import__('datetime').datetime.now().isoformat()}
        path.write_text(json.dumps(data, default=str))
        return {"saved": True, "count": len(projects)}
    except Exception as e:
        return {"saved": False, "error": str(e)}

@app.post("/portfolio/load")
async def load_portfolio(request: Request):
    try:
        payload = await request.json()
        email = str(payload.get('email', 'anon'))
        safe = hashlib.sha256(email.encode()).hexdigest()[:16]
        path = PORTFOLIO_DIR / f"p_{safe}.json"
        if not path.exists():
            return {"projects": [], "found": False}
        data = json.loads(path.read_text())
        return {"projects": data.get("projects", []), "found": True, "updated": data.get("updated")}
    except Exception as e:
        return {"projects": [], "found": False, "error": str(e)}


"""
════════════════════════════════════════════════════════════════
FINAL RENDER DEPLOYMENT CHECKLIST
════════════════════════════════════════════════════════════════

You have the £20/month Render plan — backend stays awake 24/7.
No more cold start. Every demo opens instantly.

1. COPY FILES INTO backend/ FOLDER:
   main_py_COMPLETE.py  → add all routes to main.py
   generate_board_pack.js  → stays in backend/
   package.json (from backend_package.json)  → stays in backend/

2. INSTALL DEPENDENCIES:
   pip install anthropic pdfplumber python-pptx python-docx openpyxl --break-system-packages
   npm install pptxgenjs

3. RENDER ENVIRONMENT VARIABLES (Dashboard → Service → Environment):
   ANTHROPIC_API_KEY = sk-ant-api03-...
   VITE_ANTHROPIC_KEY = same key
   NODE_VERSION = 18

4. GIT PUSH:
   cd C:/Users/jaima/897/backend
   git add main.py generate_board_pack.js package.json
   git commit -m "V290 — complete platform: board pack, challenge, actuals, workbook, memory, versions, portfolio"
   git push

   cd C:/Users/jaima/897/frontend
   git add src/App.jsx src/ProjectTimeline.jsx src/CASEY_Upgrades.jsx src/CASEY_Features.jsx src/useTimelineAI.js
   git commit -m "V290 — typewriter verdict, count-up confidence, rotating challenge, version history, portfolio strip, actualProgress wired"
   git push

════════════════════════════════════════════════════════════════
"""

# ════════════════════════════════════════════════════════════════════════════
# ═══ NEW ROUTE A: PUBLIC API — XER / P6 INGESTION ═══════════════════════
# ════════════════════════════════════════════════════════════════════════════
#
# POST /api/ingest-xer
#
# What it does:
#   Accepts a raw Primavera P6 XER file (or JSON schedule payload)
#   → Parses activities, logic ties, critical path, WBS, resources
#   → Runs CASEY intelligence engine against the real schedule
#   → Returns a full CASEY model JSON (same shape as /generate)
#   → Client can then download board pack, risk register, QCRA/QSRA
#
# This is the feature that turns CASEY into infrastructure.
# A client connects their live P6 → Sunday night cron → updated forecast.
# T&T cannot replicate a product their clients use automatically.
#
# USAGE:
#   curl -X POST https://corbit-1.onrender.com/api/ingest-xer \
#     -H "Content-Type: application/json" \
#     -d '{"xer_content": "<raw XER text>", "currency": "£", "location": "UK"}'
#
# Or from the frontend (DocumentUpload component already handles file reading):
#   POST /api/ingest-xer
#   { "xer_content": "...", "xer_filename": "HS2.xer", "currency": "£" }
#
# ════════════════════════════════════════════════════════════════════════════

@app.post("/api/ingest-xer")
async def ingest_xer(request: Request):
    """
    Accept a Primavera P6 XER file → return a full CASEY model JSON.
    The model can then be used to generate board pack, risk register,
    QCRA/QSRA, timeline — all from the client's live schedule.
    """
    import re, datetime, anthropic

    try:
        payload = await request.json()
        xer_raw = payload.get('xer_content', '')
        xer_filename = payload.get('xer_filename', 'schedule.xer')
        currency = payload.get('currency', '£')
        location = payload.get('location', '')
        client_name = payload.get('client', '')

        if not xer_raw:
            return {"error": "No XER content provided. Send xer_content as a string."}

        # ── PARSE XER ──
        # XER files are tab-delimited with %T table markers and %F field headers
        tables = {}
        current_table = None
        current_fields = []

        for line in xer_raw.split('\n'):
            line = line.rstrip('\r')
            if line.startswith('%T\t'):
                current_table = line[3:].strip()
                tables[current_table] = []
                current_fields = []
            elif line.startswith('%F\t') and current_table:
                current_fields = line[3:].split('\t')
            elif line.startswith('%R\t') and current_table and current_fields:
                values = line[3:].split('\t')
                row = dict(zip(current_fields, values))
                tables[current_table].append(row)
            elif line.startswith('%E'):
                current_table = None

        # Extract key data
        activities = tables.get('TASK', [])
        project_rows = tables.get('PROJECT', [])
        wbs_rows = tables.get('PROJWBS', [])
        relations = tables.get('TASKPRED', [])
        resources = tables.get('TASKRSRC', [])

        # Project metadata
        project = project_rows[0] if project_rows else {}
        proj_name = project.get('proj_short_name', '') or project.get('proj_id', '') or xer_filename.replace('.xer','')

        # Activity stats
        total_acts = len(activities)
        critical_acts = [a for a in activities if a.get('driving_path_flag','') == 'Y' or a.get('cstr_type','') in ['CS_ALAP']]
        
        # Parse dates
        def parse_xer_date(s):
            if not s: return None
            for fmt in ['%Y-%m-%d %H:%M', '%Y-%m-%d', '%d-%b-%y']:
                try: return datetime.datetime.strptime(s.strip(), fmt)
                except: pass
            return None

        earliest_start = None
        latest_finish = None
        for a in activities:
            s = parse_xer_date(a.get('act_start_date') or a.get('early_start_date') or a.get('target_start_date',''))
            f = parse_xer_date(a.get('act_end_date') or a.get('early_end_date') or a.get('target_end_date',''))
            if s and (earliest_start is None or s < earliest_start): earliest_start = s
            if f and (latest_finish is None or f > latest_finish): latest_finish = f

        duration_months = 0
        if earliest_start and latest_finish:
            duration_months = round((latest_finish - earliest_start).days / 30.44)

        # Logic quality
        open_ends = sum(1 for a in activities if a.get('task_type','') not in ['TT_LOE','TT_WBS'] and
                       not any(r.get('pred_task_id') == a.get('task_id') or r.get('task_id') == a.get('task_id') for r in relations))

        logic_quality = 'GOOD' if open_ends < total_acts * 0.05 else \
                       'REVIEW' if open_ends < total_acts * 0.12 else 'POOR'

        # Build XER summary for Claude
        xer_summary = f"""
PARSED XER SCHEDULE — {proj_name}
Total activities: {total_acts}
Apparent critical activities: {len(critical_acts)}
Logic ties (relationships): {len(relations)}
Open ends (missing logic): {open_ends} ({round(open_ends/max(total_acts,1)*100)}%)
Logic quality: {logic_quality}
Earliest start: {earliest_start.strftime('%b %Y') if earliest_start else 'Unknown'}
Latest finish: {latest_finish.strftime('%b %Y') if latest_finish else 'Unknown'}
Duration: {duration_months} months
WBS levels: {len(wbs_rows)}
Resources: {len(resources)}
Location: {location or 'Not specified'}
Currency: {currency}
Client: {client_name or 'Not specified'}

TOP 15 ACTIVITIES (by name):
{chr(10).join([f"  - {a.get('task_name','?')[:60]}" for a in activities[:15]])}

WBS STRUCTURE (top levels):
{chr(10).join([f"  - {w.get('wbs_name','?')[:50]}" for w in wbs_rows[:8]])}
"""

        # ── CALL CLAUDE TO GENERATE CASEY MODEL FROM XER ──
        ai_client = anthropic.Anthropic()

        system_prompt = """You are CASEY, the world's most advanced programme intelligence system.
You have been given parsed data from a real Primavera P6 XER schedule.
Generate a complete CASEY intelligence model from this data.
Return ONLY valid JSON — no markdown, no preamble, no explanation.
The JSON must be a complete CASEY model that the frontend can render."""

        user_prompt = f"""Generate a complete CASEY programme intelligence model from this real P6 XER schedule.

{xer_summary}

Return a complete JSON model with ALL of these fields populated from the schedule data:
{{
  "programme_title": "Derived from project name",
  "title": "Short title",
  "subsector": "Inferred from activity names and WBS",
  "mode": "Earth",
  "location": "{location or 'Inferred from project'}",
  "currency_symbol": "{currency}",
  "start_date": "ISO date from earliest activity",
  "schedule": "{duration_months} months",
  "schedule_months": {duration_months},
  "cost_p50": "Estimated from sector benchmarks and duration",
  "cost_p50_bn": 0.0,
  "cost_p80": "P80 cost",
  "cost_p80_bn": 0.0,
  "cost_p90": "P90 cost",
  "direct_cost": "Direct works estimate",
  "indirect_cost": "Indirect / prelims",
  "p80_reserve": "Reserve",
  "p80_reserve_pct": 18,
  "confidence_pct": 62,
  "risk": "Medium",
  "estimate_class": 3,
  "estimate_class_name": "Class 3 — Budget",
  "schedule_level": 3,
  "schedule_level_name": "Level 3",
  "oba_pct": 35,
  "governing_constraint_prominent": "Identified from critical path",
  "institutional_authority_line": "One sentence board verdict",
  "monte_carlo": {{
    "qcra": {{"p10": "...", "p50": "...", "p80": "...", "p90": "..."}},
    "qsra": {{"p10": {duration_months-20}, "p50": {duration_months}, "p80": {duration_months+25}, "p90": {duration_months+45}}}
  }},
  "xer_health": {{
    "headline": "Schedule quality assessment",
    "activity_count": {total_acts},
    "critical_count": {len(critical_acts)},
    "logic_quality": "{logic_quality}",
    "float_quality": "REVIEW",
    "critical_pct": "{round(len(critical_acts)/max(total_acts,1)*100)}",
    "open_ends": {open_ends},
    "board_flag": "Key schedule risk identified from XER"
  }},
  "risks": [
    {{"title": "Risk 1 from schedule analysis", "probability": "High", "impact": "Critical", "cause": "...", "owner": "TBC", "mitigation": "...", "cost_emv_bn": 0.1, "schedule_impact_weeks": 8}},
    {{"title": "Risk 2", "probability": "Medium", "impact": "Major", "cause": "...", "owner": "TBC", "mitigation": "...", "cost_emv_bn": 0.05, "schedule_impact_weeks": 4}},
    {{"title": "Risk 3", "probability": "High", "impact": "Major", "cause": "...", "owner": "TBC", "mitigation": "...", "cost_emv_bn": 0.08, "schedule_impact_weeks": 6}},
    {{"title": "Risk 4", "probability": "Medium", "impact": "Moderate", "cause": "...", "owner": "TBC", "mitigation": "...", "cost_emv_bn": 0.03, "schedule_impact_weeks": 2}},
    {{"title": "Risk 5", "probability": "Low", "impact": "Major", "cause": "...", "owner": "TBC", "mitigation": "...", "cost_emv_bn": 0.04, "schedule_impact_weeks": 3}}
  ],
  "schedule_detail": [
    {{"activity": "Phase name", "start": "date", "end": "date", "critical": true}}
  ],
  "benchmark_comparison": [
    {{"name": "Comparable programme", "cost_bn": 0.0, "cost_growth_pct": 25, "schedule_slip_months": 12, "lesson": "..."}}
  ],
  "board_attack_questions": [
    "Board question 1?", "Board question 2?", "Board question 3?", "Board question 4?", "Board question 5?"
  ],
  "confidence_breakdown": [
    {{"driver": "Schedule maturity", "effect": "Class 3 XER uploaded", "delta": 6}},
    {{"driver": "Logic quality", "effect": "{logic_quality}", "delta": 4}}
  ],
  "prompt": "XER ingestion: {proj_name}",
  "source": "xer_ingestion",
  "xer_filename": "{xer_filename}",
  "xer_activity_count": {total_acts},
  "xer_open_ends": {open_ends}
}}

Make all values specific and plausible for the sector, location and duration.
Populate costs from sector benchmarks (not dummy zeros).
Identify real risks from the WBS and activity names provided."""

        response = ai_client.messages.create(
            model="claude-sonnet-4-6",
            max_tokens=4000,
            system=system_prompt,
            messages=[{"role": "user", "content": user_prompt}]
        )

        raw = response.content[0].text.strip()
        if raw.startswith('```'):
            raw = raw.split('```')[1]
            if raw.startswith('json'): raw = raw[4:]
        raw = raw.strip()

        model_json = json.loads(raw)

        # Add XER health data we computed
        model_json['xer_health'] = {
            'headline': f"Schedule parsed: {total_acts} activities · {logic_quality} logic · {open_ends} open ends",
            'activity_count': total_acts,
            'critical_count': len(critical_acts),
            'logic_quality': logic_quality,
            'float_quality': 'REVIEW' if open_ends > 5 else 'GOOD',
            'critical_pct': f"{round(len(critical_acts)/max(total_acts,1)*100)}%",
            'open_ends': open_ends,
            'board_flag': f"{open_ends} open-end activities reduce schedule confidence. Review logic before board submission." if open_ends > 3 else None,
        }
        model_json['source'] = 'xer_ingestion'
        model_json['xer_filename'] = xer_filename

        return {
            "model": model_json,
            "xer_stats": {
                "activities": total_acts,
                "relations": len(relations),
                "open_ends": open_ends,
                "logic_quality": logic_quality,
                "duration_months": duration_months,
                "start": earliest_start.isoformat() if earliest_start else None,
                "finish": latest_finish.isoformat() if latest_finish else None,
            },
            "message": f"XER parsed successfully. {total_acts} activities · {duration_months} months · {logic_quality} logic quality."
        }

    except json.JSONDecodeError as e:
        return {"error": f"Claude returned invalid JSON: {str(e)}", "raw_preview": (raw or '')[:300]}
    except ImportError:
        return {"error": "anthropic not installed — pip install anthropic"}
    except Exception as e:
        import traceback
        return {"error": f"XER ingestion failed: {str(e)}", "trace": traceback.format_exc()[-600:]}


# ════════════════════════════════════════════════════════════════════════════
# ═══ NEW ROUTE B: UPGRADE challenge-document with STREAMING ════════════════
# ════════════════════════════════════════════════════════════════════════════
#
# POST /advisor/challenge-document-stream
#
# Same as /advisor/challenge-document but returns a streaming response
# so the UI can show results appearing in real time (like Claude typing).
# The frontend shows each gap as it arrives rather than waiting for all.
#
# ════════════════════════════════════════════════════════════════════════════

from fastapi.responses import StreamingResponse
import asyncio

@app.post("/advisor/challenge-document-stream")
async def challenge_document_stream(request: Request):
    """
    Streaming version of document challenge.
    Returns server-sent events (SSE) so gaps appear in real time.
    """
    import anthropic

    try:
        payload = await request.json()
        file_b64 = payload.get('file_b64', '')
        file_name = payload.get('file_name', 'document')
        model_ctx = payload.get('model', {})

        doc_text = ""
        if file_b64:
            file_bytes = base64.b64decode(file_b64)
            fname_lower = file_name.lower()
            if fname_lower.endswith('.pdf'):
                try:
                    import pdfplumber, io
                    with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
                        pages = [p.extract_text() or '' for p in pdf.pages[:20]]
                        doc_text = '\n'.join(pages)[:14000]
                except ImportError:
                    doc_text = f"[PDF: {file_name}]"
            elif fname_lower.endswith(('.pptx', '.ppt')):
                try:
                    from pptx import Presentation
                    import io
                    prs = Presentation(io.BytesIO(file_bytes))
                    doc_text = '\n'.join(shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, 'text'))[:14000]
                except ImportError:
                    doc_text = f"[PPTX: {file_name}]"
            elif fname_lower.endswith(('.docx','.doc')):
                try:
                    from docx import Document
                    import io
                    doc = Document(io.BytesIO(file_bytes))
                    doc_text = '\n'.join(p.text for p in doc.paragraphs)[:14000]
                except ImportError:
                    doc_text = f"[DOCX: {file_name}]"

        model_info = ""
        if model_ctx:
            model_info = f"""
CASEY has already run this programme:
- Title: {model_ctx.get('title') or model_ctx.get('subsector', 'Unknown')}
- P50: {model_ctx.get('cost_p50', '?')} · P80: {model_ctx.get('cost_p80', '?')}
- Confidence: {model_ctx.get('confidence_pct', '?')}% · Schedule: {model_ctx.get('schedule', '?')}
- Governing constraint: {model_ctx.get('governing_constraint_prominent', 'Not identified')}
Compare the uploaded document against these CASEY findings."""

        prompt = f"""You are CASEY — the world's most rigorous investment committee examiner.
A user has uploaded a board pack / stage gate submission.
{model_info}

DOCUMENT:
{doc_text or '[No text extracted]'}

Act as a hostile, technically expert investment committee member.
Be specific. Reference actual content from the document where you can.
Do not be polite. Do not hedge.

Structure your response exactly as:

## VERDICT
[One clear sentence: is this board-ready? Yes / Conditional / No — and why in 10 words]

## CRITICAL GAPS
[For each gap, write: CATEGORY — Finding — Recommendation]

## BOARD QUESTIONS THIS PACK CANNOT ANSWER
[Number each question. These are the exact questions a board member will ask and this pack has no answer for]

## WHAT'S MISSING
[Bullet list of specific missing elements — be technical]

## REBUILT CONFIDENCE SCORE
[State what confidence score CASEY would assign to this pack, and what it would need to reach 75%]"""

        ai_client = anthropic.Anthropic()

        async def generate():
            with ai_client.messages.stream(
                model="claude-sonnet-4-6",
                max_tokens=2500,
                messages=[{"role": "user", "content": prompt}]
            ) as stream:
                for text in stream.text_stream:
                    yield f"data: {json.dumps({'text': text})}\n\n"
            yield "data: [DONE]\n\n"

        return StreamingResponse(
            generate(),
            media_type="text/event-stream",
            headers={"Cache-Control": "no-cache", "X-Accel-Buffering": "no"}
        )

    except Exception as e:
        return {"error": str(e)}


# ════════════════════════════════════════════════════════════════════════════
# ═══ NEW ROUTE C: PUBLIC API HEALTH + DOCS ═════════════════════════════════
# ════════════════════════════════════════════════════════════════════════════

@app.get("/api")
async def api_docs():
    """Public API documentation endpoint."""
    return {
        "name": "CASEY Programme Intelligence API",
        "version": "290",
        "description": "World-class programme intelligence — cost, schedule, QCRA, QSRA, risk, board pack from one sentence or XER file.",
        "endpoints": {
            "POST /generate": "Generate full CASEY intelligence pack from a text prompt",
            "POST /api/ingest-xer": "Upload Primavera P6 XER → full CASEY model JSON",
            "POST /advisor/challenge-document": "Upload board pack PDF/PPTX → board challenge analysis",
            "POST /advisor/challenge-document-stream": "Same, streaming — gaps appear in real time",
            "POST /export/board-pack-pptx": "Generate 13-slide board pack PPTX from CASEY model",
            "POST /export/workbook-with-cover": "Generate XLSX with cover tab, risk heatmap, risk register",
            "POST /actuals/ingest": "Monthly actuals → updated model with actual_progress_t",
            "POST /ai/timeline-narrative": "AI narrative for advisor what-if events",
            "POST /advisor/memory/save": "Save advisor conversation to server",
            "POST /advisor/memory/load": "Load advisor conversation from server",
            "POST /versions/save": "Save model version (audit trail)",
            "GET  /versions/{programme_id}": "Get version history",
            "POST /portfolio/save": "Save portfolio to server",
            "POST /portfolio/load": "Load portfolio from server",
        },
        "authentication": "No auth required for demo. Contact hello@controlorbit.com for API keys.",
        "rate_limits": "50 requests/hour per IP on free tier.",
        "contact": "hello@controlorbit.com",
        "docs": "https://controlorbit.com/api-docs",
    }


# ════════════════════════════════════════════════════════════════════════════
# ═══ FRONTEND: How to call the XER ingest from the Upload component ═════════
# ════════════════════════════════════════════════════════════════════════════
#
# In your DocumentUpload component (CASEY_Features.jsx), add this case:
#
#   if (f.name.endsWith('.xer') || f.name.endsWith('.XER')) {
#     const text = await f.text();  // XER is plain text
#     const resp = await fetch(`${apiBase}/api/ingest-xer`, {
#       method: 'POST',
#       headers: { 'Content-Type': 'application/json' },
#       body: JSON.stringify({
#         xer_content: text,
#         xer_filename: f.name,
#         currency: model?.currency_symbol || '£',
#         location: model?.location || '',
#       }),
#     });
#     const data = await resp.json();
#     if (data.model) {
#       onModelFromXER(data.model);  // pass to App to setModel()
#     }
#   }
#
# ════════════════════════════════════════════════════════════════════════════
